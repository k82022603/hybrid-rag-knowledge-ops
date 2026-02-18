#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Enterprise PPT Generator Template
KERIS 스타일 기반 전문 프레젠테이션 생성기

Usage:
    python template.py [테마번호] [출력파일명]
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

# ============== 테마 정의 ==============
THEMES = {
    1: {
        'name': 'Ocean Depths (오션 뎁스)',
        'description': '기업용, 전문적, 금융, 컨설팅',
        'primary': RGBColor(0, 102, 153),
        'secondary': RGBColor(0, 153, 204),
        'accent': RGBColor(0, 51, 102),
        'dark': RGBColor(33, 33, 33),
        'light': RGBColor(204, 229, 255),
        'light2': RGBColor(230, 242, 255),
        'header_text': RGBColor(255, 255, 255),
        'table_header': RGBColor(0, 102, 153),
        'table_alt': RGBColor(230, 242, 255),
        'cover_bg': RGBColor(0, 51, 102),
        'cover_deco': RGBColor(0, 102, 153),
        'cover_title': RGBColor(255, 255, 255),      # 화이트 (배경 대비)
        'box_fill': RGBColor(0, 102, 153),
        'box_fill2': RGBColor(0, 153, 204),
        'box_fill3': RGBColor(0, 51, 102),
        'arrow': RGBColor(255, 102, 0),
        'success': RGBColor(46, 125, 50),
        'warning': RGBColor(255, 152, 0),
        'danger': RGBColor(211, 47, 47),
    },
    2: {
        'name': 'Golden Hour (골든 아워)',
        'description': '따뜻한, 환대, 호텔, F&B, 럭셔리',
        'primary': RGBColor(255, 160, 0),
        'secondary': RGBColor(255, 111, 0),
        'accent': RGBColor(121, 85, 72),
        'dark': RGBColor(62, 39, 35),
        'light': RGBColor(255, 248, 225),
        'light2': RGBColor(255, 243, 224),
        'header_text': RGBColor(255, 248, 225),
        'table_header': RGBColor(141, 110, 99),
        'table_alt': RGBColor(239, 235, 233),
        'cover_bg': RGBColor(62, 39, 35),
        'cover_deco': RGBColor(121, 85, 72),
        'cover_title': RGBColor(255, 160, 0),        # 앰버 (대비 양호, 유지)
        'box_fill': RGBColor(121, 85, 72),
        'box_fill2': RGBColor(161, 136, 127),
        'box_fill3': RGBColor(62, 39, 35),
        'arrow': RGBColor(255, 160, 0),
        'success': RGBColor(46, 125, 50),
        'warning': RGBColor(255, 152, 0),
        'danger': RGBColor(211, 47, 47),
    },
    3: {
        'name': 'Tech Innovation (테크 이노베이션)',
        'description': '기술, 스타트업, IT, 개발',
        'primary': RGBColor(46, 125, 50),
        'secondary': RGBColor(129, 199, 132),
        'accent': RGBColor(27, 94, 32),
        'dark': RGBColor(33, 33, 33),
        'light': RGBColor(232, 245, 233),
        'light2': RGBColor(241, 248, 233),
        'header_text': RGBColor(255, 255, 255),
        'table_header': RGBColor(46, 125, 50),
        'table_alt': RGBColor(232, 245, 233),
        'cover_bg': RGBColor(27, 94, 32),
        'cover_deco': RGBColor(46, 125, 50),
        'cover_title': RGBColor(255, 255, 255),      # 화이트 (배경 대비)
        'box_fill': RGBColor(46, 125, 50),
        'box_fill2': RGBColor(129, 199, 132),
        'box_fill3': RGBColor(27, 94, 32),
        'arrow': RGBColor(255, 152, 0),
        'success': RGBColor(46, 125, 50),
        'warning': RGBColor(255, 152, 0),
        'danger': RGBColor(211, 47, 47),
    },
}

FONT = {'title': 24, 'section': 14, 'subtitle': 14, 'body': 11, 'table': 10, 'small': 9, 'diagram': 10}
FONT_NAME = "맑은 고딕"


class EnterprisePPTGenerator:
    """대기업 제안서/신청서용 PPT 생성기"""

    def __init__(self, theme_id=1):
        if theme_id not in THEMES:
            raise ValueError(f"테마 번호는 1, 2, 3 중 하나여야 합니다.")
        self.theme = THEMES[theme_id]
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        print(f"선택된 테마: {self.theme['name']}")

    # ============== 기본 함수들 ==============

    def add_header(self, slide, slide_title, section_title):
        """헤더 추가"""
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(0.8))
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme['accent']
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(8), Inches(0.5))
        p = title_box.text_frame.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(FONT['title'])
        p.font.bold = True
        p.font.color.rgb = self.theme['header_text']
        p.font.name = FONT_NAME

        section_box = slide.shapes.add_textbox(Inches(8.5), Inches(0.22), Inches(4.5), Inches(0.45))
        sp = section_box.text_frame.paragraphs[0]
        sp.text = section_title
        sp.font.size = Pt(FONT['section'])
        sp.font.color.rgb = self.theme['light']
        sp.font.name = FONT_NAME
        sp.alignment = PP_ALIGN.RIGHT

    def add_subtitle_box(self, slide, subtitle, y=0.95):
        """소제목"""
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.3), Inches(0.4))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(FONT['subtitle'])
        p.font.bold = True
        p.font.color.rgb = self.theme['accent']
        p.font.name = FONT_NAME
        return y + 0.5

    def add_table(self, slide, headers, rows, start_y=1.5, col_widths=None):
        """테이블 추가"""
        num_cols = len(headers)
        num_rows = len(rows) + 1
        if col_widths is None:
            col_widths = [12.3 / num_cols] * num_cols
        table_height = 0.4 + len(rows) * 0.35
        table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(start_y),
                                       Inches(sum(col_widths)), Inches(min(table_height, 5.5))).table
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme['table_header']
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(FONT['table'])
            p.font.bold = True
            p.font.color.rgb = self.theme['header_text']
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.theme['table_alt'] if i % 2 == 1 else self.theme['light']
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(FONT['table'])
                p.font.color.rgb = self.theme['dark']
                p.font.name = FONT_NAME
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        return start_y + table_height + 0.2

    # ============== 도형 함수들 ==============

    def add_rounded_box(self, slide, left, top, width, height, text,
                        fill_color, text_color=None, font_size=10, bold=True):
        """둥근 박스"""
        if text_color is None:
            text_color = RGBColor(255, 255, 255)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = self.theme['dark']
        shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = text_color
        tf.paragraphs[0].font.name = FONT_NAME
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        return shape

    def add_arrow(self, slide, x1, y1, x2, y2, color=None, width=2):
        """화살표"""
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
            Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        connector.line.color.rgb = color or self.theme['arrow']
        connector.line.width = Pt(width)
        return connector

    def add_text_box(self, slide, left, top, width, height, text,
                     font_size=10, color=None, bold=False, align=PP_ALIGN.LEFT):
        """텍스트 박스"""
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or self.theme['dark']
        p.font.name = FONT_NAME
        p.font.bold = bold
        p.alignment = align
        return box

    def add_explanation_panel(self, slide, left, top, width, height, items, title="설명"):
        """설명 패널 (연한 녹색 배경)"""
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height))
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(226, 240, 217)
        panel.line.color.rgb = RGBColor(196, 220, 187)
        panel.line.width = Pt(2)

        title_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1),
            Inches(width - 0.3), Inches(0.35))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.theme['accent']
        p.font.name = FONT_NAME

        y = top + 0.45
        for i, item in enumerate(items, 1):
            item_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(y),
                Inches(width - 0.3), Inches(0.6))
            tf = item_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"{i}. {item}"
            p.font.size = Pt(9)
            p.font.color.rgb = self.theme['dark']
            p.font.name = FONT_NAME
            y += 0.55
        return panel

    # ============== 표지/목차 ==============

    def create_cover_slide(self, title, subtitle, date=None):
        """표지 슬라이드"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.theme['cover_bg']
        bg.line.fill.background()

        deco = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(9), Inches(4.5), Inches(5), Inches(4))
        deco.fill.solid()
        deco.fill.fore_color.rgb = self.theme['cover_deco']
        deco.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11), Inches(1.5))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.theme['cover_title']
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11), Inches(0.8))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(20)
        sp.font.color.rgb = self.theme['light']
        sp.font.name = FONT_NAME
        sp.alignment = PP_ALIGN.CENTER

        if date:
            date_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11), Inches(0.5))
            dp = date_box.text_frame.paragraphs[0]
            dp.text = date
            dp.font.size = Pt(14)
            dp.font.color.rgb = self.theme['light2']
            dp.font.name = FONT_NAME
            dp.alignment = PP_ALIGN.CENTER

        return slide

    def create_toc_slide(self, toc_items):
        """목차 슬라이드"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "목차", "")

        y = 1.3
        for item in toc_items:
            box = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(10), Inches(0.5))
            p = box.text_frame.paragraphs[0]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = self.theme['dark']
            p.font.name = FONT_NAME
            p.font.bold = True
            y += 0.55

        return slide

    # ============== 아키텍처 슬라이드 예시 ==============

    def create_architecture_slide(self, section_title):
        """MSA 아키텍처 도식화 예시"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "시스템 아키텍처", section_title)
        self.add_subtitle_box(slide, "마이크로서비스 아키텍처 구조")

        # AS-IS (왼쪽)
        self.add_text_box(slide, 0.5, 1.4, 2, 0.35, "AS-IS (현행)", 12, self.theme['danger'], True, PP_ALIGN.CENTER)
        self.add_rounded_box(slide, 0.5, 1.8, 2.5, 4.0,
            "모놀리식\n아키텍처\n\n• 단일 WAS\n• 단일 DB\n• 강결합",
            RGBColor(189, 189, 189), self.theme['dark'], 10, False)

        # 화살표
        self.add_arrow(slide, 3.2, 3.8, 4.0, 3.8, self.theme['arrow'], 3)
        self.add_text_box(slide, 3.3, 3.3, 0.8, 0.4, "전환", 10, self.theme['arrow'], True, PP_ALIGN.CENTER)

        # TO-BE (오른쪽)
        self.add_text_box(slide, 4.2, 1.4, 4.5, 0.35, "TO-BE (목표)", 12, self.theme['success'], True, PP_ALIGN.CENTER)

        # API Gateway
        self.add_rounded_box(slide, 4.2, 1.85, 4.5, 0.5, "API Gateway", self.theme['box_fill3'], None, 10)

        # Microservices
        services = [("서비스 A", 4.2), ("서비스 B", 5.35), ("서비스 C", 6.5), ("서비스 D", 7.65)]
        for name, x in services:
            self.add_rounded_box(slide, x, 2.5, 1.05, 0.9, name, self.theme['box_fill'], None, 9)

        # Database
        for name, x in services:
            self.add_rounded_box(slide, x, 3.5, 1.05, 0.55, "DB", self.theme['box_fill3'], None, 8)

        # 설명 패널
        explanations = [
            "서비스 분해: 도메인별 독립 서비스로 분리",
            "느슨한 결합: 서비스 간 의존성 최소화",
            "독립 배포: 각 서비스별 독립 배포 가능",
            "확장성: 서비스별 개별 스케일링"
        ]
        self.add_explanation_panel(slide, 9.0, 1.5, 4.0, 3.5, explanations)

        return slide

    def create_comparison_slide(self, title, section_title, subtitle, headers, rows):
        """비교표 슬라이드"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, title, section_title)
        y = self.add_subtitle_box(slide, subtitle)
        self.add_table(slide, headers, rows, start_y=y + 0.2)
        return slide

    # ============== 저장 ==============

    def save(self, output_path):
        """프레젠테이션 저장"""
        self.prs.save(output_path)
        print(f"총 슬라이드 수: {len(self.prs.slides)}개")
        print(f"저장 완료: {output_path}")


def main():
    """예제 실행"""
    theme_id = int(sys.argv[1]) if len(sys.argv) > 1 else 1
    output_name = sys.argv[2] if len(sys.argv) > 2 else "example_presentation.pptx"

    generator = EnterprisePPTGenerator(theme_id)

    # 표지
    generator.create_cover_slide(
        title="프로젝트 제안서",
        subtitle="시스템 현대화 및 클라우드 전환",
        date="2026년 1월 30일"
    )

    # 목차
    generator.create_toc_slide([
        "Ⅰ. 프로젝트 개요",
        "Ⅱ. 현황 분석",
        "Ⅲ. 목표 아키텍처",
        "Ⅳ. 추진 전략",
        "Ⅴ. 기대 효과",
    ])

    # 아키텍처
    generator.create_architecture_slide("Ⅲ. 목표 아키텍처")

    # 비교표
    generator.create_comparison_slide(
        title="역량 비교",
        section_title="Ⅱ. 현황 분석",
        subtitle="AS-IS vs TO-BE 비교",
        headers=["구분", "현행", "목표"],
        rows=[
            ["아키텍처", "모놀리식", "마이크로서비스"],
            ["배포 주기", "월 1회", "일 10회"],
            ["장애 복구", "4시간", "10분"],
            ["확장성", "수동", "자동 (HPA)"],
        ]
    )

    generator.save(output_name)


if __name__ == "__main__":
    main()
