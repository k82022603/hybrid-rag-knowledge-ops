#!/usr/bin/env python3
"""
Hybrid RAG Performance Benchmark

AI Service의 검색/업로드/인증 성능을 자동 측정하고 리포트를 생성합니다.

Usage:
    python perf_benchmark.py                     # 전체 벤치마크
    python perf_benchmark.py --only search       # 검색만 측정
    python perf_benchmark.py --only upload       # 업로드만 측정
    python perf_benchmark.py --only esknn        # ES kNN만 측정
    python perf_benchmark.py --only token        # 토큰만 측정
    python perf_benchmark.py --iterations 20     # 반복 횟수 변경
    python perf_benchmark.py --output report.md  # 결과 저장 (Markdown)
    python perf_benchmark.py --format csv --output report.csv  # CSV 출력
    python perf_benchmark.py --format json --output report.json # JSON 출력

Author: Claude (Opus 4.6)
Date: 2026-02-06
"""

import argparse
import json
import os
import statistics
import sys
import tempfile
import time
from datetime import datetime
from urllib.request import Request, urlopen
from urllib.error import HTTPError, URLError


# ─── 색상 코드 ────────────────────────────────────────────────
class C:
    R = "\033[0m"
    B = "\033[1m"
    D = "\033[2m"
    RED = "\033[91m"
    GRN = "\033[92m"
    YLW = "\033[93m"
    BLU = "\033[94m"
    CYN = "\033[96m"

    @staticmethod
    def c(text, color):
        if os.getenv("NO_COLOR"):
            return text
        return f"{color}{text}{C.R}"


# ─── HTTP 클라이언트 ─────────────────────────────────────────
class APIClient:
    def __init__(self, base_url, email="admin@example.com", password="admin1234"):
        self.base_url = base_url.rstrip("/")
        self.email = email
        self.password = password
        self.token = None

    def _req(self, method, path, data=None, headers=None, timeout=30):
        url = f"{self.base_url}{path}"
        h = {"Content-Type": "application/json"}
        if headers:
            h.update(headers)
        if self.token:
            h["Authorization"] = f"Bearer {self.token}"
        body = json.dumps(data).encode("utf-8") if data else None
        req = Request(url, data=body, headers=h, method=method)
        start = time.time()
        resp = urlopen(req, timeout=timeout)
        elapsed_ms = (time.time() - start) * 1000
        result = json.loads(resp.read().decode("utf-8"))
        return result, elapsed_ms

    def _req_raw(self, method, path, body_bytes=None, headers=None, timeout=30):
        """Raw request (for multipart upload)"""
        url = f"{self.base_url}{path}"
        h = {}
        if headers:
            h.update(headers)
        if self.token:
            h["Authorization"] = f"Bearer {self.token}"
        req = Request(url, data=body_bytes, headers=h, method=method)
        start = time.time()
        resp = urlopen(req, timeout=timeout)
        elapsed_ms = (time.time() - start) * 1000
        result = json.loads(resp.read().decode("utf-8"))
        return result, elapsed_ms

    def auth(self):
        data = {"email": self.email, "password": self.password}
        result, ms = self._req("POST", "/api/v1/auth/login", data)
        self.token = result.get("accessToken") or result.get("access_token")
        return ms

    def search(self, query, search_type="hybrid", top_k=5):
        endpoints = {
            "hybrid": "/api/v1/search/hybrid",
            "keyword": "/api/v1/search/keyword",
            "semantic": "/api/v1/search/semantic",
        }
        return self._req("POST", endpoints[search_type], {"query": query, "top_k": top_k})

    def upload_file(self, filepath):
        """Multipart file upload"""
        boundary = "----BenchmarkBoundary"
        filename = os.path.basename(filepath)
        with open(filepath, "rb") as f:
            file_data = f.read()

        body = (
            f"--{boundary}\r\n"
            f'Content-Disposition: form-data; name="file"; filename="{filename}"\r\n'
            f"Content-Type: application/octet-stream\r\n\r\n"
        ).encode("utf-8") + file_data + f"\r\n--{boundary}--\r\n".encode("utf-8")

        headers = {"Content-Type": f"multipart/form-data; boundary={boundary}"}
        return self._req_raw("POST", "/api/v1/documents/upload", body, headers)


# ─── ES 직접 쿼리 ──────────────────────────────────────────
def es_knn_search(es_url="http://localhost:9200"):
    """ES에 직접 kNN 쿼리 실행 (zero vector - 순수 검색 성능 측정용)"""
    url = f"{es_url}/knowledge_chunks/_search"
    # 실제 임베딩 대신 1024차원 작은 값 vector 사용 (성능 측정 목적)
    query = {
        "size": 5,
        "knn": {
            "field": "dense_vector",
            "query_vector": [0.01] * 1024,
            "k": 5,
            "num_candidates": 50,
        },
    }
    body = json.dumps(query).encode("utf-8")
    headers = {"Content-Type": "application/json"}
    req = Request(url, data=body, headers=headers, method="POST")
    start = time.time()
    resp = urlopen(req, timeout=10)
    elapsed_ms = (time.time() - start) * 1000
    resp.read()
    return elapsed_ms


# ─── 벤치마크 러너 ──────────────────────────────────────────
class BenchmarkRunner:
    def __init__(self, client, es_url="http://localhost:9200"):
        self.client = client
        self.es_url = es_url
        self.results = {}

    def run_search_bench(self, iterations=10, queries=None):
        """Hybrid Search 벤치마크"""
        if queries is None:
            queries = [
                "MSA 마이크로서비스 아키텍처 전환",
                "문서를 벡터로 변환하는 방법",
                "Elasticsearch 검색 성능 최적화",
                "Kubernetes deployment CI/CD",
                "RAG 파이프라인 설계",
            ]
        print(C.c(f"\n  [1/4] Hybrid Search Benchmark ({iterations} iterations)", C.BLU + C.B))
        times = []
        for i in range(iterations):
            query = queries[i % len(queries)]
            try:
                _, ms = self.client.search(query, "hybrid", 5)
                times.append(ms)
                bar = "=" * int(ms / 50)
                color = C.GRN if ms < 500 else C.YLW if ms < 1000 else C.RED
                print(f"    {i+1:2d}/{iterations}  {C.c(f'{ms:7.0f}ms', color)}  {C.c(bar, color)}")
            except Exception as e:
                print(f"    {i+1:2d}/{iterations}  {C.c(f'ERROR: {e}', C.RED)}")
        self.results["search"] = times
        if times:
            self._print_stats("Hybrid Search", times, threshold=500)

    def run_upload_bench(self, iterations=3):
        """문서 업로드 벤치마크"""
        print(C.c(f"\n  [2/4] Document Upload Benchmark ({iterations} iterations)", C.BLU + C.B))

        # 최소 PPTX 생성
        test_file = self._create_test_pptx()
        times = []
        for i in range(iterations):
            try:
                _, ms = self.client.upload_file(test_file)
                times.append(ms)
                print(f"    {i+1:2d}/{iterations}  {C.c(f'{ms:7.0f}ms', C.GRN if ms < 3000 else C.RED)}")
            except Exception as e:
                print(f"    {i+1:2d}/{iterations}  {C.c(f'ERROR: {e}', C.RED)}")
        os.unlink(test_file)
        self.results["upload"] = times
        if times:
            self._print_stats("Document Upload", times, threshold=3000)

    def run_esknn_bench(self, iterations=10):
        """ES kNN 순수 검색 벤치마크"""
        print(C.c(f"\n  [3/4] ES kNN Pure Search Benchmark ({iterations} iterations)", C.BLU + C.B))
        times = []
        for i in range(iterations):
            try:
                ms = es_knn_search(self.es_url)
                times.append(ms)
                print(f"    {i+1:2d}/{iterations}  {C.c(f'{ms:7.1f}ms', C.GRN if ms < 100 else C.RED)}")
            except Exception as e:
                print(f"    {i+1:2d}/{iterations}  {C.c(f'ERROR: {e}', C.RED)}")
        self.results["esknn"] = times
        if times:
            self._print_stats("ES kNN (pure)", times, threshold=100)

    def run_token_bench(self, iterations=5):
        """토큰 발급 벤치마크"""
        print(C.c(f"\n  [4/4] Token Auth Benchmark ({iterations} iterations)", C.BLU + C.B))
        times = []
        for i in range(iterations):
            try:
                ms = self.client.auth()
                times.append(ms)
                print(f"    {i+1:2d}/{iterations}  {C.c(f'{ms:7.0f}ms', C.GRN if ms < 2000 else C.RED)}")
            except Exception as e:
                print(f"    {i+1:2d}/{iterations}  {C.c(f'ERROR: {e}', C.RED)}")
        self.results["token"] = times
        if times:
            self._print_stats("Token Auth", times, threshold=2000)

    @staticmethod
    def _percentile(sorted_times, pct):
        """백분위수 계산"""
        if not sorted_times:
            return 0
        idx = int(len(sorted_times) * pct / 100)
        idx = min(idx, len(sorted_times) - 1)
        return sorted_times[idx]

    def _print_stats(self, name, times, threshold):
        """통계 출력"""
        avg = statistics.mean(times)
        mn = min(times)
        mx = max(times)
        med = statistics.median(times)
        s = sorted(times)
        p50 = self._percentile(s, 50)
        p95 = self._percentile(s, 95)
        p99 = self._percentile(s, 99)
        passed = avg < threshold

        status = C.c("PASS", C.GRN + C.B) if passed else C.c("FAIL", C.RED + C.B)
        print(f"\n    {C.c(name, C.B)}: avg={C.c(f'{avg:.0f}ms', C.YLW)} "
              f"min={mn:.0f} max={mx:.0f} P50={p50:.0f} P95={p95:.0f} P99={p99:.0f} "
              f"threshold=<{threshold}ms [{status}]")

    def _create_test_pptx(self):
        """벤치마크용 최소 PPTX 파일 생성"""
        try:
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Benchmark Test"
            slide.placeholders[1].text = "Performance measurement test file"
            path = tempfile.mktemp(suffix=".pptx")
            prs.save(path)
            return path
        except ImportError:
            # python-pptx 없으면 빈 파일로 대체
            path = tempfile.mktemp(suffix=".pptx")
            with open(path, "wb") as f:
                f.write(b"PK" + b"\x00" * 100)  # 최소 zip 시그니처
            return path

    def _compute_stats(self, times, threshold):
        """통계 계산 (재사용)"""
        avg = statistics.mean(times)
        mn = min(times)
        mx = max(times)
        s = sorted(times)
        return {
            "avg": avg, "min": mn, "max": mx,
            "p50": self._percentile(s, 50),
            "p95": self._percentile(s, 95),
            "p99": self._percentile(s, 99),
            "threshold": threshold,
            "passed": avg < threshold,
            "count": len(times),
            "raw": times,
        }

    def generate_report(self, output_file=None, fmt="md"):
        """리포트 생성 (md, csv, json)"""
        configs = {
            "search": ("Hybrid Search", 500),
            "upload": ("Document Upload", 3000),
            "esknn": ("ES kNN (pure)", 100),
            "token": ("Token Auth", 2000),
        }

        stats_map = {}
        for key, (name, threshold) in configs.items():
            times = self.results.get(key, [])
            if times:
                stats_map[key] = {"name": name, **self._compute_stats(times, threshold)}

        if fmt == "json":
            report = self._report_json(stats_map)
        elif fmt == "csv":
            report = self._report_csv(stats_map)
        else:
            report = self._report_md(stats_map)

        if output_file:
            with open(output_file, "w", encoding="utf-8") as f:
                f.write(report)
            print(C.c(f"\n  Report saved to: {output_file}", C.GRN))

        return report

    def _report_md(self, stats_map):
        """Markdown 리포트"""
        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        lines = [
            "# Performance Benchmark Report", "",
            f"**Date**: {now}",
            f"**AI Service**: {self.client.base_url}",
            f"**ES**: {self.es_url}", "", "---", "",
            "## Summary", "",
            "| Metric | Avg | Min | Max | P50 | P95 | P99 | Threshold | Result |",
            "|--------|-----|-----|-----|-----|-----|-----|-----------|--------|",
        ]
        total_pass = sum(1 for s in stats_map.values() if s["passed"])
        for s in stats_map.values():
            status = "**PASS**" if s["passed"] else "**FAIL**"
            lines.append(
                f"| {s['name']} | {s['avg']:.0f}ms | {s['min']:.0f}ms | {s['max']:.0f}ms "
                f"| {s['p50']:.0f}ms | {s['p95']:.0f}ms | {s['p99']:.0f}ms "
                f"| <{s['threshold']}ms | {status} |"
            )
        lines.extend(["", f"**Overall**: {total_pass}/{len(stats_map)} PASS", ""])

        for s in stats_map.values():
            lines.extend([
                f"## {s['name']} ({s['count']} iterations)", "",
                "| # | Time |", "|---|------|",
            ])
            for i, t in enumerate(s["raw"], 1):
                lines.append(f"| {i} | {t:.0f}ms |")
            lines.append("")

        return "\n".join(lines)

    def _report_csv(self, stats_map):
        """CSV 리포트"""
        lines = ["metric,avg_ms,min_ms,max_ms,p50_ms,p95_ms,p99_ms,threshold_ms,result"]
        for s in stats_map.values():
            lines.append(
                f"{s['name']},{s['avg']:.1f},{s['min']:.1f},{s['max']:.1f},"
                f"{s['p50']:.1f},{s['p95']:.1f},{s['p99']:.1f},"
                f"{s['threshold']},{'PASS' if s['passed'] else 'FAIL'}"
            )
        lines.append("")
        lines.append("# Raw data")
        lines.append("metric,iteration,time_ms")
        for s in stats_map.values():
            for i, t in enumerate(s["raw"], 1):
                lines.append(f"{s['name']},{i},{t:.1f}")
        return "\n".join(lines)

    def _report_json(self, stats_map):
        """JSON 리포트"""
        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        output = {
            "date": now,
            "ai_service_url": self.client.base_url,
            "es_url": self.es_url,
            "benchmarks": {},
        }
        for key, s in stats_map.items():
            output["benchmarks"][key] = {
                "name": s["name"],
                "avg_ms": round(s["avg"], 1),
                "min_ms": round(s["min"], 1),
                "max_ms": round(s["max"], 1),
                "p50_ms": round(s["p50"], 1),
                "p95_ms": round(s["p95"], 1),
                "p99_ms": round(s["p99"], 1),
                "threshold_ms": s["threshold"],
                "result": "PASS" if s["passed"] else "FAIL",
                "iterations": s["count"],
                "raw_times_ms": [round(t, 1) for t in s["raw"]],
            }
        total_pass = sum(1 for s in stats_map.values() if s["passed"])
        output["overall"] = f"{total_pass}/{len(stats_map)} PASS"
        return json.dumps(output, indent=2, ensure_ascii=False)


# ─── 메인 ───────────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser(
        description="Hybrid RAG Performance Benchmark",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("--url", default="http://localhost:8000", help="AI Service URL")
    parser.add_argument("--es-url", default="http://localhost:9200", help="Elasticsearch URL")
    parser.add_argument("--email", default="admin@example.com", help="Email")
    parser.add_argument("--password", default="admin1234", help="Password")
    parser.add_argument("--iterations", type=int, default=10, help="반복 횟수 (default: 10)")
    parser.add_argument("--only", choices=["search", "upload", "esknn", "token"], help="특정 항목만 측정")
    parser.add_argument("--format", default="md", choices=["md", "csv", "json"], help="출력 형식 (default: md)")
    parser.add_argument("--output", help="리포트 저장 경로")
    args = parser.parse_args()

    print(C.c("\n  Hybrid RAG Performance Benchmark", C.CYN + C.B))
    print(C.c("  " + "=" * 40, C.CYN))

    client = APIClient(args.url, args.email, args.password)

    # 연결 확인
    print(C.c("\n  Connecting...", C.D))
    try:
        ms = client.auth()
        print(C.c(f"  Connected & Authenticated ({ms:.0f}ms)", C.GRN))
    except Exception as e:
        print(C.c(f"  ERROR: {e}", C.RED))
        sys.exit(1)

    runner = BenchmarkRunner(client, args.es_url)
    iters = args.iterations

    if args.only:
        if args.only == "search":
            runner.run_search_bench(iters)
        elif args.only == "upload":
            runner.run_upload_bench(min(iters, 5))
        elif args.only == "esknn":
            runner.run_esknn_bench(iters)
        elif args.only == "token":
            runner.run_token_bench(min(iters, 5))
    else:
        runner.run_search_bench(iters)
        runner.run_upload_bench(min(iters, 5))
        runner.run_esknn_bench(iters)
        runner.run_token_bench(min(iters, 5))

    # 리포트 생성
    report = runner.generate_report(args.output, fmt=args.format)
    if not args.output:
        print(C.c(f"\n  (Use --output report.{args.format} to save report)", C.D))

    print(C.c("\n  Benchmark complete!\n", C.CYN + C.B))


if __name__ == "__main__":
    main()
