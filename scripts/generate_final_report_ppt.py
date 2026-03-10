#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
고도화 프로젝트 최종 보고서 PPT 생성
Theme 3: Tech Innovation (테크 이노베이션)
"""
import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '.claude', 'skills', 'enterprise-ppt-generator'))

from template import EnterprisePPTGenerator, FONT, FONT_NAME, PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


class FinalReportPPT(EnterprisePPTGenerator):

    def add_bullet_slide(self, title, section, subtitle, bullets, y_start=1.5):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, title, section)
        y = self.add_subtitle_box(slide, subtitle)
        y = max(y, y_start)
        for b in bullets:
            box = slide.shapes.add_textbox(Inches(1.0), Inches(y), Inches(11), Inches(0.45))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"  {b}"
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['dark']
            p.font.name = FONT_NAME
            y += 0.45
        return slide

    def add_kpi_boxes(self, slide, kpis, y=1.5):
        """KPI 카드 4개 가로 배치"""
        n = len(kpis)
        w = 2.8
        gap = 0.3
        total = n * w + (n - 1) * gap
        x_start = (13.333 - total) / 2
        for i, (label, value, sub) in enumerate(kpis):
            x = x_start + i * (w + gap)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(1.6))
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            box.line.width = Pt(2)
            # Value
            vb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.2), Inches(w - 0.2), Inches(0.6))
            p = vb.text_frame.paragraphs[0]
            p.text = value
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.CENTER
            # Label
            lb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.85), Inches(w - 0.2), Inches(0.35))
            p2 = lb.text_frame.paragraphs[0]
            p2.text = label
            p2.font.size = Pt(11)
            p2.font.bold = True
            p2.font.color.rgb = self.theme['dark']
            p2.font.name = FONT_NAME
            p2.alignment = PP_ALIGN.CENTER
            # Sub
            sb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.15), Inches(w - 0.2), Inches(0.3))
            p3 = sb.text_frame.paragraphs[0]
            p3.text = sub
            p3.font.size = Pt(9)
            p3.font.color.rgb = RGBColor(120, 120, 120)
            p3.font.name = FONT_NAME
            p3.alignment = PP_ALIGN.CENTER

    def add_retro_slide(self, name, role, did, well, learn, msg):
        """팀원 회고 슬라이드"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, f"팀원 회고 — {name}", "별첨. 팀원 회고")
        self.add_subtitle_box(slide, role)

        # Left: 기여 + 잘된점
        self.add_text_box(slide, 0.5, 1.5, 6, 0.35, "기여한 것", 12, self.theme['primary'], True)
        y = 1.9
        for item in did[:4]:
            self.add_text_box(slide, 0.7, y, 5.8, 0.35, f"  {item}", 10)
            y += 0.35

        self.add_text_box(slide, 0.5, y + 0.1, 6, 0.35, "잘된 점", 12, self.theme['success'], True)
        y += 0.5
        for item in well[:3]:
            self.add_text_box(slide, 0.7, y, 5.8, 0.35, f"  {item}", 10)
            y += 0.35

        # Right: 배운점 + 한마디
        self.add_text_box(slide, 7, 1.5, 6, 0.35, "배운 점", 12, self.theme['accent'], True)
        y2 = 1.9
        for item in learn[:3]:
            self.add_text_box(slide, 7.2, y2, 5.5, 0.35, f"  {item}", 10)
            y2 += 0.35

        # 한마디 패널
        self.add_explanation_panel(slide, 7, y2 + 0.3, 5.8, 1.5, [msg], title="팀원들에게 한마디")


def main():
    gen = FinalReportPPT(theme_id=3)  # Tech Innovation

    # ==================== 1. 표지 ====================
    gen.create_cover_slide(
        title="Hybrid RAG Knowledge Platform",
        subtitle="Graph RAG 기반 지능형 지식 검색 시스템 — 고도화 프로젝트 최종 보고서",
        date="2026년 3월 10일  |  KT DS"
    )

    # ==================== 2. 목차 ====================
    gen.create_toc_slide([
        "I.  프로젝트 개요 및 Executive Summary",
        "II.  프로젝트 진행 현황 (Sprint 이력)",
        "III. 시스템 아키텍처",
        "IV.  핵심 기능 — Hybrid RAG 4채널 융합 검색",
        "V.   품질 평가 (RAGAS)",
        "VI.  성능 테스트 및 데이터 현황",
        "VII. 프로젝트 산출물 및 비용 분석",
        "VIII. 주요 교훈 및 향후 과제",
        "별첨. 팀원 회고 (13명)"
    ])

    # ==================== 3. Executive Summary ====================
    slide = gen.prs.slides.add_slide(gen.prs.slide_layouts[6])
    gen.add_header(slide, "Executive Summary", "I. 프로젝트 개요")
    gen.add_subtitle_box(slide, '"구글처럼 검색하고, 사람처럼 답변한다"')

    gen.add_kpi_boxes(slide, [
        ("RAGAS 품질 등급", "A등급", "Mean 0.763 (역대 최고)"),
        ("처리 문서", "1,450건", "42,612 청크 / 60K 엔티티"),
        ("테스트 커버리지", "97%", "84 TC, UAT 18/18 PASS"),
        ("비용 절감", "95%", "DeepSeek vs GPT-4o"),
    ], y=1.5)

    gen.add_text_box(slide, 0.5, 3.5, 12.3, 0.4,
        "프로젝트 기간: 2026-01-10 ~ 2026-03-10 (약 2개월)  |  10 Sprints, ~291 SP, 394 커밋  |  13 AI 에이전트 가상 팀",
        12, gen.theme['dark'], False, PP_ALIGN.CENTER)

    # 진행 바
    phases = [("기획", 100), ("설계", 100), ("구현", 100), ("테스트", 100), ("배포", 100), ("고도화", 100)]
    y = 4.2
    for i, (name, pct) in enumerate(phases):
        x = 0.8 + i * 2.05
        # 라벨
        gen.add_text_box(slide, x, y, 1.8, 0.3, name, 10, gen.theme['dark'], True, PP_ALIGN.CENTER)
        # 바 배경
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y + 0.35), Inches(1.8), Inches(0.25))
        bg.fill.solid()
        bg.fill.fore_color.rgb = gen.theme['light']
        bg.line.fill.background()
        # 바 채움
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y + 0.35), Inches(1.8 * pct / 100), Inches(0.25))
        bar.fill.solid()
        bar.fill.fore_color.rgb = gen.theme['primary']
        bar.line.fill.background()
        gen.add_text_box(slide, x, y + 0.65, 1.8, 0.25, "100% DONE", 8, gen.theme['success'], True, PP_ALIGN.CENTER)

    gen.add_explanation_panel(slide, 0.5, 5.3, 12.3, 1.8, [
        "4채널 Hybrid RAG 검색: Vector + Keyword(BM25/Nori) + Sparse + Graph(Neo4j)",
        "LLM: DeepSeek V3.2 (95% 비용 절감), Embedding: BGE-M3 1024차원",
        "ETL 3-Phase: Docling 파싱 → GPU 임베딩 → 엔티티 추출 (Adaptive Gleaning)",
        "인프라: Docker Compose 18개 컨테이너, Keycloak SSO, Prometheus/Grafana"
    ], title="핵심 기술 요약")

    # ==================== 4. Sprint 이력 ====================
    gen.create_comparison_slide(
        "Sprint 이력", "II. 프로젝트 진행",
        "10개 Sprint, 약 291 Story Points 완료",
        ["Sprint", "기간", "SP", "Stories", "주요 내용"],
        [
            ["01", "01-10~01-17", "21", "5", "인프라 구축, Docker Compose 18개 컨테이너"],
            ["02", "01-17~01-24", "31", "5", "AI Service 핵심 (Embedding, Search, ETL)"],
            ["03", "01-24~01-31", "32", "5", "API Gateway, Frontend, 통합 테스트"],
            ["04~05", "01-31~02-07", "53", "9", "문서 파싱 고도화, Keycloak SSO"],
            ["06~07", "02-07~02-12", "40", "7", "기술부채 해결, Production-Ready"],
            ["08", "02-12~02-14", "34", "6", "ETL 3-Phase, GPU 임베딩 (Colab T4)"],
            ["09", "02-14~03-06", "56", "16", "UAT 18/18, RAGAS v1~v14"],
            ["10", "03-07~03-10", "24", "9", "검색 UX 고도화, 최종 성능 최적화"],
        ]
    )

    # ==================== 5. 시스템 아키텍처 ====================
    slide = gen.prs.slides.add_slide(gen.prs.slide_layouts[6])
    gen.add_header(slide, "시스템 아키텍처", "III. 아키텍처")
    gen.add_subtitle_box(slide, "Hybrid RAG Knowledge Platform — 전체 구성도")

    # User → Nginx
    gen.add_rounded_box(slide, 3.5, 1.4, 6.3, 0.5, "사용자 (Browser)", RGBColor(100,100,100))
    gen.add_arrow(slide, 6.6, 1.9, 6.6, 2.2, gen.theme['arrow'], 2)

    # Nginx
    gen.add_rounded_box(slide, 3.5, 2.2, 6.3, 0.45, "Nginx (Reverse Proxy) :80", gen.theme['box_fill3'])
    gen.add_arrow(slide, 4.5, 2.65, 4.5, 3.0)
    gen.add_arrow(slide, 8.5, 2.65, 8.5, 3.0)

    # Frontend + Gateway
    gen.add_rounded_box(slide, 1.5, 3.0, 2.8, 0.5, "Frontend\nReact 18 + Tailwind", gen.theme['box_fill2'], gen.theme['dark'], 9)
    gen.add_rounded_box(slide, 4.8, 3.0, 3.5, 0.5, "API Gateway\nSpring Cloud :8080", gen.theme['box_fill'])
    gen.add_arrow(slide, 6.5, 3.5, 6.5, 3.85)
    gen.add_arrow(slide, 5.5, 3.5, 4.0, 3.85)

    # Backend + AI Service
    gen.add_rounded_box(slide, 2.5, 3.85, 2.5, 0.55, "Backend\nSpringBoot :8081", gen.theme['box_fill2'], gen.theme['dark'], 9)
    gen.add_rounded_box(slide, 5.5, 3.85, 3.0, 0.55, "AI Service\nFastAPI+LangGraph :8000", gen.theme['box_fill'])

    # Databases
    gen.add_arrow(slide, 5.0, 4.4, 3.0, 5.0)
    gen.add_arrow(slide, 6.5, 4.4, 6.5, 5.0)
    gen.add_arrow(slide, 8.0, 4.4, 9.5, 5.0)

    gen.add_rounded_box(slide, 1.0, 5.0, 2.8, 0.9, "Elasticsearch\n42,612 chunks\nNori 한국어", gen.theme['box_fill3'], None, 9)
    gen.add_rounded_box(slide, 5.0, 5.0, 2.8, 0.9, "Neo4j\n60K Entities\nKnowledge Graph", gen.theme['box_fill3'], None, 9)
    gen.add_rounded_box(slide, 9.0, 5.0, 2.8, 0.9, "PostgreSQL\n1,450 docs\nSSOT", gen.theme['box_fill3'], None, 9)

    # Side panel
    gen.add_explanation_panel(slide, 9.5, 1.4, 3.5, 3.2, [
        "BGE-M3 1024차원 벡터 임베딩",
        "Nori 한국어 형태소 분석 BM25",
        "RRF 4채널 융합 (k=60)",
        "BGE Reranker v2 M3 재순위화",
        "DeepSeek V3.2 답변 생성"
    ], title="RAG 파이프라인")

    # ==================== 6. 핵심 기능 — Hybrid RAG ====================
    slide = gen.prs.slides.add_slide(gen.prs.slide_layouts[6])
    gen.add_header(slide, "Hybrid RAG 4채널 융합 검색", "IV. 핵심 기능")
    gen.add_subtitle_box(slide, "Vector + Keyword + Sparse + Graph → RRF 융합 → Reranker → LLM 답변")

    # 4 channels
    channels = [
        ("Vector\n(Semantic)", "BGE-M3 코사인 유사도", gen.theme['box_fill']),
        ("Keyword\n(BM25)", "Nori 한국어 형태소", gen.theme['box_fill2']),
        ("Sparse\n(Learned)", "BGE-M3 토큰 가중", gen.theme['primary']),
        ("Graph\n(Entity)", "Neo4j 관계 확장", gen.theme['box_fill3']),
    ]
    for i, (name, desc, color) in enumerate(channels):
        x = 0.5 + i * 2.4
        gen.add_rounded_box(slide, x, 1.6, 2.1, 1.0, name, color, None, 11)
        gen.add_text_box(slide, x, 2.65, 2.1, 0.3, desc, 8, gen.theme['dark'], False, PP_ALIGN.CENTER)

    # Arrows to RRF
    for i in range(4):
        x = 1.55 + i * 2.4
        gen.add_arrow(slide, x, 2.95, 5.5, 3.5, gen.theme['arrow'], 1)

    # RRF
    gen.add_rounded_box(slide, 4.0, 3.5, 3.0, 0.6, "RRF 융합 (k=60)", gen.theme['accent'])
    gen.add_arrow(slide, 5.5, 4.1, 5.5, 4.5, gen.theme['arrow'], 2)

    # Reranker
    gen.add_rounded_box(slide, 3.5, 4.5, 4.0, 0.6, "BGE Reranker v2 M3 (Cross-encoder)", gen.theme['primary'])
    gen.add_arrow(slide, 5.5, 5.1, 5.5, 5.5, gen.theme['arrow'], 2)

    # LLM
    gen.add_rounded_box(slide, 3.0, 5.5, 5.0, 0.7, "DeepSeek V3.2 답변 생성\n(RAG Context 주입)", gen.theme['box_fill3'])

    # Explanation panel
    gen.add_explanation_panel(slide, 9.0, 1.6, 4.0, 4.8, [
        "RRF: 4채널 검색 결과를 순위 기반 융합",
        "Reranker: top_k*3 후보 → Cross-encoder 재순위화",
        "Quality Gate: HIGH/PARTIAL/NONE 자동 판정",
        "Graph 효과: 60% 질문에서 유효 (+11.3% 부스트)",
        "Adaptive Gleaning: 문서 길이 기반 동적 추출",
        "캐시: Redis 임베딩 캐시 (93.8x speedup)"
    ], title="핵심 특징")

    # ==================== 7. RAGAS 품질 평가 ====================
    slide = gen.prs.slides.add_slide(gen.prs.slide_layouts[6])
    gen.add_header(slide, "RAGAS 품질 평가", "V. 품질 평가")
    gen.add_subtitle_box(slide, "18회 평가를 통한 최적 파라미터 확정 — v16 Mean 0.763 (A등급, 역대 최고)")

    gen.add_table(slide,
        ["버전", "설정 변경", "Faithfulness", "Precision", "Recall", "Mean", "등급"],
        [
            ["v7", "최초 기준선", "0.844", "0.431", "0.575", "0.617", "C+"],
            ["v11", "Chat E2E baseline", "0.935", "0.618", "0.672", "0.742", "A-"],
            ["v14", "graph_top_k 복원", "0.940", "0.682", "0.608", "0.743", "A-"],
            ["v15", "2-Pass 중복 증명", "0.907", "0.682", "0.595", "0.728", "B+"],
            ["v16", "최적 파라미터 확정", "0.859", "0.739", "0.690", "0.763", "A"],
            ["v17", "Chat API E2E", "0.800", "0.683", "0.685", "0.723", "A-"],
            ["v18", "GPT-4o Judge", "0.980", "0.000", "0.049", "0.343", "C"],
        ],
        start_y=1.6, col_widths=[0.8, 2.5, 1.5, 1.3, 1.3, 1.2, 0.8]
    )

    gen.add_explanation_panel(slide, 0.5, 5.2, 12.3, 1.8, [
        "v16 최적 파라미터: candidates_cap=50, graph_top_k=10, Reranker 1-Pass, pool=min(top_k*3, 50)",
        "Reranker 2-Pass 수학적 중복 증명: Cross-encoder는 (query, content)만 사용 → 동일 점수",
        "GPT-4o Judge: Faithfulness 관대(0.980), Context 극도로 엄격(0.000) — 평가 방법론 차이"
    ], title="핵심 인사이트")

    # ==================== 8. 성능 + 데이터 ====================
    slide = gen.prs.slides.add_slide(gen.prs.slide_layouts[6])
    gen.add_header(slide, "성능 테스트 및 데이터 현황", "VI. 성능/데이터")
    gen.add_subtitle_box(slide, "k6 Smoke Test 결과 + 처리 데이터 현황")

    # k6 results
    gen.add_table(slide,
        ["API", "성공률", "p95 Latency", "판정"],
        [
            ["Auth Login", "100%", "516ms", "PASS"],
            ["Auth /me, /refresh", "100%", "4.9ms", "PASS"],
            ["Semantic Search", "100%", "4.08s", "PASS"],
            ["Keyword Search", "100%", "1.37s", "PASS"],
            ["Hybrid Search", "60%", "timeout", "PARTIAL"],
        ],
        start_y=1.6, col_widths=[3, 1.5, 2, 1.5]
    )

    # Data stats (right side)
    gen.add_text_box(slide, 8.5, 1.5, 4.5, 0.4, "데이터 현황", 14, gen.theme['primary'], True)
    data_items = [
        ("원본 문서 (PG)", "1,450건"),
        ("청크 (ES)", "42,612건"),
        ("벡터 임베딩", "42,612건 (1024차원)"),
        ("엔티티 (Neo4j)", "~60,000건"),
        ("관계 (Neo4j)", "~54,000건"),
        ("지원 형식", "PDF, PPTX, HWP, DOCX"),
    ]
    y = 2.0
    for label, value in data_items:
        gen.add_text_box(slide, 8.7, y, 2.5, 0.35, label, 10, gen.theme['dark'])
        gen.add_text_box(slide, 11.2, y, 1.8, 0.35, value, 10, gen.theme['primary'], True)
        y += 0.38

    gen.add_explanation_panel(slide, 0.5, 4.8, 12.3, 1.2, [
        "Hybrid timeout: WSL2 14GB 메모리 제약 → Reranker ONNX CPU 리소스 고갈. Production(32GB+)에서 해소 예상.",
        "테스트 커버리지 97% (5개 모듈, 84 TC), UAT 18/18 PASS"
    ], title="Note")

    # ==================== 9. 산출물 + 비용 ====================
    slide = gen.prs.slides.add_slide(gen.prs.slide_layouts[6])
    gen.add_header(slide, "산출물 및 비용 분석", "VII. 산출물/비용")
    gen.add_subtitle_box(slide, "코드, 문서, 인프라 산출물 및 운영 비용")

    # 코드
    gen.add_table(slide,
        ["구분", "라인 수", "파일 수", "비고"],
        [
            ["Python (AI Service)", "72,584", "-", "FastAPI, LangGraph, BGE-M3"],
            ["Java (Backend/Gateway)", "~5,000", "-", "SpringBoot 3.x"],
            ["TypeScript (Frontend)", "~8,000", "-", "React 18, Tailwind"],
            ["Docker/Infra (YAML)", "3,647", "-", "18개 컨테이너"],
            ["테스트 코드", "-", "84", "Unit + Integration"],
            ["기술 문서", "-", "320", "설계서, 매뉴얼, 평가 보고서"],
            ["Git 커밋", "-", "394", "10 Sprints"],
        ],
        start_y=1.6, col_widths=[3, 1.5, 1.2, 4]
    )

    # 비용 (하단)
    gen.add_text_box(slide, 0.5, 5.0, 3, 0.35, "운영 비용 분석", 14, gen.theme['primary'], True)
    gen.add_table(slide,
        ["항목", "비용", "비고"],
        [
            ["DeepSeek API", "~$5/월", "95% 절감 (vs GPT-4o)"],
            ["OpenAI API (평가용)", "$30.65 잔여", "RAGAS Judge 전용"],
            ["LangSmith", "무료 티어", "Observability"],
            ["Docker 인프라", "자체 서버", "WSL2 14GB"],
        ],
        start_y=5.4, col_widths=[2.5, 1.5, 3.5]
    )

    # ==================== 10. 교훈 + 향후 과제 ====================
    slide = gen.prs.slides.add_slide(gen.prs.slide_layouts[6])
    gen.add_header(slide, "주요 교훈 및 향후 과제", "VIII. 교훈/과제")

    # 교훈 (좌측)
    gen.add_text_box(slide, 0.5, 1.2, 6, 0.4, "주요 교훈", 14, gen.theme['danger'], True)
    lessons = [
        'Nori 미적용 사고 (32일): "설계서에 적혀 있다고 구현된 것이 아니다"',
        "Reranker 2-Pass 중복 증명: 동일 모델 + 동일 후보 = 동일 점수",
        "능동적 대처: 질문 대신 보고, 실측 데이터 기반 즉시 실행",
        "Mock 테스트 금지: 실환경(Docker) 테스트만 유효",
    ]
    y = 1.7
    for item in lessons:
        gen.add_text_box(slide, 0.7, y, 6, 0.45, f"  {item}", 10, gen.theme['dark'])
        y += 0.45

    # 향후 과제 (우측)
    gen.add_text_box(slide, 7, 1.2, 6, 0.4, "향후 과제", 14, gen.theme['primary'], True)
    gen.add_table(slide,
        ["우선순위", "과제", "예상 효과"],
        [
            ["P1", "ONNX INT8 Reranker", "CPU 2-4x 개선"],
            ["P1", "HybridRetriever 통일", "Chat API 4%p 개선"],
            ["P2", "GPU 서빙 (T4/L4)", "임베딩 10x, 추론 5x"],
            ["P2", "의미 기반 청킹", "Context Recall 개선"],
            ["P3", "K8s 마이그레이션", "수평 확장성"],
        ],
        start_y=1.7, col_widths=[1, 2.5, 2.5]
    )

    # 결론 (하단)
    conclusion = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5))
    conclusion.fill.solid()
    conclusion.fill.fore_color.rgb = gen.theme['accent']
    conclusion.line.fill.background()
    gen.add_text_box(slide, 1.0, 5.7, 11.3, 0.4,
        "Hybrid RAG Knowledge Platform 고도화 프로젝트를 성공적으로 완료하였습니다.",
        16, RGBColor(255,255,255), True, PP_ALIGN.CENTER)
    gen.add_text_box(slide, 1.0, 6.2, 11.3, 0.4,
        "RAGAS A등급 (Mean 0.763)  |  4채널 Hybrid 검색  |  1,450문서 / 42,612청크  |  테스트 97%  |  10 Sprint / 394 커밋",
        12, gen.theme['light'], False, PP_ALIGN.CENTER)

    # ==================== 별첨: 팀원 회고 ====================
    retros = [
        ("PM", "Product Manager — Sprint/Jira/Slack 관리",
         ["10개 Sprint 관리, ~60 Stories 조율", "Jira SCRUM-10~121 이슈 관리", "Slack 커뮤니케이션 허브 운영", "팀 간 충돌 조율 및 우선순위 결정"],
         ["Agent Teams 풀가동으로 병렬 작업 극대화", "역할 분담 원칙 확립 (PM은 코딩 금지)"],
         ["가상 팀이라도 역할 분리가 품질을 높인다", "Sprint 단위 반복이 복잡한 프로젝트를 관리 가능하게 만든다"],
         "팀원 모두 고생했습니다! 2개월간 10 Sprint를 완주한 것 자체가 대단합니다."),

        ("TechLead", "기술 리더 — 아키텍처 검토, 코드 리뷰",
         ["상세 설계서 v2.4 검토/승인 (91점 A등급)", "소스코드 리뷰 72.5/100 B+", "ADR 작성 및 기술 의사결정", "Nori 미적용 사고 사후 분석"],
         ["4채널 Hybrid RAG 아키텍처 설계 안정적 운영", "Reranker 2-Pass 중복 수학적 증명"],
         ["실동작 검증의 중요성 (Nori 32일 미발견)", "Cross-encoder의 수학적 특성 이해"],
         "좋은 설계가 좋은 코드를 만듭니다. 항상 Why를 먼저 생각하세요."),

        ("Backend", "백엔드 개발자 — SpringBoot API Gateway",
         ["Spring Cloud Gateway 라우팅 구현", "Keycloak SSO 통합", "JWT 인증 필터 구현", "Resilience4j 회로 차단기"],
         ["Gateway + Keycloak + JWT 삼중 인증 안정 운영"],
         ["MSA 환경에서 인증 플로우 복잡성 관리 노하우"],
         "API Gateway는 모든 요청의 관문입니다. 견고하게 만들어서 뿌듯합니다."),

        ("Frontend", "프론트엔드 개발자 — React 18 UI",
         ["React 18 + Tailwind CSS UI 구현", "검색 결과 카드, 소스 인용 컴포넌트", "Keycloak SSO 프론트 통합"],
         ["Tailwind + Headless UI로 빠른 프로토타이핑"],
         ["Antigravity AI 디자인 도구와의 협업 경험"],
         "사용자가 직접 만지는 화면을 만드는 건 항상 설레는 일이에요!"),

        ("RAG Engineer", "ML/RAG 엔지니어 — RAG 파이프라인, AI Service",
         ["LangGraph RAG 워크플로우 설계", "HybridRetriever 4채널 융합", "BGE-M3 임베딩 + BGE Reranker", "DeepSeek LLM 통합"],
         ["4채널 RRF 융합이 단일 채널 대비 +23.6% 품질 향상", "Adaptive Gleaning으로 엔티티 추출 효율화"],
         ["Graph RAG이 60% 질문에서 유효하다는 실증 데이터 확보", "Reranker 수학적 분석으로 불필요한 연산 제거"],
         "RAG의 품질은 결국 데이터와 검색의 품질입니다. 모델만큼 파이프라인이 중요해요."),

        ("ETL Engineer", "데이터 엔지니어 — ETL 파이프라인",
         ["3-Phase ETL 파이프라인 설계/구현", "Docling 문서 파싱 (PDF/PPTX/HWP)", "GPU 임베딩 (Colab T4, 65.6 c/s)", "듀얼 API 키 배치 엔티티 추출"],
         ["3-Phase 분리로 장애 격리 및 재시작 용이", "GPU 임베딩으로 56,063건 13.5분 처리"],
         ["대용량 데이터 파이프라인에서 멱등성의 중요성", "OOM 방지를 위한 배치 크기 튜닝 경험"],
         "데이터는 AI의 연료입니다. 깨끗한 데이터가 좋은 결과를 만듭니다."),

        ("DB Designer", "데이터베이스 설계자 — 스키마 설계, 쿼리 최적화",
         ["PostgreSQL SSOT 스키마 설계", "Elasticsearch Nori 매핑 + 커스텀 이미지", "Neo4j Entity/Relation 스키마", "3-Store 정합성 보장 체계"],
         ["3개 DB 간 정합성 100% 달성"],
         ["ES 커스텀 이미지(Nori)의 필요성과 Dockerfile 관리"],
         "데이터의 진실은 하나(SSOT)여야 합니다. PG가 그 역할을 충실히 했습니다."),

        ("Infra Engineer", "인프라 엔지니어 — Docker Compose 인프라",
         ["Docker Compose 18개 컨테이너 구성", "커스텀 ES Nori 이미지 빌드", "WSL2 메모리 프로파일 관리", "네트워크/볼륨 설계"],
         ["18개 컨테이너 안정적 운영 (14GB 메모리에서!)"],
         ["메모리 제약 환경에서의 컨테이너 최적화 노하우"],
         "인프라는 보이지 않을 때 가장 잘 하고 있는 겁니다."),

        ("DevOps", "DevOps 엔지니어 — CI/CD, Observability",
         ["Prometheus/Grafana 모니터링 대시보드", "LangSmith Observability 연동", "Jaeger 분산 트레이싱", "배포 자동화 스크립트"],
         ["LangSmith로 RAG 파이프라인 실시간 모니터링 가능"],
         ["관측 가능성(Observability)이 디버깅 시간을 획기적으로 줄인다"],
         "측정할 수 없으면 개선할 수 없습니다. Observability는 선택이 아닌 필수!"),

        ("QA", "QA 엔지니어 — 테스트, RAGAS 평가",
         ["RAGAS 평가 18회 수행 (v1~v18)", "UAT 18/18 PASS 달성", "k6 성능 테스트 3종 실행", "테스트 커버리지 97% 달성"],
         ["RAGAS 자동 평가 파이프라인으로 객관적 품질 측정", "Mock 금지 원칙으로 실환경 검증 보장"],
         ["RAG 시스템의 품질 평가는 전통 SW 테스트와 다르다", "Judge LLM 선택이 평가 결과에 큰 영향"],
         "테스트는 안심을 주는 안전망입니다. 97% 커버리지가 자랑스럽습니다."),

        ("Architect", "소프트웨어 아키텍트 — 시스템/기능 상세 설계",
         ["상세 설계서 14개 작성", "Hybrid RAG 4채널 아키텍처 설계", "기술 검토 문서 (Gleaning, K8s 등)", "인프라 설계서 (18개 컨테이너)"],
         ["설계 91점 A등급으로 TechLead 승인"],
         ["설계는 한 번에 완벽할 수 없고, 구현 피드백으로 진화한다"],
         "좋은 아키텍처는 변경에 열려 있고, 핵심은 보호합니다."),

        ("Documenter", "코드 문서화 담당 — API/코드/아키텍처 문서화",
         ["API 문서화 (OpenAPI/Swagger)", "운영 매뉴얼 35개 작성 (OPS-001~035)", "개발자 가이드 3종 작성", "320개 기술 문서 관리"],
         ["OPS 매뉴얼이 트러블슈팅 시간을 크게 단축"],
         ["문서는 코드만큼 중요하고, 유지보수가 더 어렵다"],
         "코드는 기계가 읽고, 문서는 사람이 읽습니다. 둘 다 잘 써야 합니다."),

        ("Web Designer", "UI/UX 설계 — 디자인 시스템 관리",
         ["UI/UX 설계 및 디자인 가이드", "Antigravity AI 프롬프트 설계", "검색 결과 UX 최적화", "다크/라이트 테마 설계"],
         ["Antigravity + Stitch MCP로 디자인-개발 간극 축소"],
         ["AI 디자인 도구가 프로토타이핑 속도를 혁신적으로 높인다"],
         "좋은 UX는 사용자가 의식하지 못할 때 완성됩니다."),
    ]

    for name, role, did, well, learn, msg in retros:
        gen.add_retro_slide(name, role, did, well, learn, msg)

    # ==================== END 슬라이드 ====================
    slide = gen.prs.slides.add_slide(gen.prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = gen.theme['cover_bg']
    bg.line.fill.background()

    gen.add_text_box(slide, 1, 2.5, 11.3, 1, "Thank You", 48, RGBColor(255,255,255), True, PP_ALIGN.CENTER)
    gen.add_text_box(slide, 1, 3.8, 11.3, 0.6,
        "Hybrid RAG Knowledge Platform — 고도화 프로젝트 완료", 20, gen.theme['light'], False, PP_ALIGN.CENTER)
    gen.add_text_box(slide, 1, 4.5, 11.3, 0.5,
        '"구글처럼 검색하고, 사람처럼 답변한다"', 16, gen.theme['secondary'], False, PP_ALIGN.CENTER)
    gen.add_text_box(slide, 1, 5.5, 11.3, 0.4,
        "2026-03-10  |  Claude Code + 13 AI Agent Team  |  KT DS", 14, gen.theme['light2'], False, PP_ALIGN.CENTER)

    # Save
    output_path = "/mnt/d/Users/KTDS/Documents/06.과제/hybrid-rag-knowledge-ops/knowledge_service/docs/07_maintenance/고도화프로젝트_최종보고서_TechInnovation.pptx"
    gen.save(output_path)
    print(f"\n생성 완료: {output_path}")


if __name__ == "__main__":
    main()
