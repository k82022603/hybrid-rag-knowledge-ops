#!/usr/bin/env python3
"""
Frontend 상세 설계서 프레젠테이션 생성
Brown Earth Theme + 아키텍처 도식화 + 장점 부각
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os

# ========== Brown Earth Theme ==========
THEME = {
    'primary': RGBColor(121, 85, 72),        # #795548 - Brown
    'secondary': RGBColor(161, 136, 127),    # #A1887F - Light Brown
    'accent': RGBColor(255, 152, 0),         # #FF9800 - Orange
    'dark': RGBColor(62, 39, 35),            # #3E2723 - Dark Brown
    'light': RGBColor(239, 235, 233),        # #EFEBE9 - Cream
    'text': RGBColor(33, 33, 33),            # #212121 - Dark Text
    'white': RGBColor(255, 255, 255),
    # Feature Colors
    'react': RGBColor(97, 218, 251),         # #61DAFB - React Blue
    'typescript': RGBColor(49, 120, 198),    # #3178C6 - TypeScript
    'redux': RGBColor(118, 74, 188),         # #764ABC - Redux Purple
    'query': RGBColor(255, 75, 75),          # #FF4B4B - React Query Red
    'mui': RGBColor(0, 127, 255),            # #007FFF - MUI Blue
}

# ========== Helper Functions ==========

def add_title_slide(prs, title, subtitle):
    """표지 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['dark']
    bg.line.fill.background()

    # 장식 원
    deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), Inches(4.5), Inches(4.5), Inches(4))
    deco.fill.solid()
    deco.fill.fore_color.rgb = THEME['primary']
    deco.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.LEFT

    # 부제목
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.4), Inches(1))
    sf = sub_box.text_frame
    sf.word_wrap = True
    sp = sf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(24)
    sp.font.color.rgb = THEME['secondary']
    sp.font.name = "맑은 고딕"
    sp.alignment = PP_ALIGN.LEFT

    # 버전 정보
    ver_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(4), Inches(0.5))
    vf = ver_box.text_frame
    vp = vf.paragraphs[0]
    vp.text = "Version 1.0 | 2026-01-15"
    vp.font.size = Pt(14)
    vp.font.color.rgb = THEME['light']
    vp.font.name = "맑은 고딕"

    return slide


def add_content_slide(prs, title, content_items):
    """내용 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더 바
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # 내용
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = THEME['text']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(12)
        p.level = 0

    return slide


def add_tech_stack_slide(prs):
    """기술 스택 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "기술 스택"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # 카테고리별 기술 박스
    categories = [
        ("Core", ["React 18.3+", "TypeScript 5.4+", "Vite 5.x"], THEME['react']),
        ("상태 관리", ["Redux Toolkit", "React Query 5.x", "Redux Persist"], THEME['redux']),
        ("UI/UX", ["Material-UI 5.x", "Emotion", "Framer Motion"], THEME['mui']),
        ("폼/검증", ["React Hook Form", "Zod", "Resolvers"], THEME['accent']),
        ("시각화", ["Recharts", "Mermaid", "Prism.js"], THEME['query']),
        ("테스트", ["Vitest", "Playwright", "Storybook"], THEME['typescript']),
    ]

    start_x = 0.4
    start_y = 1.3
    box_width = 3.0
    box_height = 1.8
    gap_x = 0.2
    gap_y = 0.2

    for i, (cat_name, techs, color) in enumerate(categories):
        col = i % 3
        row = i // 3
        x = start_x + col * (box_width + gap_x)
        y = start_y + row * (box_height + gap_y)

        # 박스
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                     Inches(box_width), Inches(box_height))
        box.fill.solid()
        box.fill.fore_color.rgb = THEME['light']
        box.line.color.rgb = color
        box.line.width = Pt(3)

        # 카테고리 제목
        cat_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1),
                                           Inches(box_width - 0.2), Inches(0.4))
        ctf = cat_box.text_frame
        cp = ctf.paragraphs[0]
        cp.text = cat_name
        cp.font.size = Pt(14)
        cp.font.bold = True
        cp.font.color.rgb = color
        cp.font.name = "맑은 고딕"

        # 기술 목록
        tech_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.5),
                                            Inches(box_width - 0.2), Inches(1.2))
        ttf = tech_box.text_frame
        ttf.word_wrap = True
        for j, tech in enumerate(techs):
            if j == 0:
                tp = ttf.paragraphs[0]
            else:
                tp = ttf.add_paragraph()
            tp.text = f"• {tech}"
            tp.font.size = Pt(12)
            tp.font.color.rgb = THEME['text']
            tp.font.name = "맑은 고딕"
            tp.space_after = Pt(2)

    return slide


def add_component_architecture_slide(prs):
    """컴포넌트 아키텍처 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "컴포넌트 아키텍처"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # 계층 구조 다이어그램
    layers = [
        ("Pages (페이지)", "라우트와 1:1 매핑, 데이터 페칭, 레이아웃 조합",
         "DashboardPage, SearchPage, KnowledgeDetailPage", THEME['dark']),
        ("Features (기능 컴포넌트)", "비즈니스 로직 포함, 특정 도메인에 종속",
         "KnowledgeList, SearchChat, BookmarkFolder", THEME['primary']),
        ("Components (공유 컴포넌트)", "재사용 가능한 UI 컴포넌트, 비즈니스 로직 없음",
         "Button, Card, Modal, Table, Input", THEME['secondary']),
    ]

    layer_height = 1.6
    start_y = 1.3
    layer_width = 9.0

    for i, (name, desc, examples, color) in enumerate(layers):
        y = start_y + i * (layer_height + 0.3)

        # 레이어 박스
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y),
                                     Inches(layer_width), Inches(layer_height))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # 레이어 이름
        name_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.15),
                                            Inches(layer_width - 0.4), Inches(0.4))
        ntf = name_box.text_frame
        np = ntf.paragraphs[0]
        np.text = name
        np.font.size = Pt(18)
        np.font.bold = True
        np.font.color.rgb = THEME['white']
        np.font.name = "맑은 고딕"

        # 설명
        desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.55),
                                            Inches(layer_width - 0.4), Inches(0.4))
        dtf = desc_box.text_frame
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(12)
        dp.font.color.rgb = THEME['light']
        dp.font.name = "맑은 고딕"

        # 예시
        ex_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.95),
                                          Inches(layer_width - 0.4), Inches(0.5))
        etf = ex_box.text_frame
        ep = etf.paragraphs[0]
        ep.text = f"예: {examples}"
        ep.font.size = Pt(11)
        ep.font.color.rgb = THEME['accent']
        ep.font.name = "맑은 고딕"

        # 화살표 (마지막 제외)
        if i < len(layers) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                           Inches(4.75), Inches(y + layer_height),
                                           Inches(0.5), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = THEME['accent']
            arrow.line.fill.background()

    return slide


def add_state_management_slide(prs):
    """상태 관리 전략 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "상태 관리 전략"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # 3개 컬럼
    columns = [
        ("서버 상태\n(Server State)", "React Query",
         ["지식 목록", "검색 결과", "사용자 정보", "북마크", "대시보드 데이터"], THEME['query']),
        ("클라이언트 상태\n(Client State)", "Redux Toolkit",
         ["인증 상태", "UI 상태", "폼 입력", "모달 상태", "사이드바 상태"], THEME['redux']),
        ("URL 상태\n(URL State)", "React Router",
         ["페이지 번호", "필터 조건", "정렬 옵션", "검색어", "탭 선택"], THEME['mui']),
    ]

    col_width = 3.0
    start_x = 0.4
    col_gap = 0.15

    for i, (title, tool, items, color) in enumerate(columns):
        x = start_x + i * (col_width + col_gap)

        # 컬럼 헤더
        header_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Inches(x), Inches(1.2),
                                            Inches(col_width), Inches(1.0))
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = color
        header_box.line.fill.background()

        h_text = slide.shapes.add_textbox(Inches(x + 0.1), Inches(1.25),
                                          Inches(col_width - 0.2), Inches(0.9))
        htf = h_text.text_frame
        hp = htf.paragraphs[0]
        hp.text = title
        hp.font.size = Pt(14)
        hp.font.bold = True
        hp.font.color.rgb = THEME['white']
        hp.font.name = "맑은 고딕"
        hp.alignment = PP_ALIGN.CENTER

        # 도구 이름
        tool_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(x), Inches(2.3),
                                          Inches(col_width), Inches(0.5))
        tool_box.fill.solid()
        tool_box.fill.fore_color.rgb = THEME['light']
        tool_box.line.color.rgb = color
        tool_box.line.width = Pt(2)

        t_text = slide.shapes.add_textbox(Inches(x), Inches(2.35),
                                          Inches(col_width), Inches(0.4))
        ttf = t_text.text_frame
        tp = ttf.paragraphs[0]
        tp.text = tool
        tp.font.size = Pt(14)
        tp.font.bold = True
        tp.font.color.rgb = color
        tp.font.name = "맑은 고딕"
        tp.alignment = PP_ALIGN.CENTER

        # 항목 리스트
        items_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Inches(x), Inches(2.9),
                                           Inches(col_width), Inches(4.2))
        items_box.fill.solid()
        items_box.fill.fore_color.rgb = THEME['light']
        items_box.line.color.rgb = THEME['secondary']
        items_box.line.width = Pt(1)

        item_text = slide.shapes.add_textbox(Inches(x + 0.15), Inches(3.0),
                                             Inches(col_width - 0.3), Inches(4.0))
        itf = item_text.text_frame
        itf.word_wrap = True

        for j, item in enumerate(items):
            if j == 0:
                ip = itf.paragraphs[0]
            else:
                ip = itf.add_paragraph()
            ip.text = f"• {item}"
            ip.font.size = Pt(13)
            ip.font.color.rgb = THEME['text']
            ip.font.name = "맑은 고딕"
            ip.space_after = Pt(8)

    return slide


def add_project_structure_slide(prs):
    """프로젝트 구조 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "프로젝트 구조 (Feature-Based)"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # 왼쪽: 메인 폴더 구조
    folders = [
        ("src/", "소스 코드 루트", THEME['dark']),
        ("├─ app/", "앱 설정 (providers, router)", THEME['primary']),
        ("├─ components/", "공유 컴포넌트", THEME['primary']),
        ("├─ features/", "기능별 모듈 (핵심)", THEME['accent']),
        ("├─ hooks/", "공통 커스텀 훅", THEME['primary']),
        ("├─ pages/", "페이지 컴포넌트", THEME['primary']),
        ("├─ services/", "API 서비스", THEME['primary']),
        ("├─ store/", "Redux 스토어", THEME['primary']),
        ("├─ styles/", "테마, 글로벌 스타일", THEME['primary']),
        ("├─ types/", "타입 정의", THEME['primary']),
        ("└─ utils/", "유틸리티 함수", THEME['primary']),
    ]

    y_start = 1.2
    for i, (name, desc, color) in enumerate(folders):
        y = y_start + i * 0.5

        # 폴더명
        name_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(2.2), Inches(0.4))
        ntf = name_box.text_frame
        np = ntf.paragraphs[0]
        np.text = name
        np.font.size = Pt(13)
        np.font.bold = True
        np.font.color.rgb = color
        np.font.name = "Consolas"

        # 설명
        desc_box = slide.shapes.add_textbox(Inches(2.7), Inches(y), Inches(2.0), Inches(0.4))
        dtf = desc_box.text_frame
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(11)
        dp.font.color.rgb = THEME['text']
        dp.font.name = "맑은 고딕"

    # 오른쪽: Feature 모듈 상세
    feature_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(5.2), Inches(1.2),
                                         Inches(4.5), Inches(5.8))
    feature_box.fill.solid()
    feature_box.fill.fore_color.rgb = THEME['light']
    feature_box.line.color.rgb = THEME['accent']
    feature_box.line.width = Pt(3)

    # Feature 헤더
    fh_box = slide.shapes.add_textbox(Inches(5.4), Inches(1.35), Inches(4.1), Inches(0.4))
    fhtf = fh_box.text_frame
    fhp = fhtf.paragraphs[0]
    fhp.text = "Feature 모듈 구조"
    fhp.font.size = Pt(16)
    fhp.font.bold = True
    fhp.font.color.rgb = THEME['accent']
    fhp.font.name = "맑은 고딕"

    feature_folders = [
        "features/knowledge/",
        "├─ components/     # 기능 전용 컴포넌트",
        "│  ├─ KnowledgeCard/",
        "│  ├─ KnowledgeList/",
        "│  └─ KnowledgeForm/",
        "├─ hooks/          # 기능 전용 훅",
        "│  ├─ useKnowledge.ts",
        "│  └─ useKnowledgeList.ts",
        "├─ services/       # API 호출",
        "├─ store/          # 상태 관리",
        "├─ types/          # 타입 정의",
        "└─ index.ts        # Public API",
    ]

    ff_text = slide.shapes.add_textbox(Inches(5.4), Inches(1.85), Inches(4.1), Inches(5.0))
    fftf = ff_text.text_frame
    fftf.word_wrap = True

    for j, line in enumerate(feature_folders):
        if j == 0:
            fp = fftf.paragraphs[0]
        else:
            fp = fftf.add_paragraph()
        fp.text = line
        fp.font.size = Pt(11)
        if j == 0:
            fp.font.bold = True
            fp.font.color.rgb = THEME['dark']
        elif "#" in line:
            fp.font.color.rgb = THEME['secondary']
        else:
            fp.font.color.rgb = THEME['text']
        fp.font.name = "Consolas"
        fp.space_after = Pt(3)

    return slide


def add_architecture_advantages_slide(prs):
    """아키텍처 장점 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "아키텍처 장점"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # 4개 카테고리
    advantages = [
        ("Feature-Based 구조", [
            "높은 모듈성 - 기능별 독립적 개발/테스트",
            "명확한 책임 분리 - 유지보수 용이",
            "팀 협업 최적화 - 기능별 병렬 개발",
            "재사용성 - 컴포넌트 독립적 배포 가능"
        ], THEME['accent']),
        ("상태 관리 분리", [
            "Server State vs Client State 분리",
            "React Query 캐싱으로 API 호출 최소화",
            "Redux로 UI 상태 예측 가능하게 관리",
            "타입 안전성 - TypeScript 완벽 지원"
        ], THEME['redux']),
        ("MUI 디자인 시스템", [
            "일관된 UI/UX 경험 제공",
            "다크/라이트 모드 완벽 지원",
            "Pretendard 폰트로 한글 최적화",
            "반응형 디자인 내장"
        ], THEME['mui']),
        ("성능 최적화", [
            "Vite 기반 빠른 빌드 (HMR)",
            "코드 스플리팅 - 필요한 코드만 로드",
            "React.memo, useMemo 활용",
            "가상화 - 대량 데이터 효율적 렌더링"
        ], THEME['typescript']),
    ]

    col_width = 4.6
    row_height = 2.8
    start_x = 0.4
    start_y = 1.2
    gap = 0.2

    for i, (title, items, color) in enumerate(advantages):
        col = i % 2
        row = i // 2
        x = start_x + col * (col_width + gap)
        y = start_y + row * (row_height + gap)

        # 카드 배경
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(x), Inches(y),
                                      Inches(col_width), Inches(row_height))
        card.fill.solid()
        card.fill.fore_color.rgb = THEME['light']
        card.line.color.rgb = color
        card.line.width = Pt(3)

        # 카테고리 헤더
        cat_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(x), Inches(y),
                                            Inches(col_width), Inches(0.5))
        cat_header.fill.solid()
        cat_header.fill.fore_color.rgb = color
        cat_header.line.fill.background()

        cat_text = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.08),
                                            Inches(col_width - 0.3), Inches(0.4))
        ctf = cat_text.text_frame
        cp = ctf.paragraphs[0]
        cp.text = title
        cp.font.size = Pt(14)
        cp.font.bold = True
        cp.font.color.rgb = THEME['white']
        cp.font.name = "맑은 고딕"

        # 장점 리스트
        item_text = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.6),
                                             Inches(col_width - 0.3), Inches(row_height - 0.7))
        itf = item_text.text_frame
        itf.word_wrap = True

        for j, item in enumerate(items):
            if j == 0:
                ip = itf.paragraphs[0]
            else:
                ip = itf.add_paragraph()
            ip.text = f"✓ {item}"
            ip.font.size = Pt(11)
            ip.font.color.rgb = THEME['text']
            ip.font.name = "맑은 고딕"
            ip.space_after = Pt(6)

    return slide


def add_layout_diagram_slide(prs):
    """레이아웃 구조 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "레이아웃 구조"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # 레이아웃 다이어그램
    # Header
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(0.5), Inches(1.3),
                                          Inches(9.0), Inches(0.7))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = THEME['dark']
    header_shape.line.fill.background()

    header_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(4), Inches(0.5))
    htf = header_text.text_frame
    hp = htf.paragraphs[0]
    hp.text = "Header (64px) - 로고, 검색, 알림, 프로필"
    hp.font.size = Pt(12)
    hp.font.color.rgb = THEME['white']
    hp.font.name = "맑은 고딕"

    # Sidebar
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.5), Inches(2.1),
                                     Inches(2.0), Inches(4.5))
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = THEME['primary']
    sidebar.line.fill.background()

    sidebar_text = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(1.8), Inches(4.3))
    stf = sidebar_text.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = "Sidebar\n(240px)"
    sp.font.size = Pt(12)
    sp.font.bold = True
    sp.font.color.rgb = THEME['white']
    sp.font.name = "맑은 고딕"

    menu_items = ["• 대시보드", "• 검색", "• 내 지식", "• 북마크", "• 설정"]
    for item in menu_items:
        mp = stf.add_paragraph()
        mp.text = item
        mp.font.size = Pt(11)
        mp.font.color.rgb = THEME['light']
        mp.font.name = "맑은 고딕"

    # Main Content
    main = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(2.6), Inches(2.1),
                                  Inches(6.9), Inches(4.5))
    main.fill.solid()
    main.fill.fore_color.rgb = THEME['light']
    main.line.color.rgb = THEME['secondary']
    main.line.width = Pt(2)

    # Breadcrumbs
    bread = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(2.6), Inches(2.1),
                                   Inches(6.9), Inches(0.5))
    bread.fill.solid()
    bread.fill.fore_color.rgb = THEME['white']
    bread.line.color.rgb = THEME['secondary']

    bread_text = slide.shapes.add_textbox(Inches(2.75), Inches(2.18), Inches(6.5), Inches(0.35))
    btf = bread_text.text_frame
    bp = btf.paragraphs[0]
    bp.text = "Breadcrumbs: Home > Knowledge > Detail"
    bp.font.size = Pt(11)
    bp.font.color.rgb = THEME['text']
    bp.font.name = "맑은 고딕"

    # Main Content Label
    main_label = slide.shapes.add_textbox(Inches(4.5), Inches(4.0), Inches(3), Inches(0.5))
    mtf = main_label.text_frame
    mp = mtf.paragraphs[0]
    mp.text = "Main Content Area"
    mp.font.size = Pt(18)
    mp.font.bold = True
    mp.font.color.rgb = THEME['secondary']
    mp.font.name = "맑은 고딕"
    mp.alignment = PP_ALIGN.CENTER

    # Footer
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.5), Inches(6.7),
                                    Inches(9.0), Inches(0.5))
    footer.fill.solid()
    footer.fill.fore_color.rgb = THEME['secondary']
    footer.line.fill.background()

    footer_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(8.6), Inches(0.35))
    ftf = footer_text.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "Footer (48px) - Copyright © 2026"
    fp.font.size = Pt(11)
    fp.font.color.rgb = THEME['white']
    fp.font.name = "맑은 고딕"

    return slide


def add_performance_slide(prs):
    """성능 최적화 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "성능 최적화 전략"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    strategies = [
        ("코드 스플리팅", "Route-based & Component-based lazy loading으로 초기 로딩 시간 단축",
         "React.lazy + Suspense", THEME['react']),
        ("메모이제이션", "불필요한 리렌더링 방지로 성능 향상",
         "React.memo, useCallback, useMemo", THEME['typescript']),
        ("가상화", "대량 리스트 렌더링 최적화",
         "react-window, react-virtualized-auto-sizer", THEME['redux']),
        ("이미지 최적화", "Lazy loading + Intersection Observer로 필요한 시점에 로드",
         "LazyImage 컴포넌트", THEME['accent']),
        ("캐싱 전략", "React Query staleTime/gcTime으로 API 호출 최소화",
         "5분 stale, 30분 cache", THEME['query']),
        ("번들 최적화", "Vite 기반 빠른 빌드, Tree Shaking, 코드 압축",
         "ES modules, Rollup", THEME['mui']),
    ]

    start_y = 1.2
    row_height = 1.0

    for i, (title, desc, tech, color) in enumerate(strategies):
        y = start_y + i * (row_height + 0.05)

        # 아이콘 박스
        icon_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(0.4), Inches(y),
                                          Inches(0.1), Inches(row_height - 0.1))
        icon_box.fill.solid()
        icon_box.fill.fore_color.rgb = color
        icon_box.line.fill.background()

        # 제목
        title_text = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(2.3), Inches(0.4))
        ttf = title_text.text_frame
        tp = ttf.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(14)
        tp.font.bold = True
        tp.font.color.rgb = color
        tp.font.name = "맑은 고딕"

        # 설명
        desc_text = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.35), Inches(5.5), Inches(0.5))
        dtf = desc_text.text_frame
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(11)
        dp.font.color.rgb = THEME['text']
        dp.font.name = "맑은 고딕"

        # 기술
        tech_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(6.3), Inches(y + 0.15),
                                          Inches(3.3), Inches(0.5))
        tech_box.fill.solid()
        tech_box.fill.fore_color.rgb = THEME['light']
        tech_box.line.color.rgb = color
        tech_box.line.width = Pt(1)

        tech_text = slide.shapes.add_textbox(Inches(6.4), Inches(y + 0.22), Inches(3.1), Inches(0.4))
        thtf = tech_text.text_frame
        thp = thtf.paragraphs[0]
        thp.text = tech
        thp.font.size = Pt(10)
        thp.font.color.rgb = color
        thp.font.name = "Consolas"
        thp.alignment = PP_ALIGN.CENTER

    return slide


def add_qa_slide(prs):
    """Q&A 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['dark']
    bg.line.fill.background()

    # Q&A 제목
    qa_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(10), Inches(1))
    tf = qa_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Q & A"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    # 핵심 하이라이트
    highlights = [
        "React 18 + TypeScript + Vite",
        "Feature-Based 모듈화 아키텍처",
        "Server/Client 상태 관리 분리",
        "MUI 기반 한글 최적화 디자인 시스템",
        "성능 최적화 (코드 스플리팅, 캐싱, 가상화)",
    ]

    hl_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2.5))
    hltf = hl_box.text_frame
    hltf.word_wrap = True

    for i, hl in enumerate(highlights):
        if i == 0:
            hp = hltf.paragraphs[0]
        else:
            hp = hltf.add_paragraph()
        hp.text = f"✓ {hl}"
        hp.font.size = Pt(14)
        hp.font.color.rgb = THEME['accent']
        hp.font.name = "맑은 고딕"
        hp.alignment = PP_ALIGN.CENTER
        hp.space_after = Pt(8)

    return slide


def main():
    """메인 함수"""
    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 슬라이드 생성
    # 1. 표지
    add_title_slide(prs,
                    "Frontend 상세 설계서",
                    "React 기반 Knowledge Discovery Platform")

    # 2. 개요
    add_content_slide(prs, "개요", [
        "React 18 기반 SPA 애플리케이션",
        "TypeScript 전면 적용으로 타입 안정성 확보",
        "Material-UI v5 기반 디자인 시스템",
        "Redux Toolkit + React Query 상태 관리",
        "반응형 웹 디자인 (Desktop, Tablet, Mobile)",
        "핵심 기능: 대시보드, 지식 검색(RAG), 지식 관리, 개인화, 문서 변환, 관리자"
    ])

    # 3. 기술 스택
    add_tech_stack_slide(prs)

    # 4. 프로젝트 구조
    add_project_structure_slide(prs)

    # 5. 컴포넌트 아키텍처
    add_component_architecture_slide(prs)

    # 6. 상태 관리 전략
    add_state_management_slide(prs)

    # 7. 아키텍처 장점
    add_architecture_advantages_slide(prs)

    # 8. 레이아웃 구조
    add_layout_diagram_slide(prs)

    # 9. 성능 최적화
    add_performance_slide(prs)

    # 10. Q&A
    add_qa_slide(prs)

    # 저장
    output_dir = "knowledge_service/docs/02_design"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "Frontend_Detailed_Design_Brown.pptx")
    prs.save(output_path)

    print(f"프레젠테이션 생성 완료: {output_path}")
    print(f"총 슬라이드 수: {len(prs.slides)}")

    return output_path


if __name__ == "__main__":
    main()
