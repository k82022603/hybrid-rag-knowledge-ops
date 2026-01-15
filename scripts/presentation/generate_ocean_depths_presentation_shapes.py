#!/usr/bin/env python3
"""Ocean Depths 테마 프레젠테이션 생성 (PPTX 도형 버전)"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# Ocean Depths 컬러 팔레트
COLORS = {
    'primary': RGBColor(0, 102, 153),      # #006699
    'secondary': RGBColor(0, 153, 204),    # #0099CC
    'accent': RGBColor(255, 102, 0),       # #FF6600
    'dark': RGBColor(0, 51, 102),          # #003366
    'light': RGBColor(204, 229, 255),      # #CCE5FF
    'text': RGBColor(33, 33, 33),          # #212121
    'white': RGBColor(255, 255, 255)       # #FFFFFF
}

def add_arrow(slide, x1, y1, x2, y2, color=None):
    """화살표 추가"""
    if color is None:
        color = COLORS['accent']

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(4)

    # 화살표 머리 추가
    connector.line.end_arrow_type = 2  # 화살표

    return connector

def add_rounded_box(slide, left, top, width, height, text, fill_color, text_color=None, font_size=20):
    """둥근 모서리 박스 추가"""
    if text_color is None:
        text_color = COLORS['white']

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = COLORS['dark']
    shape.line.width = Pt(2)

    # 텍스트 추가
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = text_color
    paragraph.font.name = "맑은 고딕"
    paragraph.font.bold = True
    paragraph.alignment = PP_ALIGN.CENTER

    return shape

def create_architecture_diagram_slide(prs):
    """시스템 아키텍처 다이어그램 슬라이드 (도형 사용)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['secondary']
    header.line.color.rgb = COLORS['secondary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "시스템 아키텍처"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "맑은 고딕"

    # Frontend Layer
    add_rounded_box(slide, 0.5, 1.2, 9, 0.8,
                   "Frontend Layer\nReact 18 + TypeScript",
                   COLORS['secondary'], font_size=22)

    # API Gateway
    add_rounded_box(slide, 0.5, 2.3, 9, 0.7,
                   "API Gateway - Spring Cloud Gateway",
                   COLORS['primary'], font_size=20)

    # Service Layer
    services = [
        ("Knowledge\nService", 0.5),
        ("Search\nService", 2.9),
        ("User\nService", 5.3),
        ("AI Service\n(FastAPI)", 7.7)
    ]

    for service_name, x_pos in services:
        add_rounded_box(slide, x_pos, 3.4, 2.2, 1.0,
                       service_name, COLORS['accent'], font_size=16)

    # Database Layer
    databases = [
        ("PostgreSQL", 0.5),
        ("Elasticsearch", 2.9),
        ("Neo4j", 5.3),
        ("Redis", 7.7)
    ]

    for db_name, x_pos in databases:
        add_rounded_box(slide, x_pos, 4.8, 2.2, 0.9,
                       db_name, COLORS['dark'], font_size=18)

    # 화살표 추가
    arrow_positions = [1.6, 4.0, 6.4, 8.8]

    # Frontend -> API Gateway
    for x in arrow_positions:
        add_arrow(slide, x, 2.0, x, 2.3)

    # API Gateway -> Services
    for x in arrow_positions:
        add_arrow(slide, x, 3.0, x, 3.4)

    # Services -> Databases
    for x in arrow_positions:
        add_arrow(slide, x, 4.4, x, 4.8)

def create_dashboard_diagram_slide(prs):
    """대시보드 다이어그램 슬라이드 (도형 사용)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['dark']
    header.line.color.rgb = COLORS['dark']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "📊 대시보드 - 지식 검색 시스템"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    # 통계 카드 4개
    cards = [
        ("📚 전체 문서\n\n1,234", 0.5),
        ("🔍 금주 검색\n\n567", 2.7),
        ("⭐ 인기 문서\n\n89", 4.9),
        ("👥 활성 사용자\n\n45", 7.1)
    ]

    for card_text, x_pos in cards:
        box = add_rounded_box(slide, x_pos, 1.3, 2.0, 1.3,
                             card_text, COLORS['light'],
                             text_color=COLORS['text'], font_size=16)

    # 최근 활동 패널
    activity_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(3.0),
        Inches(4.3), Inches(2.8)
    )
    activity_box.fill.solid()
    activity_box.fill.fore_color.rgb = COLORS['white']
    activity_box.line.color.rgb = COLORS['primary']
    activity_box.line.width = Pt(3)

    # 최근 활동 제목
    activity_title = slide.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(4), Inches(0.4))
    tf = activity_title.text_frame
    tf.text = "최근 활동"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['text']
    p.font.name = "맑은 고딕"

    # 활동 내역
    activities = [
        "• 프로젝트 A 문서 업데이트",
        "• 신규 사용자 5명 가입",
        "• API 문서 검색 증가"
    ]

    y_pos = 3.8
    for activity in activities:
        act_text = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(4), Inches(0.4))
        tf = act_text.text_frame
        tf.text = activity
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text']
        p.font.name = "맑은 고딕"
        y_pos += 0.5

    # 인기 검색어 패널
    keyword_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(3.0),
        Inches(4.3), Inches(2.8)
    )
    keyword_box.fill.solid()
    keyword_box.fill.fore_color.rgb = COLORS['white']
    keyword_box.line.color.rgb = COLORS['secondary']
    keyword_box.line.width = Pt(3)

    # 인기 검색어 제목
    keyword_title = slide.shapes.add_textbox(Inches(5.4), Inches(3.2), Inches(4), Inches(0.4))
    tf = keyword_title.text_frame
    tf.text = "인기 검색어"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['text']
    p.font.name = "맑은 고딕"

    # 검색어 목록
    keywords = ["API 문서", "배포 가이드", "아키텍처", "인증 방법"]

    y_pos = 3.8
    for i, keyword in enumerate(keywords, 1):
        kw_text = slide.shapes.add_textbox(Inches(5.4), Inches(y_pos), Inches(4), Inches(0.4))
        tf = kw_text.text_frame
        tf.text = f"{i}. {keyword}"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text']
        p.font.name = "맑은 고딕"
        y_pos += 0.5

def create_search_modes_diagram_slide(prs):
    """검색 모드 비교 다이어그램 슬라이드 (도형 사용)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['secondary']
    header.line.color.rgb = COLORS['secondary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "검색 기능 - 채팅 vs 검색 모드"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "맑은 고딕"

    # 왼쪽: 채팅 모드 패널
    chat_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.0),
        Inches(4.6), Inches(5.8)
    )
    chat_panel.fill.solid()
    chat_panel.fill.fore_color.rgb = COLORS['light']
    chat_panel.line.color.rgb = COLORS['primary']
    chat_panel.line.width = Pt(4)

    # 채팅 모드 헤더
    chat_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.0),
        Inches(4.6), Inches(0.8)
    )
    chat_header.fill.solid()
    chat_header.fill.fore_color.rgb = COLORS['primary']
    chat_header.line.color.rgb = COLORS['primary']

    chat_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(0.4))
    tf = chat_title.text_frame
    tf.text = "💬 채팅 모드"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    # 채팅 모드 특징
    chat_features = [
        "✓ RAG 기반 답변 생성",
        "✓ 자연어 질문 처리",
        "✓ 컨텍스트 유지",
        "✓ 참조 문서 링크 제공",
        "✓ 대화형 인터페이스"
    ]

    y_pos = 2.2
    for feature in chat_features:
        feat_text = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos), Inches(4), Inches(0.4))
        tf = feat_text.text_frame
        tf.text = feature
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text']
        p.font.name = "맑은 고딕"
        y_pos += 0.65

    # 채팅 모드 적합 용도
    chat_usage = slide.shapes.add_textbox(Inches(0.6), Inches(6.0), Inches(4), Inches(0.5))
    tf = chat_usage.text_frame
    tf.text = "최적: 복잡한 질문"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    # 오른쪽: 검색 모드 패널
    search_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.1), Inches(1.0),
        Inches(4.6), Inches(5.8)
    )
    search_panel.fill.solid()
    search_panel.fill.fore_color.rgb = COLORS['light']
    search_panel.line.color.rgb = COLORS['secondary']
    search_panel.line.width = Pt(4)

    # 검색 모드 헤더
    search_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.1), Inches(1.0),
        Inches(4.6), Inches(0.8)
    )
    search_header.fill.solid()
    search_header.fill.fore_color.rgb = COLORS['secondary']
    search_header.line.color.rgb = COLORS['secondary']

    search_title = slide.shapes.add_textbox(Inches(5.3), Inches(1.2), Inches(4.2), Inches(0.4))
    tf = search_title.text_frame
    tf.text = "🔍 검색 모드"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    # 검색 모드 특징
    search_features = [
        "✓ 하이브리드 검색 (BM25+벡터)",
        "✓ 필터링 (날짜, 타입 등)",
        "✓ 정렬 옵션",
        "✓ 페이지네이션",
        "✓ 북마크 기능"
    ]

    y_pos = 2.2
    for feature in search_features:
        feat_text = slide.shapes.add_textbox(Inches(5.4), Inches(y_pos), Inches(4), Inches(0.4))
        tf = feat_text.text_frame
        tf.text = feature
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text']
        p.font.name = "맑은 고딕"
        y_pos += 0.65

    # 검색 모드 적합 용도
    search_usage = slide.shapes.add_textbox(Inches(5.4), Inches(6.0), Inches(4), Inches(0.5))
    tf = search_usage.text_frame
    tf.text = "최적: 정확한 키워드"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

def create_title_slide(prs, title, subtitle):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['dark']
    bg.line.color.rgb = COLORS['dark']

    # 장식
    wave = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(6), Inches(5),
        Inches(5), Inches(3)
    )
    wave.fill.solid()
    wave.fill.fore_color.rgb = COLORS['primary']
    wave.line.color.rgb = COLORS['primary']

    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    # 부제목
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['light']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

def create_content_slide(prs, title, content_lines):
    """컨텐츠 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    header.line.color.rgb = COLORS['primary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "맑은 고딕"

    # 내용
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for line in content_lines:
        p = content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['text']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(12)

def create_two_column_slide(prs, title, left_items, right_items):
    """2단 컬럼 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    header.line.color.rgb = COLORS['primary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "맑은 고딕"

    # 왼쪽 패널
    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5),
        Inches(4.25), Inches(5.5)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = COLORS['light']
    left_panel.line.color.rgb = COLORS['primary']
    left_panel.line.width = Pt(3)

    left_text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(3.7), Inches(5))
    left_frame = left_text_box.text_frame
    left_frame.word_wrap = True

    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(10)

    # 오른쪽 패널
    right_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.25), Inches(1.5),
        Inches(4.25), Inches(5.5)
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = COLORS['light']
    right_panel.line.color.rgb = COLORS['secondary']
    right_panel.line.width = Pt(3)

    right_text_box = slide.shapes.add_textbox(Inches(5.55), Inches(1.8), Inches(3.7), Inches(5))
    right_frame = right_text_box.text_frame
    right_frame.word_wrap = True

    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(10)

def main():
    output_dir = r"knowledge_service\docs\02_design\ui_storyboard"
    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("프레젠테이션 생성 중...")

    # 1. 표지
    create_title_slide(prs, "사내 지식 검색 시스템", "UI 스토리보드 발표")

    # 2. 개요
    create_content_slide(prs, "프로젝트 개요", [
        "📌 목표",
        "  • Graph RAG 기반 지능형 지식 검색",
        "  • 사내 문서의 효율적 관리 및 활용",
        "",
        "💡 핵심 기능",
        "  • AI 기반 자연어 질의응답",
        "  • 하이브리드 검색 (키워드 + 벡터)",
        "  • 엔티티 기반 관계 그래프",
        "  • 시간 기반 지식 버전 관리"
    ])

    # 3. 시스템 아키텍처 (도형)
    create_architecture_diagram_slide(prs)

    # 4. 로그인 & 대시보드 (도형)
    create_dashboard_diagram_slide(prs)

    # 5. 검색 기능 (도형)
    create_search_modes_diagram_slide(prs)

    # 6. 지식 관리
    create_two_column_slide(prs, "지식 관리", [
        "📋 지식 목록",
        "• 카드 뷰 / 리스트 뷰",
        "• 필터링 (프로젝트, 타입, 날짜)",
        "• 정렬 (최신순, 인기순)",
        "• 무한 스크롤",
        "",
        "📝 지식 작성/수정",
        "• 마크다운 에디터",
        "• 파일 업로드 (PDF, DOCX)",
        "• 메타데이터 자동 추출",
        "• 미리보기 기능"
    ], [
        "📄 지식 상세",
        "• 풀스크린 뷰어",
        "• 목차 네비게이션",
        "• 관련 문서 추천",
        "• 참조 그래프 시각화",
        "",
        "🔖 부가 기능",
        "• 북마크",
        "• 댓글/피드백",
        "• 버전 히스토리",
        "• 공유 링크 생성"
    ])

    # 7. 프로필 & 관리자
    create_two_column_slide(prs, "프로필 & 관리자", [
        "👤 프로필 관리",
        "• 개인정보 수정",
        "• 비밀번호 변경",
        "• 알림 설정",
        "• 테마 설정 (라이트/다크)",
        "",
        "⭐ 북마크",
        "• 즐겨찾기 문서 목록",
        "• 태그별 분류",
        "• 검색 히스토리"
    ], [
        "🔧 관리자 페이지",
        "• 사용자 관리 (CRUD)",
        "• 권한 관리 (역할 기반)",
        "• 지식 승인/거부",
        "• 통계 대시보드",
        "",
        "📊 시스템 모니터링",
        "• 검색 분석",
        "• 성능 지표",
        "• 사용량 리포트"
    ])

    # 8. 기술 스택
    create_two_column_slide(prs, "기술 스택", [
        "🎨 Frontend",
        "• React 18 + TypeScript",
        "• Material-UI (MUI)",
        "• TanStack Query (데이터 관리)",
        "• Recharts (차트)",
        "",
        "🔧 Backend",
        "• Spring Boot 3.x (Java 17)",
        "• Spring Security (JWT)",
        "• FastAPI (AI Service)"
    ], [
        "🗄️ Database",
        "• PostgreSQL (마스터 DB)",
        "• Elasticsearch (검색/벡터)",
        "• Neo4j (그래프)",
        "• Redis (캐시)",
        "",
        "🤖 AI/ML",
        "• DeepSeek-V3.2 (LLM)",
        "• BGE-M3 (임베딩)",
        "• LangGraph (워크플로우)"
    ])

    # 9. 기대 효과
    create_content_slide(prs, "기대 효과", [
        "⏰ 검색 시간 단축",
        "  • 기존 대비 70% 이상 시간 절감",
        "  • AI 기반 즉시 답변 제공",
        "",
        "📈 지식 활용도 증대",
        "  • 관련 문서 자동 추천",
        "  • 시간 기반 정확한 정보 제공",
        "",
        "💰 비용 효율",
        "  • DeepSeek 활용으로 95% 비용 절감",
        "  • 오픈소스 기반 확장 가능"
    ])

    # 10. Q&A
    create_title_slide(prs, "Q & A", "질문해 주세요")

    # 저장
    output_path = os.path.join(output_dir, "UI스토리보드_발표자료_OceanDepths.pptx")
    prs.save(output_path)

    print(f"\n프레젠테이션 생성 완료!")
    print(f"파일 위치: {output_path}")
    print(f"총 슬라이드 수: {len(prs.slides)}")
    print(f"Ocean Depths 테마 적용")
    print(f"다이어그램: PPTX 도형으로 제작 (편집 가능)")

if __name__ == "__main__":
    main()
