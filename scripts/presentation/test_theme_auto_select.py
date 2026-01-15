#!/usr/bin/env python3
"""테마 자동 선택 테스트"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os

# ========== 테마 팔레트 정의 ==========

OCEAN_DEPTHS = {
    'primary': RGBColor(0, 102, 153),
    'secondary': RGBColor(0, 153, 204),
    'accent': RGBColor(255, 102, 0),
    'dark': RGBColor(0, 51, 102),
    'light': RGBColor(204, 229, 255),
    'text': RGBColor(33, 33, 33),
    'white': RGBColor(255, 255, 255)
}

TECH_INNOVATION = {
    'primary': RGBColor(46, 125, 50),
    'secondary': RGBColor(129, 199, 132),
    'accent': RGBColor(255, 152, 0),
    'dark': RGBColor(27, 94, 32),
    'light': RGBColor(232, 245, 233),
    'text': RGBColor(33, 33, 33),
    'white': RGBColor(255, 255, 255)
}

FOREST_CANOPY = {
    'primary': RGBColor(56, 142, 60),
    'secondary': RGBColor(139, 195, 74),
    'accent': RGBColor(121, 85, 72),
    'dark': RGBColor(27, 94, 32),
    'light': RGBColor(241, 248, 233),
    'text': RGBColor(33, 33, 33),
    'white': RGBColor(255, 255, 255)
}

ARCTIC_FROST = {
    'primary': RGBColor(3, 169, 244),
    'secondary': RGBColor(79, 195, 247),
    'accent': RGBColor(0, 188, 212),
    'dark': RGBColor(1, 87, 155),
    'light': RGBColor(225, 245, 254),
    'text': RGBColor(33, 33, 33),
    'white': RGBColor(255, 255, 255)
}

SUNSET_BOULEVARD = {
    'primary': RGBColor(255, 87, 34),
    'secondary': RGBColor(255, 193, 7),
    'accent': RGBColor(156, 39, 176),
    'dark': RGBColor(191, 54, 12),
    'light': RGBColor(255, 243, 224),
    'text': RGBColor(33, 33, 33),
    'white': RGBColor(255, 255, 255)
}

MIDNIGHT_GALAXY = {
    'primary': RGBColor(63, 81, 181),
    'secondary': RGBColor(121, 134, 203),
    'accent': RGBColor(255, 64, 129),
    'dark': RGBColor(26, 35, 126),
    'light': RGBColor(232, 234, 246),
    'text': RGBColor(255, 255, 255),
    'white': RGBColor(255, 255, 255)
}

# 모든 테마 팔레트
THEME_PALETTES = {
    'tech_innovation': TECH_INNOVATION,
    'ocean_depths': OCEAN_DEPTHS,
    'forest_canopy': FOREST_CANOPY,
    'arctic_frost': ARCTIC_FROST,
    'sunset_boulevard': SUNSET_BOULEVARD,
    'midnight_galaxy': MIDNIGHT_GALAXY
}

# ========== 테마 자동 선택 함수 ==========

def auto_select_theme(content_text, title=""):
    """
    콘텐츠와 제목을 분석하여 적절한 테마를 자동 선택
    """
    text = (content_text + " " + title).lower()

    theme_keywords = {
        'tech_innovation': [
            'ai', '인공지능', 'ml', '머신러닝', 'api', '개발', 'developer',
            '코드', 'code', '프로그래밍', 'software', '소프트웨어', 'devops',
            '클라우드', 'cloud', '데이터', 'data', 'tech', '기술', 'it',
            'rag', 'llm', '딥러닝', 'deep learning', 'python', 'java',
            '시스템', 'system', '아키텍처', 'architecture', '플랫폼', 'platform'
        ],
        'ocean_depths': [
            '기업', 'enterprise', '비즈니스', 'business', '전략', 'strategy',
            '금융', 'finance', '투자', 'investment', '컨설팅', 'consulting',
            '경영', 'management', '조직', 'organization', '리더십', 'leadership',
            '분기', 'quarter', '실적', 'performance', '보고', 'report'
        ],
        'forest_canopy': [
            '환경', 'environment', '지속가능', 'sustainable', 'esg', '친환경',
            'green', '탄소', 'carbon', '에너지', 'energy', '재활용', 'recycle',
            '자연', 'nature', '생태', 'ecology', '기후', 'climate'
        ],
        'arctic_frost': [
            '의료', 'medical', '헬스케어', 'healthcare', '병원', 'hospital',
            '제약', 'pharmaceutical', '바이오', 'bio', '임상', 'clinical',
            '과학', 'science', '연구', 'research', '실험', 'experiment'
        ],
        'sunset_boulevard': [
            '마케팅', 'marketing', '광고', 'advertising', '캠페인', 'campaign',
            '브랜드', 'brand', '이벤트', 'event', '프로모션', 'promotion',
            '소셜', 'social', '미디어', 'media', '콘텐츠', 'content'
        ],
        'midnight_galaxy': [
            '게임', 'game', '엔터테인먼트', 'entertainment',
            '영화', 'movie', '음악', 'music', '스트리밍', 'streaming',
            'vr', 'ar', '메타버스', 'metaverse'
        ]
    }

    scores = {theme: 0 for theme in theme_keywords}

    for theme, keywords in theme_keywords.items():
        for keyword in keywords:
            if keyword in text:
                scores[theme] += 1

    best_theme = max(scores, key=scores.get)

    if scores[best_theme] == 0:
        best_theme = 'ocean_depths'

    return best_theme, THEME_PALETTES[best_theme], scores

# ========== 슬라이드 생성 함수 ==========

def create_title_slide(prs, title, subtitle, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme['dark']
    bg.line.color.rgb = theme['dark']

    deco = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(6), Inches(5), Inches(5), Inches(3)
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = theme['primary']
    deco.line.color.rgb = theme['primary']

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = theme['white']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    sf = sub_box.text_frame
    sf.text = subtitle
    sp = sf.paragraphs[0]
    sp.font.size = Pt(24)
    sp.font.color.rgb = theme['light']
    sp.font.name = "맑은 고딕"
    sp.alignment = PP_ALIGN.CENTER

    return slide

def create_content_slide(prs, title, content_lines, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = theme['primary']
    header.line.color.rgb = theme['primary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = theme['white']
    p.font.name = "맑은 고딕"

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    cf = content_box.text_frame
    cf.word_wrap = True

    for line in content_lines:
        p = cf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = theme['text']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(12)

    return slide

# ========== 테스트 케이스 ==========

test_cases = [
    {
        "title": "AI 기반 RAG 시스템 구축",
        "subtitle": "Hybrid RAG Knowledge Platform",
        "content": [
            "LLM과 벡터 데이터베이스 통합",
            "DeepSeek-V3.2 API 활용",
            "시스템 아키텍처 설계",
            "Python FastAPI 개발"
        ],
        "expected": "tech_innovation"
    },
    {
        "title": "2024년 4분기 경영 실적 보고",
        "subtitle": "비즈니스 전략 및 성과 분석",
        "content": [
            "매출 성장률 분석",
            "투자 수익률 (ROI)",
            "조직 구조 개편 현황",
            "리더십 프로그램 성과"
        ],
        "expected": "ocean_depths"
    },
    {
        "title": "ESG 지속가능경영 보고서",
        "subtitle": "탄소중립 실현을 위한 전략",
        "content": [
            "환경 영향 저감 활동",
            "친환경 에너지 전환",
            "재활용률 개선",
            "기후변화 대응 계획"
        ],
        "expected": "forest_canopy"
    },
    {
        "title": "신약 임상시험 결과 발표",
        "subtitle": "바이오 의약품 연구 성과",
        "content": [
            "임상 3상 실험 결과",
            "헬스케어 혁신",
            "과학적 검증 데이터",
            "제약 산업 동향"
        ],
        "expected": "arctic_frost"
    },
    {
        "title": "신제품 마케팅 캠페인",
        "subtitle": "브랜드 인지도 향상 전략",
        "content": [
            "소셜 미디어 광고",
            "인플루언서 마케팅",
            "이벤트 프로모션",
            "콘텐츠 마케팅"
        ],
        "expected": "sunset_boulevard"
    },
    {
        "title": "메타버스 게임 출시 발표",
        "subtitle": "VR/AR 엔터테인먼트 플랫폼",
        "content": [
            "게임 플레이 소개",
            "VR 기술 적용",
            "스트리밍 서비스",
            "엔터테인먼트 콘텐츠"
        ],
        "expected": "midnight_galaxy"
    }
]

# ========== 테스트 실행 ==========

def run_tests():
    print("=" * 60)
    print("테마 자동 선택 테스트")
    print("=" * 60)

    results = []

    for i, case in enumerate(test_cases, 1):
        content_text = case["title"] + " " + case["subtitle"] + " " + " ".join(case["content"])
        selected_theme, palette, scores = auto_select_theme(content_text, case["title"])

        is_correct = selected_theme == case["expected"]
        results.append(is_correct)

        status = "PASS" if is_correct else "FAIL"

        print(f"\n[Test {i}] {status}")
        print(f"  제목: {case['title']}")
        print(f"  예상 테마: {case['expected']}")
        print(f"  선택된 테마: {selected_theme}")
        print(f"  점수: {scores}")

    print("\n" + "=" * 60)
    passed = sum(results)
    total = len(results)
    print(f"테스트 결과: {passed}/{total} 통과 ({passed/total*100:.0f}%)")
    print("=" * 60)

    return results

def create_sample_presentations():
    """샘플 프레젠테이션 생성 (처음 2개만)"""
    output_dir = "knowledge_service/docs/02_design/ui_storyboard"
    os.makedirs(output_dir, exist_ok=True)

    print("\n샘플 프레젠테이션 생성...")

    for i, case in enumerate(test_cases[:2], 1):
        content_text = case["title"] + " " + case["subtitle"] + " " + " ".join(case["content"])
        theme_name, theme, _ = auto_select_theme(content_text, case["title"])

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # 표지
        create_title_slide(prs, case["title"], case["subtitle"], theme)

        # 내용
        create_content_slide(prs, "주요 내용", case["content"], theme)

        # 저장
        filename = f"test_theme_{theme_name}.pptx"
        output_path = os.path.join(output_dir, filename)
        prs.save(output_path)

        print(f"  [{i}] {filename} (테마: {theme_name})")

    print("\n샘플 프레젠테이션 생성 완료!")

if __name__ == "__main__":
    run_tests()
    create_sample_presentations()
