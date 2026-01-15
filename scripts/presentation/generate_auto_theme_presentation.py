#!/usr/bin/env python3
"""테마 자동 선택 프레젠테이션 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os

# ========== 테마 팔레트 정의 ==========

OCEAN_DEPTHS = {
    'primary': RGBColor(0, 102, 153),
    'secondary': RGBColor(0, 153, 204),
    'accent': RGBColor(255, 102, 0),
    'dark': RGBColor(0, 51, 102),
    'light': RGBColor(204, 229, 255),
    'text': RGBColor(33, 33, 33),
    'white': RGBColor(255, 255, 255)
}

TECH_INNOVATION = {
    'primary': RGBColor(46, 125, 50),
    'secondary': RGBColor(129, 199, 132),
    'accent': RGBColor(255, 152, 0),
    'dark': RGBColor(27, 94, 32),
    'light': RGBColor(232, 245, 233),
    'text': RGBColor(33, 33, 33),
    'white': RGBColor(255, 255, 255)
}

FOREST_CANOPY = {
    'primary': RGBColor(56, 142, 60),
    'secondary': RGBColor(139, 195, 74),
    'accent': RGBColor(121, 85, 72),
    'dark': RGBColor(27, 94, 32),
    'light': RGBColor(241, 248, 233),
    'text': RGBColor(33, 33, 33),
    'white': RGBColor(255, 255, 255)
}

ARCTIC_FROST = {
    'primary': RGBColor(3, 169, 244),
    'secondary': RGBColor(79, 195, 247),
    'accent': RGBColor(0, 188, 212),
    'dark': RGBColor(1, 87, 155),
    'light': RGBColor(225, 245, 254),
    'text': RGBColor(33, 33, 33),
    'white': RGBColor(255, 255, 255)
}

SUNSET_BOULEVARD = {
    'primary': RGBColor(255, 87, 34),
    'secondary': RGBColor(255, 193, 7),
    'accent': RGBColor(156, 39, 176),
    'dark': RGBColor(191, 54, 12),
    'light': RGBColor(255, 243, 224),
    'text': RGBColor(33, 33, 33),
    'white': RGBColor(255, 255, 255)
}

MIDNIGHT_GALAXY = {
    'primary': RGBColor(63, 81, 181),
    'secondary': RGBColor(121, 134, 203),
    'accent': RGBColor(255, 64, 129),
    'dark': RGBColor(26, 35, 126),
    'light': RGBColor(232, 234, 246),
    'text': RGBColor(255, 255, 255),
    'white': RGBColor(255, 255, 255)
}

THEME_PALETTES = {
    'tech_innovation': TECH_INNOVATION,
    'ocean_depths': OCEAN_DEPTHS,
    'forest_canopy': FOREST_CANOPY,
    'arctic_frost': ARCTIC_FROST,
    'sunset_boulevard': SUNSET_BOULEVARD,
    'midnight_galaxy': MIDNIGHT_GALAXY
}

THEME_NAMES_KR = {
    'tech_innovation': 'Tech Innovation (녹색 계열)',
    'ocean_depths': 'Ocean Depths (파란색 계열)',
    'forest_canopy': 'Forest Canopy (자연 녹색)',
    'arctic_frost': 'Arctic Frost (하늘색 계열)',
    'sunset_boulevard': 'Sunset Boulevard (오렌지 계열)',
    'midnight_galaxy': 'Midnight Galaxy (인디고 계열)'
}

# ========== 테마 자동 선택 ==========

def auto_select_theme(content_text, title=""):
    """콘텐츠 분석 기반 테마 자동 선택"""
    text = (content_text + " " + title).lower()

    theme_keywords = {
        'tech_innovation': [
            'ai', '인공지능', 'ml', '머신러닝', 'api', '개발', 'developer',
            '코드', 'code', '프로그래밍', 'software', '소프트웨어', 'devops',
            '클라우드', 'cloud', '데이터', 'data', 'tech', '기술', 'it',
            'rag', 'llm', '딥러닝', 'deep learning', 'python', 'java',
            '시스템', 'system', '아키텍처', 'architecture', '플랫폼', 'platform',
            'elasticsearch', 'neo4j', 'postgresql', 'fastapi', 'spring'
        ],
        'ocean_depths': [
            '기업', 'enterprise', '비즈니스', 'business', '전략', 'strategy',
            '금융', 'finance', '투자', 'investment', '컨설팅', 'consulting',
            '경영', 'management', '조직', 'organization', '리더십', 'leadership',
            '분기', 'quarter', '실적', 'performance', '보고', 'report'
        ],
        'forest_canopy': [
            '환경', 'environment', '지속가능', 'sustainable', 'esg', '친환경',
            'green', '탄소', 'carbon', '에너지', 'energy', '재활용', 'recycle',
            '자연', 'nature', '생태', 'ecology', '기후', 'climate'
        ],
        'arctic_frost': [
            '의료', 'medical', '헬스케어', 'healthcare', '병원', 'hospital',
            '제약', 'pharmaceutical', '바이오', 'bio', '임상', 'clinical',
            '과학', 'science', '연구', 'research', '실험', 'experiment'
        ],
        'sunset_boulevard': [
            '마케팅', 'marketing', '광고', 'advertising', '캠페인', 'campaign',
            '브랜드', 'brand', '이벤트', 'event', '프로모션', 'promotion',
            '소셜', 'social', '미디어', 'media', '콘텐츠', 'content'
        ],
        'midnight_galaxy': [
            '게임', 'game', '엔터테인먼트', 'entertainment',
            '영화', 'movie', '음악', 'music', '스트리밍', 'streaming',
            'vr', 'ar', '메타버스', 'metaverse'
        ]
    }

    scores = {theme: 0 for theme in theme_keywords}

    for theme, keywords in theme_keywords.items():
        for keyword in keywords:
            if keyword in text:
                scores[theme] += 1

    best_theme = max(scores, key=scores.get)

    if scores[best_theme] == 0:
        best_theme = 'ocean_depths'

    return best_theme, THEME_PALETTES[best_theme], scores

# ========== 슬라이드 생성 함수 ==========

def add_arrow(slide, x1, y1, x2, y2, color):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(4)
    connector.line.end_arrow_type = 2
    return connector

def add_rounded_box(slide, left, top, width, height, text, fill_color, text_color=None, font_size=20):
    if text_color is None:
        text_color = RGBColor(255, 255, 255)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(33, 33, 33)
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = text_color
    paragraph.font.name = "맑은 고딕"
    paragraph.font.bold = True
    paragraph.alignment = PP_ALIGN.CENTER

    return shape

def create_title_slide(prs, title, subtitle, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme['dark']
    bg.line.color.rgb = theme['dark']

    deco = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(6), Inches(5), Inches(5), Inches(3)
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = theme['primary']
    deco.line.color.rgb = theme['primary']

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = theme['white']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    sf = sub_box.text_frame
    sf.text = subtitle
    sp = sf.paragraphs[0]
    sp.font.size = Pt(24)
    sp.font.color.rgb = theme['light']
    sp.font.name = "맑은 고딕"
    sp.alignment = PP_ALIGN.CENTER

def create_content_slide(prs, title, content_lines, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = theme['primary']
    header.line.color.rgb = theme['primary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = theme['white']
    p.font.name = "맑은 고딕"

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    cf = content_box.text_frame
    cf.word_wrap = True

    for line in content_lines:
        p = cf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = theme['text']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(12)

def create_architecture_slide(prs, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = theme['secondary']
    header.line.color.rgb = theme['secondary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "시스템 아키텍처"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme['white']
    p.font.name = "맑은 고딕"

    # Frontend
    add_rounded_box(slide, 0.5, 1.1, 9, 0.7,
                   "Frontend - React 18 + TypeScript",
                   theme['secondary'], font_size=18)

    # API Gateway
    add_rounded_box(slide, 0.5, 2.0, 9, 0.6,
                   "API Gateway - Spring Cloud",
                   theme['primary'], font_size=16)

    # Services
    services = [
        ("Knowledge\nService", 0.5),
        ("Search\nService", 2.9),
        ("User\nService", 5.3),
        ("AI Service\n(FastAPI)", 7.7)
    ]
    for name, x in services:
        add_rounded_box(slide, x, 2.9, 2.2, 0.9, name, theme['accent'], font_size=14)

    # Databases
    databases = [
        ("PostgreSQL", 0.5),
        ("Elasticsearch", 2.9),
        ("Neo4j", 5.3),
        ("Redis", 7.7)
    ]
    for name, x in databases:
        add_rounded_box(slide, x, 4.1, 2.2, 0.7, name, theme['dark'], font_size=16)

    # Arrows
    arrow_x = [1.6, 4.0, 6.4, 8.8]
    for x in arrow_x:
        add_arrow(slide, x, 1.8, x, 2.0, theme['accent'])
        add_arrow(slide, x, 2.6, x, 2.9, theme['accent'])
        add_arrow(slide, x, 3.8, x, 4.1, theme['accent'])

def create_two_column_slide(prs, title, left_items, right_items, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = theme['primary']
    header.line.color.rgb = theme['primary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = theme['white']
    p.font.name = "맑은 고딕"

    # Left panel
    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(4.25), Inches(5.5)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = theme['light']
    left_panel.line.color.rgb = theme['primary']
    left_panel.line.width = Pt(3)

    left_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(3.7), Inches(5))
    ltf = left_text.text_frame
    ltf.word_wrap = True
    for item in left_items:
        p = ltf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = theme['text']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(10)

    # Right panel
    right_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.25), Inches(1.5), Inches(4.25), Inches(5.5)
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = theme['light']
    right_panel.line.color.rgb = theme['secondary']
    right_panel.line.width = Pt(3)

    right_text = slide.shapes.add_textbox(Inches(5.55), Inches(1.8), Inches(3.7), Inches(5))
    rtf = right_text.text_frame
    rtf.word_wrap = True
    for item in right_items:
        p = rtf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = theme['text']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(10)

# ========== 메인 ==========

def main():
    # 프레젠테이션 콘텐츠 정의 (테마 지정 없음!)
    presentation = {
        "title": "Hybrid RAG 지식 검색 플랫폼",
        "subtitle": "AI 기반 사내 지식 관리 시스템",
        "slides": [
            {
                "type": "content",
                "title": "프로젝트 개요",
                "content": [
                    "목표: Graph RAG 기반 지능형 지식 검색",
                    "",
                    "핵심 기능:",
                    "  - AI 기반 자연어 질의응답 (RAG)",
                    "  - 하이브리드 검색 (BM25 + 벡터)",
                    "  - 엔티티 기반 관계 그래프",
                    "  - 시간 기반 지식 버전 관리",
                    "",
                    "기대 효과: 검색 시간 70% 단축"
                ]
            },
            {
                "type": "architecture"
            },
            {
                "type": "two_column",
                "title": "기술 스택",
                "left": [
                    "Frontend",
                    "  - React 18 + TypeScript",
                    "  - Material-UI",
                    "  - TanStack Query",
                    "",
                    "Backend",
                    "  - Spring Boot 3.x",
                    "  - FastAPI (AI Service)",
                    "  - Spring Security"
                ],
                "right": [
                    "Database",
                    "  - PostgreSQL (마스터)",
                    "  - Elasticsearch (검색)",
                    "  - Neo4j (그래프)",
                    "",
                    "AI/ML",
                    "  - DeepSeek-V3.2 (LLM)",
                    "  - BGE-M3 (임베딩)",
                    "  - LangGraph"
                ]
            },
            {
                "type": "content",
                "title": "비용 효율성",
                "content": [
                    "DeepSeek-V3.2 활용으로 95% 비용 절감",
                    "",
                    "비용 비교 (1,000문서 기준):",
                    "  - GPT-4: $45.50",
                    "  - DeepSeek-V3.2: $2.26",
                    "",
                    "오픈소스 기반 확장 가능",
                    "  - Elasticsearch (Apache 2.0)",
                    "  - Neo4j Community",
                    "  - LangGraph"
                ]
            }
        ]
    }

    # 콘텐츠에서 텍스트 추출
    content_text = presentation["title"] + " " + presentation["subtitle"]
    for slide in presentation["slides"]:
        if slide["type"] == "content":
            content_text += " " + " ".join(slide["content"])
        elif slide["type"] == "two_column":
            content_text += " " + " ".join(slide["left"]) + " " + " ".join(slide["right"])

    # 테마 자동 선택!
    theme_name, theme, scores = auto_select_theme(content_text, presentation["title"])

    print("=" * 60)
    print("테마 자동 선택 결과")
    print("=" * 60)
    print(f"선택된 테마: {THEME_NAMES_KR[theme_name]}")
    print(f"키워드 매칭 점수: {scores}")
    print("=" * 60)

    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 표지
    create_title_slide(prs, presentation["title"], presentation["subtitle"], theme)

    # 슬라이드 생성
    for slide_data in presentation["slides"]:
        if slide_data["type"] == "content":
            create_content_slide(prs, slide_data["title"], slide_data["content"], theme)
        elif slide_data["type"] == "architecture":
            create_architecture_slide(prs, theme)
        elif slide_data["type"] == "two_column":
            create_two_column_slide(prs, slide_data["title"],
                                   slide_data["left"], slide_data["right"], theme)

    # Q&A 슬라이드
    create_title_slide(prs, "Q & A", "질문해 주세요", theme)

    # 저장
    output_dir = "knowledge_service/docs/02_design/ui_storyboard"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "HybridRAG_자동테마.pptx")
    prs.save(output_path)

    print(f"\n프레젠테이션 생성 완료!")
    print(f"파일: {output_path}")
    print(f"슬라이드 수: {len(prs.slides)}")

if __name__ == "__main__":
    main()
