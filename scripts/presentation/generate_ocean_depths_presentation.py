"""
사내 지식 검색 시스템 UI 스토리보드 프레젠테이션 생성
Ocean Depths 테마 적용
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Ocean Depths 테마 색상 팔레트
OCEAN_DEPTHS_COLORS = {
    'primary': RGBColor(0, 102, 153),      # 진한 청록색 #006699
    'secondary': RGBColor(0, 153, 204),    # 밝은 청록색 #0099CC
    'accent': RGBColor(255, 102, 0),       # 오렌지 #FF6600
    'dark': RGBColor(0, 51, 102),          # 어두운 파란색 #003366
    'light': RGBColor(204, 229, 255),      # 연한 파란색 #CCE5FF
    'text': RGBColor(33, 33, 33),          # 진한 회색 #212121
    'bg': RGBColor(255, 255, 255)          # 흰색 #FFFFFF
}

def create_svg_architecture():
    """시스템 아키텍처 SVG 생성"""
    svg_content = """<?xml version="1.0" encoding="UTF-8"?>
<svg width="800" height="400" xmlns="http://www.w3.org/2000/svg">
  <defs>
    <style>
      .box { fill: #0099CC; stroke: #003366; stroke-width: 2; }
      .text { fill: white; font-family: 맑은 고딕; font-size: 16px; text-anchor: middle; }
      .arrow { stroke: #FF6600; stroke-width: 3; fill: none; marker-end: url(#arrowhead); }
    </style>
    <marker id="arrowhead" markerWidth="10" markerHeight="10" refX="9" refY="3" orient="auto">
      <polygon points="0 0, 10 3, 0 6" fill="#FF6600" />
    </marker>
  </defs>

  <!-- Frontend Layer -->
  <rect class="box" x="50" y="30" width="700" height="60" rx="5"/>
  <text class="text" x="400" y="65">Frontend (React 18 + TypeScript)</text>

  <!-- API Gateway -->
  <rect class="box" x="50" y="120" width="700" height="60" rx="5"/>
  <text class="text" x="400" y="155">API Gateway (Spring Cloud)</text>

  <!-- Services Layer -->
  <rect class="box" x="50" y="210" width="160" height="60" rx="5"/>
  <text class="text" x="130" y="245">Knowledge Service</text>

  <rect class="box" x="230" y="210" width="160" height="60" rx="5"/>
  <text class="text" x="310" y="245">Search Service</text>

  <rect class="box" x="410" y="210" width="160" height="60" rx="5"/>
  <text class="text" x="490" y="245">User Service</text>

  <rect class="box" x="590" y="210" width="160" height="60" rx="5"/>
  <text class="text" x="670" y="245">AI Service</text>

  <!-- Database Layer -->
  <rect class="box" x="50" y="310" width="160" height="60" rx="5"/>
  <text class="text" x="130" y="345">PostgreSQL</text>

  <rect class="box" x="230" y="310" width="160" height="60" rx="5"/>
  <text class="text" x="310" y="345">Elasticsearch</text>

  <rect class="box" x="410" y="310" width="160" height="60" rx="5"/>
  <text class="text" x="490" y="345">Neo4j</text>

  <rect class="box" x="590" y="310" width="160" height="60" rx="5"/>
  <text class="text" x="670" y="345">Redis</text>

  <!-- Arrows -->
  <path class="arrow" d="M 130 90 L 130 120"/>
  <path class="arrow" d="M 310 90 L 310 120"/>
  <path class="arrow" d="M 490 90 L 490 120"/>
  <path class="arrow" d="M 670 90 L 670 120"/>

  <path class="arrow" d="M 130 180 L 130 210"/>
  <path class="arrow" d="M 310 180 L 310 210"/>
  <path class="arrow" d="M 490 180 L 490 210"/>
  <path class="arrow" d="M 670 180 L 670 210"/>

  <path class="arrow" d="M 130 270 L 130 310"/>
  <path class="arrow" d="M 310 270 L 310 310"/>
  <path class="arrow" d="M 490 270 L 490 310"/>
  <path class="arrow" d="M 670 270 L 670 310"/>
</svg>"""
    return svg_content

def create_svg_login_dashboard():
    """로그인 & 대시보드 SVG 생성"""
    svg_content = """<?xml version="1.0" encoding="UTF-8"?>
<svg width="800" height="400" xmlns="http://www.w3.org/2000/svg">
  <defs>
    <style>
      .container { fill: white; stroke: #0099CC; stroke-width: 2; }
      .header { fill: #006699; }
      .card { fill: #CCE5FF; stroke: #0099CC; stroke-width: 2; }
      .text-title { fill: white; font-family: 맑은 고딕; font-size: 20px; font-weight: bold; text-anchor: middle; }
      .text-label { fill: #003366; font-family: 맑은 고딕; font-size: 14px; text-anchor: middle; }
      .text-value { fill: #FF6600; font-family: 맑은 고딕; font-size: 24px; font-weight: bold; text-anchor: middle; }
    </style>
  </defs>

  <!-- Main Container -->
  <rect class="container" x="50" y="30" width="700" height="340" rx="10"/>

  <!-- Header -->
  <rect class="header" x="50" y="30" width="700" height="50" rx="10"/>
  <text class="text-title" x="400" y="62">Knowledge Hub - 대시보드</text>

  <!-- Statistics Cards -->
  <rect class="card" x="80" y="110" width="150" height="100" rx="5"/>
  <text class="text-label" x="155" y="135">전체 지식</text>
  <text class="text-value" x="155" y="175">1,234</text>

  <rect class="card" x="260" y="110" width="150" height="100" rx="5"/>
  <text class="text-label" x="335" y="135">오늘 등록</text>
  <text class="text-value" x="335" y="175">12</text>

  <rect class="card" x="440" y="110" width="150" height="100" rx="5"/>
  <text class="text-label" x="515" y="135">내 지식</text>
  <text class="text-value" x="515" y="175">45</text>

  <rect class="card" x="620" y="110" width="110" height="100" rx="5"/>
  <text class="text-label" x="675" y="135">북마크</text>
  <text class="text-value" x="675" y="175">23</text>

  <!-- Activity Section -->
  <rect class="card" x="80" y="240" width="310" height="110" rx="5"/>
  <text class="text-label" x="235" y="265">🔥 인기 지식 TOP 5</text>
  <text class="text-label" x="120" y="290" text-anchor="start">1. React 컴포넌트 설계</text>
  <text class="text-label" x="120" y="310" text-anchor="start">2. Spring Boot 가이드</text>
  <text class="text-label" x="120" y="330" text-anchor="start">3. 코드 리뷰 체크리스트</text>

  <rect class="card" x="420" y="240" width="310" height="110" rx="5"/>
  <text class="text-label" x="575" y="265">📝 최신 지식</text>
  <text class="text-label" x="460" y="290" text-anchor="start">• API 설계 규칙</text>
  <text class="text-label" x="460" y="310" text-anchor="start">• Docker 컨테이너 관리</text>
  <text class="text-label" x="460" y="330" text-anchor="start">• 신규 보안 정책</text>
</svg>"""
    return svg_content

def create_svg_search_modes():
    """검색 모드 비교 SVG 생성"""
    svg_content = """<?xml version="1.0" encoding="UTF-8"?>
<svg width="800" height="400" xmlns="http://www.w3.org/2000/svg">
  <defs>
    <style>
      .panel { fill: #CCE5FF; stroke: #0099CC; stroke-width: 2; }
      .header { fill: #006699; }
      .text-header { fill: white; font-family: 맑은 고딕; font-size: 18px; font-weight: bold; text-anchor: middle; }
      .text-body { fill: #003366; font-family: 맑은 고딕; font-size: 14px; }
      .icon { fill: #FF6600; font-size: 24px; }
    </style>
  </defs>

  <!-- Chat Mode Panel -->
  <rect class="panel" x="50" y="50" width="340" height="320" rx="10"/>
  <rect class="header" x="50" y="50" width="340" height="50" rx="10"/>
  <text class="text-header" x="220" y="82">💬 채팅 모드</text>

  <text class="text-body" x="80" y="130">• RAG 기반 대화형 검색</text>
  <text class="text-body" x="80" y="160">• 자연어 질문 가능</text>
  <text class="text-body" x="80" y="190">• 컨텍스트 유지</text>
  <text class="text-body" x="80" y="220">• 실시간 스트리밍 답변</text>
  <text class="text-body" x="80" y="250">• 출처 문서 표시</text>
  <text class="text-body" x="80" y="280">• 피드백 기능 (👍/👎)</text>

  <text class="text-body" x="80" y="320" font-weight="bold">사용: 복잡한 질문</text>

  <!-- Keyword Search Panel -->
  <rect class="panel" x="410" y="50" width="340" height="320" rx="10"/>
  <rect class="header" x="410" y="50" width="340" height="50" rx="10"/>
  <text class="text-header" x="580" y="82">🔍 키워드 검색</text>

  <text class="text-body" x="440" y="130">• 전통적 검색 방식</text>
  <text class="text-body" x="440" y="160">• 필터링 (카테고리, 태그)</text>
  <text class="text-body" x="440" y="190">• 정렬 (관련도, 최신, 인기)</text>
  <text class="text-body" x="440" y="220">• 페이지네이션</text>
  <text class="text-body" x="440" y="250">• 북마크 기능</text>
  <text class="text-body" x="440" y="280">• 일괄 작업</text>

  <text class="text-body" x="440" y="320" font-weight="bold">사용: 정확한 키워드</text>
</svg>"""
    return svg_content

def save_svg_to_png(svg_content, png_path):
    """SVG를 PNG로 변환 (svglib + reportlab 사용)"""
    import io
    from svglib.svglib import svg2rlg
    from reportlab.graphics import renderPM

    try:
        # SVG 콘텐츠를 임시 파일로 저장
        svg_path = png_path.replace('.png', '.svg')
        with open(svg_path, 'w', encoding='utf-8') as f:
            f.write(svg_content)

        # svglib로 변환
        drawing = svg2rlg(svg_path)
        if drawing:
            renderPM.drawToFile(drawing, png_path, fmt='PNG', dpi=144)
            print(f"SVG converted to PNG: {png_path}")
            return True
        else:
            print(f"Failed to convert SVG. SVG saved to: {svg_path}")
            return False
    except Exception as e:
        print(f"Error converting SVG to PNG: {e}")
        svg_path = png_path.replace('.png', '.svg')
        with open(svg_path, 'w', encoding='utf-8') as f:
            f.write(svg_content)
        print(f"SVG saved to: {svg_path}")
        return False

def create_title_slide(prs, title, subtitle):
    """표지 슬라이드 - Ocean Depths 테마"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

    # 배경 그라데이션 효과 (어두운 파란색)
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['dark']
    bg_shape.line.color.rgb = OCEAN_DEPTHS_COLORS['dark']

    # 장식 요소 (파도 느낌)
    wave_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(6), Inches(5),
        Inches(5), Inches(3)
    )
    wave_shape.fill.solid()
    wave_shape.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['primary']
    wave_shape.line.color.rgb = OCEAN_DEPTHS_COLORS['primary']

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.2)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.name = '맑은 고딕'
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 부제목
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4), Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.name = '맑은 고딕'
    p.font.size = Pt(28)
    p.font.color.rgb = OCEAN_DEPTHS_COLORS['light']
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_content_slide_ocean(prs, title, content_items):
    """콘텐츠 슬라이드 - Ocean Depths 테마"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg_shape.line.color.rgb = RGBColor(255, 255, 255)

    # 헤더 바
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['primary']
    header_shape.line.color.rgb = OCEAN_DEPTHS_COLORS['primary']

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.name = '맑은 고딕'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 콘텐츠
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(5.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for idx, item in enumerate(content_items):
        if idx > 0:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]

        p.text = item
        p.font.name = '맑은 고딕'
        p.font.size = Pt(18)
        p.font.color.rgb = OCEAN_DEPTHS_COLORS['text']
        p.space_before = Pt(12)

        # 불릿 포인트 색상
        if item.startswith('•') or item.startswith('-'):
            p.font.color.rgb = OCEAN_DEPTHS_COLORS['dark']

    return slide

def create_image_slide_ocean(prs, title, image_path):
    """이미지 슬라이드 - Ocean Depths 테마"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg_shape.line.color.rgb = RGBColor(255, 255, 255)

    # 헤더 바
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['primary']
    header_shape.line.color.rgb = OCEAN_DEPTHS_COLORS['primary']

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.name = '맑은 고딕'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 이미지
    if os.path.exists(image_path):
        slide.shapes.add_picture(
            image_path,
            Inches(1), Inches(1.5),
            width=Inches(8)
        )

    return slide

def create_two_column_slide_ocean(prs, title, left_items, right_items):
    """2열 레이아웃 - Ocean Depths 테마"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg_shape.line.color.rgb = RGBColor(255, 255, 255)

    # 헤더 바
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['primary']
    header_shape.line.color.rgb = OCEAN_DEPTHS_COLORS['primary']

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.name = '맑은 고딕'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 왼쪽 패널
    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.3),
        Inches(4.4), Inches(5.7)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['light']
    left_panel.line.color.rgb = OCEAN_DEPTHS_COLORS['secondary']
    left_panel.line.width = Pt(2)

    # 오른쪽 패널
    right_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.1), Inches(1.3),
        Inches(4.4), Inches(5.7)
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['light']
    right_panel.line.color.rgb = OCEAN_DEPTHS_COLORS['secondary']
    right_panel.line.width = Pt(2)

    # 왼쪽 콘텐츠
    left_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(4), Inches(5.3)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for idx, item in enumerate(left_items):
        if idx > 0:
            p = left_frame.add_paragraph()
        else:
            p = left_frame.paragraphs[0]
        p.text = item
        p.font.name = '맑은 고딕'
        p.font.size = Pt(16)
        p.font.color.rgb = OCEAN_DEPTHS_COLORS['text']
        p.space_before = Pt(8)

    # 오른쪽 콘텐츠
    right_box = slide.shapes.add_textbox(
        Inches(5.3), Inches(1.5), Inches(4), Inches(5.3)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for idx, item in enumerate(right_items):
        if idx > 0:
            p = right_frame.add_paragraph()
        else:
            p = right_frame.paragraphs[0]
        p.text = item
        p.font.name = '맑은 고딕'
        p.font.size = Pt(16)
        p.font.color.rgb = OCEAN_DEPTHS_COLORS['text']
        p.space_before = Pt(8)

    return slide

# 프레젠테이션 생성
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# SVG 이미지 경로 설정
svg_dir = r"D:\Users\KTDS\Documents\06.과제\hybrid-rag-knowledge-ops\knowledge_service\docs\02_design\ui_storyboard"
arch_png = os.path.join(svg_dir, "architecture_diagram.png")
dashboard_png = os.path.join(svg_dir, "dashboard_diagram.png")
search_png = os.path.join(svg_dir, "search_modes_diagram.png")

# SVG 생성 및 PNG 변환
print("SVG 다이어그램 생성 중...")
save_svg_to_png(create_svg_architecture(), arch_png)
save_svg_to_png(create_svg_login_dashboard(), dashboard_png)
save_svg_to_png(create_svg_search_modes(), search_png)

# 1. 표지
create_title_slide(
    prs,
    "사내 지식 검색 시스템",
    "UI/UX 스토리보드\nHybrid RAG Knowledge Platform"
)

# 2. 프로젝트 개요
create_content_slide_ocean(
    prs,
    "프로젝트 개요",
    [
        "목표: 사내 지식을 통합 검색하고 지능형 답변 제공",
        "",
        "핵심 가치:",
        "  • Graph RAG 기반 지능형 검색",
        "  • 자연어 대화형 인터페이스",
        "  • 체계적인 지식 관리 (등록, 수정, 버전 관리)",
        "  • 개인화 (북마크, 검색 기록, AI 추천)",
        "",
        "사용자: 전 직원 (500명)",
        "기대 효과: 검색 시간 70% 단축, 정확도 95%"
    ]
)

# 3. 시스템 아키텍처 (SVG)
if os.path.exists(arch_png):
    create_image_slide_ocean(prs, "시스템 아키텍처", arch_png)
else:
    # SVG를 변환할 수 없으면 텍스트로 대체
    create_content_slide_ocean(
        prs,
        "시스템 아키텍처",
        [
            "Frontend: React 18 + TypeScript + Vite",
            "  ↓",
            "API Gateway: Spring Cloud Gateway",
            "  ↓",
            "Services:",
            "  • Knowledge Service  • Search Service",
            "  • User Service       • AI Service (Python)",
            "  ↓",
            "Database:",
            "  • PostgreSQL  • Elasticsearch  • Neo4j  • Redis"
        ]
    )

# 4. 로그인 & 대시보드 (SVG)
if os.path.exists(dashboard_png):
    create_image_slide_ocean(prs, "로그인 & 메인 대시보드", dashboard_png)
else:
    create_content_slide_ocean(
        prs,
        "로그인 & 메인 대시보드",
        [
            "로그인",
            "  • OAuth 2.0 + PKCE 인증",
            "  • SSO 연동 (Keycloak)",
            "  • 안전한 회사 계정 로그인",
            "",
            "대시보드 위젯",
            "  • 통계 카드 (전체 지식, 오늘 등록, 내 지식, 북마크)",
            "  • 인기 지식 TOP 5",
            "  • 최신 지식 목록",
            "  • AI 추천 지식",
            "  • 검색 트렌드 (워드 클라우드)"
        ]
    )

# 5. 검색 화면 (SVG)
if os.path.exists(search_png):
    create_image_slide_ocean(prs, "검색 화면 - 2가지 모드", search_png)
else:
    create_two_column_slide_ocean(
        prs,
        "검색 화면 - 2가지 모드",
        [
            "💬 채팅 모드",
            "━━━━━━━━━━━",
            "• RAG 기반 대화형 검색",
            "• 자연어 질문",
            "• 컨텍스트 유지",
            "• 실시간 스트리밍 답변",
            "• 출처 문서 표시",
            "• 피드백 (👍/👎)",
            "",
            "사용: 복잡한 질문, 맥락 이해"
        ],
        [
            "🔍 키워드 검색",
            "━━━━━━━━━━━",
            "• 전통적 검색 방식",
            "• 필터링 (카테고리, 태그)",
            "• 정렬 (관련도, 최신, 인기)",
            "• 페이지네이션",
            "• 북마크 기능",
            "• 일괄 작업",
            "",
            "사용: 정확한 키워드 검색"
        ]
    )

# 6. 지식 관리
create_content_slide_ocean(
    prs,
    "지식 관리",
    [
        "지식 목록",
        "  • 전체 / 내 지식 / 북마크 / 임시저장 탭",
        "  • 검색, 필터, 정렬 기능",
        "  • 체크박스로 일괄 작업",
        "",
        "지식 상세",
        "  • 마크다운 완전 지원 (코드, 다이어그램, 수식)",
        "  • 댓글 및 답글, 버전 이력 관리",
        "",
        "지식 작성/수정",
        "  • 마크다운 에디터 (실시간 미리보기)",
        "  • 자동 태그 추천, 임시저장 (10초마다)"
    ]
)

# 7. 프로필 & 관리자
create_two_column_slide_ocean(
    prs,
    "프로필 & 관리자",
    [
        "👤 프로필",
        "━━━━━━━━",
        "• 활동 통계",
        "  - 작성 지식: 45개",
        "  - 조회수: 8,234",
        "  - 받은 댓글: 156",
        "  - 받은 북마크: 234",
        "",
        "• 최근 작성 지식",
        "• 활동 그래프 (30일)",
        "",
        "⭐ 북마크",
        "• 폴더별 관리"
    ],
    [
        "⚙️ 관리자",
        "━━━━━━━━",
        "• 시스템 대시보드",
        "  - 사용자/트래픽 통계",
        "  - 시스템 알림",
        "",
        "• 사용자 관리",
        "  - 역할 (USER/MANAGER/ADMIN)",
        "  - 계정 상태 관리",
        "",
        "• 지식 승인/반려",
        "• 카테고리 관리",
        "• 시스템 설정"
    ]
)

# 8. 기술 스택
create_content_slide_ocean(
    prs,
    "기술 스택",
    [
        "Frontend: React 18 + TypeScript + Vite",
        "  • Material-UI v5, Redux Toolkit, React Query",
        "",
        "Backend: Spring Boot 3.x + Spring Cloud Gateway",
        "",
        "AI Service: Python FastAPI",
        "  • LangGraph, DeepSeek V3.2, BGE-M3",
        "",
        "Database",
        "  • PostgreSQL (마스터)  • Elasticsearch (검색)",
        "  • Neo4j (그래프)       • Redis (캐시)"
    ]
)

# 9. 기대 효과
create_content_slide_ocean(
    prs,
    "주요 기능 및 기대 효과",
    [
        "주요 기능",
        "  • Graph RAG 기반 지능형 검색",
        "  • 대화형 인터페이스 (채팅 모드)",
        "  • 메타데이터 자동 추출 (DeepSeek V3.2)",
        "  • 버전 관리 및 협업 기능",
        "",
        "기대 효과",
        "  • 검색 시간 70% 단축 (15분 → 2분)",
        "  • 검색 정확도 95%",
        "  • 지식 재사용 증가, 중복 작업 감소",
        "  • 비용 절감: DeepSeek 활용 95% 비용 절감"
    ]
)

# 10. Q&A
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 배경 그라데이션
bg_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(10), Inches(7.5)
)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['dark']
bg_shape.line.color.rgb = OCEAN_DEPTHS_COLORS['dark']

# 장식
wave_shape = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(-2), Inches(5),
    Inches(6), Inches(4)
)
wave_shape.fill.solid()
wave_shape.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['primary']
wave_shape.line.color.rgb = OCEAN_DEPTHS_COLORS['primary']

title_box = slide.shapes.add_textbox(
    Inches(2), Inches(2.5), Inches(6), Inches(1)
)
title_frame = title_box.text_frame
title_frame.text = "Q & A"
p = title_frame.paragraphs[0]
p.font.name = '맑은 고딕'
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(
    Inches(2), Inches(4), Inches(6), Inches(0.8)
)
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "감사합니다"
p = subtitle_frame.paragraphs[0]
p.font.name = '맑은 고딕'
p.font.size = Pt(36)
p.font.color.rgb = OCEAN_DEPTHS_COLORS['light']
p.alignment = PP_ALIGN.CENTER

date_box = slide.shapes.add_textbox(
    Inches(2), Inches(6), Inches(6), Inches(0.5)
)
date_frame = date_box.text_frame
date_frame.text = "Hybrid RAG Knowledge Platform | 2026-01-15"
p = date_frame.paragraphs[0]
p.font.name = '맑은 고딕'
p.font.size = Pt(16)
p.font.color.rgb = OCEAN_DEPTHS_COLORS['light']
p.alignment = PP_ALIGN.CENTER

# 저장
output_path = os.path.join(svg_dir, "UI스토리보드_발표자료_OceanDepths.pptx")
prs.save(output_path)
print(f"\n프레젠테이션 생성 완료!")
print(f"파일 위치: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
print(f"\nSVG 다이어그램:")
print(f"  - architecture_diagram.svg/png")
print(f"  - dashboard_diagram.svg/png")
print(f"  - search_modes_diagram.svg/png")
