#!/usr/bin/env python3
"""
인증/권한 관리 상세 설계서 프레젠테이션 생성
Brown Earth Theme 적용
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os

# ========== Brown Earth Theme ==========
BROWN_EARTH = {
    'primary': RGBColor(121, 85, 72),      # #795548 Brown
    'secondary': RGBColor(161, 136, 127),  # #A1887F Light Brown
    'accent': RGBColor(255, 152, 0),       # #FF9800 Orange
    'dark': RGBColor(62, 39, 35),          # #3E2723 Dark Brown
    'light': RGBColor(239, 235, 233),      # #EFEBE9 Light Beige
    'text': RGBColor(33, 33, 33),          # #212121 Dark Text
    'white': RGBColor(255, 255, 255),
    # Component colors
    'keycloak': RGBColor(0, 150, 136),     # Teal for Keycloak
    'jwt': RGBColor(255, 87, 34),          # Deep Orange for JWT
    'spring': RGBColor(76, 175, 80),       # Green for Spring
    'redis': RGBColor(220, 53, 69),        # Red for Redis
    'frontend': RGBColor(33, 150, 243),    # Blue for Frontend
}

def create_title_slide(prs, title, subtitle, theme):
    """표지 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme['dark']
    bg.line.fill.background()

    # 장식 원
    deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6), Inches(4.5), Inches(5), Inches(4))
    deco.fill.solid()
    deco.fill.fore_color.rgb = theme['primary']
    deco.line.fill.background()

    # 상단 액센트 바
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = theme['accent']
    accent_bar.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = theme['white']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.LEFT

    # 부제목
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(7), Inches(0.8))
    sf = sub_box.text_frame
    sf.word_wrap = True
    sp = sf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(22)
    sp.font.color.rgb = theme['light']
    sp.font.name = "맑은 고딕"
    sp.alignment = PP_ALIGN.LEFT

    # 날짜
    date_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(4), Inches(0.5))
    df = date_box.text_frame
    dp = df.paragraphs[0]
    dp.text = "2026-01-15"
    dp.font.size = Pt(14)
    dp.font.color.rgb = theme['secondary']
    dp.font.name = "맑은 고딕"

    return slide

def create_section_header(prs, title, theme):
    """섹션 헤더 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme['primary']
    bg.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = theme['white']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_slide_header(slide, title, theme):
    """슬라이드 헤더 추가"""
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = theme['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme['white']
    p.font.name = "맑은 고딕"

def add_shape_with_text(slide, shape_type, left, top, width, height, fill_color, text, font_size=12, text_color=None, bold=False):
    """텍스트가 포함된 도형 추가"""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(100, 100, 100)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = text_color or RGBColor(255, 255, 255)
    tf.paragraphs[0].font.name = "맑은 고딕"
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.anchor = MSO_ANCHOR.MIDDLE

    return shape

def add_arrow(slide, start_x, start_y, end_x, end_y, color):
    """화살표 추가"""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(start_x), Inches(start_y),
        Inches(end_x), Inches(end_y)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector

def create_overview_slide(prs, theme):
    """아키텍처 개요 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "인증/권한 관리 아키텍처 개요", theme)

    # 설명 텍스트
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "OAuth 2.0 + PKCE 기반 보안 인증 시스템 | Keycloak IdP 통합"
    p.font.size = Pt(16)
    p.font.color.rgb = theme['text']
    p.font.name = "맑은 고딕"

    # 4개 핵심 컴포넌트 박스
    components = [
        ("Keycloak IdP", "OAuth 2.0 Provider\n사용자 인증/관리\nSSO 지원", theme['keycloak']),
        ("JWT Token", "Access Token (15분)\nRefresh Token (7일)\n토큰 회전", theme['jwt']),
        ("Spring Security", "JWT 검증 필터\nRole 기반 인가\nCSRF/XSS 방어", theme['spring']),
        ("Redis Blacklist", "토큰 무효화\n세션 관리\n실시간 동기화", theme['redis']),
    ]

    start_x = 0.5
    for i, (title, desc, color) in enumerate(components):
        x = start_x + i * 2.4

        # 컴포넌트 박스
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(2.2), Inches(2.8))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # 제목
        title_box = slide.shapes.add_textbox(Inches(x), Inches(2.3), Inches(2.2), Inches(0.5))
        ttf = title_box.text_frame
        tp = ttf.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(14)
        tp.font.bold = True
        tp.font.color.rgb = theme['white']
        tp.font.name = "맑은 고딕"
        tp.alignment = PP_ALIGN.CENTER

        # 설명
        desc_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.9), Inches(2), Inches(2))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        for line in desc.split('\n'):
            dp = dtf.add_paragraph()
            dp.text = "• " + line
            dp.font.size = Pt(11)
            dp.font.color.rgb = theme['white']
            dp.font.name = "맑은 고딕"

    # 하단 요약 박스
    summary_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.3), Inches(9), Inches(1.8))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = theme['light']
    summary_box.line.color.rgb = theme['secondary']

    summary_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(8.6), Inches(1.5))
    stf = summary_text.text_frame
    stf.word_wrap = True

    summaries = [
        "✓ Keycloak: 오픈소스 IdP로 비용 효율적인 인증 인프라 구축",
        "✓ PKCE: Authorization Code 탈취 공격 방어 (Native/SPA 앱 보안)",
        "✓ Token Rotation: Refresh Token 사용 시마다 새 토큰 발급으로 보안 강화",
        "✓ Redis Blacklist: 즉각적인 토큰 무효화로 보안 사고 대응 가능"
    ]

    for summary in summaries:
        sp = stf.add_paragraph()
        sp.text = summary
        sp.font.size = Pt(12)
        sp.font.color.rgb = theme['text']
        sp.font.name = "맑은 고딕"
        sp.space_after = Pt(4)

    return slide

def create_oauth_flow_slide(prs, theme):
    """OAuth 2.0 + PKCE 인증 흐름 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "OAuth 2.0 Authorization Code + PKCE Flow", theme)

    # 설명
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(9), Inches(0.4))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SPA/Native 앱을 위한 보안 강화 인증 흐름 | Code Verifier + Code Challenge"
    p.font.size = Pt(13)
    p.font.color.rgb = theme['secondary']
    p.font.name = "맑은 고딕"

    # 액터들
    actors = [
        ("User\n(Browser)", 0.3, theme['frontend']),
        ("Frontend\n(SPA)", 2.4, theme['frontend']),
        ("API Gateway", 4.5, theme['spring']),
        ("Keycloak\n(IdP)", 6.6, theme['keycloak']),
        ("Backend\nServices", 8.7, theme['primary']),
    ]

    for name, x, color in actors:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.7), Inches(1.3), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        text_box = slide.shapes.add_textbox(Inches(x), Inches(1.75), Inches(1.3), Inches(0.8))
        ttf = text_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = name
        tp.font.size = Pt(10)
        tp.font.bold = True
        tp.font.color.rgb = theme['white']
        tp.font.name = "맑은 고딕"
        tp.alignment = PP_ALIGN.CENTER

    # 흐름 단계
    steps = [
        ("1", "로그인 버튼 클릭", 0.9, 3.0, 1.5, theme['text']),
        ("2", "Generate code_verifier\n+ code_challenge", 2.4, 3.0, 1.4, theme['frontend']),
        ("3", "Authorization Request\n(+code_challenge)", 3.3, 3.6, 3.0, theme['accent']),
        ("4", "사용자 인증\n(ID/PW 입력)", 6.6, 3.6, 1.3, theme['keycloak']),
        ("5", "Authorization Code\n반환", 4.3, 4.3, 2.0, theme['keycloak']),
        ("6", "Token Request\n(code + code_verifier)", 3.3, 5.0, 3.0, theme['accent']),
        ("7", "Verify & Issue\nAccess + Refresh Token", 6.6, 5.0, 1.4, theme['keycloak']),
        ("8", "API 요청\n(Bearer Token)", 4.5, 5.7, 1.5, theme['spring']),
        ("9", "Protected\nResource", 8.5, 5.7, 1.2, theme['primary']),
    ]

    for step_num, text, x, y, width, color in steps:
        # 단계 번호 원
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.15), Inches(y - 0.15), Inches(0.3), Inches(0.3))
        circle.fill.solid()
        circle.fill.fore_color.rgb = theme['accent']
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(x - 0.15), Inches(y - 0.12), Inches(0.3), Inches(0.3))
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        np.text = step_num
        np.font.size = Pt(9)
        np.font.bold = True
        np.font.color.rgb = theme['white']
        np.font.name = "맑은 고딕"
        np.alignment = PP_ALIGN.CENTER

        # 설명 텍스트
        desc = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y - 0.1), Inches(width), Inches(0.6))
        dtf = desc.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = text
        dp.font.size = Pt(9)
        dp.font.color.rgb = color
        dp.font.name = "맑은 고딕"

    # PKCE 설명 박스
    pkce_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(6.5), Inches(9.4), Inches(0.8))
    pkce_box.fill.solid()
    pkce_box.fill.fore_color.rgb = theme['light']
    pkce_box.line.color.rgb = theme['accent']
    pkce_box.line.width = Pt(2)

    pkce_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.55), Inches(9), Inches(0.7))
    ptf = pkce_text.text_frame
    ptf.word_wrap = True
    pp = ptf.paragraphs[0]
    pp.text = "🔐 PKCE (Proof Key for Code Exchange): code_verifier는 클라이언트에만 저장, code_challenge만 서버로 전송"
    pp.font.size = Pt(12)
    pp.font.color.rgb = theme['text']
    pp.font.name = "맑은 고딕"
    pp2 = ptf.add_paragraph()
    pp2.text = "→ Authorization Code가 탈취되어도 code_verifier 없이는 토큰 교환 불가 (Man-in-the-Middle 공격 방어)"
    pp2.font.size = Pt(11)
    pp2.font.color.rgb = theme['primary']
    pp2.font.name = "맑은 고딕"

    return slide

def create_jwt_strategy_slide(prs, theme):
    """JWT 토큰 전략 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "JWT Token 전략", theme)

    # Access Token vs Refresh Token 비교
    # Access Token 박스
    at_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3), Inches(4.5), Inches(2.8))
    at_box.fill.solid()
    at_box.fill.fore_color.rgb = theme['jwt']
    at_box.line.fill.background()

    at_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.3), Inches(0.5))
    attf = at_title.text_frame
    atp = attf.paragraphs[0]
    atp.text = "🔑 Access Token"
    atp.font.size = Pt(18)
    atp.font.bold = True
    atp.font.color.rgb = theme['white']
    atp.font.name = "맑은 고딕"

    at_content = slide.shapes.add_textbox(Inches(0.6), Inches(1.95), Inches(4.2), Inches(2))
    actf = at_content.text_frame
    actf.word_wrap = True
    at_items = [
        "만료 시간: 15분 (단기)",
        "저장 위치: 메모리 (XSS 방어)",
        "용도: API 요청 인증",
        "Payload: user_id, roles, permissions",
        "갱신: Refresh Token으로 재발급"
    ]
    for item in at_items:
        p = actf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(12)
        p.font.color.rgb = theme['white']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(6)

    # Refresh Token 박스
    rt_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.3), Inches(4.5), Inches(2.8))
    rt_box.fill.solid()
    rt_box.fill.fore_color.rgb = theme['keycloak']
    rt_box.line.fill.background()

    rt_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.5))
    rttf = rt_title.text_frame
    rtp = rttf.paragraphs[0]
    rtp.text = "🔄 Refresh Token"
    rtp.font.size = Pt(18)
    rtp.font.bold = True
    rtp.font.color.rgb = theme['white']
    rtp.font.name = "맑은 고딕"

    rt_content = slide.shapes.add_textbox(Inches(5.3), Inches(1.95), Inches(4.2), Inches(2))
    rctf = rt_content.text_frame
    rctf.word_wrap = True
    rt_items = [
        "만료 시간: 7일 (장기)",
        "저장 위치: HttpOnly Cookie (CSRF 방어)",
        "용도: Access Token 갱신",
        "보안: Secure, SameSite=Strict",
        "정책: Token Rotation (사용시 재발급)"
    ]
    for item in rt_items:
        p = rctf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(12)
        p.font.color.rgb = theme['white']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(6)

    # Token Rotation 다이어그램
    rotation_title = slide.shapes.add_textbox(Inches(0.4), Inches(4.3), Inches(9), Inches(0.4))
    rottf = rotation_title.text_frame
    rotp = rottf.paragraphs[0]
    rotp.text = "Token Rotation 전략"
    rotp.font.size = Pt(16)
    rotp.font.bold = True
    rotp.font.color.rgb = theme['text']
    rotp.font.name = "맑은 고딕"

    # 흐름도
    steps = [
        ("RT 사용", 0.5, 5.0, theme['keycloak']),
        ("검증", 2.3, 5.0, theme['spring']),
        ("새 AT+RT\n발급", 4.1, 5.0, theme['jwt']),
        ("기존 RT\n무효화", 5.9, 5.0, theme['redis']),
        ("Redis\n저장", 7.7, 5.0, theme['redis']),
    ]

    for i, (text, x, y, color) in enumerate(steps):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.5), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        text_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.15), Inches(1.5), Inches(0.7))
        ttf = text_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = text
        tp.font.size = Pt(11)
        tp.font.bold = True
        tp.font.color.rgb = theme['white']
        tp.font.name = "맑은 고딕"
        tp.alignment = PP_ALIGN.CENTER

        # 화살표 (마지막 제외)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.55), Inches(y + 0.35), Inches(0.6), Inches(0.25))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = theme['accent']
            arrow.line.fill.background()

    # 보안 이점 설명
    benefit_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(6.1), Inches(9.2), Inches(1.2))
    benefit_box.fill.solid()
    benefit_box.fill.fore_color.rgb = theme['light']
    benefit_box.line.color.rgb = theme['primary']

    benefit_text = slide.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(8.8), Inches(1))
    btf = benefit_text.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.text = "🛡️ Token Rotation 보안 이점"
    bp.font.size = Pt(13)
    bp.font.bold = True
    bp.font.color.rgb = theme['primary']
    bp.font.name = "맑은 고딕"

    bp2 = btf.add_paragraph()
    bp2.text = "• RT 탈취 시 단 1회만 사용 가능 → 탈취 탐지 후 즉시 무효화 | • Replay Attack 방어 | • 세션 하이재킹 위험 최소화"
    bp2.font.size = Pt(11)
    bp2.font.color.rgb = theme['text']
    bp2.font.name = "맑은 고딕"

    return slide

def create_rbac_slide(prs, theme):
    """RBAC 권한 관리 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "RBAC (Role-Based Access Control)", theme)

    # 역할 계층 구조
    # ADMIN (최상위)
    admin_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.75), Inches(1.4), Inches(2.5), Inches(0.8))
    admin_box.fill.solid()
    admin_box.fill.fore_color.rgb = theme['redis']
    admin_box.line.fill.background()

    admin_text = slide.shapes.add_textbox(Inches(3.75), Inches(1.5), Inches(2.5), Inches(0.6))
    atf = admin_text.text_frame
    ap = atf.paragraphs[0]
    ap.text = "🔴 ADMIN"
    ap.font.size = Pt(16)
    ap.font.bold = True
    ap.font.color.rgb = theme['white']
    ap.font.name = "맑은 고딕"
    ap.alignment = PP_ALIGN.CENTER

    # KNOWLEDGE_MANAGER (중간)
    km_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.25), Inches(2.6), Inches(3.5), Inches(0.8))
    km_box.fill.solid()
    km_box.fill.fore_color.rgb = theme['accent']
    km_box.line.fill.background()

    km_text = slide.shapes.add_textbox(Inches(3.25), Inches(2.7), Inches(3.5), Inches(0.6))
    kmtf = km_text.text_frame
    kmp = kmtf.paragraphs[0]
    kmp.text = "🟠 KNOWLEDGE_MANAGER"
    kmp.font.size = Pt(16)
    kmp.font.bold = True
    kmp.font.color.rgb = theme['white']
    kmp.font.name = "맑은 고딕"
    kmp.alignment = PP_ALIGN.CENTER

    # USER (기본)
    user_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.75), Inches(3.8), Inches(4.5), Inches(0.8))
    user_box.fill.solid()
    user_box.fill.fore_color.rgb = theme['spring']
    user_box.line.fill.background()

    user_text = slide.shapes.add_textbox(Inches(2.75), Inches(3.9), Inches(4.5), Inches(0.6))
    utf = user_text.text_frame
    up = utf.paragraphs[0]
    up.text = "🟢 USER (기본 역할)"
    up.font.size = Pt(16)
    up.font.bold = True
    up.font.color.rgb = theme['white']
    up.font.name = "맑은 고딕"
    up.alignment = PP_ALIGN.CENTER

    # 연결선 (위에서 아래로)
    line1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.85), Inches(2.2), Inches(0.3), Inches(0.35))
    line1.fill.solid()
    line1.fill.fore_color.rgb = theme['secondary']
    line1.line.fill.background()

    line2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.85), Inches(3.4), Inches(0.3), Inches(0.35))
    line2.fill.solid()
    line2.fill.fore_color.rgb = theme['secondary']
    line2.line.fill.background()

    # 권한 테이블
    table_title = slide.shapes.add_textbox(Inches(0.4), Inches(4.9), Inches(9), Inches(0.4))
    tttf = table_title.text_frame
    ttp = tttf.paragraphs[0]
    ttp.text = "권한 매트릭스"
    ttp.font.size = Pt(14)
    ttp.font.bold = True
    ttp.font.color.rgb = theme['text']
    ttp.font.name = "맑은 고딕"

    # 테이블 헤더
    headers = ["권한", "USER", "KNOWLEDGE_MANAGER", "ADMIN"]
    header_widths = [2.8, 1.5, 2.4, 1.5]
    x_pos = 0.6
    for i, (header, width) in enumerate(zip(headers, header_widths)):
        hbox = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(5.3), Inches(width), Inches(0.4))
        hbox.fill.solid()
        hbox.fill.fore_color.rgb = theme['primary']
        hbox.line.color.rgb = theme['dark']

        htext = slide.shapes.add_textbox(Inches(x_pos), Inches(5.35), Inches(width), Inches(0.35))
        htf = htext.text_frame
        hp = htf.paragraphs[0]
        hp.text = header
        hp.font.size = Pt(11)
        hp.font.bold = True
        hp.font.color.rgb = theme['white']
        hp.font.name = "맑은 고딕"
        hp.alignment = PP_ALIGN.CENTER

        x_pos += width

    # 테이블 데이터
    rows = [
        ("지식 검색 (SEARCH)", "✅", "✅", "✅"),
        ("지식 등록/수정 (WRITE)", "❌", "✅", "✅"),
        ("시스템 설정 (ADMIN)", "❌", "❌", "✅"),
    ]

    y_pos = 5.7
    for perm, user_perm, km_perm, admin_perm in rows:
        x_pos = 0.6
        values = [perm, user_perm, km_perm, admin_perm]
        for i, (value, width) in enumerate(zip(values, header_widths)):
            rbox = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(y_pos), Inches(width), Inches(0.35))
            rbox.fill.solid()
            rbox.fill.fore_color.rgb = theme['light'] if i == 0 else theme['white']
            rbox.line.color.rgb = theme['secondary']

            rtext = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 0.05), Inches(width), Inches(0.3))
            rtf = rtext.text_frame
            rp = rtf.paragraphs[0]
            rp.text = value
            rp.font.size = Pt(10)
            rp.font.color.rgb = theme['text']
            rp.font.name = "맑은 고딕"
            rp.alignment = PP_ALIGN.CENTER

            x_pos += width
        y_pos += 0.35

    # 설명
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.4))
    ntf = note_box.text_frame
    np = ntf.paragraphs[0]
    np.text = "💡 상위 역할은 하위 역할의 모든 권한을 자동 상속 (계층적 RBAC)"
    np.font.size = Pt(11)
    np.font.color.rgb = theme['primary']
    np.font.name = "맑은 고딕"

    return slide

def create_spring_security_slide(prs, theme):
    """Spring Security 구성 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Spring Security + API Gateway 구성", theme)

    # 왼쪽: API Gateway
    gw_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.4), Inches(4.5), Inches(3.2))
    gw_box.fill.solid()
    gw_box.fill.fore_color.rgb = theme['spring']
    gw_box.line.fill.background()

    gw_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(0.5))
    gwtf = gw_title.text_frame
    gwp = gwtf.paragraphs[0]
    gwp.text = "🌐 API Gateway (Spring Cloud)"
    gwp.font.size = Pt(16)
    gwp.font.bold = True
    gwp.font.color.rgb = theme['white']
    gwp.font.name = "맑은 고딕"

    gw_content = slide.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(4.2), Inches(2.4))
    gwctf = gw_content.text_frame
    gwctf.word_wrap = True

    gw_items = [
        "JwtAuthenticationFilter",
        "  → Bearer 토큰 추출 및 검증",
        "  → SecurityContext 설정",
        "RateLimitFilter",
        "  → 요청 속도 제한 (분당 100회)",
        "  → IP/User별 제한",
        "CORS Configuration",
        "  → 허용 도메인 설정"
    ]

    for item in gw_items:
        p = gwctf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = theme['white']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(3)

    # 오른쪽: Backend Services
    be_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.4), Inches(4.5), Inches(3.2))
    be_box.fill.solid()
    be_box.fill.fore_color.rgb = theme['primary']
    be_box.line.fill.background()

    be_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(0.5))
    betf = be_title.text_frame
    bep = betf.paragraphs[0]
    bep.text = "🔐 Backend Services (Spring Security)"
    bep.font.size = Pt(16)
    bep.font.bold = True
    bep.font.color.rgb = theme['white']
    bep.font.name = "맑은 고딕"

    be_content = slide.shapes.add_textbox(Inches(5.3), Inches(2.1), Inches(4.2), Inches(2.4))
    bectf = be_content.text_frame
    bectf.word_wrap = True

    be_items = [
        "@PreAuthorize 어노테이션",
        "  → 메서드 레벨 권한 체크",
        "  → SpEL 표현식 지원",
        "SecurityFilterChain",
        "  → Stateless 세션 정책",
        "  → CSRF 비활성화 (JWT)",
        "ExceptionHandler",
        "  → 인증/인가 에러 처리"
    ]

    for item in be_items:
        p = bectf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = theme['white']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(3)

    # 화살표 (Gateway → Backend)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(2.8), Inches(0.6), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = theme['accent']
    arrow.line.fill.background()

    # 하단: 보안 계층 다이어그램
    layer_title = slide.shapes.add_textbox(Inches(0.4), Inches(4.8), Inches(9), Inches(0.4))
    ltf = layer_title.text_frame
    lp = ltf.paragraphs[0]
    lp.text = "보안 필터 체인 순서"
    lp.font.size = Pt(14)
    lp.font.bold = True
    lp.font.color.rgb = theme['text']
    lp.font.name = "맑은 고딕"

    layers = [
        ("CORS", theme['secondary']),
        ("Rate\nLimit", theme['accent']),
        ("JWT\nAuth", theme['jwt']),
        ("Role\nCheck", theme['spring']),
        ("Method\nSecurity", theme['primary']),
        ("Business\nLogic", theme['keycloak']),
    ]

    x = 0.5
    for i, (name, color) in enumerate(layers):
        lbox = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(5.3), Inches(1.5), Inches(0.9))
        lbox.fill.solid()
        lbox.fill.fore_color.rgb = color
        lbox.line.fill.background()

        ltext = slide.shapes.add_textbox(Inches(x + 0.2), Inches(5.4), Inches(1.1), Inches(0.7))
        lttf = ltext.text_frame
        lttf.word_wrap = True
        ltp = lttf.paragraphs[0]
        ltp.text = name
        ltp.font.size = Pt(10)
        ltp.font.bold = True
        ltp.font.color.rgb = theme['white']
        ltp.font.name = "맑은 고딕"
        ltp.alignment = PP_ALIGN.CENTER

        x += 1.55

    # 설명
    note = slide.shapes.add_textbox(Inches(0.4), Inches(6.4), Inches(9.2), Inches(0.9))
    ntf = note.text_frame
    ntf.word_wrap = True
    np = ntf.paragraphs[0]
    np.text = "✓ 요청은 순차적으로 모든 필터를 통과 | 각 단계에서 검증 실패 시 즉시 거부"
    np.font.size = Pt(11)
    np.font.color.rgb = theme['text']
    np.font.name = "맑은 고딕"

    return slide

def create_architecture_advantages_slide(prs, theme):
    """아키텍처 장점 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "인증/권한 아키텍처 장점", theme)

    advantages = [
        {
            "icon": "🏠",
            "title": "Keycloak (오픈소스 IdP)",
            "items": [
                "상용 IdP 대비 95% 비용 절감",
                "SSO, LDAP, Social 로그인 기본 지원",
                "대시보드로 사용자/세션 통합 관리",
                "SAML, OAuth 2.0, OIDC 표준 준수"
            ],
            "color": theme['keycloak']
        },
        {
            "icon": "🛡️",
            "title": "PKCE 인증 보안",
            "items": [
                "Authorization Code 탈취 공격 방어",
                "Native/SPA 앱 필수 보안 요구사항 충족",
                "code_verifier 클라이언트 로컬 저장",
                "MITM (Man-in-the-Middle) 공격 차단"
            ],
            "color": theme['jwt']
        },
        {
            "icon": "🔄",
            "title": "Token Rotation 전략",
            "items": [
                "Refresh Token 재사용 탐지",
                "탈취 시 1회만 사용 가능",
                "자동 토큰 갱신으로 UX 향상",
                "세션 하이재킹 위험 최소화"
            ],
            "color": theme['spring']
        },
        {
            "icon": "⚡",
            "title": "Redis 기반 Blacklist",
            "items": [
                "즉각적인 토큰 무효화 (<1ms)",
                "보안 사고 시 실시간 대응",
                "분산 환경 동기화 지원",
                "TTL 기반 자동 정리"
            ],
            "color": theme['redis']
        }
    ]

    positions = [(0.3, 1.3), (5.0, 1.3), (0.3, 4.4), (5.0, 4.4)]

    for (x, y), adv in zip(positions, advantages):
        # 카드 배경
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.7), Inches(2.9))
        card.fill.solid()
        card.fill.fore_color.rgb = adv['color']
        card.line.fill.background()

        # 아이콘과 제목
        title_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.15), Inches(4.4), Inches(0.5))
        ttf = title_box.text_frame
        tp = ttf.paragraphs[0]
        tp.text = f"{adv['icon']} {adv['title']}"
        tp.font.size = Pt(15)
        tp.font.bold = True
        tp.font.color.rgb = theme['white']
        tp.font.name = "맑은 고딕"

        # 항목들
        content_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.7), Inches(4.3), Inches(2.1))
        ctf = content_box.text_frame
        ctf.word_wrap = True

        for item in adv['items']:
            p = ctf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(11)
            p.font.color.rgb = theme['white']
            p.font.name = "맑은 고딕"
            p.space_after = Pt(6)

    return slide

def create_security_highlights_slide(prs, theme):
    """보안 하이라이트 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "보안 체크리스트 하이라이트", theme)

    # 3열 구조
    categories = [
        {
            "title": "🔐 인증 보안",
            "items": [
                "PKCE 필수 적용",
                "HTTPS Only (TLS 1.3)",
                "Secure Cookie 설정",
                "Password 해싱 (bcrypt)",
                "브루트포스 방어"
            ],
            "color": theme['keycloak']
        },
        {
            "title": "🎫 토큰 보안",
            "items": [
                "단기 Access Token (15분)",
                "HttpOnly Refresh Token",
                "Token Rotation 활성화",
                "Redis Blacklist 사용",
                "JWT 서명 검증 (RS256)"
            ],
            "color": theme['jwt']
        },
        {
            "title": "🛡️ 인프라 보안",
            "items": [
                "Rate Limiting (100/분)",
                "CORS 화이트리스트",
                "CSRF 방어 (SameSite)",
                "XSS 방어 (CSP)",
                "SQL Injection 방어"
            ],
            "color": theme['spring']
        }
    ]

    x = 0.4
    for cat in categories:
        # 카드
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(3.0), Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = cat['color']
        card.line.fill.background()

        # 제목
        title_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(1.4), Inches(2.8), Inches(0.5))
        ttf = title_box.text_frame
        tp = ttf.paragraphs[0]
        tp.text = cat['title']
        tp.font.size = Pt(16)
        tp.font.bold = True
        tp.font.color.rgb = theme['white']
        tp.font.name = "맑은 고딕"
        tp.alignment = PP_ALIGN.CENTER

        # 항목
        content_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.0), Inches(2.7), Inches(3.4))
        ctf = content_box.text_frame
        ctf.word_wrap = True

        for item in cat['items']:
            p = ctf.add_paragraph()
            p.text = "✓ " + item
            p.font.size = Pt(12)
            p.font.color.rgb = theme['white']
            p.font.name = "맑은 고딕"
            p.space_after = Pt(12)

        x += 3.2

    # 하단 메시지
    footer_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(5.8), Inches(9.2), Inches(1.4))
    footer_box.fill.solid()
    footer_box.fill.fore_color.rgb = theme['light']
    footer_box.line.color.rgb = theme['primary']
    footer_box.line.width = Pt(2)

    footer_text = slide.shapes.add_textbox(Inches(0.6), Inches(5.95), Inches(8.8), Inches(1.2))
    ftf = footer_text.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.text = "🎯 보안 원칙: Defense in Depth (심층 방어)"
    fp.font.size = Pt(14)
    fp.font.bold = True
    fp.font.color.rgb = theme['primary']
    fp.font.name = "맑은 고딕"

    fp2 = ftf.add_paragraph()
    fp2.text = "단일 보안 계층에 의존하지 않고, 여러 계층의 보안 메커니즘을 중첩 적용하여"
    fp2.font.size = Pt(12)
    fp2.font.color.rgb = theme['text']
    fp2.font.name = "맑은 고딕"

    fp3 = ftf.add_paragraph()
    fp3.text = "하나의 계층이 뚫려도 다른 계층에서 방어할 수 있는 구조 구현"
    fp3.font.size = Pt(12)
    fp3.font.color.rgb = theme['text']
    fp3.font.name = "맑은 고딕"

    return slide

def create_qna_slide(prs, theme):
    """Q&A 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme['dark']
    bg.line.fill.background()

    # Q&A 텍스트
    qa_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(10), Inches(1.5))
    qtf = qa_box.text_frame
    qp = qtf.paragraphs[0]
    qp.text = "Q & A"
    qp.font.size = Pt(72)
    qp.font.bold = True
    qp.font.color.rgb = theme['white']
    qp.font.name = "맑은 고딕"
    qp.alignment = PP_ALIGN.CENTER

    # 하이라이트 요약
    highlights = [
        "🔐 OAuth 2.0 + PKCE 보안 인증",
        "🔄 Token Rotation 자동 갱신",
        "⚡ Redis Blacklist 실시간 무효화",
        "🛡️ 심층 방어 (Defense in Depth)"
    ]

    y = 4.3
    for hl in highlights:
        hl_box = slide.shapes.add_textbox(Inches(2), Inches(y), Inches(6), Inches(0.4))
        htf = hl_box.text_frame
        hp = htf.paragraphs[0]
        hp.text = hl
        hp.font.size = Pt(16)
        hp.font.color.rgb = theme['secondary']
        hp.font.name = "맑은 고딕"
        hp.alignment = PP_ALIGN.CENTER
        y += 0.5

    return slide

def main():
    """메인 함수 - 프레젠테이션 생성"""
    print("=" * 60)
    print("인증/권한 관리 상세 설계서 프레젠테이션 생성")
    print("테마: Brown Earth")
    print("=" * 60)

    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    theme = BROWN_EARTH

    # 슬라이드 생성
    print("\n슬라이드 생성 중...")

    # 1. 표지
    create_title_slide(
        prs,
        "인증/권한 관리\n상세 설계서",
        "OAuth 2.0 + PKCE | JWT Token Strategy | RBAC",
        theme
    )
    print("  [1/9] 표지 생성 완료")

    # 2. 아키텍처 개요
    create_overview_slide(prs, theme)
    print("  [2/9] 아키텍처 개요 생성 완료")

    # 3. OAuth 2.0 + PKCE Flow
    create_oauth_flow_slide(prs, theme)
    print("  [3/9] OAuth 2.0 + PKCE Flow 생성 완료")

    # 4. JWT Token 전략
    create_jwt_strategy_slide(prs, theme)
    print("  [4/9] JWT Token 전략 생성 완료")

    # 5. RBAC 권한 관리
    create_rbac_slide(prs, theme)
    print("  [5/9] RBAC 권한 관리 생성 완료")

    # 6. Spring Security + API Gateway
    create_spring_security_slide(prs, theme)
    print("  [6/9] Spring Security 구성 생성 완료")

    # 7. 아키텍처 장점
    create_architecture_advantages_slide(prs, theme)
    print("  [7/9] 아키텍처 장점 생성 완료")

    # 8. 보안 하이라이트
    create_security_highlights_slide(prs, theme)
    print("  [8/9] 보안 하이라이트 생성 완료")

    # 9. Q&A
    create_qna_slide(prs, theme)
    print("  [9/9] Q&A 생성 완료")

    # 저장
    output_dir = "knowledge_service/docs/02_design"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "Auth_Authorization_Design_Brown.pptx")
    prs.save(output_path)

    print("\n" + "=" * 60)
    print(f"프레젠테이션 생성 완료!")
    print(f"파일 위치: {output_path}")
    print(f"총 슬라이드: 9장")
    print(f"파일 크기: {os.path.getsize(output_path):,} bytes")
    print("=" * 60)

    return output_path

if __name__ == "__main__":
    main()
