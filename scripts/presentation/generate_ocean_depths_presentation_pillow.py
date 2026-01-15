#!/usr/bin/env python3
"""Ocean Depths 테마 프레젠테이션 생성 (Pillow 버전)"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont
import os

# Ocean Depths 컬러 팔레트
OCEAN_DEPTHS_COLORS = {
    'primary': RGBColor(0, 102, 153),      # #006699
    'secondary': RGBColor(0, 153, 204),    # #0099CC
    'accent': RGBColor(255, 102, 0),       # #FF6600
    'dark': RGBColor(0, 51, 102),          # #003366
    'light': RGBColor(204, 229, 255),      # #CCE5FF
    'text': RGBColor(33, 33, 33),          # #212121
    'bg': RGBColor(255, 255, 255)          # #FFFFFF
}

# PIL용 컬러 (튜플 형식)
PIL_COLORS = {
    'primary': (0, 102, 153),
    'secondary': (0, 153, 204),
    'accent': (255, 102, 0),
    'dark': (0, 51, 102),
    'light': (204, 229, 255),
    'text': (33, 33, 33),
    'bg': (255, 255, 255),
    'white': (255, 255, 255)
}

def create_architecture_diagram(output_path):
    """시스템 아키텍처 다이어그램 생성 (Pillow)"""
    # 이미지 크기
    width, height = 1600, 800
    img = Image.new('RGB', (width, height), PIL_COLORS['bg'])
    draw = ImageDraw.Draw(img)

    try:
        # 맑은 고딕 폰트 로드
        font_title = ImageFont.truetype("malgun.ttf", 32)
        font_text = ImageFont.truetype("malgun.ttf", 28)
    except:
        font_title = ImageFont.load_default()
        font_text = ImageFont.load_default()

    # 레이어 높이
    y_positions = [60, 240, 420, 600]

    # Frontend Layer
    draw.rounded_rectangle([100, y_positions[0], 1500, y_positions[0]+100],
                          radius=10, fill=PIL_COLORS['secondary'], outline=PIL_COLORS['dark'], width=3)
    draw.text((800, y_positions[0]+50), "Frontend Layer", fill=PIL_COLORS['white'],
             font=font_title, anchor="mm")
    draw.text((800, y_positions[0]+85), "React 18 + TypeScript", fill=PIL_COLORS['white'],
             font=font_text, anchor="mm")

    # API Gateway
    draw.rounded_rectangle([100, y_positions[1], 1500, y_positions[1]+100],
                          radius=10, fill=PIL_COLORS['primary'], outline=PIL_COLORS['dark'], width=3)
    draw.text((800, y_positions[1]+50), "API Gateway", fill=PIL_COLORS['white'],
             font=font_title, anchor="mm")
    draw.text((800, y_positions[1]+85), "Spring Cloud Gateway", fill=PIL_COLORS['white'],
             font=font_text, anchor="mm")

    # Service Layer (4개)
    services = [
        ("Knowledge\nService", 200),
        ("Search\nService", 500),
        ("User\nService", 800),
        ("AI Service\n(FastAPI)", 1100)
    ]

    for service_name, x_pos in services:
        draw.rounded_rectangle([x_pos, y_positions[2], x_pos+280, y_positions[2]+100],
                              radius=10, fill=PIL_COLORS['accent'], outline=PIL_COLORS['dark'], width=3)
        draw.text((x_pos+140, y_positions[2]+50), service_name, fill=PIL_COLORS['white'],
                 font=font_text, anchor="mm", align="center")

    # Database Layer (4개)
    databases = [
        ("PostgreSQL", 200),
        ("Elasticsearch", 500),
        ("Neo4j", 800),
        ("Redis", 1100)
    ]

    for db_name, x_pos in databases:
        draw.rounded_rectangle([x_pos, y_positions[3], x_pos+280, y_positions[3]+100],
                              radius=10, fill=PIL_COLORS['dark'], outline=PIL_COLORS['dark'], width=3)
        draw.text((x_pos+140, y_positions[3]+50), db_name, fill=PIL_COLORS['white'],
                 font=font_text, anchor="mm")

    # 화살표 그리기 (수직선)
    arrow_x_positions = [340, 640, 940, 1240]
    for x in arrow_x_positions:
        for i in range(3):
            y_start = y_positions[i] + 100
            y_end = y_positions[i+1]
            # 선
            draw.line([(x, y_start), (x, y_end)], fill=PIL_COLORS['accent'], width=5)
            # 화살표 머리
            draw.polygon([(x, y_end), (x-10, y_end-15), (x+10, y_end-15)],
                        fill=PIL_COLORS['accent'])

    img.save(output_path, 'PNG')
    print(f"Architecture diagram saved: {output_path}")
    return True

def create_dashboard_diagram(output_path):
    """대시보드 레이아웃 다이어그램 생성 (Pillow)"""
    width, height = 1600, 800
    img = Image.new('RGB', (width, height), PIL_COLORS['bg'])
    draw = ImageDraw.Draw(img)

    try:
        font_title = ImageFont.truetype("malgun.ttf", 36)
        font_text = ImageFont.truetype("malgun.ttf", 28)
        font_small = ImageFont.truetype("malgun.ttf", 24)
    except:
        font_title = ImageFont.load_default()
        font_text = ImageFont.load_default()
        font_small = ImageFont.load_default()

    # 헤더 영역
    draw.rectangle([0, 0, 1600, 100], fill=PIL_COLORS['dark'])
    draw.text((800, 50), "📊 대시보드 - 지식 검색 시스템", fill=PIL_COLORS['white'],
             font=font_title, anchor="mm")

    # 통계 카드 4개
    cards = [
        ("📚 전체 문서", "1,234", 150, 180),
        ("🔍 금주 검색", "567", 490, 180),
        ("⭐ 인기 문서", "89", 830, 180),
        ("👥 활성 사용자", "45", 1170, 180)
    ]

    for title, value, x, y in cards:
        # 카드 배경
        draw.rounded_rectangle([x, y, x+300, y+180],
                              radius=15, fill=PIL_COLORS['light'], outline=PIL_COLORS['primary'], width=3)
        # 제목
        draw.text((x+150, y+50), title, fill=PIL_COLORS['text'],
                 font=font_text, anchor="mm")
        # 값
        draw.text((x+150, y+120), value, fill=PIL_COLORS['accent'],
                 font=font_title, anchor="mm")

    # 최근 활동 영역
    draw.rounded_rectangle([100, 420, 750, 720],
                          radius=15, fill=PIL_COLORS['white'], outline=PIL_COLORS['primary'], width=3)
    draw.text((425, 450), "최근 활동", fill=PIL_COLORS['text'],
             font=font_title, anchor="mm")

    activities = [
        "• 프로젝트 A 문서 업데이트",
        "• 신규 사용자 5명 가입",
        "• API 문서 검색 증가"
    ]

    y_pos = 520
    for activity in activities:
        draw.text((120, y_pos), activity, fill=PIL_COLORS['text'], font=font_small)
        y_pos += 50

    # 인기 검색어
    draw.rounded_rectangle([850, 420, 1500, 720],
                          radius=15, fill=PIL_COLORS['white'], outline=PIL_COLORS['secondary'], width=3)
    draw.text((1175, 450), "인기 검색어", fill=PIL_COLORS['text'],
             font=font_title, anchor="mm")

    keywords = ["API 문서", "배포 가이드", "아키텍처", "인증 방법"]
    y_pos = 520
    for i, keyword in enumerate(keywords, 1):
        draw.text((870, y_pos), f"{i}. {keyword}", fill=PIL_COLORS['text'], font=font_small)
        y_pos += 50

    img.save(output_path, 'PNG')
    print(f"Dashboard diagram saved: {output_path}")
    return True

def create_search_modes_diagram(output_path):
    """검색 모드 비교 다이어그램 생성 (Pillow)"""
    width, height = 1600, 800
    img = Image.new('RGB', (width, height), PIL_COLORS['bg'])
    draw = ImageDraw.Draw(img)

    try:
        font_title = ImageFont.truetype("malgun.ttf", 40)
        font_subtitle = ImageFont.truetype("malgun.ttf", 32)
        font_text = ImageFont.truetype("malgun.ttf", 26)
    except:
        font_title = ImageFont.load_default()
        font_subtitle = ImageFont.load_default()
        font_text = ImageFont.load_default()

    # 왼쪽: 채팅 모드
    draw.rounded_rectangle([50, 50, 750, 750],
                          radius=20, fill=PIL_COLORS['light'], outline=PIL_COLORS['primary'], width=4)
    draw.rectangle([50, 50, 750, 150], fill=PIL_COLORS['primary'])
    draw.text((400, 100), "💬 채팅 모드", fill=PIL_COLORS['white'],
             font=font_title, anchor="mm")

    chat_features = [
        "RAG 기반 답변 생성",
        "자연어 질문 처리",
        "컨텍스트 유지",
        "참조 문서 링크 제공",
        "대화형 인터페이스"
    ]

    y_pos = 220
    for feature in chat_features:
        draw.text((100, y_pos), f"✓ {feature}", fill=PIL_COLORS['text'], font=font_text)
        y_pos += 80

    draw.text((400, 670), "최적: 복잡한 질문", fill=PIL_COLORS['accent'],
             font=font_subtitle, anchor="mm")

    # 오른쪽: 검색 모드
    draw.rounded_rectangle([850, 50, 1550, 750],
                          radius=20, fill=PIL_COLORS['light'], outline=PIL_COLORS['secondary'], width=4)
    draw.rectangle([850, 50, 1550, 150], fill=PIL_COLORS['secondary'])
    draw.text((1200, 100), "🔍 검색 모드", fill=PIL_COLORS['white'],
             font=font_title, anchor="mm")

    search_features = [
        "하이브리드 검색 (BM25+벡터)",
        "필터링 (날짜, 타입 등)",
        "정렬 옵션",
        "페이지네이션",
        "북마크 기능"
    ]

    y_pos = 220
    for feature in search_features:
        draw.text((900, y_pos), f"✓ {feature}", fill=PIL_COLORS['text'], font=font_text)
        y_pos += 80

    draw.text((1200, 670), "최적: 정확한 키워드", fill=PIL_COLORS['accent'],
             font=font_subtitle, anchor="mm")

    img.save(output_path, 'PNG')
    print(f"Search modes diagram saved: {output_path}")
    return True

def create_title_slide(prs, title, subtitle):
    """표지 슬라이드 - Ocean Depths 테마"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['dark']
    bg_shape.line.color.rgb = OCEAN_DEPTHS_COLORS['dark']

    # 장식 요소
    wave_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(6), Inches(5),
        Inches(5), Inches(3)
    )
    wave_shape.fill.solid()
    wave_shape.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['primary']
    wave_shape.line.color.rgb = OCEAN_DEPTHS_COLORS['primary']

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.2)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.font.name = "맑은 고딕"
    title_para.alignment = PP_ALIGN.CENTER

    # 부제목
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4), Inches(8), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = OCEAN_DEPTHS_COLORS['light']
    subtitle_para.font.name = "맑은 고딕"
    subtitle_para.alignment = PP_ALIGN.CENTER

def create_content_slide_ocean(prs, title, content_lines):
    """일반 컨텐츠 슬라이드 - Ocean Depths 테마"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더 바
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['primary']
    header_shape.line.color.rgb = OCEAN_DEPTHS_COLORS['primary']

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.font.name = "맑은 고딕"

    # 내용
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for line in content_lines:
        p = content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = OCEAN_DEPTHS_COLORS['text']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(12)

def create_image_slide_ocean(prs, title, image_path):
    """이미지 슬라이드 - Ocean Depths 테마"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더 바
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['secondary']
    header_shape.line.color.rgb = OCEAN_DEPTHS_COLORS['secondary']

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15), Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.font.name = "맑은 고딕"

    # 이미지
    if os.path.exists(image_path):
        slide.shapes.add_picture(
            image_path,
            Inches(0.5), Inches(1.2),
            width=Inches(9)
        )

def create_two_column_slide_ocean(prs, title, left_items, right_items):
    """2단 컬럼 슬라이드 - Ocean Depths 테마"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더 바
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['primary']
    header_shape.line.color.rgb = OCEAN_DEPTHS_COLORS['primary']

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.font.name = "맑은 고딕"

    # 왼쪽 패널
    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5),
        Inches(4.25), Inches(5.5)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['light']
    left_panel.line.color.rgb = OCEAN_DEPTHS_COLORS['primary']
    left_panel.line.width = Pt(3)

    left_text_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.8), Inches(3.7), Inches(5)
    )
    left_frame = left_text_box.text_frame
    left_frame.word_wrap = True

    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = OCEAN_DEPTHS_COLORS['text']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(10)

    # 오른쪽 패널
    right_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.25), Inches(1.5),
        Inches(4.25), Inches(5.5)
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = OCEAN_DEPTHS_COLORS['light']
    right_panel.line.color.rgb = OCEAN_DEPTHS_COLORS['secondary']
    right_panel.line.width = Pt(3)

    right_text_box = slide.shapes.add_textbox(
        Inches(5.55), Inches(1.8), Inches(3.7), Inches(5)
    )
    right_frame = right_text_box.text_frame
    right_frame.word_wrap = True

    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = OCEAN_DEPTHS_COLORS['text']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(10)

def main():
    # 출력 디렉토리
    output_dir = r"knowledge_service\docs\02_design\ui_storyboard"
    os.makedirs(output_dir, exist_ok=True)

    # Pillow로 다이어그램 생성
    print("다이어그램 생성 중...")
    arch_png = os.path.join(output_dir, "architecture_diagram.png")
    dash_png = os.path.join(output_dir, "dashboard_diagram.png")
    search_png = os.path.join(output_dir, "search_modes_diagram.png")

    create_architecture_diagram(arch_png)
    create_dashboard_diagram(dash_png)
    create_search_modes_diagram(search_png)

    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 표지
    create_title_slide(prs,
                      "사내 지식 검색 시스템",
                      "UI 스토리보드 발표")

    # 2. 개요
    create_content_slide_ocean(prs, "프로젝트 개요", [
        "📌 목표",
        "  • Graph RAG 기반 지능형 지식 검색",
        "  • 사내 문서의 효율적 관리 및 활용",
        "",
        "💡 핵심 기능",
        "  • AI 기반 자연어 질의응답",
        "  • 하이브리드 검색 (키워드 + 벡터)",
        "  • 엔티티 기반 관계 그래프",
        "  • 시간 기반 지식 버전 관리"
    ])

    # 3. 시스템 아키텍처
    create_image_slide_ocean(prs, "시스템 아키텍처", arch_png)

    # 4. 로그인 & 대시보드
    create_image_slide_ocean(prs, "로그인 & 대시보드", dash_png)

    # 5. 검색 기능
    create_image_slide_ocean(prs, "검색 기능 - 채팅 vs 검색 모드", search_png)

    # 6. 지식 관리
    create_two_column_slide_ocean(prs, "지식 관리", [
        "📋 지식 목록",
        "• 카드 뷰 / 리스트 뷰",
        "• 필터링 (프로젝트, 타입, 날짜)",
        "• 정렬 (최신순, 인기순)",
        "• 무한 스크롤",
        "",
        "📝 지식 작성/수정",
        "• 마크다운 에디터",
        "• 파일 업로드 (PDF, DOCX)",
        "• 메타데이터 자동 추출",
        "• 미리보기 기능"
    ], [
        "📄 지식 상세",
        "• 풀스크린 뷰어",
        "• 목차 네비게이션",
        "• 관련 문서 추천",
        "• 참조 그래프 시각화",
        "",
        "🔖 부가 기능",
        "• 북마크",
        "• 댓글/피드백",
        "• 버전 히스토리",
        "• 공유 링크 생성"
    ])

    # 7. 프로필 & 관리자
    create_two_column_slide_ocean(prs, "프로필 & 관리자", [
        "👤 프로필 관리",
        "• 개인정보 수정",
        "• 비밀번호 변경",
        "• 알림 설정",
        "• 테마 설정 (라이트/다크)",
        "",
        "⭐ 북마크",
        "• 즐겨찾기 문서 목록",
        "• 태그별 분류",
        "• 검색 히스토리"
    ], [
        "🔧 관리자 페이지",
        "• 사용자 관리 (CRUD)",
        "• 권한 관리 (역할 기반)",
        "• 지식 승인/거부",
        "• 통계 대시보드",
        "",
        "📊 시스템 모니터링",
        "• 검색 분석",
        "• 성능 지표",
        "• 사용량 리포트"
    ])

    # 8. 기술 스택
    create_two_column_slide_ocean(prs, "기술 스택", [
        "🎨 Frontend",
        "• React 18 + TypeScript",
        "• Material-UI (MUI)",
        "• TanStack Query (데이터 관리)",
        "• Recharts (차트)",
        "",
        "🔧 Backend",
        "• Spring Boot 3.x (Java 17)",
        "• Spring Security (JWT)",
        "• FastAPI (AI Service)"
    ], [
        "🗄️ Database",
        "• PostgreSQL (마스터 DB)",
        "• Elasticsearch (검색/벡터)",
        "• Neo4j (그래프)",
        "• Redis (캐시)",
        "",
        "🤖 AI/ML",
        "• DeepSeek-V3.2 (LLM)",
        "• BGE-M3 (임베딩)",
        "• LangGraph (워크플로우)"
    ])

    # 9. 기대 효과
    create_content_slide_ocean(prs, "기대 효과", [
        "⏰ 검색 시간 단축",
        "  • 기존 대비 70% 이상 시간 절감",
        "  • AI 기반 즉시 답변 제공",
        "",
        "📈 지식 활용도 증대",
        "  • 관련 문서 자동 추천",
        "  • 시간 기반 정확한 정보 제공",
        "",
        "💰 비용 효율",
        "  • DeepSeek 활용으로 95% 비용 절감",
        "  • 오픈소스 기반 확장 가능"
    ])

    # 10. Q&A
    create_title_slide(prs, "Q & A", "질문해 주세요")

    # 저장
    output_path = os.path.join(output_dir, "UI스토리보드_발표자료_OceanDepths.pptx")
    prs.save(output_path)

    print(f"\n프레젠테이션 생성 완료!")
    print(f"파일 위치: {output_path}")
    print(f"총 슬라이드 수: {len(prs.slides)}")
    print(f"\n다이어그램:")
    print(f"  - architecture_diagram.png")
    print(f"  - dashboard_diagram.png")
    print(f"  - search_modes_diagram.png")

if __name__ == "__main__":
    main()
