#!/usr/bin/env python3
"""
Hybrid RAG Platform Detailed Design Presentation Generator
- Based on hybrid_rag_platform_detailed_design.md v2.3
- Tech Innovation Theme
- PPT Shape-based Diagrams with detailed comments
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os

# ========== Tech Innovation Theme ==========
THEME = {
    'primary': RGBColor(46, 125, 50),       # #2E7D32 - Green
    'secondary': RGBColor(129, 199, 132),   # #81C784 - Light Green
    'accent': RGBColor(255, 152, 0),        # #FF9800 - Orange
    'dark': RGBColor(27, 94, 32),           # #1B5E20 - Dark Green
    'light': RGBColor(232, 245, 233),       # #E8F5E9 - Mint
    'text': RGBColor(33, 33, 33),           # #212121
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(158, 158, 158),        # #9E9E9E
    'lightgray': RGBColor(238, 238, 238),   # #EEEEEE
    'red': RGBColor(244, 67, 54),           # #F44336
    'blue': RGBColor(33, 150, 243),         # #2196F3
}

FONT_NAME = "맑은 고딕"


def add_shape_with_text(slide, shape_type, left, top, width, height,
                        fill_color, text="", font_size=12, font_color=None,
                        bold=False, border_color=None, border_width=Pt(1)):
    """Add a shape with text"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.color.rgb = fill_color

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.name = FONT_NAME
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].font.color.rgb = font_color or THEME['text']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 font_color=None, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = FONT_NAME
    p.font.bold = bold
    p.font.color.rgb = font_color or THEME['text']
    p.alignment = alignment
    return textbox


def create_section_header(slide, title, subtitle=""):
    """Add section header to slide"""
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.color.rgb = THEME['primary']

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
                 title, font_size=32, font_color=THEME['white'], bold=True)

    if subtitle:
        add_text_box(slide, Inches(10.5), Inches(0.4), Inches(2.5), Inches(0.5),
                     subtitle, font_size=14, font_color=THEME['light'])


def create_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['dark']
    bg.line.color.rgb = THEME['dark']

    # Decorative elements
    deco1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(4.5), Inches(5), Inches(4))
    deco1.fill.solid()
    deco1.fill.fore_color.rgb = THEME['primary']
    deco1.line.color.rgb = THEME['primary']

    deco2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(-1), Inches(4), Inches(3))
    deco2.fill.solid()
    deco2.fill.fore_color.rgb = THEME['secondary']
    deco2.line.color.rgb = THEME['secondary']

    # Title
    add_text_box(slide, Inches(0.5), Inches(2), Inches(12), Inches(1.2),
                 "Neo4j Graph RAG 기반", font_size=28, font_color=THEME['secondary'],
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(2.8), Inches(12), Inches(1.5),
                 "Hybrid 지식 플랫폼 상세 설계서", font_size=44, font_color=THEME['white'],
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Version and date
    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(0.6),
                 "Version 2.3 | Review 완료 (코드 검증됨)", font_size=18, font_color=THEME['light'],
                 alignment=PP_ALIGN.CENTER)

    # Key points
    points = [
        "Vector + Graph Hybrid Search",
        "DeepSeek V3.2 LLM 통합",
        "95% 비용 절감"
    ]
    for i, point in enumerate(points):
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.5 + i * 2.8), Inches(5.5), Inches(2.6), Inches(0.5),
            THEME['accent'], point, font_size=12, font_color=THEME['white'], bold=True
        )

    # Date
    add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                 "2026-01-14", font_size=14, font_color=THEME['secondary'],
                 alignment=PP_ALIGN.CENTER)

    return slide


def create_overview_slide(prs):
    """Slide 2: System Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "System Overview", "시스템 개요")

    # Description
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.5), Inches(12.7), Inches(0.8),
        THEME['light'],
        "Vector Search와 Graph Search를 결합하여 기업 내부 지식을 효과적으로 검색하는 Hybrid RAG 시스템",
        font_size=16, font_color=THEME['dark'], bold=True
    )

    # Core principles (4 boxes)
    principles = [
        ("비용 효율성", "LLM 비용 95% 절감\nDeepSeek 단일 모델"),
        ("제로 조인", "DB 간 조인 제거\nES 메타데이터 비정규화"),
        ("단일 진실 공급원", "데이터 일관성 보장\nPostgreSQL SSOT"),
        ("16GB RAM 운영", "제한된 리소스 최적화\nSlim Graph 전략"),
    ]

    for i, (title, desc) in enumerate(principles):
        x = 0.3 + i * 3.2
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2.5), Inches(3), Inches(1.5),
            THEME['white'], border_color=THEME['primary']
        )
        add_text_box(slide, Inches(x + 0.1), Inches(2.6), Inches(2.8), Inches(0.4),
                     title, font_size=14, font_color=THEME['primary'],
                     bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.1), Inches(3.0), Inches(2.8), Inches(0.9),
                     desc, font_size=11, font_color=THEME['text'], alignment=PP_ALIGN.CENTER)

    # Hybrid Search Flow Diagram
    add_text_box(slide, Inches(0.3), Inches(4.2), Inches(4), Inches(0.4),
                 "Hybrid Search 흐름", font_size=14, font_color=THEME['primary'], bold=True)

    # Flow: Query -> Hybrid Engine -> Vector/Graph -> Integration -> LLM -> Response
    flow_items = [
        ("사용자\n질의", THEME['light']),
        ("Hybrid\nEngine", THEME['primary']),
        ("Vector\n(ES)", THEME['secondary']),
        ("Graph\n(Neo4j)", THEME['secondary']),
        ("결과\n통합", THEME['primary']),
        ("DeepSeek\n답변합성", THEME['accent']),
        ("최종\n응답", THEME['light']),
    ]

    # Draw flow
    x_start = 0.5
    for i, (text, color) in enumerate(flow_items):
        x = x_start + i * 1.8
        y = 4.7

        # Special handling for Vector/Graph parallel
        if i == 2:  # Vector
            add_shape_with_text(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y - 0.5), Inches(1.5), Inches(0.8),
                color, text, font_size=10,
                font_color=THEME['white'] if color == THEME['primary'] else THEME['dark']
            )
        elif i == 3:  # Graph
            add_shape_with_text(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x - 1.8), Inches(y + 0.5), Inches(1.5), Inches(0.8),
                color, text, font_size=10,
                font_color=THEME['white'] if color == THEME['primary'] else THEME['dark']
            )
        elif i > 3:  # After parallel
            add_shape_with_text(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x - 1.8), Inches(y), Inches(1.5), Inches(0.8),
                color, text, font_size=10,
                font_color=THEME['white'] if color in [THEME['primary'], THEME['accent']] else THEME['dark']
            )
        else:
            add_shape_with_text(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(1.5), Inches(0.8),
                color, text, font_size=10,
                font_color=THEME['white'] if color == THEME['primary'] else THEME['dark']
            )

        # Arrows
        if i < len(flow_items) - 1 and i != 2 and i != 3:
            add_shape_with_text(
                slide, MSO_SHAPE.RIGHT_ARROW,
                Inches(x + 1.5), Inches(y + 0.3), Inches(0.25), Inches(0.2),
                THEME['gray']
            )

    # Key terms
    add_text_box(slide, Inches(0.3), Inches(5.9), Inches(4), Inches(0.4),
                 "핵심 용어", font_size=12, font_color=THEME['primary'], bold=True)

    terms = [
        ("Hybrid RAG", "Vector + Graph 결합 검색"),
        ("VIP 아키텍처", "Value-Intelligent-Planning 3단계"),
        ("RRF", "Reciprocal Rank Fusion"),
        ("Slim Graph", "Neo4j 최소 속성 저장"),
    ]

    for i, (term, desc) in enumerate(terms):
        x = 0.3 + (i % 2) * 6.5
        y = 6.2 + (i // 2) * 0.45
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(1.5), Inches(0.4),
            THEME['secondary'], term, font_size=9, font_color=THEME['dark'], bold=True
        )
        add_text_box(slide, Inches(x + 1.6), Inches(y + 0.05), Inches(4.5), Inches(0.35),
                     desc, font_size=10, font_color=THEME['text'])

    return slide


def create_tech_stack_slide(prs):
    """Slide 3: Tech Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "Tech Stack", "기술 스택")

    # Categories
    categories = [
        ("문서 파싱", [("Docling", "2.x", "PDF/DOCX 텍스트, 테이블 추출")]),
        ("LLM", [
            ("DeepSeek-Chat", "V3.2", "엔티티 추출, 답변 합성"),
            ("DeepSeek-Reasoner", "V3.2", "복잡한 시계열 추론"),
        ]),
        ("임베딩", [("BGE-M3", "-", "Dense + Sparse 벡터 생성")]),
        ("데이터베이스", [
            ("Neo4j", "5.x", "지식 그래프 저장"),
            ("Elasticsearch", "8.x", "벡터 검색 + 메타데이터"),
            ("PostgreSQL", "16+", "마스터 레코드 (SSOT)"),
        ]),
        ("오케스트레이션", [
            ("LangGraph", "1.0+", "AI 워크플로우 관리"),
            ("LangChain", "1.2+", "LLM 통합"),
        ]),
    ]

    y_pos = 1.5
    for cat_name, items in categories:
        # Category header
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3), Inches(y_pos), Inches(2.2), Inches(0.45),
            THEME['primary'], cat_name, font_size=12, font_color=THEME['white'], bold=True
        )

        # Items
        for i, (tech, ver, desc) in enumerate(items):
            x = 2.7 + i * 3.5
            if i >= 3:  # Wrap to next line
                continue

            add_shape_with_text(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y_pos), Inches(3.3), Inches(0.45),
                THEME['light'], border_color=THEME['secondary']
            )
            add_text_box(slide, Inches(x + 0.1), Inches(y_pos + 0.05), Inches(1.2), Inches(0.35),
                         f"{tech} {ver}", font_size=10, font_color=THEME['primary'], bold=True)
            add_text_box(slide, Inches(x + 1.4), Inches(y_pos + 0.05), Inches(1.8), Inches(0.35),
                         desc, font_size=9, font_color=THEME['gray'])

        y_pos += 0.55
        if len(items) > 3:
            for i, (tech, ver, desc) in enumerate(items[3:]):
                x = 2.7 + i * 3.5
                add_shape_with_text(
                    slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x), Inches(y_pos), Inches(3.3), Inches(0.45),
                    THEME['light'], border_color=THEME['secondary']
                )
                add_text_box(slide, Inches(x + 0.1), Inches(y_pos + 0.05), Inches(1.2), Inches(0.35),
                             f"{tech} {ver}", font_size=10, font_color=THEME['primary'], bold=True)
                add_text_box(slide, Inches(x + 1.4), Inches(y_pos + 0.05), Inches(1.8), Inches(0.35),
                             desc, font_size=9, font_color=THEME['gray'])
            y_pos += 0.55

    # Service Architecture
    add_text_box(slide, Inches(0.3), Inches(4.5), Inches(4), Inches(0.4),
                 "서비스 분리 아키텍처", font_size=14, font_color=THEME['primary'], bold=True)

    # SpringBoot
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(5.0), Inches(4), Inches(1.8),
        THEME['white'], border_color=THEME['primary'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(0.4), Inches(5.05), Inches(3.8), Inches(0.35),
                 "SpringBoot Backend", font_size=12, font_color=THEME['primary'], bold=True)
    add_text_box(slide, Inches(0.4), Inches(5.4), Inches(3.8), Inches(1.3),
                 "Java 17+, Spring Boot 3.x\n\n- 비즈니스 로직, CRUD\n- 사용자 인증/인가 (JWT)\n- 트랜잭션 관리\n- AI Service 연동 (WebClient)",
                 font_size=9, font_color=THEME['text'])

    # Arrow
    add_shape_with_text(
        slide, MSO_SHAPE.RIGHT_ARROW,
        Inches(4.4), Inches(5.7), Inches(0.5), Inches(0.4),
        THEME['accent']
    )
    add_text_box(slide, Inches(4.4), Inches(6.15), Inches(0.5), Inches(0.3),
                 "REST", font_size=8, font_color=THEME['gray'], alignment=PP_ALIGN.CENTER)

    # AI Service
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.0), Inches(5.0), Inches(4), Inches(1.8),
        THEME['white'], border_color=THEME['accent'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(5.1), Inches(5.05), Inches(3.8), Inches(0.35),
                 "AI Service (Python)", font_size=12, font_color=THEME['accent'], bold=True)
    add_text_box(slide, Inches(5.1), Inches(5.4), Inches(3.8), Inches(1.3),
                 "Python 3.11+, FastAPI\n\n- Hybrid 검색 (Vector+Graph)\n- 엔티티/메타데이터 추출\n- 임베딩 생성 (BGE-M3)\n- 답변 합성 (VIP Pipeline)",
                 font_size=9, font_color=THEME['text'])

    # Databases
    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(9.3), Inches(5.0), Inches(1.2), Inches(0.8),
        THEME['dark'], "PG", font_size=10, font_color=THEME['white']
    )
    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(10.6), Inches(5.0), Inches(1.2), Inches(0.8),
        THEME['dark'], "ES", font_size=10, font_color=THEME['white']
    )
    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(11.9), Inches(5.0), Inches(1.2), Inches(0.8),
        THEME['dark'], "Neo4j", font_size=10, font_color=THEME['white']
    )

    # Key point
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.3), Inches(6.0), Inches(3.8), Inches(0.8),
        THEME['light'], border_color=THEME['secondary']
    )
    add_text_box(slide, Inches(9.4), Inches(6.05), Inches(3.6), Inches(0.7),
                 "SpringBoot는 AI 모델과\n직접 연동하지 않음!\n-> AI Service REST API 호출",
                 font_size=9, font_color=THEME['primary'], bold=True, alignment=PP_ALIGN.CENTER)

    return slide


def create_vip_architecture_slide(prs):
    """Slide 4: VIP 3-Stage Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "VIP 3-Stage LLM Architecture", "3단계 LLM 처리 파이프라인")

    # Stage 1: Value
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.5), Inches(4), Inches(2.5),
        THEME['light'], border_color=THEME['primary'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(0.4), Inches(1.55), Inches(3.8), Inches(0.4),
                 "Stage 1: Value (엔티티 채굴)", font_size=14, font_color=THEME['primary'], bold=True)

    stage1_items = [
        ("문서 입력", "청크 텍스트 (512 토큰)"),
        ("DeepSeek-Chat", "단순 엔티티/관계 추출"),
        ("DeepSeek-Reasoner", "복잡한 관계 추론"),
        ("출력", "엔티티, 관계, 메타데이터"),
    ]
    for i, (item, desc) in enumerate(stage1_items):
        y = 2.0 + i * 0.5
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y), Inches(1.4), Inches(0.4),
            THEME['secondary'] if i < 3 else THEME['accent'],
            item, font_size=9, font_color=THEME['dark']
        )
        add_text_box(slide, Inches(2.0), Inches(y + 0.05), Inches(2.2), Inches(0.35),
                     desc, font_size=9, font_color=THEME['text'])

    # Stage 2: Intelligent
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.6), Inches(1.5), Inches(4), Inches(2.5),
        THEME['light'], border_color=THEME['blue'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(4.7), Inches(1.55), Inches(3.8), Inches(0.4),
                 "Stage 2: Intelligent (오케스트레이션)", font_size=14, font_color=THEME['blue'], bold=True)

    stage2_items = [
        ("사용자 질의", "자연어 질문 입력"),
        ("의도 분석", "Reasoner: 복잡도 판단"),
        ("검색 전략", "es_only/neo4j_only/hybrid"),
        ("쿼리 실행", "Chat: Vector+Graph 검색"),
    ]
    for i, (item, desc) in enumerate(stage2_items):
        y = 2.0 + i * 0.5
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(4.8), Inches(y), Inches(1.4), Inches(0.4),
            THEME['secondary'] if i < 3 else THEME['accent'],
            item, font_size=9, font_color=THEME['dark']
        )
        add_text_box(slide, Inches(6.3), Inches(y + 0.05), Inches(2.2), Inches(0.35),
                     desc, font_size=9, font_color=THEME['text'])

    # Stage 3: Planning
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.9), Inches(1.5), Inches(4.1), Inches(2.5),
        THEME['light'], border_color=THEME['accent'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(9.0), Inches(1.55), Inches(3.9), Inches(0.4),
                 "Stage 3: Planning (답변 합성)", font_size=14, font_color=THEME['accent'], bold=True)

    stage3_items = [
        ("검색 결과", "Vector + Graph 통합"),
        ("컨텍스트 구성", "RRF 재순위화"),
        ("답변 생성", "DeepSeek-Chat"),
        ("최종 응답", "출처 포함 자연어"),
    ]
    for i, (item, desc) in enumerate(stage3_items):
        y = 2.0 + i * 0.5
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.1), Inches(y), Inches(1.4), Inches(0.4),
            THEME['secondary'] if i < 3 else THEME['accent'],
            item, font_size=9, font_color=THEME['dark']
        )
        add_text_box(slide, Inches(10.6), Inches(y + 0.05), Inches(2.2), Inches(0.35),
                     desc, font_size=9, font_color=THEME['text'])

    # Arrows between stages
    add_shape_with_text(
        slide, MSO_SHAPE.RIGHT_ARROW,
        Inches(4.35), Inches(2.6), Inches(0.25), Inches(0.3),
        THEME['gray']
    )
    add_shape_with_text(
        slide, MSO_SHAPE.RIGHT_ARROW,
        Inches(8.65), Inches(2.6), Inches(0.25), Inches(0.3),
        THEME['gray']
    )

    # Model usage table
    add_text_box(slide, Inches(0.3), Inches(4.2), Inches(4), Inches(0.4),
                 "Stage별 모델 사용", font_size=12, font_color=THEME['primary'], bold=True)

    # Table header
    headers = ["Stage", "모델", "Temperature", "용도"]
    col_widths = [2.5, 2.5, 1.5, 3.5]
    x = 0.3
    for i, header in enumerate(headers):
        add_shape_with_text(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(4.6), Inches(col_widths[i]), Inches(0.4),
            THEME['primary'], header, font_size=10, font_color=THEME['white'], bold=True
        )
        x += col_widths[i]

    # Table rows
    rows = [
        ("Stage 1 (단순)", "deepseek-chat", "0", "엔티티/관계 추출"),
        ("Stage 1 (복잡)", "deepseek-reasoner", "1", "복잡한 관계 추론"),
        ("Stage 2 (분석)", "deepseek-reasoner", "1", "의도 분석, 전략 수립"),
        ("Stage 2 (실행)", "deepseek-chat", "0", "검색 쿼리 실행"),
        ("Stage 3", "deepseek-chat", "0", "답변 합성"),
    ]

    for row_idx, row in enumerate(rows):
        x = 0.3
        y = 5.0 + row_idx * 0.35
        for col_idx, cell in enumerate(row):
            add_shape_with_text(
                slide, MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(col_widths[col_idx]), Inches(0.35),
                THEME['white'] if row_idx % 2 == 0 else THEME['lightgray'],
                cell, font_size=9, font_color=THEME['text'],
                border_color=THEME['gray']
            )
            x += col_widths[col_idx]

    # Key note
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(6.8), Inches(12.7), Inches(0.5),
        THEME['accent'],
        "핵심: Reasoner는 Thinking Mode (temperature=1) 필수 | Chat은 Non-thinking (temperature=0)",
        font_size=11, font_color=THEME['white'], bold=True
    )

    return slide


def create_data_model_slide(prs):
    """Slide 5: Data Model"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "Data Model Design", "데이터 모델 설계")

    # PostgreSQL
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.5), Inches(4), Inches(3.3),
        THEME['white'], border_color=THEME['dark'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(0.4), Inches(1.55), Inches(3.8), Inches(0.4),
                 "PostgreSQL (SSOT)", font_size=14, font_color=THEME['dark'], bold=True)

    pg_tables = [
        "documents - 문서 마스터",
        "chunks - 청크 저장",
        "entities - 엔티티 마스터",
        "projects - 프로젝트",
        "persons - 인물 정보",
        "document_entities - 문서-엔티티 연결",
        "entity_relationships - 엔티티 관계",
    ]
    for i, table in enumerate(pg_tables):
        add_text_box(slide, Inches(0.5), Inches(2.0 + i * 0.35), Inches(3.6), Inches(0.35),
                     f"- {table}", font_size=10, font_color=THEME['text'])

    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(4.5), Inches(3.6), Inches(0.25),
        THEME['secondary'], "마스터 레코드 - 다른 DB에서 직접 수정 금지",
        font_size=8, font_color=THEME['dark']
    )

    # Elasticsearch
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.6), Inches(1.5), Inches(4), Inches(3.3),
        THEME['white'], border_color=THEME['blue'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(4.7), Inches(1.55), Inches(3.8), Inches(0.4),
                 "Elasticsearch (Vector + Meta)", font_size=14, font_color=THEME['blue'], bold=True)

    es_fields = [
        ("chunk_id", "청크 식별자"),
        ("document_id", "문서 ID (PG 참조)"),
        ("content", "청크 텍스트"),
        ("dense_vector", "1024차원 Dense"),
        ("sparse_vector", "Sparse (BM25)"),
        ("metadata.*", "비정규화된 메타"),
    ]
    for i, (field, desc) in enumerate(es_fields):
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(4.8), Inches(2.0 + i * 0.4), Inches(1.4), Inches(0.35),
            THEME['light'], field, font_size=9, font_color=THEME['text']
        )
        add_text_box(slide, Inches(6.3), Inches(2.0 + i * 0.4), Inches(2.2), Inches(0.35),
                     desc, font_size=9, font_color=THEME['gray'])

    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.8), Inches(4.5), Inches(3.6), Inches(0.25),
        THEME['secondary'], "제로 조인 - 메타데이터 비정규화",
        font_size=8, font_color=THEME['dark']
    )

    # Neo4j
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.9), Inches(1.5), Inches(4.1), Inches(3.3),
        THEME['white'], border_color=THEME['primary'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(9.0), Inches(1.55), Inches(3.9), Inches(0.4),
                 "Neo4j (Knowledge Graph)", font_size=14, font_color=THEME['primary'], bold=True)

    neo4j_items = [
        ("Entity", "id, name, type, description"),
        ("TextUnit", "id, document_id, chunk_index"),
        ("Community", "id, title, summary, level"),
        ("Document", "id, title, type"),
    ]
    add_text_box(slide, Inches(9.1), Inches(1.95), Inches(3.8), Inches(0.3),
                 "Nodes:", font_size=10, font_color=THEME['primary'], bold=True)
    for i, (node, desc) in enumerate(neo4j_items):
        add_shape_with_text(
            slide, MSO_SHAPE.OVAL,
            Inches(9.2), Inches(2.25 + i * 0.45), Inches(1.2), Inches(0.4),
            THEME['secondary'], node, font_size=9, font_color=THEME['dark']
        )
        add_text_box(slide, Inches(10.5), Inches(2.25 + i * 0.45), Inches(2.4), Inches(0.4),
                     desc, font_size=8, font_color=THEME['gray'])

    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.1), Inches(4.5), Inches(3.7), Inches(0.25),
        THEME['secondary'], "Slim Graph - 최소 속성만 저장",
        font_size=8, font_color=THEME['dark']
    )

    # Sync diagram
    add_text_box(slide, Inches(0.3), Inches(5.0), Inches(4), Inches(0.4),
                 "데이터 동기화 흐름", font_size=12, font_color=THEME['primary'], bold=True)

    # Flow: PG -> ES, PG -> Neo4j
    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(1.5), Inches(5.5), Inches(2), Inches(0.8),
        THEME['dark'], "PostgreSQL\n(SSOT)", font_size=10, font_color=THEME['white']
    )

    # Arrows
    add_shape_with_text(
        slide, MSO_SHAPE.RIGHT_ARROW,
        Inches(3.6), Inches(5.6), Inches(0.8), Inches(0.3),
        THEME['accent']
    )
    add_shape_with_text(
        slide, MSO_SHAPE.RIGHT_ARROW,
        Inches(3.6), Inches(6.0), Inches(0.8), Inches(0.3),
        THEME['accent']
    )

    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(4.5), Inches(5.3), Inches(1.8), Inches(0.6),
        THEME['blue'], "ES", font_size=10, font_color=THEME['white']
    )
    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(4.5), Inches(6.0), Inches(1.8), Inches(0.6),
        THEME['primary'], "Neo4j", font_size=10, font_color=THEME['white']
    )

    add_text_box(slide, Inches(6.5), Inches(5.5), Inches(6), Inches(0.8),
                 "동기화 규칙:\n- 3개 DB 동시 저장 (asyncio.gather)\n- PostgreSQL이 실패하면 전체 롤백\n- ES/Neo4j 실패시 재시도 큐에 등록",
                 font_size=10, font_color=THEME['text'])

    return slide


def create_hybrid_search_slide(prs):
    """Slide 6: Hybrid Search"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "Hybrid Search Engine", "하이브리드 검색 엔진")

    # Vector Search
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.5), Inches(6), Inches(2.5),
        THEME['white'], border_color=THEME['blue'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(0.4), Inches(1.55), Inches(5.8), Inches(0.4),
                 "Vector Search (Elasticsearch)", font_size=14, font_color=THEME['blue'], bold=True)

    vector_features = [
        ("Dense Vector", "BGE-M3 1024차원\n의미론적 유사도 검색"),
        ("Sparse Vector", "BM25 스타일\n키워드 가중치 매칭"),
        ("메타데이터 필터", "project_name, date\n제로 조인 쿼리"),
    ]

    for i, (title, desc) in enumerate(vector_features):
        x = 0.5 + i * 1.9
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2.0), Inches(1.75), Inches(1.0),
            THEME['light'], border_color=THEME['secondary']
        )
        add_text_box(slide, Inches(x + 0.1), Inches(2.05), Inches(1.55), Inches(0.3),
                     title, font_size=10, font_color=THEME['blue'], bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.1), Inches(2.35), Inches(1.55), Inches(0.6),
                     desc, font_size=8, font_color=THEME['text'], alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.1), Inches(5.6), Inches(0.8),
                 "ES knn 쿼리: num_candidates=100, k=20\n필터와 함께 단일 쿼리로 검색 (제로 조인)",
                 font_size=9, font_color=THEME['gray'])

    # Graph Search
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.6), Inches(1.5), Inches(6.4), Inches(2.5),
        THEME['white'], border_color=THEME['primary'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(6.7), Inches(1.55), Inches(6.2), Inches(0.4),
                 "Graph Search (Neo4j)", font_size=14, font_color=THEME['primary'], bold=True)

    graph_features = [
        ("엔티티 검색", "FULLTEXT INDEX\n엔티티명 매칭"),
        ("관계 탐색", "2-hop 확장\n관련 문서 발견"),
        ("커뮤니티", "계층적 그룹\n주제별 탐색"),
    ]

    for i, (title, desc) in enumerate(graph_features):
        x = 6.8 + i * 2.0
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2.0), Inches(1.85), Inches(1.0),
            THEME['light'], border_color=THEME['secondary']
        )
        add_text_box(slide, Inches(x + 0.1), Inches(2.05), Inches(1.65), Inches(0.3),
                     title, font_size=10, font_color=THEME['primary'], bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.1), Inches(2.35), Inches(1.65), Inches(0.6),
                     desc, font_size=8, font_color=THEME['text'], alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(6.8), Inches(3.1), Inches(6.0), Inches(0.8),
                 "Cypher: 엔티티 -> MENTIONED_IN -> TextUnit -> document_id\n관계 그래프로 숨겨진 연결 발견",
                 font_size=9, font_color=THEME['gray'])

    # RRF Fusion
    add_text_box(slide, Inches(0.3), Inches(4.2), Inches(6), Inches(0.4),
                 "RRF (Reciprocal Rank Fusion)", font_size=14, font_color=THEME['primary'], bold=True)

    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(4.6), Inches(12.7), Inches(1.5),
        THEME['light'], border_color=THEME['accent'], border_width=Pt(2)
    )

    # RRF formula
    add_text_box(slide, Inches(0.5), Inches(4.7), Inches(4), Inches(0.4),
                 "융합 공식", font_size=11, font_color=THEME['accent'], bold=True)
    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(4), Inches(0.8),
                 "score(d) = SUM( 1 / (k + rank_i(d)) )\n\nk = 60 (상수)\nrank_i = 각 검색에서의 순위",
                 font_size=10, font_color=THEME['text'])

    # Implementation note
    add_text_box(slide, Inches(4.8), Inches(4.7), Inches(4), Inches(0.4),
                 "구현 방법", font_size=11, font_color=THEME['accent'], bold=True)
    add_text_box(slide, Inches(4.8), Inches(5.1), Inches(4), Inches(0.8),
                 "ranx 라이브러리 사용\n(ES 내장 RRF는 Platinum 라이선스)\n\nfrom ranx import fuse\nfuse(runs, method='rrf')",
                 font_size=10, font_color=THEME['text'])

    # Weights
    add_text_box(slide, Inches(9.2), Inches(4.7), Inches(3.5), Inches(0.4),
                 "가중치 설정", font_size=11, font_color=THEME['accent'], bold=True)
    weights = [
        ("Vector Dense", "0.4"),
        ("Vector Sparse", "0.3"),
        ("Graph", "0.3"),
    ]
    for i, (name, weight) in enumerate(weights):
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.3), Inches(5.1 + i * 0.35), Inches(2), Inches(0.3),
            THEME['secondary'], name, font_size=9, font_color=THEME['dark']
        )
        add_text_box(slide, Inches(11.4), Inches(5.1 + i * 0.35), Inches(0.8), Inches(0.3),
                     weight, font_size=10, font_color=THEME['primary'], bold=True)

    # Performance target
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.5),
        THEME['primary'],
        "성능 목표: 검색 응답 시간 < 2초 | Top-10 Recall > 85% | MRR > 0.7",
        font_size=12, font_color=THEME['white'], bold=True
    )

    return slide


def create_cost_analysis_slide(prs):
    """Slide 7: Cost Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "Cost Analysis", "비용 분석")

    # Cost comparison table
    add_text_box(slide, Inches(0.3), Inches(1.4), Inches(6), Inches(0.4),
                 "LLM 비용 비교 (1,000문서 기준)", font_size=14, font_color=THEME['primary'], bold=True)

    # Table
    headers = ["작업", "Claude Sonnet", "GPT-4o", "DeepSeek", "절감률"]
    col_widths = [2.5, 2, 2, 2, 1.5]

    x = 0.3
    for i, header in enumerate(headers):
        add_shape_with_text(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(1.8), Inches(col_widths[i]), Inches(0.45),
            THEME['primary'], header, font_size=11, font_color=THEME['white'], bold=True
        )
        x += col_widths[i]

    rows = [
        ("엔티티 추출", "$10.50", "$3.00", "$0.56", "94.7%"),
        ("관계 추론", "$15.00", "$5.00", "$1.20", "92.0%"),
        ("오케스트레이션", "$8.00", "$3.00", "$0.08", "99.0%"),
        ("답변 합성", "$12.00", "$3.50", "$0.42", "96.5%"),
        ("총계", "$45.50", "$14.50", "$2.26", "95.0%"),
    ]

    for row_idx, row in enumerate(rows):
        x = 0.3
        y = 2.25 + row_idx * 0.45
        is_total = row_idx == len(rows) - 1
        for col_idx, cell in enumerate(row):
            bg = THEME['accent'] if (col_idx == 3 and not is_total) else (THEME['secondary'] if is_total else THEME['white'])
            fc = THEME['white'] if col_idx == 3 and not is_total else (THEME['dark'] if is_total else THEME['text'])
            add_shape_with_text(
                slide, MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(col_widths[col_idx]), Inches(0.45),
                bg, cell, font_size=11 if is_total else 10,
                font_color=fc, bold=is_total, border_color=THEME['gray']
            )
            x += col_widths[col_idx]

    # DeepSeek pricing
    add_text_box(slide, Inches(0.3), Inches(4.7), Inches(5), Inches(0.4),
                 "DeepSeek 모델별 비용", font_size=14, font_color=THEME['primary'], bold=True)

    pricing = [
        ("deepseek-chat", "$0.28/1M", "$1.10/1M", "$0.028/1M"),
        ("deepseek-reasoner", "$2.19/1M", "$8.98/1M", "-"),
    ]

    add_shape_with_text(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0.3), Inches(5.1), Inches(2), Inches(0.35),
        THEME['primary'], "모델", font_size=10, font_color=THEME['white'], bold=True
    )
    add_shape_with_text(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(2.3), Inches(5.1), Inches(1.5), Inches(0.35),
        THEME['primary'], "입력", font_size=10, font_color=THEME['white'], bold=True
    )
    add_shape_with_text(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(3.8), Inches(5.1), Inches(1.5), Inches(0.35),
        THEME['primary'], "출력", font_size=10, font_color=THEME['white'], bold=True
    )
    add_shape_with_text(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(5.3), Inches(5.1), Inches(1.5), Inches(0.35),
        THEME['primary'], "캐시 히트", font_size=10, font_color=THEME['white'], bold=True
    )

    for i, (model, input_cost, output_cost, cache) in enumerate(pricing):
        y = 5.45 + i * 0.35
        add_shape_with_text(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(0.3), Inches(y), Inches(2), Inches(0.35),
            THEME['light'], model, font_size=9, font_color=THEME['text'], border_color=THEME['gray']
        )
        add_shape_with_text(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(2.3), Inches(y), Inches(1.5), Inches(0.35),
            THEME['white'], input_cost, font_size=9, font_color=THEME['text'], border_color=THEME['gray']
        )
        add_shape_with_text(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(3.8), Inches(y), Inches(1.5), Inches(0.35),
            THEME['white'], output_cost, font_size=9, font_color=THEME['text'], border_color=THEME['gray']
        )
        add_shape_with_text(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(5.3), Inches(y), Inches(1.5), Inches(0.35),
            THEME['white'], cache, font_size=9, font_color=THEME['text'], border_color=THEME['gray']
        )

    # Monthly cost summary
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.5), Inches(4.7), Inches(5.5), Inches(2.3),
        THEME['light'], border_color=THEME['accent'], border_width=Pt(2)
    )

    add_text_box(slide, Inches(7.7), Inches(4.8), Inches(5.1), Inches(0.4),
                 "월간 예상 비용 (1,000문서)", font_size=14, font_color=THEME['accent'], bold=True)

    add_text_box(slide, Inches(7.7), Inches(5.2), Inches(5.1), Inches(1.6),
                 "문서 처리: $2.26/월\n검색 쿼리 (1,000회): $0.50/월\n─────────────────\n총 예상 비용: $2.76/월\n연간 비용: $33.12/년\n\nvs Claude: $546/년\n절감액: $513/년 (94% 절감)",
                 font_size=11, font_color=THEME['text'])

    return slide


def create_do_dont_slide(prs):
    """Slide 8: Do/Don't Guidelines"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "Do / Don't Guidelines", "필수 준수사항 및 금지사항")

    # DO section
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.4), Inches(6.2), Inches(3.0),
        THEME['light'], border_color=THEME['primary'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.4),
                 "DO (필수 준수사항)", font_size=16, font_color=THEME['primary'], bold=True)

    do_items = [
        "LLM: DeepSeek-Chat/Reasoner만 사용",
        "문서 파싱: Docling 사용 (온프레미스)",
        "청킹: HybridChunker + 512 토큰 제한",
        "임베딩: FlagEmbedding BGE-M3",
        "Dense + Sparse 동시 생성",
        "PostgreSQL = SSOT (Single Source)",
        "ES 메타데이터 비정규화 (제로 조인)",
        "Neo4j Slim Graph (최소 속성)",
        "3개 DB 동시 저장 (asyncio.gather)",
    ]

    for i, item in enumerate(do_items):
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.95 + i * 0.27), Inches(0.25), Inches(0.22),
            THEME['primary'], "V", font_size=10, font_color=THEME['white'], bold=True
        )
        add_text_box(slide, Inches(0.85), Inches(1.95 + i * 0.27), Inches(5.5), Inches(0.25),
                     item, font_size=9, font_color=THEME['text'])

    # DON'T section
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.4), Inches(6.2), Inches(3.0),
        THEME['lightgray'], border_color=THEME['red'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.4),
                 "DON'T (금지사항)", font_size=16, font_color=THEME['red'], bold=True)

    dont_items = [
        "LlamaParse 클라우드 (데이터 외부 전송)",
        "PyPDF2/pdfminer 직접 사용 (구조 손실)",
        "OpenAI/Claude 모델 사용",
        "LangChain HuggingFaceEmbeddings",
        "PostgreSQL 조인으로 검색",
        "Neo4j에 텍스트 본문 저장",
        "ES 내장 RRF (Platinum 라이선스)",
        "ColBERT 벡터 생성 (메모리 과다)",
        "동기적 DB 저장 (성능 저하)",
    ]

    for i, item in enumerate(dont_items):
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.0), Inches(1.95 + i * 0.27), Inches(0.25), Inches(0.22),
            THEME['red'], "X", font_size=10, font_color=THEME['white'], bold=True
        )
        add_text_box(slide, Inches(7.35), Inches(1.95 + i * 0.27), Inches(5.5), Inches(0.25),
                     item, font_size=9, font_color=THEME['text'])

    # Performance checklist
    add_text_box(slide, Inches(0.3), Inches(4.6), Inches(6), Inches(0.4),
                 "성능 최적화 체크리스트", font_size=14, font_color=THEME['primary'], bold=True)

    checklist = [
        "Docling + HybridChunker 설정",
        "DeepSeek API 키 + Temperature 설정",
        "BGE-M3 배치 크기 8 이하",
        "PostgreSQL/ES/Neo4j 인덱스 생성",
        "메모리 사용량 85% 이하",
        "응답 시간 2초 이내",
    ]

    for i, item in enumerate(checklist):
        x = 0.3 + (i % 2) * 6.5
        y = 5.0 + (i // 2) * 0.4
        add_shape_with_text(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(0.25), Inches(0.25),
            THEME['white'], "", border_color=THEME['gray']
        )
        add_text_box(slide, Inches(x + 0.35), Inches(y), Inches(5.5), Inches(0.35),
                     item, font_size=10, font_color=THEME['text'])

    # Key warning
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.5),
        THEME['red'],
        "절대 금지: 민감 데이터 외부 전송 (LlamaParse 등 클라우드 서비스 사용 금지)",
        font_size=12, font_color=THEME['white'], bold=True
    )

    return slide


def create_deployment_slide(prs):
    """Slide 9: Deployment Guide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "Deployment & Testing", "배포 및 테스트 전략")

    # Docker deployment
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.4), Inches(6.2), Inches(2.6),
        THEME['white'], border_color=THEME['primary'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.4),
                 "Docker 기반 배포", font_size=14, font_color=THEME['primary'], bold=True)

    containers = [
        ("postgres:16", "5432", "마스터 DB"),
        ("elasticsearch:8.x", "9200", "벡터 검색"),
        ("neo4j:5.x", "7474/7687", "그래프 DB"),
        ("ai-service", "8000", "FastAPI"),
        ("backend", "8080", "SpringBoot"),
    ]

    for i, (name, port, desc) in enumerate(containers):
        y = 1.95 + i * 0.4
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y), Inches(2.2), Inches(0.35),
            THEME['secondary'], name, font_size=9, font_color=THEME['dark']
        )
        add_text_box(slide, Inches(2.8), Inches(y + 0.02), Inches(1.2), Inches(0.3),
                     f":{port}", font_size=9, font_color=THEME['gray'])
        add_text_box(slide, Inches(4.1), Inches(y + 0.02), Inches(2), Inches(0.3),
                     desc, font_size=9, font_color=THEME['text'])

    # Memory allocation
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.4), Inches(6.2), Inches(2.6),
        THEME['white'], border_color=THEME['accent'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.4),
                 "메모리 할당 (16GB 환경)", font_size=14, font_color=THEME['accent'], bold=True)

    memory = [
        ("Elasticsearch", "4GB", "JVM heap"),
        ("Neo4j", "2GB", "JVM heap"),
        ("PostgreSQL", "1GB", "shared_buffers"),
        ("AI Service", "4GB", "BGE-M3 모델"),
        ("Backend", "1GB", "Spring Boot"),
        ("예비", "4GB", "OS/캐시"),
    ]

    for i, (name, size, note) in enumerate(memory):
        y = 1.95 + i * 0.38
        add_text_box(slide, Inches(7.0), Inches(y), Inches(2.2), Inches(0.35),
                     name, font_size=10, font_color=THEME['text'])
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.3), Inches(y), Inches(0.8), Inches(0.3),
            THEME['secondary'], size, font_size=9, font_color=THEME['dark'], bold=True
        )
        add_text_box(slide, Inches(10.2), Inches(y), Inches(2.5), Inches(0.35),
                     note, font_size=9, font_color=THEME['gray'])

    # Test strategy
    add_text_box(slide, Inches(0.3), Inches(4.2), Inches(4), Inches(0.4),
                 "테스트 전략", font_size=14, font_color=THEME['primary'], bold=True)

    test_types = [
        ("Unit Test", "개별 함수/클래스\npytest, JUnit"),
        ("Integration", "서비스 간 연동\nTestcontainers"),
        ("E2E", "전체 플로우\nPlaywright"),
        ("Performance", "부하 테스트\nLocust"),
    ]

    for i, (name, desc) in enumerate(test_types):
        x = 0.3 + i * 3.2
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(4.6), Inches(3), Inches(1.2),
            THEME['light'], border_color=THEME['secondary']
        )
        add_text_box(slide, Inches(x + 0.1), Inches(4.65), Inches(2.8), Inches(0.35),
                     name, font_size=11, font_color=THEME['primary'], bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.1), Inches(5.0), Inches(2.8), Inches(0.7),
                     desc, font_size=9, font_color=THEME['text'], alignment=PP_ALIGN.CENTER)

    # CI/CD
    add_text_box(slide, Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.7),
                 "CI/CD: GitLab CI → Docker Build → Container Registry → K8s Deploy\n목표: 커버리지 80%+ | 빌드 시간 < 10분 | 무중단 배포",
                 font_size=11, font_color=THEME['gray'])

    return slide


def create_closing_slide(prs):
    """Slide 10: Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['dark']
    bg.line.color.rgb = THEME['dark']

    # Decorative
    deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(4.5), Inches(5), Inches(4))
    deco.fill.solid()
    deco.fill.fore_color.rgb = THEME['primary']
    deco.line.color.rgb = THEME['primary']

    # Title
    add_text_box(slide, Inches(0.5), Inches(2), Inches(12), Inches(1),
                 "Thank You", font_size=56, font_color=THEME['white'],
                 bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(12), Inches(0.7),
                 "Questions & Answers", font_size=28, font_color=THEME['light'],
                 alignment=PP_ALIGN.CENTER)

    # Key highlights
    highlights = [
        "Hybrid RAG = Vector Search + Graph Search",
        "VIP 3-Stage LLM Architecture (DeepSeek V3.2)",
        "95% 비용 절감 ($45.50 -> $2.26/1,000문서)",
        "제로 조인 아키텍처 (ES 메타데이터 비정규화)",
        "서비스 분리: SpringBoot (비즈니스) + AI Service (Python)"
    ]

    for i, highlight in enumerate(highlights):
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(4.3 + i * 0.45), Inches(9), Inches(0.4),
            THEME['primary'], highlight, font_size=11, font_color=THEME['white']
        )

    # Footer
    add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                 "Neo4j Graph RAG Hybrid Knowledge Platform v2.3 | 2026-01-14",
                 font_size=14, font_color=THEME['secondary'], alignment=PP_ALIGN.CENTER)

    return slide


def main():
    """Generate presentation"""
    print("=" * 60)
    print("Hybrid RAG Platform Detailed Design Presentation Generator")
    print("Theme: Tech Innovation")
    print("=" * 60)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("\n[1/10] Creating title slide...")
    create_title_slide(prs)

    print("[2/10] Creating overview slide...")
    create_overview_slide(prs)

    print("[3/10] Creating tech stack slide...")
    create_tech_stack_slide(prs)

    print("[4/10] Creating VIP architecture slide...")
    create_vip_architecture_slide(prs)

    print("[5/10] Creating data model slide...")
    create_data_model_slide(prs)

    print("[6/10] Creating hybrid search slide...")
    create_hybrid_search_slide(prs)

    print("[7/10] Creating cost analysis slide...")
    create_cost_analysis_slide(prs)

    print("[8/10] Creating do/don't slide...")
    create_do_dont_slide(prs)

    print("[9/10] Creating deployment slide...")
    create_deployment_slide(prs)

    print("[10/10] Creating closing slide...")
    create_closing_slide(prs)

    # Save
    output_dir = "knowledge_service/docs/02_design"
    os.makedirs(output_dir, exist_ok=True)

    output_path = os.path.join(output_dir, "Hybrid_RAG_Platform_Design_v2.3.pptx")
    prs.save(output_path)

    print("\n" + "=" * 60)
    print(f"Presentation saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"Slide size: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches")
    print("=" * 60)


if __name__ == "__main__":
    main()
