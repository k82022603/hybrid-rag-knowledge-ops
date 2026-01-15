#!/usr/bin/env python3
"""
UI Storyboard 프레젠테이션 생성
Brown Earth Theme 적용
각 페이지마다 스토리보드 상세 설명 및 개발자 주의사항 포함
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# ========== Brown Earth Theme ==========
BROWN_EARTH = {
    'primary': RGBColor(121, 85, 72),      # #795548 Brown
    'secondary': RGBColor(161, 136, 127),  # #A1887F Light Brown
    'accent': RGBColor(255, 152, 0),       # #FF9800 Orange
    'dark': RGBColor(62, 39, 35),          # #3E2723 Dark Brown
    'light': RGBColor(239, 235, 233),      # #EFEBE9 Light Beige
    'text': RGBColor(33, 33, 33),          # #212121 Dark Text
    'white': RGBColor(255, 255, 255),
    # UI Component colors
    'header': RGBColor(63, 81, 181),       # Indigo
    'sidebar': RGBColor(55, 71, 79),       # Blue Grey
    'content': RGBColor(250, 250, 250),    # Light Grey
    'success': RGBColor(76, 175, 80),      # Green
    'warning': RGBColor(255, 193, 7),      # Amber
    'error': RGBColor(244, 67, 54),        # Red
    'info': RGBColor(33, 150, 243),        # Blue
}

def create_title_slide(prs, title, subtitle, theme):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme['dark']
    bg.line.fill.background()

    # 장식
    deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6), Inches(4.5), Inches(5), Inches(4))
    deco.fill.solid()
    deco.fill.fore_color.rgb = theme['primary']
    deco.line.fill.background()

    # 상단 액센트 바
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = theme['accent']
    accent_bar.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = theme['white']
    p.font.name = "맑은 고딕"

    # 부제목
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(7), Inches(0.8))
    sf = sub_box.text_frame
    sf.word_wrap = True
    sp = sf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(22)
    sp.font.color.rgb = theme['light']
    sp.font.name = "맑은 고딕"

    # 날짜
    date_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(4), Inches(0.5))
    df = date_box.text_frame
    dp = df.paragraphs[0]
    dp.text = "2026-01-15 | Knowledge Hub UI Design v1.0"
    dp.font.size = Pt(14)
    dp.font.color.rgb = theme['secondary']
    dp.font.name = "맑은 고딕"

    return slide

def add_slide_header(slide, title, theme):
    """슬라이드 헤더"""
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = theme['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = theme['white']
    p.font.name = "맑은 고딕"

def add_description_box(slide, x, y, width, height, title, items, theme, bg_color=None):
    """설명 박스 추가"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color or theme['light']
    box.line.color.rgb = theme['secondary']
    box.line.width = Pt(1)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1), Inches(width - 0.2), Inches(0.4))
    ttf = title_box.text_frame
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(12)
    tp.font.bold = True
    tp.font.color.rgb = theme['primary']
    tp.font.name = "맑은 고딕"

    # 항목들
    content_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.45), Inches(width - 0.2), Inches(height - 0.55))
    ctf = content_box.text_frame
    ctf.word_wrap = True

    for item in items:
        p = ctf.add_paragraph()
        p.text = item
        p.font.size = Pt(9)
        p.font.color.rgb = theme['text']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(3)

def create_login_slide(prs, theme):
    """로그인 화면 스토리보드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "1. 로그인 화면 스토리보드", theme)

    # 왼쪽: 와이어프레임
    login_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.2), Inches(5.2), Inches(5.5))
    login_bg.fill.solid()
    login_bg.fill.fore_color.rgb = theme['light']
    login_bg.line.color.rgb = theme['secondary']

    # 로고 영역
    logo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(1.8), Inches(2.4), Inches(1))
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = theme['primary']
    logo_box.line.fill.background()

    logo_text = slide.shapes.add_textbox(Inches(1.8), Inches(2.0), Inches(2.4), Inches(0.6))
    ltf = logo_text.text_frame
    lp = ltf.paragraphs[0]
    lp.text = "📚 Knowledge Hub"
    lp.font.size = Pt(16)
    lp.font.bold = True
    lp.font.color.rgb = theme['white']
    lp.font.name = "맑은 고딕"
    lp.alignment = PP_ALIGN.CENTER

    # 부제목
    sub_text = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(4), Inches(0.4))
    stf = sub_text.text_frame
    sp = stf.paragraphs[0]
    sp.text = "사내 지식 검색 시스템"
    sp.font.size = Pt(14)
    sp.font.color.rgb = theme['text']
    sp.font.name = "맑은 고딕"
    sp.alignment = PP_ALIGN.CENTER

    # SSO 버튼
    sso_btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3.8), Inches(3), Inches(0.7))
    sso_btn.fill.solid()
    sso_btn.fill.fore_color.rgb = theme['accent']
    sso_btn.line.fill.background()

    sso_text = slide.shapes.add_textbox(Inches(1.5), Inches(3.95), Inches(3), Inches(0.5))
    btf = sso_text.text_frame
    bp = btf.paragraphs[0]
    bp.text = "🔐 SSO 로그인"
    bp.font.size = Pt(14)
    bp.font.bold = True
    bp.font.color.rgb = theme['white']
    bp.font.name = "맑은 고딕"
    bp.alignment = PP_ALIGN.CENTER

    # OAuth 설명
    oauth_text = slide.shapes.add_textbox(Inches(1), Inches(4.7), Inches(4), Inches(0.3))
    otf = oauth_text.text_frame
    op = otf.paragraphs[0]
    op.text = "안전한 OAuth 2.0 + PKCE 인증"
    op.font.size = Pt(10)
    op.font.color.rgb = theme['secondary']
    op.font.name = "맑은 고딕"
    op.alignment = PP_ALIGN.CENTER

    # 푸터
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(5.5), Inches(5.2), Inches(0.5))
    footer.fill.solid()
    footer.fill.fore_color.rgb = theme['secondary']
    footer.line.fill.background()

    footer_text = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(5), Inches(0.3))
    ftf = footer_text.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "© 2026 Company | 도움말 | 개인정보처리방침"
    fp.font.size = Pt(9)
    fp.font.color.rgb = theme['white']
    fp.font.name = "맑은 고딕"
    fp.alignment = PP_ALIGN.CENTER

    # 상태 표시
    states_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(5), Inches(0.5))
    states_tf = states_text.text_frame
    states_p = states_tf.paragraphs[0]
    states_p.text = "상태: 기본 | 로딩중 ⏳ | 에러 ⚠️"
    states_p.font.size = Pt(10)
    states_p.font.color.rgb = theme['text']
    states_p.font.name = "맑은 고딕"

    # 오른쪽: 스토리보드 상세 설명
    add_description_box(slide, 5.8, 1.2, 3.8, 2.4, "📋 스토리보드 상세", [
        "• 화면 목적: 사용자 인증 진입점",
        "• 인증 방식: OAuth 2.0 + PKCE (SSO)",
        "• 주요 컴포넌트:",
        "  - 로고 및 서비스명",
        "  - SSO 로그인 버튼 (CTA)",
        "  - 푸터 링크 (도움말, 정책)",
        "• 상태별 화면: 기본/로딩/에러",
        "• 에러 시 '다시 시도' 버튼 표시"
    ], theme)

    add_description_box(slide, 5.8, 3.8, 3.8, 2.9, "⚠️ 개발자 주의사항", [
        "🔐 보안:",
        "• PKCE code_verifier 클라이언트만 저장",
        "• Redirect URI 화이트리스트 검증",
        "• state 파라미터로 CSRF 방어",
        "",
        "🎨 UI/UX:",
        "• 로딩 중 버튼 비활성화 + 스피너",
        "• 에러 메시지 사용자 친화적 표현",
        "• 키보드 Enter로 로그인 가능",
        "",
        "📱 반응형: 모바일 시 전체화면"
    ], theme, theme['white'])

    return slide

def create_dashboard_slide(prs, theme):
    """대시보드 화면 스토리보드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "2. 대시보드 스토리보드", theme)

    # 왼쪽: 와이어프레임
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.1), Inches(5.2), Inches(0.5))
    header.fill.solid()
    header.fill.fore_color.rgb = theme['header']
    header.line.fill.background()

    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(0.3))
    htf = header_text.text_frame
    hp = htf.paragraphs[0]
    hp.text = "📚 LOGO    🔍 검색...               🔔 👤"
    hp.font.size = Pt(10)
    hp.font.color.rgb = theme['white']
    hp.font.name = "맑은 고딕"

    # Sidebar
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.6), Inches(1), Inches(5))
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = theme['sidebar']
    sidebar.line.fill.background()

    sidebar_text = slide.shapes.add_textbox(Inches(0.45), Inches(1.8), Inches(0.9), Inches(4.5))
    stf = sidebar_text.text_frame
    stf.word_wrap = True
    menus = ["📊 대시보드", "🔍 검색", "📚 지식관리", "⭐ 북마크", "👤 프로필", "⚙️ 설정"]
    for menu in menus:
        p = stf.add_paragraph()
        p.text = menu
        p.font.size = Pt(8)
        p.font.color.rgb = theme['white']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(8)

    # Main Content Area
    content = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.4), Inches(1.6), Inches(4.2), Inches(5))
    content.fill.solid()
    content.fill.fore_color.rgb = theme['content']
    content.line.color.rgb = theme['secondary']

    # 통계 카드 4개
    stats = [("📚", "1,234"), ("📝", "12"), ("📖", "45"), ("⭐", "23")]
    for i, (icon, num) in enumerate(stats):
        x = 1.5 + i * 1.0
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(0.9), Inches(0.7))
        card.fill.solid()
        card.fill.fore_color.rgb = theme['white']
        card.line.color.rgb = theme['secondary']

        card_text = slide.shapes.add_textbox(Inches(x), Inches(1.9), Inches(0.9), Inches(0.5))
        ctf = card_text.text_frame
        cp = ctf.paragraphs[0]
        cp.text = f"{icon}\n{num}"
        cp.font.size = Pt(9)
        cp.font.color.rgb = theme['text']
        cp.font.name = "맑은 고딕"
        cp.alignment = PP_ALIGN.CENTER

    # 위젯 영역
    widget1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.7), Inches(1.95), Inches(1.8))
    widget1.fill.solid()
    widget1.fill.fore_color.rgb = theme['white']
    widget1.line.color.rgb = theme['secondary']

    w1_text = slide.shapes.add_textbox(Inches(1.55), Inches(2.8), Inches(1.85), Inches(0.3))
    w1tf = w1_text.text_frame
    w1p = w1tf.paragraphs[0]
    w1p.text = "🔥 인기 지식 TOP 5"
    w1p.font.size = Pt(9)
    w1p.font.bold = True
    w1p.font.color.rgb = theme['text']
    w1p.font.name = "맑은 고딕"

    widget2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.55), Inches(2.7), Inches(1.95), Inches(1.8))
    widget2.fill.solid()
    widget2.fill.fore_color.rgb = theme['white']
    widget2.line.color.rgb = theme['secondary']

    w2_text = slide.shapes.add_textbox(Inches(3.6), Inches(2.8), Inches(1.85), Inches(0.3))
    w2tf = w2_text.text_frame
    w2p = w2tf.paragraphs[0]
    w2p.text = "📝 최신 지식"
    w2p.font.size = Pt(9)
    w2p.font.bold = True
    w2p.font.color.rgb = theme['text']
    w2p.font.name = "맑은 고딕"

    widget3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(4.6), Inches(1.95), Inches(1.5))
    widget3.fill.solid()
    widget3.fill.fore_color.rgb = theme['white']
    widget3.line.color.rgb = theme['secondary']

    w3_text = slide.shapes.add_textbox(Inches(1.55), Inches(4.7), Inches(1.85), Inches(0.3))
    w3tf = w3_text.text_frame
    w3p = w3tf.paragraphs[0]
    w3p.text = "💡 AI 추천"
    w3p.font.size = Pt(9)
    w3p.font.bold = True
    w3p.font.color.rgb = theme['text']
    w3p.font.name = "맑은 고딕"

    widget4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.55), Inches(4.6), Inches(1.95), Inches(1.5))
    widget4.fill.solid()
    widget4.fill.fore_color.rgb = theme['white']
    widget4.line.color.rgb = theme['secondary']

    w4_text = slide.shapes.add_textbox(Inches(3.6), Inches(4.7), Inches(1.85), Inches(0.3))
    w4tf = w4_text.text_frame
    w4p = w4tf.paragraphs[0]
    w4p.text = "🔎 검색 트렌드"
    w4p.font.size = Pt(9)
    w4p.font.bold = True
    w4p.font.color.rgb = theme['text']
    w4p.font.name = "맑은 고딕"

    # 오른쪽: 스토리보드 상세 설명
    add_description_box(slide, 5.8, 1.1, 3.8, 2.8, "📋 스토리보드 상세", [
        "• 레이아웃: Header(64px) + Sidebar(240px) + Content",
        "• 통계 카드 4종:",
        "  - 전체 지식 수 (증감률 표시)",
        "  - 오늘 등록 지식",
        "  - 내 지식 수",
        "  - 북마크 수",
        "• 위젯 4종:",
        "  - 인기 지식 TOP 5 (조회수 기준)",
        "  - 최신 지식 목록",
        "  - AI 추천 (관심사 기반)",
        "  - 검색 트렌드 (워드 클라우드)"
    ], theme)

    add_description_box(slide, 5.8, 4.1, 3.8, 2.6, "⚠️ 개발자 주의사항", [
        "📡 API:",
        "• GET /api/v1/dashboard/stats",
        "• 통계 데이터 5분 캐싱 권장",
        "• 인기 지식 실시간 조회수 반영",
        "",
        "🎨 UI/UX:",
        "• Skeleton 로딩 UI 적용",
        "• 통계 카드 클릭 → 상세 목록 이동",
        "• 위젯 더보기 → 해당 페이지 이동",
        "",
        "📱 반응형: Tablet 2열, Mobile 1열"
    ], theme, theme['white'])

    return slide

def create_search_chat_slide(prs, theme):
    """검색 - 채팅 모드 스토리보드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "3. 검색 - 채팅 모드 스토리보드", theme)

    # 왼쪽: 와이어프레임
    chat_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.1), Inches(5.2), Inches(5.5))
    chat_bg.fill.solid()
    chat_bg.fill.fore_color.rgb = theme['content']
    chat_bg.line.color.rgb = theme['secondary']

    # 탭
    tab1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.3), Inches(1.5), Inches(0.4))
    tab1.fill.solid()
    tab1.fill.fore_color.rgb = theme['primary']
    tab1.line.fill.background()

    tab1_text = slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(1.5), Inches(0.3))
    t1tf = tab1_text.text_frame
    t1p = t1tf.paragraphs[0]
    t1p.text = "💬 채팅 모드"
    t1p.font.size = Pt(10)
    t1p.font.bold = True
    t1p.font.color.rgb = theme['white']
    t1p.font.name = "맑은 고딕"
    t1p.alignment = PP_ALIGN.CENTER

    tab2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.2), Inches(1.3), Inches(1.5), Inches(0.4))
    tab2.fill.solid()
    tab2.fill.fore_color.rgb = theme['light']
    tab2.line.color.rgb = theme['secondary']

    tab2_text = slide.shapes.add_textbox(Inches(2.2), Inches(1.35), Inches(1.5), Inches(0.3))
    t2tf = tab2_text.text_frame
    t2p = t2tf.paragraphs[0]
    t2p.text = "🔍 키워드 검색"
    t2p.font.size = Pt(10)
    t2p.font.color.rgb = theme['text']
    t2p.font.name = "맑은 고딕"
    t2p.alignment = PP_ALIGN.CENTER

    # 채팅 영역
    chat_area = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.9), Inches(4.8), Inches(4))
    chat_area.fill.solid()
    chat_area.fill.fore_color.rgb = theme['white']
    chat_area.line.color.rgb = theme['secondary']

    # 사용자 메시지
    user_msg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(2.1), Inches(2.7), Inches(0.5))
    user_msg.fill.solid()
    user_msg.fill.fore_color.rgb = theme['info']
    user_msg.line.fill.background()

    user_text = slide.shapes.add_textbox(Inches(2.5), Inches(2.2), Inches(2.7), Inches(0.3))
    utf = user_text.text_frame
    up = utf.paragraphs[0]
    up.text = "👤 React 설계 방법 알려주세요"
    up.font.size = Pt(9)
    up.font.color.rgb = theme['white']
    up.font.name = "맑은 고딕"

    # AI 응답
    ai_msg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.8), Inches(3.5), Inches(2))
    ai_msg.fill.solid()
    ai_msg.fill.fore_color.rgb = theme['light']
    ai_msg.line.color.rgb = theme['secondary']

    ai_text = slide.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(3.3), Inches(1.8))
    atf = ai_text.text_frame
    atf.word_wrap = True
    ap = atf.paragraphs[0]
    ap.text = "🤖 React 컴포넌트 설계 가이드\n\n• Presentational Components\n• Container Components\n\n📚 참조: React 설계 가이드\n    관련도: 95%"
    ap.font.size = Pt(8)
    ap.font.color.rgb = theme['text']
    ap.font.name = "맑은 고딕"

    # 피드백 버튼
    feedback_text = slide.shapes.add_textbox(Inches(0.9), Inches(4.7), Inches(3), Inches(0.2))
    ftf = feedback_text.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "👍 👎  📋 복사  🔗 공유"
    fp.font.size = Pt(8)
    fp.font.color.rgb = theme['secondary']
    fp.font.name = "맑은 고딕"

    # 입력창
    input_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.0), Inches(4.8), Inches(0.5))
    input_box.fill.solid()
    input_box.fill.fore_color.rgb = theme['white']
    input_box.line.color.rgb = theme['primary']

    input_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.1), Inches(4.5), Inches(0.3))
    itf = input_text.text_frame
    ip = itf.paragraphs[0]
    ip.text = "💬 질문을 입력하세요...                    📎 ➤"
    ip.font.size = Pt(10)
    ip.font.color.rgb = theme['secondary']
    ip.font.name = "맑은 고딕"

    # 오른쪽: 스토리보드 상세 설명
    add_description_box(slide, 5.8, 1.1, 3.8, 2.6, "📋 스토리보드 상세", [
        "• 검색 모드: 채팅 모드 (RAG 기반)",
        "• 사용 케이스: 복잡한 질문, 맥락 이해 필요",
        "• 메시지 구성:",
        "  - 사용자 메시지 (우측 정렬, 파란색)",
        "  - AI 응답 (좌측 정렬, 마크다운 지원)",
        "  - 참조 문서 카드 (클릭 시 상세 이동)",
        "• 피드백: 👍👎 버튼으로 답변 품질 평가",
        "• 하단: 새 대화, 대화 내역 버튼"
    ], theme)

    add_description_box(slide, 5.8, 3.9, 3.8, 2.8, "⚠️ 개발자 주의사항", [
        "🤖 RAG 연동:",
        "• POST /api/v1/chat/message",
        "• 스트리밍 응답 처리 (SSE)",
        "• 참조 문서 관련도 점수 표시",
        "",
        "🎨 UI/UX:",
        "• 타이핑 인디케이터 애니메이션",
        "• 마크다운 렌더링 (코드 하이라이트)",
        "• 응답 생성 중 입력 비활성화",
        "",
        "💾 대화 내역 localStorage 저장",
        "📱 모바일: 전체화면 채팅"
    ], theme, theme['white'])

    return slide

def create_search_keyword_slide(prs, theme):
    """검색 - 키워드 검색 스토리보드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "4. 검색 - 키워드 검색 스토리보드", theme)

    # 왼쪽: 와이어프레임
    search_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.1), Inches(5.2), Inches(5.5))
    search_bg.fill.solid()
    search_bg.fill.fore_color.rgb = theme['content']
    search_bg.line.color.rgb = theme['secondary']

    # 검색창
    search_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.3), Inches(4.8), Inches(0.45))
    search_box.fill.solid()
    search_box.fill.fore_color.rgb = theme['white']
    search_box.line.color.rgb = theme['primary']
    search_box.line.width = Pt(2)

    search_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.38), Inches(4.5), Inches(0.3))
    stf = search_text.text_frame
    sp = stf.paragraphs[0]
    sp.text = "🔍 React                                    🔍"
    sp.font.size = Pt(11)
    sp.font.color.rgb = theme['text']
    sp.font.name = "맑은 고딕"

    # 필터
    filter_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.85), Inches(4.8), Inches(0.4))
    filter_box.fill.solid()
    filter_box.fill.fore_color.rgb = theme['light']
    filter_box.line.fill.background()

    filter_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.92), Inches(4.5), Inches(0.25))
    ftf = filter_text.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "📁 카테고리 ▼  │  📅 기간 ▼  │  🏷️ 태그 ▼  │  ⚙️ 더보기"
    fp.font.size = Pt(9)
    fp.font.color.rgb = theme['text']
    fp.font.name = "맑은 고딕"

    # 결과 헤더
    result_header = slide.shapes.add_textbox(Inches(0.6), Inches(2.35), Inches(4.8), Inches(0.25))
    rhtf = result_header.text_frame
    rhp = rhtf.paragraphs[0]
    rhp.text = '"React" 검색 결과 (총 156건)             정렬: 관련도순 ▼'
    rhp.font.size = Pt(9)
    rhp.font.color.rgb = theme['text']
    rhp.font.name = "맑은 고딕"

    # 검색 결과 카드 3개
    for i in range(3):
        y = 2.7 + i * 1.3
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y), Inches(4.8), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = theme['white']
        card.line.color.rgb = theme['secondary']

        titles = ["📄 React 컴포넌트 설계 가이드", "📄 React Hooks 완벽 가이드", "📄 React + TypeScript 설정"]
        card_title = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.1), Inches(4.5), Inches(0.25))
        cttf = card_title.text_frame
        ctp = cttf.paragraphs[0]
        ctp.text = titles[i] + "                    ⭐ ⋮"
        ctp.font.size = Pt(10)
        ctp.font.bold = True
        ctp.font.color.rgb = theme['text']
        ctp.font.name = "맑은 고딕"

        card_desc = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.4), Inches(4.5), Inches(0.3))
        cdtf = card_desc.text_frame
        cdp = cdtf.paragraphs[0]
        cdp.text = "검색어가 포함된 본문 미리보기..."
        cdp.font.size = Pt(8)
        cdp.font.color.rgb = theme['secondary']
        cdp.font.name = "맑은 고딕"

        card_meta = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.7), Inches(4.5), Inches(0.4))
        cmtf = card_meta.text_frame
        cmp = cmtf.paragraphs[0]
        cmp.text = "📁 프론트엔드 │ 👤 김개발 │ 📅 2026-01-10\n🏷️ React  컴포넌트    👁 234  💬 12"
        cmp.font.size = Pt(7)
        cmp.font.color.rgb = theme['secondary']
        cmp.font.name = "맑은 고딕"

    # 페이지네이션
    page_text = slide.shapes.add_textbox(Inches(1.5), Inches(6.4), Inches(3), Inches(0.25))
    ptf = page_text.text_frame
    pp = ptf.paragraphs[0]
    pp.text = "< 1  2  3  4  5 ... 16 >"
    pp.font.size = Pt(10)
    pp.font.color.rgb = theme['primary']
    pp.font.name = "맑은 고딕"
    pp.alignment = PP_ALIGN.CENTER

    # 오른쪽: 스토리보드 상세 설명
    add_description_box(slide, 5.8, 1.1, 3.8, 2.6, "📋 스토리보드 상세", [
        "• 검색 모드: 키워드 검색 (전통적 방식)",
        "• 필터 옵션:",
        "  - 카테고리 (프론트엔드, 백엔드...)",
        "  - 기간 (전체, 7일, 30일, 직접입력)",
        "  - 태그 (다중 선택 가능)",
        "  - 고급 필터 (문서유형, 작성자, 조회수)",
        "• 정렬: 관련도순, 최신순, 조회수순",
        "• 결과 카드: 제목, 요약, 메타정보, 태그",
        "• 검색어 하이라이트 처리"
    ], theme)

    add_description_box(slide, 5.8, 3.9, 3.8, 2.8, "⚠️ 개발자 주의사항", [
        "🔍 검색 기능:",
        "• GET /api/v1/search?q=&filters=",
        "• Elasticsearch Hybrid Search 연동",
        "• 자동완성: 최근검색 + 인기검색어",
        "• 검색어 디바운싱 (300ms)",
        "",
        "🎨 UI/UX:",
        "• 결과 없음 시 추천 키워드 표시",
        "• Skeleton 로딩 카드 표시",
        "• 무한 스크롤 또는 페이지네이션",
        "",
        "📱 모바일: 필터 모달 처리"
    ], theme, theme['white'])

    return slide

def create_knowledge_slide(prs, theme):
    """지식 관리 스토리보드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "5. 지식 관리 스토리보드", theme)

    # 왼쪽: 와이어프레임 - 목록
    list_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.1), Inches(5.2), Inches(5.5))
    list_bg.fill.solid()
    list_bg.fill.fore_color.rgb = theme['content']
    list_bg.line.color.rgb = theme['secondary']

    # 탭 + 버튼
    tabs_area = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.9), Inches(0.3))
    tatf = tabs_area.text_frame
    tap = tatf.paragraphs[0]
    tap.text = "전체(1,234) │ 내지식(45) │ 북마크(23) │ 임시(3)    ➕새 지식"
    tap.font.size = Pt(8)
    tap.font.color.rgb = theme['text']
    tap.font.name = "맑은 고딕"

    # 필터 + 검색
    filter_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.55), Inches(4.9), Inches(0.35))
    filter_bar.fill.solid()
    filter_bar.fill.fore_color.rgb = theme['light']
    filter_bar.line.fill.background()

    filter_text = slide.shapes.add_textbox(Inches(0.55), Inches(1.6), Inches(4.8), Inches(0.25))
    ftf = filter_text.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "🔍 검색...     │  📁 카테고리 ▼  │  📅 기간 ▼  │  정렬 ▼"
    fp.font.size = Pt(8)
    fp.font.color.rgb = theme['text']
    fp.font.name = "맑은 고딕"

    # 목록 아이템
    items = [
        "📄 React 컴포넌트 설계 가이드  │  김개발  │  01-10",
        "📄 Spring Boot 가이드  │  이스프링  │  01-09",
        "📄 코드 리뷰 체크리스트  │  박리뷰  │  01-08",
        "📄 배포 프로세스 안내  │  최배포  │  01-07",
    ]

    for i, item in enumerate(items):
        y = 2.0 + i * 0.85
        item_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(4.9), Inches(0.75))
        item_box.fill.solid()
        item_box.fill.fore_color.rgb = theme['white']
        item_box.line.color.rgb = theme['secondary']

        item_text = slide.shapes.add_textbox(Inches(0.55), Inches(y + 0.1), Inches(4.8), Inches(0.55))
        itf = item_text.text_frame
        itf.word_wrap = True
        ip = itf.paragraphs[0]
        ip.text = f"☐ {item}\n    요약 텍스트...  👁 234  💬 12        ⭐ ⋮"
        ip.font.size = Pt(8)
        ip.font.color.rgb = theme['text']
        ip.font.name = "맑은 고딕"

    # 하단 액션
    action_text = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(4.9), Inches(0.25))
    atf = action_text.text_frame
    ap = atf.paragraphs[0]
    ap.text = "☐ 전체선택    선택: 0개    [🗑️ 삭제]          < 1 2 3 ... >"
    ap.font.size = Pt(8)
    ap.font.color.rgb = theme['text']
    ap.font.name = "맑은 고딕"

    # 에디터 미니 프리뷰
    editor_label = slide.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(2), Inches(0.25))
    eltf = editor_label.text_frame
    elp = eltf.paragraphs[0]
    elp.text = "📝 에디터 주요 기능:"
    elp.font.size = Pt(9)
    elp.font.bold = True
    elp.font.color.rgb = theme['primary']
    elp.font.name = "맑은 고딕"

    editor_features = slide.shapes.add_textbox(Inches(0.5), Inches(6.15), Inches(5), Inches(0.5))
    eftf = editor_features.text_frame
    efp = eftf.paragraphs[0]
    efp.text = "마크다운 WYSIWYG │ 실시간 미리보기 │ 파일 드래그 업로드 │ 자동 임시저장"
    efp.font.size = Pt(8)
    efp.font.color.rgb = theme['secondary']
    efp.font.name = "맑은 고딕"

    # 오른쪽: 스토리보드 상세 설명
    add_description_box(slide, 5.8, 1.1, 3.8, 2.5, "📋 스토리보드 상세", [
        "• 탭: 전체, 내 지식, 북마크, 임시저장",
        "• 목록 기능:",
        "  - 체크박스 다중 선택",
        "  - 카테고리/기간 필터",
        "  - 정렬 (최신순, 조회순)",
        "• 더보기 메뉴: 수정, 복사, 공유, 삭제",
        "• 에디터:",
        "  - 마크다운 WYSIWYG",
        "  - 실시간 미리보기",
        "  - 코드 하이라이팅"
    ], theme)

    add_description_box(slide, 5.8, 3.8, 3.8, 2.9, "⚠️ 개발자 주의사항", [
        "📡 API:",
        "• GET /api/v1/knowledge (목록)",
        "• POST /api/v1/knowledge (생성)",
        "• PUT /api/v1/knowledge/:id (수정)",
        "• DELETE /api/v1/knowledge/:id",
        "",
        "📝 에디터:",
        "• Tiptap 또는 Toast UI Editor",
        "• 이미지 S3 업로드 연동",
        "• XSS 방지: sanitize-html",
        "",
        "🔐 권한: 본인 지식만 수정/삭제"
    ], theme, theme['white'])

    return slide

def create_profile_admin_slide(prs, theme):
    """프로필 & 관리자 스토리보드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "6. 프로필 & 관리자 스토리보드", theme)

    # 왼쪽 상단: 프로필
    profile_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.1), Inches(2.5), Inches(2.8))
    profile_bg.fill.solid()
    profile_bg.fill.fore_color.rgb = theme['white']
    profile_bg.line.color.rgb = theme['secondary']

    profile_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2.3), Inches(0.25))
    ptf = profile_title.text_frame
    pp = ptf.paragraphs[0]
    pp.text = "👤 프로필"
    pp.font.size = Pt(10)
    pp.font.bold = True
    pp.font.color.rgb = theme['primary']
    pp.font.name = "맑은 고딕"

    profile_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2.3), Inches(2.2))
    pctf = profile_content.text_frame
    pctf.word_wrap = True
    pcp = pctf.paragraphs[0]
    pcp.text = "┌───┐\n│ 👤 │  김개발\n└───┘  Frontend Dev\n       개발1팀\n\n📊 활동 통계\n📝 45  👁 8,234\n💬 156  ⭐ 234\n\n[✏️ 프로필 수정]"
    pcp.font.size = Pt(8)
    pcp.font.color.rgb = theme['text']
    pcp.font.name = "맑은 고딕"

    # 왼쪽 하단: 설정
    settings_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(4.1), Inches(2.5), Inches(2.5))
    settings_bg.fill.solid()
    settings_bg.fill.fore_color.rgb = theme['white']
    settings_bg.line.color.rgb = theme['secondary']

    settings_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(2.3), Inches(0.25))
    stf = settings_title.text_frame
    sp = stf.paragraphs[0]
    sp.text = "⚙️ 설정"
    sp.font.size = Pt(10)
    sp.font.bold = True
    sp.font.color.rgb = theme['primary']
    sp.font.name = "맑은 고딕"

    settings_content = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(2.3), Inches(1.9))
    sctf = settings_content.text_frame
    sctf.word_wrap = True
    scp = sctf.paragraphs[0]
    scp.text = "🔔 알림 설정\n  ☑️ 이메일 알림\n  ☑️ 푸시 알림\n\n🎨 화면 설정\n  ○ 라이트  ● 다크\n\n🔒 보안 설정\n  ☑️ 2FA 활성화"
    scp.font.size = Pt(8)
    scp.font.color.rgb = theme['text']
    scp.font.name = "맑은 고딕"

    # 오른쪽 상단: 관리자
    admin_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.1), Inches(1.1), Inches(2.5), Inches(5.5))
    admin_bg.fill.solid()
    admin_bg.fill.fore_color.rgb = theme['dark']
    admin_bg.line.fill.background()

    admin_title = slide.shapes.add_textbox(Inches(3.2), Inches(1.2), Inches(2.3), Inches(0.25))
    atf = admin_title.text_frame
    ap = atf.paragraphs[0]
    ap.text = "⚙️ 관리자 (ADMIN)"
    ap.font.size = Pt(10)
    ap.font.bold = True
    ap.font.color.rgb = theme['accent']
    ap.font.name = "맑은 고딕"

    admin_content = slide.shapes.add_textbox(Inches(3.2), Inches(1.5), Inches(2.3), Inches(4.9))
    actf = admin_content.text_frame
    actf.word_wrap = True
    acp = actf.paragraphs[0]
    acp.text = """📊 시스템 현황
┌────┐ ┌────┐
│👥256│ │📚1234│
└────┘ └────┘
┌────┐ ┌────┐
│🔍567│ │💬 89│
└────┘ └────┘

📈 트래픽 그래프
[최근 7일 차트]

⚠️ 시스템 알림
• 저장공간 80%
• 백업 완료

👥 사용자 관리
📚 지식 관리
📁 카테고리 관리
⚙️ 시스템 설정"""
    acp.font.size = Pt(7)
    acp.font.color.rgb = theme['white']
    acp.font.name = "맑은 고딕"

    # 오른쪽: 스토리보드 상세 설명
    add_description_box(slide, 5.8, 1.1, 3.8, 2.4, "📋 스토리보드 상세", [
        "• 프로필: 사진, 정보, 통계, 활동 그래프",
        "• 설정 메뉴:",
        "  - 알림 (이메일, 푸시)",
        "  - 화면 (테마, 폰트)",
        "  - 보안 (2FA, 세션)",
        "  - 데이터 (내보내기)",
        "• 관리자 전용:",
        "  - 사용자 관리 (역할, 상태)",
        "  - 지식 관리 (승인, 반려)",
        "  - 카테고리 관리"
    ], theme)

    add_description_box(slide, 5.8, 3.7, 3.8, 3.0, "⚠️ 개발자 주의사항", [
        "🔐 권한 체계:",
        "• USER: 기본 기능",
        "• KNOWLEDGE_MANAGER: 지식 승인",
        "• ADMIN: 전체 관리",
        "",
        "📡 API:",
        "• GET/PUT /api/v1/users/me",
        "• GET /api/v1/admin/stats",
        "• PUT /api/v1/admin/users/:id/role",
        "",
        "🎨 UI/UX:",
        "• 권한 부족 시 403 페이지",
        "• 감사 로그 기록 필수"
    ], theme, theme['white'])

    return slide

def create_developer_notes_slide(prs, theme):
    """웹코더 & 개발자 가이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "7. 웹코더 & 개발자 구현 가이드", theme)

    # 3열 구조
    # 열 1: 기술 스택
    col1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.1), Inches(3.1), Inches(2.8))
    col1.fill.solid()
    col1.fill.fore_color.rgb = theme['white']
    col1.line.color.rgb = theme['secondary']

    col1_title = slide.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(2.9), Inches(0.3))
    c1tf = col1_title.text_frame
    c1p = c1tf.paragraphs[0]
    c1p.text = "🛠️ 기술 스택"
    c1p.font.size = Pt(12)
    c1p.font.bold = True
    c1p.font.color.rgb = theme['primary']
    c1p.font.name = "맑은 고딕"

    col1_content = slide.shapes.add_textbox(Inches(0.4), Inches(1.55), Inches(2.9), Inches(2.2))
    c1ctf = col1_content.text_frame
    c1ctf.word_wrap = True
    c1cp = c1ctf.paragraphs[0]
    c1cp.text = """• React 18 + TypeScript 5.0
• TailwindCSS
• Zustand (클라이언트 상태)
• TanStack Query (서버 상태)
• React Router v6

📦 라이브러리:
• Tiptap (에디터)
• react-virtualized
• react-hook-form
• axios + interceptor"""
    c1cp.font.size = Pt(9)
    c1cp.font.color.rgb = theme['text']
    c1cp.font.name = "맑은 고딕"

    # 열 2: 반응형 & 접근성
    col2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(1.1), Inches(3.1), Inches(2.8))
    col2.fill.solid()
    col2.fill.fore_color.rgb = theme['white']
    col2.line.color.rgb = theme['secondary']

    col2_title = slide.shapes.add_textbox(Inches(3.6), Inches(1.2), Inches(2.9), Inches(0.3))
    c2tf = col2_title.text_frame
    c2p = c2tf.paragraphs[0]
    c2p.text = "📱 반응형 & 접근성"
    c2p.font.size = Pt(12)
    c2p.font.bold = True
    c2p.font.color.rgb = theme['primary']
    c2p.font.name = "맑은 고딕"

    col2_content = slide.shapes.add_textbox(Inches(3.6), Inches(1.55), Inches(2.9), Inches(2.2))
    c2ctf = col2_content.text_frame
    c2ctf.word_wrap = True
    c2cp = c2ctf.paragraphs[0]
    c2cp.text = """📐 브레이크포인트:
• Desktop: 1280px+
• Tablet: 768px ~ 1279px
• Mobile: < 768px

♿ 접근성:
• aria-label 필수
• 키보드 네비게이션
• 색상 대비 4.5:1
• 포커스 표시
• 스크린리더 테스트"""
    c2cp.font.size = Pt(9)
    c2cp.font.color.rgb = theme['text']
    c2cp.font.name = "맑은 고딕"

    # 열 3: 성능 & 보안
    col3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.1), Inches(3.0), Inches(2.8))
    col3.fill.solid()
    col3.fill.fore_color.rgb = theme['white']
    col3.line.color.rgb = theme['secondary']

    col3_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(2.8), Inches(0.3))
    c3tf = col3_title.text_frame
    c3p = c3tf.paragraphs[0]
    c3p.text = "⚡ 성능 & 보안"
    c3p.font.size = Pt(12)
    c3p.font.bold = True
    c3p.font.color.rgb = theme['primary']
    c3p.font.name = "맑은 고딕"

    col3_content = slide.shapes.add_textbox(Inches(6.8), Inches(1.55), Inches(2.8), Inches(2.2))
    c3ctf = col3_content.text_frame
    c3ctf.word_wrap = True
    c3cp = c3ctf.paragraphs[0]
    c3cp.text = """🚀 최적화:
• 코드 스플리팅 (lazy)
• 이미지 WebP
• 가상화 목록
• useMemo

🔒 보안:
• XSS: sanitize-html
• CSRF: SameSite
• 토큰: HttpOnly
• 인증: PKCE + JWT"""
    c3cp.font.size = Pt(9)
    c3cp.font.color.rgb = theme['text']
    c3cp.font.name = "맑은 고딕"

    # 하단: API & 테스트
    api_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(4.1), Inches(4.65), Inches(2.5))
    api_box.fill.solid()
    api_box.fill.fore_color.rgb = theme['light']
    api_box.line.color.rgb = theme['accent']
    api_box.line.width = Pt(2)

    api_title = slide.shapes.add_textbox(Inches(0.4), Inches(4.2), Inches(4.45), Inches(0.3))
    atf = api_title.text_frame
    ap = atf.paragraphs[0]
    ap.text = "📡 API 연동 주의사항"
    ap.font.size = Pt(12)
    ap.font.bold = True
    ap.font.color.rgb = theme['accent']
    ap.font.name = "맑은 고딕"

    api_content = slide.shapes.add_textbox(Inches(0.4), Inches(4.55), Inches(4.45), Inches(1.9))
    actf = api_content.text_frame
    actf.word_wrap = True
    acp = actf.paragraphs[0]
    acp.text = """• Base URL: /api/v1
• 인증: Bearer Token
• 에러: { code, message, details }
• 페이지네이션: page, size, sort

🔄 Axios Interceptor:
• 401 → 토큰 갱신
• 403 → 권한 없음 페이지
• 500 → 에러 토스트
• 네트워크 에러 → 재시도"""
    acp.font.size = Pt(9)
    acp.font.color.rgb = theme['text']
    acp.font.name = "맑은 고딕"

    test_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.05), Inches(4.1), Inches(4.65), Inches(2.5))
    test_box.fill.solid()
    test_box.fill.fore_color.rgb = theme['light']
    test_box.line.color.rgb = theme['success']
    test_box.line.width = Pt(2)

    test_title = slide.shapes.add_textbox(Inches(5.15), Inches(4.2), Inches(4.45), Inches(0.3))
    ttf = test_title.text_frame
    tp = ttf.paragraphs[0]
    tp.text = "🧪 테스트 & 품질"
    tp.font.size = Pt(12)
    tp.font.bold = True
    tp.font.color.rgb = theme['success']
    tp.font.name = "맑은 고딕"

    test_content = slide.shapes.add_textbox(Inches(5.15), Inches(4.55), Inches(4.45), Inches(1.9))
    tctf = test_content.text_frame
    tctf.word_wrap = True
    tcp = tctf.paragraphs[0]
    tcp.text = """• 단위: Vitest (커버리지 80%+)
• 컴포넌트: Testing Library
• E2E: Playwright
• 모킹: MSW

✅ 품질 기준:
• TypeScript strict
• ESLint + Prettier
• Lighthouse 90+
• Chrome, Firefox, Safari, Edge"""
    tcp.font.size = Pt(9)
    tcp.font.color.rgb = theme['text']
    tcp.font.name = "맑은 고딕"

    return slide

def create_qna_slide(prs, theme):
    """Q&A 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme['dark']
    bg.line.fill.background()

    # Q&A
    qa_box = slide.shapes.add_textbox(Inches(0), Inches(2), Inches(10), Inches(1.5))
    qtf = qa_box.text_frame
    qp = qtf.paragraphs[0]
    qp.text = "Q & A"
    qp.font.size = Pt(72)
    qp.font.bold = True
    qp.font.color.rgb = theme['white']
    qp.font.name = "맑은 고딕"
    qp.alignment = PP_ALIGN.CENTER

    # 하이라이트
    highlights = [
        "📱 반응형 디자인 (Desktop, Tablet, Mobile)",
        "🔍 검색: 채팅 모드 (RAG) + 키워드 검색",
        "📝 마크다운 에디터 + 실시간 미리보기",
        "🔐 역할 기반 접근 제어 (USER, MANAGER, ADMIN)",
        "♿ 웹 접근성 WCAG 2.1 AA 준수"
    ]

    y = 3.8
    for hl in highlights:
        hl_box = slide.shapes.add_textbox(Inches(2), Inches(y), Inches(6), Inches(0.4))
        htf = hl_box.text_frame
        hp = htf.paragraphs[0]
        hp.text = hl
        hp.font.size = Pt(16)
        hp.font.color.rgb = theme['secondary']
        hp.font.name = "맑은 고딕"
        hp.alignment = PP_ALIGN.CENTER
        y += 0.5

    return slide

def main():
    """메인 함수"""
    print("=" * 60)
    print("UI Storyboard 프레젠테이션 생성")
    print("테마: Brown Earth")
    print("=" * 60)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    theme = BROWN_EARTH

    print("\n슬라이드 생성 중...")

    # 1. 표지
    create_title_slide(
        prs,
        "UI Storyboard\n상세 설계서",
        "Knowledge Hub | 사내 지식 검색 시스템 UI/UX 설계",
        theme
    )
    print("  [1/8] 표지 생성 완료")

    # 2. 로그인
    create_login_slide(prs, theme)
    print("  [2/8] 로그인 화면 생성 완료")

    # 3. 대시보드
    create_dashboard_slide(prs, theme)
    print("  [3/8] 대시보드 생성 완료")

    # 4. 검색 - 채팅
    create_search_chat_slide(prs, theme)
    print("  [4/8] 검색 채팅 모드 생성 완료")

    # 5. 검색 - 키워드
    create_search_keyword_slide(prs, theme)
    print("  [5/8] 검색 키워드 모드 생성 완료")

    # 6. 지식 관리
    create_knowledge_slide(prs, theme)
    print("  [6/8] 지식 관리 생성 완료")

    # 7. 프로필 & 관리자
    create_profile_admin_slide(prs, theme)
    print("  [7/8] 프로필 & 관리자 생성 완료")

    # 8. 개발자 가이드
    create_developer_notes_slide(prs, theme)
    print("  [8/8] 개발자 가이드 생성 완료")

    # 9. Q&A
    create_qna_slide(prs, theme)
    print("  [9/9] Q&A 생성 완료")

    # 저장
    output_dir = "knowledge_service/docs/02_design/ui_storyboard"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "UI_Storyboard_Brown.pptx")
    prs.save(output_path)

    print("\n" + "=" * 60)
    print(f"프레젠테이션 생성 완료!")
    print(f"파일 위치: {output_path}")
    print(f"총 슬라이드: 9장")
    print(f"파일 크기: {os.path.getsize(output_path):,} bytes")
    print("=" * 60)

    return output_path

if __name__ == "__main__":
    main()
