#!/usr/bin/env python3
"""
Hybrid RAG Platform Detailed Design Presentation Generator
- Brown Theme (Earth Tones)
- Detailed Architecture Diagrams
- Storage Advantages Highlight
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os

# ========== Brown Earth Theme ==========
THEME = {
    'primary': RGBColor(121, 85, 72),        # #795548 - Brown
    'secondary': RGBColor(161, 136, 127),    # #A1887F - Light Brown
    'accent': RGBColor(255, 152, 0),         # #FF9800 - Orange
    'dark': RGBColor(62, 39, 35),            # #3E2723 - Dark Brown
    'light': RGBColor(239, 235, 233),        # #EFEBE9 - Cream
    'text': RGBColor(33, 33, 33),            # #212121
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(158, 158, 158),         # #9E9E9E
    'lightgray': RGBColor(245, 245, 245),    # #F5F5F5
    'gold': RGBColor(255, 193, 7),           # #FFC107 - Gold
    'teal': RGBColor(0, 150, 136),           # #009688 - Teal (for ES)
    'blue': RGBColor(63, 81, 181),           # #3F51B5 - Indigo (for Neo4j)
    'green': RGBColor(76, 175, 80),          # #4CAF50 - Green (for PG)
}

FONT_NAME = "맑은 고딕"


def add_shape_with_text(slide, shape_type, left, top, width, height,
                        fill_color, text="", font_size=12, font_color=None,
                        bold=False, border_color=None, border_width=Pt(1)):
    """Add a shape with text"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.color.rgb = fill_color

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.name = FONT_NAME
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].font.color.rgb = font_color or THEME['text']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 font_color=None, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = FONT_NAME
    p.font.bold = bold
    p.font.color.rgb = font_color or THEME['text']
    p.alignment = alignment
    return textbox


def create_section_header(slide, title, subtitle=""):
    """Add section header to slide"""
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.color.rgb = THEME['primary']

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
                 title, font_size=32, font_color=THEME['white'], bold=True)

    if subtitle:
        add_text_box(slide, Inches(10.5), Inches(0.4), Inches(2.5), Inches(0.5),
                     subtitle, font_size=14, font_color=THEME['light'])


def create_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['dark']
    bg.line.color.rgb = THEME['dark']

    # Decorative elements
    deco1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(4), Inches(5), Inches(4.5))
    deco1.fill.solid()
    deco1.fill.fore_color.rgb = THEME['primary']
    deco1.line.color.rgb = THEME['primary']

    deco2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-1.5), Inches(5), Inches(4))
    deco2.fill.solid()
    deco2.fill.fore_color.rgb = THEME['secondary']
    deco2.line.color.rgb = THEME['secondary']

    # Title
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.8),
                 "Neo4j Graph RAG 기반", font_size=24, font_color=THEME['secondary'],
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(1.5),
                 "Hybrid 지식 플랫폼", font_size=52, font_color=THEME['white'],
                 bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.7), Inches(12), Inches(0.8),
                 "상세 설계서 v2.3", font_size=28, font_color=THEME['gold'],
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Key features
    features = [
        "Vector + Graph\nHybrid Search",
        "DeepSeek V3.2\n95% 비용 절감",
        "3-Store\n아키텍처"
    ]
    for i, feat in enumerate(features):
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2 + i * 3.2), Inches(5), Inches(2.8), Inches(1),
            THEME['accent'], feat, font_size=12, font_color=THEME['white'], bold=True
        )

    # Date
    add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                 "2026-01-14 | Review 완료 (코드 검증됨)", font_size=14, font_color=THEME['secondary'],
                 alignment=PP_ALIGN.CENTER)

    return slide


def create_architecture_overview_slide(prs):
    """Slide 2: Architecture Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "System Architecture Overview", "전체 시스템 구조")

    # Description
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.4), Inches(12.7), Inches(0.6),
        THEME['light'],
        "5-Layer Architecture: 사용자 → API Gateway → 처리 계층 → 검색 계층 → 저장 계층",
        font_size=14, font_color=THEME['dark'], bold=True
    )

    # ===== Layer 1: User =====
    add_text_box(slide, Inches(0.3), Inches(2.1), Inches(2), Inches(0.3),
                 "사용자 계층", font_size=10, font_color=THEME['gray'], bold=True)

    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(2.4), Inches(2), Inches(1.2),
        THEME['light'], "사용자\n\nWeb UI\nAPI Client",
        font_size=10, font_color=THEME['text'], border_color=THEME['secondary']
    )
    add_text_box(slide, Inches(0.3), Inches(3.65), Inches(2), Inches(0.3),
                 "React SPA\nTypeScript", font_size=8, font_color=THEME['gray'], alignment=PP_ALIGN.CENTER)

    # Arrow
    add_shape_with_text(
        slide, MSO_SHAPE.RIGHT_ARROW,
        Inches(2.35), Inches(2.85), Inches(0.35), Inches(0.35),
        THEME['accent']
    )

    # ===== Layer 2: API Gateway =====
    add_text_box(slide, Inches(2.8), Inches(2.1), Inches(2), Inches(0.3),
                 "API Gateway", font_size=10, font_color=THEME['gray'], bold=True)

    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.8), Inches(2.4), Inches(2), Inches(1.2),
        THEME['primary'], "Spring Boot\n\nREST API\nJWT 인증",
        font_size=10, font_color=THEME['white'], bold=True
    )
    add_text_box(slide, Inches(2.8), Inches(3.65), Inches(2), Inches(0.3),
                 "인증/인가\n라우팅", font_size=8, font_color=THEME['gray'], alignment=PP_ALIGN.CENTER)

    # Arrow
    add_shape_with_text(
        slide, MSO_SHAPE.RIGHT_ARROW,
        Inches(4.85), Inches(2.85), Inches(0.35), Inches(0.35),
        THEME['accent']
    )

    # ===== Layer 3: Processing (VIP) =====
    add_text_box(slide, Inches(5.3), Inches(2.1), Inches(3.5), Inches(0.3),
                 "처리 계층 (VIP Pipeline)", font_size=10, font_color=THEME['gray'], bold=True)

    # VIP Box
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.3), Inches(2.4), Inches(3.5), Inches(2.5),
        THEME['light'], border_color=THEME['primary'], border_width=Pt(2)
    )

    add_text_box(slide, Inches(5.4), Inches(2.45), Inches(3.3), Inches(0.3),
                 "LangGraph Orchestrator", font_size=10, font_color=THEME['primary'], bold=True, alignment=PP_ALIGN.CENTER)

    # Stage boxes
    stages = [
        ("Stage 1", "엔티티 추출", THEME['green']),
        ("Stage 2", "오케스트레이션", THEME['teal']),
        ("Stage 3", "답변 합성", THEME['accent']),
    ]
    for i, (stage, desc, color) in enumerate(stages):
        y = 2.8 + i * 0.65
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.5), Inches(y), Inches(3.1), Inches(0.55),
            color, f"{stage}: {desc}", font_size=9, font_color=THEME['white']
        )

    add_text_box(slide, Inches(5.3), Inches(4.95), Inches(3.5), Inches(0.3),
                 "DeepSeek-Chat / Reasoner", font_size=8, font_color=THEME['gray'], alignment=PP_ALIGN.CENTER)

    # Arrow
    add_shape_with_text(
        slide, MSO_SHAPE.RIGHT_ARROW,
        Inches(8.85), Inches(3.5), Inches(0.35), Inches(0.35),
        THEME['accent']
    )

    # ===== Layer 4: Search =====
    add_text_box(slide, Inches(9.3), Inches(2.1), Inches(2), Inches(0.3),
                 "검색 계층", font_size=10, font_color=THEME['gray'], bold=True)

    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.3), Inches(2.4), Inches(2), Inches(1),
        THEME['teal'], "Vector Search\nElasticsearch",
        font_size=10, font_color=THEME['white']
    )
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.3), Inches(3.5), Inches(2), Inches(1),
        THEME['blue'], "Graph Search\nNeo4j",
        font_size=10, font_color=THEME['white']
    )
    add_text_box(slide, Inches(9.3), Inches(4.55), Inches(2), Inches(0.3),
                 "RRF Fusion", font_size=9, font_color=THEME['accent'], bold=True, alignment=PP_ALIGN.CENTER)

    # Arrow
    add_shape_with_text(
        slide, MSO_SHAPE.RIGHT_ARROW,
        Inches(11.35), Inches(3.2), Inches(0.35), Inches(0.35),
        THEME['accent']
    )

    # ===== Layer 5: Storage =====
    add_text_box(slide, Inches(11.8), Inches(2.1), Inches(1.4), Inches(0.3),
                 "저장 계층", font_size=10, font_color=THEME['gray'], bold=True)

    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(11.8), Inches(2.4), Inches(1.4), Inches(0.7),
        THEME['green'], "PG", font_size=10, font_color=THEME['white'], bold=True
    )
    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(11.8), Inches(3.2), Inches(1.4), Inches(0.7),
        THEME['teal'], "ES", font_size=10, font_color=THEME['white'], bold=True
    )
    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(11.8), Inches(4.0), Inches(1.4), Inches(0.7),
        THEME['blue'], "Neo4j", font_size=10, font_color=THEME['white'], bold=True
    )

    # Key architecture points
    add_text_box(slide, Inches(0.3), Inches(5.3), Inches(12.7), Inches(0.3),
                 "핵심 아키텍처 특징", font_size=12, font_color=THEME['primary'], bold=True)

    points = [
        ("마이크로서비스", "SpringBoot(비즈니스) + AI Service(Python) 분리"),
        ("VIP Pipeline", "Value → Intelligent → Planning 3단계 LLM 처리"),
        ("Hybrid Search", "Vector(의미) + Graph(관계) 검색 융합"),
        ("3-Store 전략", "PG(SSOT) + ES(검색) + Neo4j(그래프) 역할 분리"),
    ]

    for i, (title, desc) in enumerate(points):
        x = 0.3 + (i % 2) * 6.5
        y = 5.65 + (i // 2) * 0.55
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(1.8), Inches(0.45),
            THEME['secondary'], title, font_size=10, font_color=THEME['white'], bold=True
        )
        add_text_box(slide, Inches(x + 1.9), Inches(y + 0.05), Inches(4.4), Inches(0.4),
                     desc, font_size=10, font_color=THEME['text'])

    return slide


def create_architecture_detail_slide(prs):
    """Slide 3: Detailed Architecture Diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "Service Architecture Detail", "서비스 분리 아키텍처 상세")

    # ===== Frontend =====
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.5), Inches(2.5), Inches(1.5),
        THEME['light'], border_color=THEME['secondary'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(0.4), Inches(1.55), Inches(2.3), Inches(0.35),
                 "Frontend", font_size=12, font_color=THEME['secondary'], bold=True)
    add_text_box(slide, Inches(0.4), Inches(1.9), Inches(2.3), Inches(1.0),
                 "React 18 + TypeScript\nTailwindCSS\nReact Query\n\n반응형 SPA",
                 font_size=9, font_color=THEME['text'], alignment=PP_ALIGN.CENTER)

    # Arrow
    add_shape_with_text(
        slide, MSO_SHAPE.DOWN_ARROW,
        Inches(1.4), Inches(3.05), Inches(0.3), Inches(0.35),
        THEME['accent']
    )

    # ===== SpringBoot Backend =====
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(3.5), Inches(5.5), Inches(3.3),
        THEME['white'], border_color=THEME['primary'], border_width=Pt(3)
    )
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(3.5), Inches(5.5), Inches(0.5),
        THEME['primary'], "SpringBoot Backend (Java 17+)", font_size=14, font_color=THEME['white'], bold=True
    )

    # SpringBoot components
    sb_components = [
        ("REST API Controller", "엔드포인트 정의\n요청/응답 처리"),
        ("Business Services", "비즈니스 로직\n트랜잭션 관리"),
        ("JPA Repositories", "데이터 접근\nCRUD 처리"),
        ("WebClient", "AI Service 호출\n비동기 통신"),
        ("Spring Security", "JWT 인증\n권한 관리"),
    ]

    for i, (name, desc) in enumerate(sb_components):
        x = 0.5 + (i % 3) * 1.8
        y = 4.1 + (i // 3) * 1.3
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(1.65), Inches(1.15),
            THEME['light'], border_color=THEME['secondary']
        )
        add_text_box(slide, Inches(x + 0.05), Inches(y + 0.05), Inches(1.55), Inches(0.35),
                     name, font_size=8, font_color=THEME['primary'], bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.05), Inches(y + 0.4), Inches(1.55), Inches(0.7),
                     desc, font_size=7, font_color=THEME['gray'], alignment=PP_ALIGN.CENTER)

    # Arrow to AI Service
    add_shape_with_text(
        slide, MSO_SHAPE.RIGHT_ARROW,
        Inches(5.9), Inches(4.8), Inches(0.5), Inches(0.5),
        THEME['accent']
    )
    add_text_box(slide, Inches(5.9), Inches(5.35), Inches(0.5), Inches(0.3),
                 "REST", font_size=8, font_color=THEME['accent'], bold=True, alignment=PP_ALIGN.CENTER)

    # ===== AI Service =====
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.5), Inches(1.5), Inches(6.5), Inches(5.3),
        THEME['white'], border_color=THEME['accent'], border_width=Pt(3)
    )
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.5), Inches(1.5), Inches(6.5), Inches(0.5),
        THEME['accent'], "AI Service (Python 3.11+ / FastAPI)", font_size=14, font_color=THEME['white'], bold=True
    )

    # AI Service components
    ai_components = [
        ("FastAPI Router", "API 엔드포인트\n/search, /extract, /embed"),
        ("VIP Pipeline", "3단계 LLM 처리\nLangGraph 오케스트레이션"),
        ("Search Service", "Hybrid 검색\nVector + Graph 융합"),
        ("Extract Service", "엔티티 추출\n메타데이터 생성"),
        ("Embed Service", "임베딩 생성\nBGE-M3 모델"),
        ("LLM Client", "DeepSeek API\nChat/Reasoner"),
    ]

    for i, (name, desc) in enumerate(ai_components):
        x = 6.7 + (i % 3) * 2.1
        y = 2.1 + (i // 3) * 1.5
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(1.95), Inches(1.35),
            THEME['light'], border_color=THEME['gold']
        )
        add_text_box(slide, Inches(x + 0.05), Inches(y + 0.05), Inches(1.85), Inches(0.4),
                     name, font_size=9, font_color=THEME['accent'], bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.05), Inches(y + 0.45), Inches(1.85), Inches(0.85),
                     desc, font_size=8, font_color=THEME['text'], alignment=PP_ALIGN.CENTER)

    # External APIs
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.7), Inches(5.2), Inches(3), Inches(0.5),
        THEME['dark'], "DeepSeek API (External)", font_size=10, font_color=THEME['white']
    )
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.9), Inches(5.2), Inches(3), Inches(0.5),
        THEME['secondary'], "BGE-M3 (Local Model)", font_size=10, font_color=THEME['white']
    )

    # Database connections
    add_shape_with_text(
        slide, MSO_SHAPE.DOWN_ARROW,
        Inches(2.4), Inches(6.85), Inches(0.3), Inches(0.35),
        THEME['green']
    )
    add_shape_with_text(
        slide, MSO_SHAPE.DOWN_ARROW,
        Inches(8.5), Inches(5.75), Inches(0.3), Inches(0.35),
        THEME['teal']
    )
    add_shape_with_text(
        slide, MSO_SHAPE.DOWN_ARROW,
        Inches(11), Inches(5.75), Inches(0.3), Inches(0.35),
        THEME['blue']
    )

    return slide


def create_architecture_advantages_slide(prs):
    """Slide 4: Architecture Advantages"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "Architecture Advantages", "아키텍처 장점")

    advantages = [
        {
            "title": "서비스 분리의 장점",
            "color": THEME['primary'],
            "items": [
                ("독립 배포", "SpringBoot와 AI Service 독립 배포 가능\n장애 격리 및 롤백 용이"),
                ("기술 최적화", "Java: 엔터프라이즈 트랜잭션\nPython: AI/ML 생태계 활용"),
                ("확장성", "AI 부하 증가시 AI Service만 스케일 아웃\n비용 효율적 리소스 관리"),
                ("팀 분리", "백엔드팀 / AI팀 독립 개발\n병렬 개발로 생산성 향상"),
            ]
        },
        {
            "title": "VIP 파이프라인의 장점",
            "color": THEME['accent'],
            "items": [
                ("단계별 최적화", "각 Stage에 최적 모델 배치\nChat: 단순작업, Reasoner: 복잡추론"),
                ("비용 절감", "95% 비용 절감 ($45.50 → $2.26)\n캐시 히트로 추가 절감"),
                ("품질 향상", "단계별 검증으로 오류 조기 발견\n중간 결과 재사용 가능"),
                ("유연한 확장", "새로운 Stage 추가 용이\nA/B 테스트 지원"),
            ]
        },
        {
            "title": "Hybrid Search의 장점",
            "color": THEME['teal'],
            "items": [
                ("높은 정확도", "Vector: 의미적 유사도\nGraph: 관계 기반 확장"),
                ("숨겨진 연결 발견", "문서 간 직접 연결 없어도\n엔티티 관계로 연결 발견"),
                ("컨텍스트 풍부화", "단순 키워드 매칭 넘어\n관련 문서 자동 확장"),
                ("RRF 융합", "여러 검색 결과 최적 결합\n단일 검색 대비 성능 향상"),
            ]
        },
    ]

    y_start = 1.5
    for adv in advantages:
        # Section title
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3), Inches(y_start), Inches(12.7), Inches(0.45),
            adv["color"], adv["title"], font_size=14, font_color=THEME['white'], bold=True
        )

        # Items (4 columns)
        for i, (title, desc) in enumerate(adv["items"]):
            x = 0.3 + i * 3.2
            y = y_start + 0.55

            add_shape_with_text(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(3), Inches(1.35),
                THEME['light'], border_color=adv["color"]
            )
            add_text_box(slide, Inches(x + 0.1), Inches(y + 0.05), Inches(2.8), Inches(0.35),
                         title, font_size=11, font_color=adv["color"], bold=True)
            add_text_box(slide, Inches(x + 0.1), Inches(y + 0.4), Inches(2.8), Inches(0.9),
                         desc, font_size=9, font_color=THEME['text'])

        y_start += 2.0

    return slide


def create_postgresql_slide(prs):
    """Slide 5: PostgreSQL Advantages"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "PostgreSQL - SSOT", "단일 진실 공급원 (마스터 DB)")

    # Main description
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.4), Inches(12.7), Inches(0.7),
        THEME['green'],
        "PostgreSQL은 모든 데이터의 마스터 레코드 (Single Source of Truth)로서 데이터 일관성을 보장합니다",
        font_size=16, font_color=THEME['white'], bold=True
    )

    # Icon and title
    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(0.5), Inches(2.3), Inches(2), Inches(1.2),
        THEME['green'], "PostgreSQL\n16+", font_size=18, font_color=THEME['white'], bold=True
    )

    # Key advantages
    advantages = [
        {
            "title": "ACID 트랜잭션",
            "desc": "완벽한 트랜잭션 보장\n\n- Atomicity: 원자성\n- Consistency: 일관성\n- Isolation: 격리성\n- Durability: 지속성\n\n금융급 데이터 신뢰성",
            "icon": "A"
        },
        {
            "title": "관계형 무결성",
            "desc": "참조 무결성 보장\n\n- FK 제약조건\n- CASCADE 삭제\n- CHECK 제약조건\n- UNIQUE 보장\n\n데이터 품질 자동 유지",
            "icon": "R"
        },
        {
            "title": "복잡한 쿼리",
            "desc": "고급 SQL 지원\n\n- JOIN 최적화\n- Window Functions\n- CTE (WITH절)\n- JSONB 지원\n\n복잡한 비즈니스 로직",
            "icon": "Q"
        },
        {
            "title": "확장성 & 성능",
            "desc": "엔터프라이즈 검증\n\n- 파티셔닝\n- 병렬 쿼리\n- VACUUM 자동화\n- WAL 복제\n\n대용량 데이터 처리",
            "icon": "S"
        },
    ]

    for i, adv in enumerate(advantages):
        x = 2.8 + i * 2.6
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2.3), Inches(2.4), Inches(3.2),
            THEME['light'], border_color=THEME['green'], border_width=Pt(2)
        )
        add_shape_with_text(
            slide, MSO_SHAPE.OVAL,
            Inches(x + 0.9), Inches(2.4), Inches(0.6), Inches(0.6),
            THEME['green'], adv["icon"], font_size=16, font_color=THEME['white'], bold=True
        )
        add_text_box(slide, Inches(x + 0.1), Inches(3.1), Inches(2.2), Inches(0.35),
                     adv["title"], font_size=12, font_color=THEME['green'], bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.1), Inches(3.5), Inches(2.2), Inches(1.9),
                     adv["desc"], font_size=9, font_color=THEME['text'])

    # Data tables
    add_text_box(slide, Inches(0.3), Inches(5.7), Inches(4), Inches(0.35),
                 "주요 테이블", font_size=12, font_color=THEME['green'], bold=True)

    tables = [
        ("documents", "문서 마스터"),
        ("chunks", "청크 저장"),
        ("entities", "엔티티 마스터"),
        ("projects", "프로젝트 정보"),
        ("persons", "인물 정보"),
        ("entity_relationships", "엔티티 관계"),
    ]

    for i, (table, desc) in enumerate(tables):
        x = 0.3 + (i % 3) * 4.3
        y = 6.05 + (i // 3) * 0.45
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(1.5), Inches(0.4),
            THEME['green'], table, font_size=9, font_color=THEME['white']
        )
        add_text_box(slide, Inches(x + 1.6), Inches(y + 0.05), Inches(2.5), Inches(0.35),
                     desc, font_size=9, font_color=THEME['text'])

    return slide


def create_elasticsearch_slide(prs):
    """Slide 6: Elasticsearch Advantages"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "Elasticsearch - Vector Search", "벡터 검색 + 메타데이터 통합")

    # Main description
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.4), Inches(12.7), Inches(0.7),
        THEME['teal'],
        "Elasticsearch는 Dense/Sparse 벡터 검색과 메타데이터 필터링을 단일 쿼리로 처리합니다 (제로 조인)",
        font_size=16, font_color=THEME['white'], bold=True
    )

    # Icon
    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(0.5), Inches(2.3), Inches(2), Inches(1.2),
        THEME['teal'], "Elastic\nsearch 8.x", font_size=16, font_color=THEME['white'], bold=True
    )

    # Key advantages
    advantages = [
        {
            "title": "Dense Vector 검색",
            "desc": "의미론적 유사도\n\n- BGE-M3 1024차원\n- kNN 근접 검색\n- 자연어 이해\n- 동의어/유사어 처리\n\n질문과 문서 의미 매칭",
            "icon": "D"
        },
        {
            "title": "Sparse Vector 검색",
            "desc": "키워드 가중치\n\n- BM25 스타일\n- 정확한 용어 매칭\n- 전문 용어 검색\n- TF-IDF 기반\n\n정확한 키워드 검색",
            "icon": "S"
        },
        {
            "title": "제로 조인 쿼리",
            "desc": "메타데이터 비정규화\n\n- project_name 내장\n- valid_date 내장\n- categories 내장\n- 단일 쿼리 완결\n\nPostgreSQL 조인 불필요",
            "icon": "Z"
        },
        {
            "title": "실시간 검색",
            "desc": "Near Real-Time\n\n- 1초 이내 인덱싱\n- refresh_interval 설정\n- 분산 검색\n- 자동 샤딩\n\n빠른 검색 응답",
            "icon": "R"
        },
    ]

    for i, adv in enumerate(advantages):
        x = 2.8 + i * 2.6
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2.3), Inches(2.4), Inches(3.2),
            THEME['light'], border_color=THEME['teal'], border_width=Pt(2)
        )
        add_shape_with_text(
            slide, MSO_SHAPE.OVAL,
            Inches(x + 0.9), Inches(2.4), Inches(0.6), Inches(0.6),
            THEME['teal'], adv["icon"], font_size=16, font_color=THEME['white'], bold=True
        )
        add_text_box(slide, Inches(x + 0.1), Inches(3.1), Inches(2.2), Inches(0.35),
                     adv["title"], font_size=12, font_color=THEME['teal'], bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.1), Inches(3.5), Inches(2.2), Inches(1.9),
                     adv["desc"], font_size=9, font_color=THEME['text'])

    # Index mapping
    add_text_box(slide, Inches(0.3), Inches(5.7), Inches(4), Inches(0.35),
                 "인덱스 필드 구조", font_size=12, font_color=THEME['teal'], bold=True)

    fields = [
        ("dense_vector", "1024차원 Dense"),
        ("sparse_vector", "Sparse (lexical)"),
        ("content", "청크 텍스트"),
        ("metadata.project", "프로젝트명"),
        ("metadata.dates", "유효 기간"),
        ("metadata.categories", "계층 분류"),
    ]

    for i, (field, desc) in enumerate(fields):
        x = 0.3 + (i % 3) * 4.3
        y = 6.05 + (i // 3) * 0.45
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(1.7), Inches(0.4),
            THEME['teal'], field, font_size=9, font_color=THEME['white']
        )
        add_text_box(slide, Inches(x + 1.8), Inches(y + 0.05), Inches(2.3), Inches(0.35),
                     desc, font_size=9, font_color=THEME['text'])

    return slide


def create_neo4j_slide(prs):
    """Slide 7: Neo4j Advantages"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "Neo4j - Knowledge Graph", "지식 그래프 기반 관계 탐색")

    # Main description
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.4), Inches(12.7), Inches(0.7),
        THEME['blue'],
        "Neo4j는 엔티티 간 관계를 그래프로 저장하여 숨겨진 연결을 발견하고 지식을 확장합니다",
        font_size=16, font_color=THEME['white'], bold=True
    )

    # Icon
    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(0.5), Inches(2.3), Inches(2), Inches(1.2),
        THEME['blue'], "Neo4j\n5.x", font_size=18, font_color=THEME['white'], bold=True
    )

    # Key advantages
    advantages = [
        {
            "title": "관계 탐색",
            "desc": "연결 기반 검색\n\n- 2-hop 확장 탐색\n- 경로 발견\n- 패턴 매칭\n- 관계 필터링\n\n숨겨진 연결 발견",
            "icon": "R"
        },
        {
            "title": "커뮤니티 탐지",
            "desc": "계층적 그룹화\n\n- Louvain 알고리즘\n- 주제별 클러스터\n- 자동 분류\n- 계층 구조 생성\n\n관련 문서 그룹화",
            "icon": "C"
        },
        {
            "title": "Slim Graph 전략",
            "desc": "메모리 최적화\n\n- 최소 속성만 저장\n- ID + 이름 + 타입\n- 본문은 ES 참조\n- 16GB RAM 운영\n\n효율적 메모리 사용",
            "icon": "S"
        },
        {
            "title": "Cypher 쿼리",
            "desc": "직관적 그래프 쿼리\n\n- 패턴 매칭 문법\n- 경로 표현\n- 집계 함수\n- FULLTEXT 검색\n\n복잡한 관계 쿼리",
            "icon": "Q"
        },
    ]

    for i, adv in enumerate(advantages):
        x = 2.8 + i * 2.6
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2.3), Inches(2.4), Inches(3.2),
            THEME['light'], border_color=THEME['blue'], border_width=Pt(2)
        )
        add_shape_with_text(
            slide, MSO_SHAPE.OVAL,
            Inches(x + 0.9), Inches(2.4), Inches(0.6), Inches(0.6),
            THEME['blue'], adv["icon"], font_size=16, font_color=THEME['white'], bold=True
        )
        add_text_box(slide, Inches(x + 0.1), Inches(3.1), Inches(2.2), Inches(0.35),
                     adv["title"], font_size=12, font_color=THEME['blue'], bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.1), Inches(3.5), Inches(2.2), Inches(1.9),
                     adv["desc"], font_size=9, font_color=THEME['text'])

    # Graph model
    add_text_box(slide, Inches(0.3), Inches(5.7), Inches(2), Inches(0.35),
                 "노드 타입", font_size=12, font_color=THEME['blue'], bold=True)

    nodes = ["Entity", "TextUnit", "Community", "Document"]
    for i, node in enumerate(nodes):
        add_shape_with_text(
            slide, MSO_SHAPE.OVAL,
            Inches(0.3 + i * 1.6), Inches(6.1), Inches(1.4), Inches(0.6),
            THEME['blue'], node, font_size=10, font_color=THEME['white']
        )

    add_text_box(slide, Inches(6.8), Inches(5.7), Inches(2), Inches(0.35),
                 "관계 타입", font_size=12, font_color=THEME['blue'], bold=True)

    rels = ["RELATED_TO", "MENTIONED_IN", "BELONGS_TO", "PART_OF"]
    for i, rel in enumerate(rels):
        add_shape_with_text(
            slide, MSO_SHAPE.RIGHT_ARROW,
            Inches(6.8 + i * 1.6), Inches(6.1), Inches(1.4), Inches(0.6),
            THEME['secondary'], rel, font_size=8, font_color=THEME['white']
        )

    return slide


def create_3store_comparison_slide(prs):
    """Slide 8: 3-Store Strategy Comparison"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "3-Store Strategy", "역할별 저장소 분리 전략")

    # Comparison table
    add_text_box(slide, Inches(0.3), Inches(1.4), Inches(4), Inches(0.4),
                 "저장소별 역할 비교", font_size=14, font_color=THEME['primary'], bold=True)

    # Headers
    headers = [("", 1.8), ("PostgreSQL", 3.2), ("Elasticsearch", 3.2), ("Neo4j", 3.2)]
    x = 0.3
    for header, width in headers:
        color = THEME['primary'] if header == "" else (THEME['green'] if header == "PostgreSQL" else (THEME['teal'] if header == "Elasticsearch" else THEME['blue']))
        add_shape_with_text(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(1.8), Inches(width), Inches(0.5),
            color, header, font_size=12, font_color=THEME['white'], bold=True
        )
        x += width

    # Rows
    rows = [
        ("주요 역할", "마스터 데이터 (SSOT)", "벡터 검색 + 메타필터", "관계 그래프 탐색"),
        ("데이터 유형", "정형 데이터, FK 관계", "벡터, 비정규화 메타", "노드, 엣지, 관계"),
        ("쿼리 특성", "복잡한 조인, 트랜잭션", "kNN, BM25, 필터", "패턴 매칭, 경로 탐색"),
        ("동기화 방향", "→ ES, Neo4j 전파", "← PG 수신 전용", "← PG 수신 전용"),
        ("수정 권한", "쓰기/읽기 (Primary)", "읽기 전용 (검색용)", "읽기 전용 (탐색용)"),
        ("메모리 할당", "1GB shared_buffers", "4GB JVM heap", "2GB JVM heap"),
    ]

    for row_idx, row in enumerate(rows):
        x = 0.3
        y = 2.3 + row_idx * 0.55
        widths = [1.8, 3.2, 3.2, 3.2]
        for col_idx, (cell, width) in enumerate(zip(row, widths)):
            bg = THEME['secondary'] if col_idx == 0 else (THEME['lightgray'] if row_idx % 2 == 0 else THEME['white'])
            fc = THEME['white'] if col_idx == 0 else THEME['text']
            add_shape_with_text(
                slide, MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(width), Inches(0.55),
                bg, cell, font_size=10, font_color=fc,
                bold=(col_idx == 0), border_color=THEME['gray']
            )
            x += width

    # Data flow diagram
    add_text_box(slide, Inches(0.3), Inches(5.7), Inches(4), Inches(0.4),
                 "데이터 동기화 흐름", font_size=14, font_color=THEME['primary'], bold=True)

    # PG -> ES, Neo4j
    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(2), Inches(6.2), Inches(2.5), Inches(1),
        THEME['green'], "PostgreSQL\n(SSOT - 마스터)", font_size=12, font_color=THEME['white'], bold=True
    )

    add_shape_with_text(
        slide, MSO_SHAPE.RIGHT_ARROW,
        Inches(4.6), Inches(6.3), Inches(0.8), Inches(0.4),
        THEME['accent']
    )
    add_shape_with_text(
        slide, MSO_SHAPE.RIGHT_ARROW,
        Inches(4.6), Inches(6.8), Inches(0.8), Inches(0.4),
        THEME['accent']
    )

    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(5.5), Inches(6.0), Inches(2.5), Inches(0.7),
        THEME['teal'], "Elasticsearch", font_size=11, font_color=THEME['white']
    )
    add_shape_with_text(
        slide, MSO_SHAPE.FLOWCHART_DATA,
        Inches(5.5), Inches(6.8), Inches(2.5), Inches(0.7),
        THEME['blue'], "Neo4j", font_size=11, font_color=THEME['white']
    )

    # Key principle
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(6.1), Inches(4.5), Inches(1.1),
        THEME['light'], border_color=THEME['accent'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(8.6), Inches(6.15), Inches(4.3), Inches(1),
                 "핵심 원칙:\nPostgreSQL 실패 → 전체 롤백\nES/Neo4j 실패 → 재시도 큐 등록\n\nasyncio.gather로 동시 저장",
                 font_size=10, font_color=THEME['text'])

    return slide


def create_cost_slide(prs):
    """Slide 9: Cost Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "Cost Efficiency", "비용 효율성 분석")

    # Main highlight
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.4), Inches(12.7), Inches(1),
        THEME['accent'],
        "DeepSeek V3.2 활용으로 95% 비용 절감: $45.50 → $2.26 / 1,000문서",
        font_size=24, font_color=THEME['white'], bold=True
    )

    # Cost comparison
    add_text_box(slide, Inches(0.3), Inches(2.6), Inches(6), Inches(0.4),
                 "LLM 비용 비교 (1,000문서 기준)", font_size=14, font_color=THEME['primary'], bold=True)

    # Chart-like visual
    models = [
        ("Claude Sonnet", 45.50, THEME['gray']),
        ("GPT-4o", 14.50, THEME['gray']),
        ("DeepSeek", 2.26, THEME['accent']),
    ]

    max_cost = 50
    for i, (name, cost, color) in enumerate(models):
        y = 3.0 + i * 1.0
        bar_width = (cost / max_cost) * 8

        add_text_box(slide, Inches(0.3), Inches(y + 0.1), Inches(2), Inches(0.4),
                     name, font_size=12, font_color=THEME['text'], bold=True)

        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.5), Inches(y), Inches(bar_width), Inches(0.6),
            color, f"${cost:.2f}", font_size=12, font_color=THEME['white'], bold=True
        )

        if name == "DeepSeek":
            add_shape_with_text(
                slide, MSO_SHAPE.STAR_5_POINT,
                Inches(2.5 + bar_width + 0.2), Inches(y + 0.1), Inches(0.4), Inches(0.4),
                THEME['gold']
            )
            add_text_box(slide, Inches(2.5 + bar_width + 0.7), Inches(y + 0.1), Inches(2), Inches(0.4),
                         "95% 절감!", font_size=12, font_color=THEME['accent'], bold=True)

    # Monthly projection
    add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(6.0), Inches(6), Inches(1.2),
        THEME['light'], border_color=THEME['primary'], border_width=Pt(2)
    )
    add_text_box(slide, Inches(0.5), Inches(6.1), Inches(5.6), Inches(1),
                 "월간 예상 비용\n\n문서 처리: $2.26 + 검색 (1,000회): $0.50\n= 총 $2.76/월 | 연간 $33.12",
                 font_size=12, font_color=THEME['text'])

    # DeepSeek pricing detail
    add_text_box(slide, Inches(6.8), Inches(2.6), Inches(6), Inches(0.4),
                 "DeepSeek 모델별 요금", font_size=14, font_color=THEME['primary'], bold=True)

    pricing = [
        ("deepseek-chat", "$0.28/1M", "$1.10/1M", "엔티티 추출, 답변 합성"),
        ("deepseek-reasoner", "$2.19/1M", "$8.98/1M", "복잡한 추론"),
    ]

    headers = ["모델", "입력", "출력", "용도"]
    x = 6.8
    for header in headers:
        w = 1.8 if header == "모델" else (1.2 if header != "용도" else 2.3)
        add_shape_with_text(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(3.0), Inches(w), Inches(0.4),
            THEME['primary'], header, font_size=10, font_color=THEME['white'], bold=True
        )
        x += w

    for row_idx, row in enumerate(pricing):
        x = 6.8
        y = 3.4 + row_idx * 0.45
        widths = [1.8, 1.2, 1.2, 2.3]
        for col_idx, (cell, w) in enumerate(zip(row, widths)):
            add_shape_with_text(
                slide, MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(0.45),
                THEME['lightgray'] if row_idx % 2 == 0 else THEME['white'],
                cell, font_size=9, font_color=THEME['text'], border_color=THEME['gray']
            )
            x += w

    # Cost optimization tips
    add_text_box(slide, Inches(6.8), Inches(4.5), Inches(6), Inches(0.4),
                 "비용 최적화 전략", font_size=14, font_color=THEME['primary'], bold=True)

    tips = [
        ("캐시 히트 활용", "시스템 프롬프트 고정으로 90% 캐시 히트"),
        ("모델 선택 최적화", "단순 작업은 Chat, 복잡한 추론만 Reasoner"),
        ("배치 처리", "대량 문서는 배치로 처리하여 오버헤드 감소"),
        ("API 호출 최소화", "불필요한 재호출 방지, 결과 재사용"),
    ]

    for i, (title, desc) in enumerate(tips):
        x = 6.8 + (i % 2) * 3.2
        y = 4.9 + (i // 2) * 0.9
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(3), Inches(0.8),
            THEME['light'], border_color=THEME['secondary']
        )
        add_text_box(slide, Inches(x + 0.1), Inches(y + 0.05), Inches(2.8), Inches(0.3),
                     title, font_size=10, font_color=THEME['accent'], bold=True)
        add_text_box(slide, Inches(x + 0.1), Inches(y + 0.35), Inches(2.8), Inches(0.4),
                     desc, font_size=8, font_color=THEME['text'])

    return slide


def create_closing_slide(prs):
    """Slide 10: Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['dark']
    bg.line.color.rgb = THEME['dark']

    # Decorative
    deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(4), Inches(5.5), Inches(4.5))
    deco.fill.solid()
    deco.fill.fore_color.rgb = THEME['primary']
    deco.line.color.rgb = THEME['primary']

    deco2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-1), Inches(4), Inches(3))
    deco2.fill.solid()
    deco2.fill.fore_color.rgb = THEME['secondary']
    deco2.line.color.rgb = THEME['secondary']

    # Title
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(1),
                 "Thank You", font_size=56, font_color=THEME['white'],
                 bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3), Inches(12), Inches(0.7),
                 "Questions & Answers", font_size=28, font_color=THEME['light'],
                 alignment=PP_ALIGN.CENTER)

    # Key highlights
    highlights = [
        ("아키텍처", "SpringBoot + AI Service 분리 | VIP 3-Stage Pipeline"),
        ("저장소 전략", "PG(SSOT) + ES(Vector) + Neo4j(Graph) 역할 분리"),
        ("비용 효율", "DeepSeek V3.2 → 95% 절감 ($45.50 → $2.26)"),
        ("검색 품질", "Hybrid Search = Vector + Graph + RRF Fusion"),
    ]

    for i, (title, desc) in enumerate(highlights):
        add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(4 + i * 0.6), Inches(2), Inches(0.5),
            THEME['accent'], title, font_size=11, font_color=THEME['white'], bold=True
        )
        add_text_box(slide, Inches(3.6), Inches(4.05 + i * 0.6), Inches(8), Inches(0.5),
                     desc, font_size=12, font_color=THEME['light'])

    # Footer
    add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                 "Hybrid RAG Knowledge Platform v2.3 | 2026-01-14",
                 font_size=14, font_color=THEME['secondary'], alignment=PP_ALIGN.CENTER)

    return slide


def main():
    """Generate presentation"""
    print("=" * 60)
    print("Hybrid RAG Platform - Brown Theme Presentation")
    print("=" * 60)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("\n[1/10] Creating title slide...")
    create_title_slide(prs)

    print("[2/10] Creating architecture overview slide...")
    create_architecture_overview_slide(prs)

    print("[3/10] Creating architecture detail slide...")
    create_architecture_detail_slide(prs)

    print("[4/10] Creating architecture advantages slide...")
    create_architecture_advantages_slide(prs)

    print("[5/10] Creating PostgreSQL advantages slide...")
    create_postgresql_slide(prs)

    print("[6/10] Creating Elasticsearch advantages slide...")
    create_elasticsearch_slide(prs)

    print("[7/10] Creating Neo4j advantages slide...")
    create_neo4j_slide(prs)

    print("[8/10] Creating 3-store comparison slide...")
    create_3store_comparison_slide(prs)

    print("[9/10] Creating cost analysis slide...")
    create_cost_slide(prs)

    print("[10/10] Creating closing slide...")
    create_closing_slide(prs)

    # Save
    output_dir = "knowledge_service/docs/02_design"
    os.makedirs(output_dir, exist_ok=True)

    output_path = os.path.join(output_dir, "Hybrid_RAG_Platform_Design_Brown.pptx")
    prs.save(output_path)

    print("\n" + "=" * 60)
    print(f"Presentation saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    print("Theme: Brown Earth Tones")
    print("=" * 60)


if __name__ == "__main__":
    main()
