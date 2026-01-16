#!/usr/bin/env python3
"""
Hybrid RAG Knowledge Platform 통합 설계서 프레젠테이션 생성기 v3
Brown Theme - 아키텍처 도식화 개선, 설명 상세화
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os

# ============================================================
# Brown Theme Color Palette
# ============================================================
THEME = {
    'primary': RGBColor(101, 67, 33),
    'secondary': RGBColor(139, 90, 43),
    'accent': RGBColor(205, 133, 63),
    'dark': RGBColor(62, 39, 35),
    'light': RGBColor(245, 235, 220),
    'text': RGBColor(62, 39, 35),
    'white': RGBColor(255, 255, 255),
    'highlight': RGBColor(255, 152, 0),
    'success': RGBColor(76, 175, 80),
    'info': RGBColor(33, 150, 243),
    'danger': RGBColor(244, 67, 54),
    'desc_bg': RGBColor(232, 245, 233),
    'layer_bg': RGBColor(250, 245, 240),
    'layer_border': RGBColor(180, 150, 120),
}

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)
HEADER_HEIGHT = Cm(2.2)
CONTENT_TOP = Cm(2.8)


def add_header(slide, title):
    """공통 헤더"""
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, HEADER_HEIGHT
    )
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(30), Cm(1.2))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"


def add_box(slide, left, top, width, height, text, fill_color,
            text_color=None, font_size=10, bold=False, align_center=True):
    """박스 추가"""
    if text_color is None:
        text_color = THEME['white']

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(left), Cm(top), Cm(width), Cm(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = "맑은 고딕"
    p.font.bold = bold
    if align_center:
        p.alignment = PP_ALIGN.CENTER

    return shape


def add_layer_box(slide, left, top, width, height, label):
    """레이어 구분 박스 (배경)"""
    # 배경 박스
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(left), Cm(top), Cm(width), Cm(height)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['layer_bg']
    bg.line.color.rgb = THEME['layer_border']
    bg.line.width = Pt(1.5)

    # 레이블
    label_box = slide.shapes.add_textbox(Cm(left + 0.2), Cm(top + 0.1), Cm(width - 0.4), Cm(0.5))
    lf = label_box.text_frame
    lf.text = label
    lp = lf.paragraphs[0]
    lp.font.size = Pt(9)
    lp.font.bold = True
    lp.font.color.rgb = THEME['dark']
    lp.font.name = "맑은 고딕"

    return bg


def add_detailed_description(slide, descriptions):
    """상세 설명 박스 (우측) - 더 풍부한 내용"""
    # 배경
    desc_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(21.5), Cm(2.8), Cm(11.8), Cm(14.5)
    )
    desc_bg.fill.solid()
    desc_bg.fill.fore_color.rgb = THEME['desc_bg']
    desc_bg.line.color.rgb = RGBColor(200, 230, 200)
    desc_bg.line.width = Pt(1)

    # 텍스트
    desc_text = slide.shapes.add_textbox(Cm(21.8), Cm(3), Cm(11.2), Cm(14))
    tf = desc_text.text_frame
    tf.word_wrap = True

    first = True
    for title, bullets in descriptions:
        # 제목
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(8)

        p.text = title
        p.font.bold = True
        p.font.size = Pt(9)
        p.font.color.rgb = THEME['dark']
        p.font.name = "맑은 고딕"

        # 불렛 포인트들
        for bullet in bullets:
            bp = tf.add_paragraph()
            bp.text = f"  • {bullet}"
            bp.font.size = Pt(8)
            bp.font.color.rgb = THEME['text']
            bp.font.name = "맑은 고딕"
            bp.space_after = Pt(2)


def add_arrow_text(slide, x, y, text="→"):
    """화살표 텍스트"""
    arr = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(1), Cm(0.5))
    af = arr.text_frame
    af.text = text
    ap = af.paragraphs[0]
    ap.font.size = Pt(12)
    ap.font.color.rgb = THEME['accent']
    ap.font.bold = True


def create_title_slide(prs):
    """표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['dark']
    bg.line.fill.background()

    deco1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(5), Inches(5)
    )
    deco1.fill.solid()
    deco1.fill.fore_color.rgb = THEME['primary']
    deco1.line.fill.background()

    deco2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(9), Inches(5.5), Inches(5), Inches(2.5)
    )
    deco2.fill.solid()
    deco2.fill.fore_color.rgb = THEME['secondary']
    deco2.line.fill.background()

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.8), Inches(0.08), Inches(1.8)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = THEME['accent']
    line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Hybrid RAG Knowledge Platform"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    p2 = tf.add_paragraph()
    p2.text = "통합 상세 설계서"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = THEME['accent']
    p2.font.name = "맑은 고딕"

    sub_box = slide.shapes.add_textbox(Inches(1.1), Inches(4.5), Inches(8), Inches(0.8))
    sf = sub_box.text_frame
    sf.text = "Vector + Graph 기반 차세대 지식 검색 시스템"
    sp = sf.paragraphs[0]
    sp.font.size = Pt(16)
    sp.font.color.rgb = THEME['light']
    sp.font.name = "맑은 고딕"

    date_box = slide.shapes.add_textbox(Inches(1.1), Inches(6.5), Inches(5), Inches(0.4))
    df = date_box.text_frame
    df.text = "2026-01-16  |  Version 1.0  |  Final Draft"
    dp = df.paragraphs[0]
    dp.font.size = Pt(11)
    dp.font.color.rgb = THEME['light']
    dp.font.name = "맑은 고딕"

    return slide


def create_section_slide(prs, section_number, section_title):
    """섹션 구분"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['primary']
    bg.line.fill.background()

    num_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(2), Inches(1))
    nf = num_box.text_frame
    nf.text = f"{section_number:02d}"
    np = nf.paragraphs[0]
    np.font.size = Pt(60)
    np.font.bold = True
    np.font.color.rgb = THEME['accent']
    np.font.name = "맑은 고딕"

    title_box = slide.shapes.add_textbox(Inches(3.2), Inches(3), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = section_title
    tp = tf.paragraphs[0]
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = THEME['white']
    tp.font.name = "맑은 고딕"

    return slide


def create_executive_summary_slide(prs):
    """Executive Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Executive Summary - 프로젝트 비전")

    # 비전 문구
    vision_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(3.2), Cm(20), Cm(1.8)
    )
    vision_box.fill.solid()
    vision_box.fill.fore_color.rgb = THEME['secondary']
    vision_box.line.fill.background()

    vision_text = slide.shapes.add_textbox(Cm(1.5), Cm(3.5), Cm(19), Cm(1.3))
    vf = vision_text.text_frame
    vf.text = '"기업 지식의 80%가 잠들어 있는 문서에서, AI가 즉시 답을 찾아드립니다"'
    vp = vf.paragraphs[0]
    vp.font.size = Pt(14)
    vp.font.bold = True
    vp.font.color.rgb = THEME['white']
    vp.font.name = "맑은 고딕"
    vp.alignment = PP_ALIGN.CENTER

    # 핵심 가치
    values = [
        ("💰 비용 효율성", "DeepSeek 단일 모델\nLLM 비용 95% 절감\n연간 ~$14,000", THEME['success']),
        ("🎯 정확한 검색", "Vector + Graph 융합\nRRF 재순위 알고리즘\nPrecision@5 > 85%", THEME['info']),
        ("🔐 기업 보안", "OAuth 2.0 + PKCE\nAES-256-GCM 암호화\nRBAC 권한 관리", RGBColor(156, 39, 176)),
        ("📈 확장 가능성", "Docker → K8s 경로\n마이크로서비스 구조\nAPI 기반 통합", THEME['highlight']),
    ]

    x_pos = 1
    for title, desc, color in values:
        add_box(slide, x_pos, 5.8, 4.8, 1, title, color, font_size=10, bold=True)

        desc_box = slide.shapes.add_textbox(Cm(x_pos), Cm(7), Cm(4.8), Cm(2.5))
        df = desc_box.text_frame
        df.word_wrap = True
        for i, line in enumerate(desc.split('\n')):
            p = df.paragraphs[0] if i == 0 else df.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(8)
            p.font.color.rgb = THEME['text']
            p.font.name = "맑은 고딕"

        x_pos += 5

    # 우측 상세 설명
    add_detailed_description(slide, [
        ("1. Hybrid RAG 기술이란?", [
            "Vector Search: 의미적 유사도 기반 검색 (BGE-M3 임베딩)",
            "Graph Search: 엔티티 관계 기반 확장 검색 (Neo4j)",
            "RRF 융합: 두 검색 결과를 Reciprocal Rank로 통합",
            "기존 키워드 검색 대비 정확도 40% 이상 향상",
        ]),
        ("2. 비용 절감 원리", [
            "DeepSeek V3.2: GPT-4 대비 토큰당 95% 저렴",
            "동일 품질 유지하면서 연간 $98,000 절감",
            "단일 모델로 엔티티 추출 + 답변 합성 모두 처리",
            "추가 Fine-tuning 불필요 → 운영 비용 최소화",
        ]),
        ("3. 기업 보안 강화 방안", [
            "Keycloak 기반 OAuth 2.0 + PKCE 인증 적용",
            "AES-256-GCM으로 민감 데이터 필드 암호화",
            "RBAC 3단계 역할: Admin/Manager/User",
            "TLS 1.3 전송 암호화 + HashiCorp Vault 키 관리",
        ]),
        ("4. 확장 전략 로드맵", [
            "Phase 1: Docker Compose 단일 서버 ($14K/년)",
            "Phase 2: Redis Sentinel 캐시 이중화 ($20K/년)",
            "Phase 3: Docker Swarm 3노드 클러스터 ($40K/년)",
            "Phase 4: Kubernetes 풀 클러스터 ($100K/년)",
        ]),
    ])

    return slide


def create_key_metrics_slide(prs):
    """핵심 수치"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "핵심 목표 수치 (Key Metrics)")

    metrics = [
        ("응답 시간", "< 3초", "(P95)", THEME['info']),
        ("검색 정확도", "> 85%", "(Precision@5)", THEME['success']),
        ("LLM 비용", "95%↓", "(GPT-4 대비)", THEME['highlight']),
        ("가용성", "99.5%", "(SLA)", THEME['secondary']),
        ("동시 사용자", "100명", "(1단계)", THEME['primary']),
    ]

    x_pos = 1
    for title, value, sub, color in metrics:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x_pos), Cm(3.3), Cm(3.6), Cm(3.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()

        val_box = slide.shapes.add_textbox(Cm(x_pos), Cm(3.7), Cm(3.6), Cm(1.2))
        vf = val_box.text_frame
        vf.text = value
        vp = vf.paragraphs[0]
        vp.font.size = Pt(20)
        vp.font.bold = True
        vp.font.color.rgb = THEME['white']
        vp.alignment = PP_ALIGN.CENTER

        title_box = slide.shapes.add_textbox(Cm(x_pos), Cm(5), Cm(3.6), Cm(0.6))
        ttf = title_box.text_frame
        ttf.text = title
        ttp = ttf.paragraphs[0]
        ttp.font.size = Pt(9)
        ttp.font.color.rgb = THEME['white']
        ttp.alignment = PP_ALIGN.CENTER

        sub_box = slide.shapes.add_textbox(Cm(x_pos), Cm(5.6), Cm(3.6), Cm(0.5))
        sf = sub_box.text_frame
        sf.text = sub
        sp = sf.paragraphs[0]
        sp.font.size = Pt(7)
        sp.font.color.rgb = RGBColor(220, 220, 220)
        sp.alignment = PP_ALIGN.CENTER

        x_pos += 3.9

    # 비용 비교 테이블
    table_title = slide.shapes.add_textbox(Cm(1), Cm(7.5), Cm(10), Cm(0.5))
    ttf = table_title.text_frame
    ttf.text = "📊 연간 비용 비교"
    ttp = ttf.paragraphs[0]
    ttp.font.size = Pt(10)
    ttp.font.bold = True
    ttp.font.color.rgb = THEME['dark']

    table = slide.shapes.add_table(3, 4, Cm(1), Cm(8.2), Cm(19), Cm(2.5)).table

    headers = ["항목", "GPT-4 사용 시", "DeepSeek 사용 시", "절감액"]
    data = [
        ["LLM API 비용", "$100,000", "$1,800", "$98,200 (98%)"],
        ["총 운영 비용", "$114,000", "$14,040", "$99,960 (88%)"],
    ]

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME['secondary']
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = THEME['white']
        p.alignment = PP_ALIGN.CENTER

    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = THEME['light']
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.font.color.rgb = THEME['text']
            p.alignment = PP_ALIGN.CENTER

    # 우측 상세 설명
    add_detailed_description(slide, [
        ("1. 응답 시간 < 3초 (P95)", [
            "Hybrid 검색(ES + Neo4j 병렬) 실행 시간: ~1초",
            "DeepSeek 답변 합성 시간: ~1.5초",
            "네트워크 + 기타 오버헤드: ~0.5초",
            "전체 P95 기준 3초 이내 완료 목표",
        ]),
        ("2. 검색 정확도 > 85% (Precision@5)", [
            "상위 5개 검색 결과 중 관련 문서 비율",
            "Vector Search만 사용 시: ~70%",
            "Graph Search 융합 후: 85%+ 목표",
            "평가 방법: 수동 레이블링 테스트셋 200건",
        ]),
        ("3. LLM 비용 95% 절감", [
            "GPT-4: $30/1M input + $60/1M output 토큰",
            "DeepSeek: $0.14/1M input + $0.28/1M output",
            "동일 사용량 기준 연간 $98,200 절감",
            "품질 비교: MMLU 벤치마크 동등 수준",
        ]),
        ("4. 가용성 99.5% (SLA)", [
            "연간 허용 다운타임: 43.8시간",
            "월간 허용 다운타임: 3.65시간",
            "헬스체크 주기: 30초",
            "장애 복구 목표 시간(RTO): 30분",
        ]),
    ])

    return slide


def create_architecture_slide(prs):
    """시스템 아키텍처 - 레이어 구분 명확화"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "시스템 아키텍처 (System Architecture)")

    # ===== External Layer =====
    add_layer_box(slide, 0.5, 2.8, 20.5, 2, "🌐 External Layer (외부 접점)")
    add_box(slide, 1, 3.5, 2.8, 1, "👤 사용자\nWeb Browser", THEME['accent'], font_size=8)
    add_arrow_text(slide, 3.8, 3.7, "→")
    add_box(slide, 4.5, 3.5, 3, 1, "🚪 Nginx\nSSL Termination\nReverse Proxy", THEME['secondary'], font_size=8)
    add_arrow_text(slide, 7.5, 3.7, "→")
    add_box(slide, 8.2, 3.5, 3.5, 1, "🔀 API Gateway\nJWT 검증\nRate Limiting", THEME['secondary'], font_size=8)
    add_arrow_text(slide, 11.7, 3.7, "→")
    add_box(slide, 12.5, 3.5, 3, 1, "🔐 Keycloak\nOAuth 2.0 IdP\nPKCE 인증", RGBColor(183, 28, 28), font_size=8)

    # DeepSeek (외부 API)
    add_box(slide, 17, 3.5, 3.5, 1, "🤖 DeepSeek API\nLLM 호출\n(External)", RGBColor(26, 115, 232), font_size=8)

    # ===== Application Layer =====
    add_layer_box(slide, 0.5, 5, 20.5, 2.3, "⚙️ Application Layer (애플리케이션 계층)")
    add_box(slide, 1, 5.8, 3.5, 1.2, "🖥️ Frontend\nReact 18 + TypeScript\nMUI v5 / Redux", RGBColor(97, 218, 251), THEME['dark'], font_size=8)
    add_arrow_text(slide, 4.5, 6.1, "→")
    add_box(slide, 5.3, 5.8, 4, 1.2, "⚙️ Backend\nSpring Boot 3.x\nSecurity / JPA\nResilience4j", RGBColor(109, 179, 63), THEME['white'], font_size=8)
    add_arrow_text(slide, 9.3, 6.1, "→")
    add_box(slide, 10, 5.8, 4.5, 1.2, "🧠 AI Service\nFastAPI + LangGraph\nBGE-M3 Embedding\nHybrid Search", RGBColor(0, 150, 136), THEME['white'], font_size=8)
    add_arrow_text(slide, 14.5, 6.1, "→")

    # LLM 연결 표시
    llm_conn = slide.shapes.add_textbox(Cm(15), Cm(5.5), Cm(5), Cm(0.4))
    lcf = llm_conn.text_frame
    lcf.text = "→ DeepSeek API 호출 ↗"
    lcp = lcf.paragraphs[0]
    lcp.font.size = Pt(7)
    lcp.font.color.rgb = THEME['info']

    # ===== Data Layer =====
    add_layer_box(slide, 0.5, 7.5, 20.5, 2.5, "💾 Data Layer (데이터 계층)")
    add_box(slide, 1, 8.3, 3, 1.3, "🐘 PostgreSQL 16\nSSOT (원본 데이터)\n사용자/문서/설정", RGBColor(51, 103, 145), font_size=8)
    add_box(slide, 4.3, 8.3, 3.5, 1.3, "🔍 Elasticsearch 8\nVector Search\n1024 dims 임베딩", RGBColor(0, 175, 170), font_size=8)
    add_box(slide, 8, 8.3, 3, 1.3, "🕸️ Neo4j 5\nGraph DB\n엔티티/관계 저장", RGBColor(0, 110, 170), font_size=8)
    add_box(slide, 11.2, 8.3, 2.5, 1.3, "💾 Redis 7\nCache\nSession 관리", RGBColor(220, 63, 66), font_size=8)
    add_box(slide, 14, 8.3, 2.8, 1.3, "📦 MinIO\nFile Storage\nS3 호환 API", RGBColor(198, 40, 40), font_size=8)
    add_box(slide, 17, 8.3, 3.5, 1.3, "🔑 Vault\nKey Management\n암호화 키 관리", RGBColor(100, 100, 100), font_size=8)

    # ===== Monitoring Layer =====
    add_layer_box(slide, 0.5, 10.2, 14, 2, "📊 Monitoring Layer (모니터링 계층)")
    add_box(slide, 1, 10.9, 3, 1, "📊 Prometheus\nMetrics 수집\nAlertmanager", RGBColor(230, 84, 49), font_size=8)
    add_box(slide, 4.3, 10.9, 3, 1, "📈 Grafana\nDashboard\n시각화", RGBColor(241, 144, 13), THEME['dark'], font_size=8)
    add_box(slide, 7.6, 10.9, 3, 1, "📝 Loki\nLog 집계\n중앙화 로깅", RGBColor(140, 86, 75), font_size=8)
    add_box(slide, 10.9, 10.9, 2.8, 1, "🔍 Promtail\nLog 수집\n에이전트", RGBColor(170, 116, 95), font_size=8)

    # 우측 상세 설명
    add_detailed_description(slide, [
        ("1. External Layer (외부 접점)", [
            "Nginx: SSL 인증서 관리, 정적 파일 서빙, 리버스 프록시",
            "API Gateway: JWT 토큰 검증, Rate Limiting(분당 100회)",
            "Keycloak: OAuth 2.0 + PKCE 인증, 사용자 관리, SSO 지원",
            "DeepSeek API: 외부 LLM 서비스, AI Service에서만 호출",
        ]),
        ("2. Application Layer (애플리케이션)", [
            "Frontend: React SPA, MUI 컴포넌트, Redux 상태관리",
            "Backend: 비즈니스 로직, 트랜잭션 관리, AI Service 연동",
            "AI Service: 임베딩 생성, Hybrid 검색, 답변 합성",
            "서비스 분리: SpringBoot ↔ AI Service REST API 통신",
        ]),
        ("3. Data Layer (데이터 계층)", [
            "PostgreSQL: SSOT(Single Source of Truth), 마스터 데이터",
            "Elasticsearch: Vector Search용 비정규화 인덱스",
            "Neo4j: 엔티티-관계 그래프, Graph Search용",
            "Redis: 세션 캐시, API 응답 캐시, 분산 락",
            "MinIO: 문서 파일 저장소, S3 호환 API 제공",
        ]),
        ("4. Monitoring Layer (모니터링)", [
            "Prometheus: 메트릭 수집 (CPU, Memory, API Latency)",
            "Grafana: 대시보드 시각화, 알림 설정",
            "Loki: 중앙화 로그 집계, 로그 검색",
            "Alertmanager: 임계치 초과 시 Slack/Email 알림",
        ]),
        ("5. 핵심 설계 원칙", [
            "Zero-Join: 검색 시 DB 조인 없이 단일 인덱스 조회",
            "비정규화: PostgreSQL → ES/Neo4j 데이터 복제",
            "서비스 분리: SpringBoot는 LLM 직접 호출 금지",
            "장애 격리: Circuit Breaker로 AI Service 장애 전파 방지",
        ]),
    ])

    return slide


def create_vip_architecture_slide(prs):
    """VIP 3단계 LLM 아키텍처"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "VIP 3단계 LLM 아키텍처 (Value-Intelligent-Planning)")

    # Stage 배경 박스들
    add_layer_box(slide, 0.5, 3, 6, 8, "Stage 1: Value (가치 추출)")
    add_layer_box(slide, 7, 3, 6, 8, "Stage 2: Intelligent (지능형 검색)")
    add_layer_box(slide, 13.5, 3, 6, 8, "Stage 3: Planning (답변 계획)")

    # Stage 1 내용
    add_box(slide, 1, 4, 5, 0.8, "📄 문서 업로드", THEME['secondary'], font_size=9)
    add_arrow_text(slide, 3, 4.9, "↓")
    add_box(slide, 1, 5.3, 5, 0.8, "🔧 Docling 파싱", RGBColor(100, 149, 237), font_size=9)
    add_arrow_text(slide, 3, 6.2, "↓")
    add_box(slide, 1, 6.6, 5, 0.8, "✂️ 청킹 (500자)", RGBColor(100, 149, 237), font_size=9)
    add_arrow_text(slide, 3, 7.5, "↓")
    add_box(slide, 1, 7.9, 5, 0.8, "🤖 DeepSeek 엔티티 추출", THEME['info'], font_size=9)
    add_arrow_text(slide, 3, 8.8, "↓")
    add_box(slide, 1, 9.2, 5, 0.8, "💎 지식 그래프 구축", THEME['success'], font_size=9)

    # Stage 2 내용
    add_box(slide, 7.5, 4, 5, 0.8, "❓ 사용자 질문 입력", THEME['secondary'], font_size=9)
    add_arrow_text(slide, 9.5, 4.9, "↓")
    add_box(slide, 7.5, 5.3, 5, 0.8, "🧠 질의 의도 분석", RGBColor(100, 149, 237), font_size=9)
    add_arrow_text(slide, 9.5, 6.2, "↓")
    add_box(slide, 7.5, 6.6, 5, 0.8, "🎯 검색 전략 수립", RGBColor(100, 149, 237), font_size=9)
    add_arrow_text(slide, 9.5, 7.5, "↓")
    add_box(slide, 7.5, 7.9, 5, 0.8, "🔍 Vector Search", RGBColor(249, 183, 22), THEME['dark'], font_size=9)
    add_box(slide, 7.5, 8.8, 5, 0.8, "🕸️ Graph Search", RGBColor(1, 139, 255), font_size=9)
    add_box(slide, 7.5, 9.7, 5, 0.5, "⚡ 병렬 실행", THEME['highlight'], THEME['dark'], font_size=8)

    # Stage 3 내용
    add_box(slide, 14, 4, 5, 0.8, "📊 검색 결과 수신", THEME['secondary'], font_size=9)
    add_arrow_text(slide, 16, 4.9, "↓")
    add_box(slide, 14, 5.3, 5, 0.8, "⚡ RRF 융합", RGBColor(255, 107, 107), font_size=9)
    add_arrow_text(slide, 16, 6.2, "↓")
    add_box(slide, 14, 6.6, 5, 0.8, "📝 컨텍스트 구성", RGBColor(100, 149, 237), font_size=9)
    add_arrow_text(slide, 16, 7.5, "↓")
    add_box(slide, 14, 7.9, 5, 0.8, "🤖 DeepSeek 답변 합성", THEME['info'], font_size=9)
    add_arrow_text(slide, 16, 8.8, "↓")
    add_box(slide, 14, 9.2, 5, 0.8, "💬 최종 답변 + 출처", THEME['success'], font_size=9)

    # 화살표 (Stage 간)
    arr1 = slide.shapes.add_textbox(Cm(6), Cm(6.5), Cm(1.2), Cm(0.5))
    a1f = arr1.text_frame
    a1f.text = "→"
    a1p = a1f.paragraphs[0]
    a1p.font.size = Pt(20)
    a1p.font.bold = True
    a1p.font.color.rgb = THEME['accent']

    arr2 = slide.shapes.add_textbox(Cm(12.5), Cm(6.5), Cm(1.2), Cm(0.5))
    a2f = arr2.text_frame
    a2f.text = "→"
    a2p = a2f.paragraphs[0]
    a2p.font.size = Pt(20)
    a2p.font.bold = True
    a2p.font.color.rgb = THEME['accent']

    # 우측 상세 설명
    add_detailed_description(slide, [
        ("1. Value 단계 (오프라인 처리)", [
            "문서 업로드 시 1회 실행되는 오프라인 파이프라인",
            "Docling: PDF/DOCX/PPT 등 다양한 포맷 파싱",
            "청킹: 500자 단위 + 50자 오버랩으로 분할",
            "DeepSeek 추출: 엔티티(Person, Project, Tech) + 관계",
            "결과 저장: Elasticsearch(벡터), Neo4j(그래프)",
        ]),
        ("2. Intelligent 단계 (온라인 처리)", [
            "사용자 질문 입력 시 실시간 실행",
            "의도 분석: 사실 질의 vs 비교 질의 vs 요약 질의 분류",
            "검색 전략: Vector/Graph 가중치 동적 결정",
            "Vector Search: BGE-M3 임베딩으로 의미 유사도 검색",
            "Graph Search: Cypher 쿼리로 엔티티 관계 탐색",
            "병렬 실행: ES + Neo4j 동시 조회로 응답 시간 최소화",
        ]),
        ("3. Planning 단계 (답변 생성)", [
            "RRF 융합: Reciprocal Rank Fusion으로 결과 통합",
            "Top-10 선정: 최종 컨텍스트 구성용 문서 선별",
            "프롬프트 구성: 시스템 지시 + 컨텍스트 + 질문",
            "DeepSeek 호출: 답변 합성 요청 (최대 4K 토큰)",
            "출처 첨부: 참조 문서 목록과 함께 응답",
        ]),
        ("4. 비용 최적화 포인트", [
            "단일 모델(DeepSeek)로 추출 + 합성 모두 처리",
            "오프라인 추출: 업로드 시 1회만 실행",
            "캐싱: 동일 질의 결과 Redis 캐시 (TTL 1시간)",
            "결과: GPT-4 대비 95% 비용 절감",
        ]),
    ])

    return slide


def create_hybrid_search_slide(prs):
    """Hybrid Search 융합"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Hybrid Search 융합 (Vector + Graph)")

    # 사용자 질의
    add_box(slide, 7, 3, 6, 1.2, "❓ 사용자 질의 입력\n\"프로젝트 A의 담당자가 작성한 기술 문서는?\"", THEME['dark'], font_size=9)

    # 분기
    add_arrow_text(slide, 5.5, 4.3, "↙")
    add_arrow_text(slide, 12.5, 4.3, "↘")

    # Vector Search 박스
    add_layer_box(slide, 0.5, 5, 9, 4, "🔍 Vector Search (의미 기반)")
    add_box(slide, 1, 5.8, 8, 0.7, "1. BGE-M3 임베딩 생성 (1024차원)", RGBColor(249, 183, 22), THEME['dark'], font_size=8)
    add_box(slide, 1, 6.7, 8, 0.7, "2. Elasticsearch knn 검색", RGBColor(249, 183, 22), THEME['dark'], font_size=8)
    add_box(slide, 1, 7.6, 8, 0.7, "3. 의미적 유사도 Top-20 반환", RGBColor(249, 183, 22), THEME['dark'], font_size=8)

    # Graph Search 박스
    add_layer_box(slide, 10, 5, 9.5, 4, "🕸️ Graph Search (관계 기반)")
    add_box(slide, 10.5, 5.8, 8.5, 0.7, "1. 엔티티 추출 (프로젝트 A, 담당자)", RGBColor(1, 139, 255), font_size=8)
    add_box(slide, 10.5, 6.7, 8.5, 0.7, "2. Neo4j Cypher 관계 탐색", RGBColor(1, 139, 255), font_size=8)
    add_box(slide, 10.5, 7.6, 8.5, 0.7, "3. 연관 문서 Top-20 반환", RGBColor(1, 139, 255), font_size=8)

    # 병렬 표시
    parallel = slide.shapes.add_textbox(Cm(8), Cm(6.5), Cm(4), Cm(0.5))
    pf = parallel.text_frame
    pf.text = "⚡ 병렬 실행"
    pp = pf.paragraphs[0]
    pp.font.size = Pt(11)
    pp.font.bold = True
    pp.font.color.rgb = THEME['highlight']
    pp.alignment = PP_ALIGN.CENTER

    # 수렴
    add_arrow_text(slide, 5.5, 9, "↘")
    add_arrow_text(slide, 12.5, 9, "↙")

    # RRF 융합
    add_box(slide, 5.5, 9.8, 9, 1.3, "⚡ RRF (Reciprocal Rank Fusion)\nscore = Σ 1/(k+rank), k=60", RGBColor(255, 107, 107), font_size=10)

    # 결과
    add_box(slide, 5.5, 11.3, 9, 0.9, "📊 통합 검색 결과 (Top-10) → 답변 생성", THEME['success'], font_size=10, bold=True)

    # 우측 상세 설명
    add_detailed_description(slide, [
        ("1. Vector Search 특징", [
            "BGE-M3: 다국어 지원 임베딩 모델 (1024차원)",
            "Elasticsearch knn: 근사 최근접 이웃 검색 (HNSW)",
            "장점: 동의어, 유사 표현, 의미적 유사 문서 검색",
            "약점: '프로젝트 A 담당자' 같은 관계 질의에 부족",
            "결과: 의미적으로 유사한 문서 Top-20 반환",
        ]),
        ("2. Graph Search 특징", [
            "엔티티 추출: 질문에서 Person, Project, Tech 식별",
            "Neo4j Cypher: 엔티티 간 관계 탐색 쿼리 실행",
            "장점: '담당자가 작성한' 같은 관계 기반 검색 강점",
            "약점: 유사 표현, 동의어 검색에 제한",
            "결과: 관계로 연결된 문서 Top-20 반환",
        ]),
        ("3. RRF 융합 알고리즘", [
            "공식: score(d) = Σ 1/(k + rank(d))",
            "k=60: 순위 차이의 민감도 조절 파라미터",
            "예: Vector 1위 + Graph 3위 = 1/61 + 1/63 = 0.032",
            "효과: 두 검색에서 모두 상위인 문서 우선 선정",
            "출력: 융합 점수 기준 Top-10 최종 선정",
        ]),
        ("4. 융합의 효과", [
            "Vector만 사용: Precision@5 약 70%",
            "Graph만 사용: Precision@5 약 65%",
            "Hybrid 융합: Precision@5 85%+ 목표",
            "상호 보완: Vector가 놓친 관계 문서를 Graph가 보완",
            "정확도 향상: 단일 검색 대비 20%+ 정확도 개선",
        ]),
    ])

    return slide


def create_data_flow_slide(prs):
    """데이터 흐름"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "검색 데이터 흐름 (Search Flow)")

    # 상단 흐름
    steps = [
        ("① 사용자", "질문 입력", THEME['accent'], 1),
        ("② Frontend", "API 호출", RGBColor(97, 218, 251), 4.2),
        ("③ Gateway", "JWT 검증", THEME['secondary'], 7.4),
        ("④ Backend", "요청 전달", RGBColor(109, 179, 63), 10.6),
        ("⑤ AI Service", "검색 실행", RGBColor(0, 150, 136), 13.8),
        ("⑥ DeepSeek", "답변 생성", RGBColor(26, 115, 232), 17),
    ]

    for title, desc, color, x in steps:
        add_box(slide, x, 3.2, 2.8, 1.8, f"{title}\n\n{desc}", color, font_size=8)

    for i in range(5):
        add_arrow_text(slide, 3.8 + i * 3.2, 3.8, "→")

    # 병렬 검색 박스
    add_layer_box(slide, 0.5, 5.5, 20, 5.2, "⚡ AI Service 내부 처리 (병렬 검색 + RRF 융합 + 답변 합성)")

    # ES
    add_box(slide, 1, 6.5, 4.5, 1.8, "🔍 Elasticsearch\n\nVector Search (knn)\nBGE-M3 유사도\nTop-20 반환", RGBColor(249, 183, 22), THEME['dark'], font_size=8, align_center=False)

    # Neo4j
    add_box(slide, 6, 6.5, 4.5, 1.8, "🕸️ Neo4j\n\nGraph Search\nCypher 관계 탐색\nTop-20 반환", RGBColor(1, 139, 255), font_size=8, align_center=False)

    # 병렬 표시
    parallel = slide.shapes.add_textbox(Cm(4), Cm(8.5), Cm(4), Cm(0.5))
    pf = parallel.text_frame
    pf.text = "병렬 실행"
    pp = pf.paragraphs[0]
    pp.font.size = Pt(8)
    pp.font.bold = True
    pp.font.color.rgb = THEME['highlight']

    add_arrow_text(slide, 10.5, 7, "→")

    # RRF
    add_box(slide, 11, 6.5, 4, 1.8, "⚡ RRF 융합\n\nReciprocal Rank\nTop-10 선정\n컨텍스트 구성", RGBColor(255, 107, 107), font_size=8, align_center=False)

    add_arrow_text(slide, 15, 7, "→")

    # DeepSeek
    add_box(slide, 15.5, 6.5, 4.5, 1.8, "🤖 DeepSeek\n\n답변 합성\n출처 첨부\n최종 응답", RGBColor(76, 175, 80), font_size=8, align_center=False)

    # 응답 흐름
    resp_label = slide.shapes.add_textbox(Cm(1), Cm(9), Cm(18), Cm(0.5))
    rlf = resp_label.text_frame
    rlf.text = "← 응답 흐름: DeepSeek → AI Service → Backend → Gateway → Frontend → 사용자"
    rlp = rlf.paragraphs[0]
    rlp.font.size = Pt(9)
    rlp.font.color.rgb = THEME['secondary']

    # 우측 상세 설명
    add_detailed_description(slide, [
        ("1. 요청 흐름 상세", [
            "사용자: 브라우저에서 질문 입력 및 전송",
            "Frontend: React Query로 POST /api/v1/search/chat",
            "Gateway: JWT 토큰 검증, Rate Limiting 체크",
            "Backend: 요청 로깅, AI Service로 전달",
            "AI Service: Hybrid 검색 및 답변 생성 수행",
        ]),
        ("2. Gateway 역할 상세", [
            "JWT 검증: RS256 서명 확인, 만료 체크",
            "Rate Limiting: 사용자당 분당 100회 제한",
            "요청 라우팅: /api/* → Backend, /auth/* → Keycloak",
            "로깅: 요청/응답 메타데이터 기록",
        ]),
        ("3. 병렬 검색 상세", [
            "Elasticsearch: knn 검색 (약 300ms)",
            "Neo4j: Cypher 쿼리 (약 400ms)",
            "병렬 실행: 총 소요 시간 max(300, 400) = 400ms",
            "순차 실행 대비 약 40% 시간 절감",
        ]),
        ("4. RRF 융합 상세", [
            "입력: ES Top-20 + Neo4j Top-20",
            "점수 계산: 1/(60+rank) 합산",
            "중복 제거: 동일 문서 점수 합산",
            "출력: 융합 점수 상위 Top-10 선정",
        ]),
        ("5. 답변 생성 상세", [
            "컨텍스트: Top-10 문서 텍스트 (최대 8K 토큰)",
            "프롬프트: 시스템 지시 + 컨텍스트 + 질문",
            "DeepSeek 호출: 약 1.5초 소요",
            "후처리: 출처 문서 목록 첨부",
        ]),
    ])

    return slide


def create_advantages_slide(prs):
    """아키텍처 장점"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "아키텍처 장점 (Architecture Advantages)")

    advantages = [
        ("💰 비용 효율성 (95% 절감)", [
            "DeepSeek V3.2: GPT-4 대비 토큰당 95% 저렴",
            "연간 LLM 비용: $100,000 → $1,800 절감",
            "인프라 비용: K8s 대비 86% 절감 ($100K → $14K)",
            "단일 모델로 엔티티 추출 + 답변 합성 처리",
        ], THEME['success']),
        ("🎯 검색 정확도 향상", [
            "Vector Search: 의미적 유사도 기반 검색",
            "Graph Search: 엔티티 관계 기반 확장 검색",
            "RRF 융합: 두 검색 결과 최적 통합",
            "목표: Precision@5 > 85%, 기존 대비 40%↑",
        ], THEME['info']),
        ("🔐 기업 보안 강화", [
            "Keycloak: OAuth 2.0 + PKCE 인증",
            "AES-256-GCM: 민감 데이터 필드 암호화",
            "RBAC: 역할 기반 접근 제어 (Admin/Manager/User)",
            "TLS 1.3 + Vault: 전송 암호화 + 키 관리",
        ], RGBColor(156, 39, 176)),
        ("📈 점진적 확장 전략", [
            "Phase 1: Docker Compose 단일 서버 ($14K/년)",
            "Phase 2: Redis Sentinel 캐시 이중화 ($20K/년)",
            "Phase 3: Docker Swarm 3노드 ($40K/년)",
            "Phase 4: Kubernetes 풀 클러스터 ($100K/년)",
        ], THEME['highlight']),
    ]

    y = 3
    for title, items, color in advantages:
        add_box(slide, 0.5, y, 5.5, 0.8, title, color, font_size=9, bold=True)

        for i, item in enumerate(items):
            item_box = slide.shapes.add_textbox(Cm(6.3), Cm(y + 0.1 + i * 0.5), Cm(13), Cm(0.5))
            itf = item_box.text_frame
            itf.text = f"• {item}"
            itp = itf.paragraphs[0]
            itp.font.size = Pt(8)
            itp.font.color.rgb = THEME['text']

        y += 2.5

    # 우측 상세 설명
    add_detailed_description(slide, [
        ("1. 비용 효율성 핵심", [
            "DeepSeek V3.2: $0.14/1M input, $0.28/1M output",
            "GPT-4: $30/1M input, $60/1M output",
            "동일 사용량 기준 연간 $98,200 절감",
            "MMLU 벤치마크: DeepSeek ≈ GPT-4 품질",
        ]),
        ("2. Hybrid 검색 장점", [
            "Vector만: 관계 질의 ('A의 담당자') 검색 불가",
            "Graph만: 유사 표현 ('책임자' ≈ '담당자') 검색 불가",
            "융합 효과: 상호 보완으로 정확도 20%+ 향상",
            "RRF: 두 검색에서 모두 상위인 문서 우선 선정",
        ]),
        ("3. 보안 설계 원칙", [
            "Zero Trust: 모든 요청에 인증/인가 검증",
            "PKCE: Authorization Code 탈취 공격 방지",
            "계층 암호화: 전송(TLS) + 필드(AES) + 저장(TDE)",
            "키 관리: HashiCorp Vault 중앙 집중 관리",
        ]),
        ("4. 확장 트리거 조건", [
            "Phase 2: 동시 사용자 100+, 일일 요청 10,000+",
            "Phase 3: 동시 사용자 500+, 일일 요청 50,000+",
            "Phase 4: 동시 사용자 1,000+, 일일 요청 200,000+",
            "원칙: 실제 트래픽 기반 점진적 투자",
        ]),
    ])

    return slide


def create_security_slide(prs):
    """보안 아키텍처"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "인증 및 보안 설계 (Security Architecture)")

    # 인증
    add_layer_box(slide, 0.5, 3, 6.3, 3.5, "🔐 인증 (Authentication)")
    items = ["OAuth 2.0 Authorization Code Flow", "PKCE (Proof Key for Code Exchange)", "Keycloak IdP 연동", "JWT Access Token (RS256)"]
    for i, item in enumerate(items):
        add_box(slide, 1, 3.8 + i * 0.7, 5.3, 0.55, f"• {item}", THEME['secondary'], font_size=7, align_center=False)

    # 권한
    add_layer_box(slide, 7, 3, 6.3, 3.5, "👤 권한 (Authorization)")
    items = ["RBAC (Role-Based Access Control)", "3단계: ADMIN / MANAGER / USER", "API 레벨 권한 검증", "리소스 기반 접근 제어"]
    for i, item in enumerate(items):
        add_box(slide, 7.5, 3.8 + i * 0.7, 5.3, 0.55, f"• {item}", RGBColor(156, 39, 176), font_size=7, align_center=False)

    # 암호화
    add_layer_box(slide, 13.5, 3, 6, 3.5, "🔒 암호화 (Encryption)")
    items = ["TLS 1.3 (전송 구간)", "AES-256-GCM (필드)", "bcrypt (비밀번호)", "Vault (키 관리)"]
    for i, item in enumerate(items):
        add_box(slide, 14, 3.8 + i * 0.7, 5, 0.55, f"• {item}", RGBColor(198, 40, 40), font_size=7, align_center=False)

    # JWT 전략
    add_layer_box(slide, 0.5, 6.8, 9, 2.5, "🎫 JWT 토큰 전략")
    jwt_items = [
        "Access Token: 15분, 메모리 저장, 만료 2분 전 자동 갱신",
        "Refresh Token: 7일, HttpOnly Cookie, Rotation 적용",
        "토큰 구조: Header.Payload.Signature (RS256 서명)",
    ]
    for i, item in enumerate(jwt_items):
        tb = slide.shapes.add_textbox(Cm(1), Cm(7.6 + i * 0.5), Cm(8.5), Cm(0.5))
        tf = tb.text_frame
        tf.text = f"• {item}"
        tp = tf.paragraphs[0]
        tp.font.size = Pt(7)
        tp.font.color.rgb = THEME['text']

    # 데이터 분류
    add_layer_box(slide, 10, 6.8, 9.5, 2.5, "📊 데이터 분류 체계")
    data_items = [
        "Level 4 (극비): 암호화 키 → HSM/Vault",
        "Level 3 (비밀): API 키, 토큰 → DB 암호화",
        "Level 2 (대외비): 이름, 이메일 → 필드 암호화",
    ]
    for i, item in enumerate(data_items):
        tb = slide.shapes.add_textbox(Cm(10.5), Cm(7.6 + i * 0.5), Cm(8.5), Cm(0.5))
        tf = tb.text_frame
        tf.text = f"• {item}"
        tp = tf.paragraphs[0]
        tp.font.size = Pt(7)
        tp.font.color.rgb = THEME['text']

    # RBAC 테이블
    table = slide.shapes.add_table(4, 5, Cm(0.5), Cm(9.8), Cm(19), Cm(2.3)).table
    headers = ["역할", "사용자 관리", "지식 CRUD", "지식 조회", "북마크"]
    data = [["ADMIN", "✅", "✅", "✅", "✅"], ["MANAGER", "❌", "✅", "✅", "✅"], ["USER", "❌", "❌", "✅", "✅"]]

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME['secondary']
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = THEME['white']
        p.alignment = PP_ALIGN.CENTER

    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.font.color.rgb = THEME['text']
            p.alignment = PP_ALIGN.CENTER

    # 우측 상세 설명
    add_detailed_description(slide, [
        ("1. OAuth 2.0 + PKCE 상세", [
            "Authorization Code Flow: 가장 안전한 OAuth 2.0 방식",
            "PKCE: code_verifier 해시로 인증 코드 탈취 방지",
            "Keycloak: 오픈소스 IdP, 사용자/역할 관리 UI 제공",
            "적용: 로그인 시 PKCE 파라미터 자동 생성/검증",
        ]),
        ("2. JWT 토큰 보안 상세", [
            "Access Token 메모리 저장: XSS로 탈취 시 자동 만료",
            "Refresh Token HttpOnly: JavaScript 접근 불가",
            "Rotation: Refresh 사용 시 새 토큰 발급, 기존 무효화",
            "RS256: 공개키로 검증, 비밀키는 Keycloak만 보유",
        ]),
        ("3. 암호화 계층 상세", [
            "TLS 1.3: 모든 네트워크 통신 암호화",
            "AES-256-GCM: 인증 태그 포함, 무결성 보장",
            "bcrypt: cost factor 12, 레인보우 테이블 공격 방어",
            "Vault: 암호화 키 중앙 관리, 자동 키 로테이션",
        ]),
        ("4. RBAC 역할 상세", [
            "ADMIN: 시스템 전체 관리, 사용자 생성/삭제",
            "MANAGER: 지식 CRUD, 부서별 문서 관리",
            "USER: 검색/조회만 가능, 북마크 개인 관리",
            "원칙: 최소 권한 부여, 필요 시 역할 승격",
        ]),
    ])

    return slide


def create_infra_slide(prs):
    """인프라 구성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "인프라 구성 (Docker Compose 기반)")

    # Docker Host
    add_layer_box(slide, 0.5, 2.8, 19.5, 9.5, "🐳 Docker Host (Production: 32C / 128GB RAM / 1TB SSD)")

    # Application
    add_layer_box(slide, 1, 3.5, 18.5, 2, "⚙️ Application Layer (6개)")
    apps = [("nginx", 1.3), ("frontend", 4.3), ("gateway", 7.3), ("backend", 10.3), ("ai-service", 13.3), ("keycloak", 16.3)]
    for name, x in apps:
        add_box(slide, x, 4.2, 2.7, 0.8, name, RGBColor(129, 236, 236), THEME['dark'], font_size=8)

    # Database
    add_layer_box(slide, 1, 5.7, 18.5, 2, "💾 Data Layer (6개)")
    dbs = [("postgresql", 1.3), ("elasticsearch", 4.3), ("neo4j", 7.3), ("redis", 10.3), ("minio", 13.3), ("vault", 16.3)]
    for name, x in dbs:
        add_box(slide, x, 6.4, 2.7, 0.8, name, RGBColor(162, 155, 254), font_size=8)

    # Monitoring
    add_layer_box(slide, 1, 7.9, 13, 2, "📊 Monitoring Layer (4개)")
    mons = [("prometheus", 1.3), ("grafana", 4.3), ("loki", 7.3), ("promtail", 10.3)]
    for name, x in mons:
        add_box(slide, x, 8.6, 2.7, 0.8, name, RGBColor(253, 121, 168), font_size=8)

    # 사양 요약
    spec_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(10.3), Cm(18.5), Cm(1.8)
    )
    spec_box.fill.solid()
    spec_box.fill.fore_color.rgb = THEME['secondary']
    spec_box.line.fill.background()

    spec_text = slide.shapes.add_textbox(Cm(1.3), Cm(10.5), Cm(18), Cm(1.5))
    stf = spec_text.text_frame
    stf.text = "📦 총 16개 컨테이너 | Production: 32C/128GB ($800/월) | Staging: 16C/64GB ($200/월) | 연간 총비용: $14,040"
    stp = stf.paragraphs[0]
    stp.font.size = Pt(10)
    stp.font.color.rgb = THEME['white']
    stp.alignment = PP_ALIGN.CENTER

    # 우측 상세 설명
    add_detailed_description(slide, [
        ("1. Application Layer 상세", [
            "nginx: SSL 종단, 정적 파일 서빙, 리버스 프록시",
            "frontend: React SPA 빌드 결과물 서빙",
            "gateway: Spring Cloud Gateway, JWT 검증",
            "backend: Spring Boot 3.x, 비즈니스 로직",
            "ai-service: FastAPI, Hybrid 검색, 답변 생성",
            "keycloak: OAuth 2.0 IdP, 사용자 관리",
        ]),
        ("2. Data Layer 상세", [
            "postgresql: SSOT, 사용자/문서/설정 데이터",
            "elasticsearch: Vector Search, 1024차원 임베딩",
            "neo4j: Graph DB, 엔티티-관계 저장",
            "redis: 세션 캐시, API 응답 캐시",
            "minio: 문서 파일 저장, S3 호환",
            "vault: 암호화 키 관리, 비밀 저장",
        ]),
        ("3. Monitoring Layer 상세", [
            "prometheus: 메트릭 수집 (15초 주기)",
            "grafana: 대시보드 (CPU, Memory, Latency)",
            "loki: 중앙화 로그 집계, LogQL 검색",
            "promtail: 컨테이너 로그 수집 에이전트",
        ]),
        ("4. 리소스 할당", [
            "postgresql: 4GB RAM, 100GB 디스크",
            "elasticsearch: 8GB RAM, 200GB 디스크",
            "neo4j: 4GB RAM, 50GB 디스크",
            "ai-service: 16GB RAM (임베딩 모델 로드)",
        ]),
        ("5. 확장 경로", [
            "Phase 2: Redis Sentinel 3노드 캐시 이중화",
            "Phase 3: Docker Swarm 3노드 클러스터",
            "Phase 4: Kubernetes + Helm 차트 배포",
        ]),
    ])

    return slide


def create_cicd_slide(prs):
    """CI/CD 파이프라인"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "CI/CD 파이프라인 (GitLab)")

    # CI
    add_layer_box(slide, 0.5, 3, 19, 2.3, "🔄 CI (Continuous Integration)")
    ci_steps = [("Git Push", THEME['accent']), ("Build", THEME['info']), ("Unit Test", THEME['success']), ("Security Scan", RGBColor(156, 39, 176)), ("SonarQube", THEME['highlight'])]
    x = 1
    for name, color in ci_steps:
        add_box(slide, x, 3.8, 3.3, 1, name, color, font_size=9)
        if x < 15:
            add_arrow_text(slide, x + 3.3, 4, "→")
        x += 3.7

    # CD
    add_layer_box(slide, 0.5, 5.5, 19, 2.3, "🚀 CD (Continuous Deployment)")
    cd_steps = [("Docker Image", THEME['info']), ("Push Registry", THEME['secondary']), ("Staging", THEME['success']), ("Approve", THEME['highlight']), ("Production", THEME['danger'])]
    x = 1
    for name, color in cd_steps:
        add_box(slide, x, 6.3, 3.3, 1, name, color, font_size=9)
        if x < 15:
            add_arrow_text(slide, x + 3.3, 6.5, "→")
        x += 3.7

    # 품질 게이트
    add_layer_box(slide, 0.5, 8, 19, 1.5, "🚦 품질 게이트")
    gate_text = slide.shapes.add_textbox(Cm(1), Cm(8.4), Cm(18.5), Cm(0.8))
    gtf = gate_text.text_frame
    gtf.text = "테스트 커버리지 80%+ | 보안 취약점 High 0건 | SonarQube Quality Gate Pass | 코드 리뷰 승인"
    gtp = gtf.paragraphs[0]
    gtp.font.size = Pt(9)
    gtp.font.color.rgb = THEME['dark']

    # 브랜치 전략
    add_layer_box(slide, 0.5, 9.7, 19, 1.5, "🌿 브랜치 전략")
    branch_text = slide.shapes.add_textbox(Cm(1), Cm(10.1), Cm(18.5), Cm(0.8))
    btf = branch_text.text_frame
    btf.text = "main (프로덕션) | develop (개발 통합) | feature/* (새 기능) | fix/* (버그 수정) | hotfix/* (긴급 수정)"
    btp = btf.paragraphs[0]
    btp.font.size = Pt(9)
    btp.font.color.rgb = THEME['dark']

    # 우측 상세 설명
    add_detailed_description(slide, [
        ("1. CI 단계 상세", [
            "Git Push: GitLab 트리거, MR 생성 시 자동 실행",
            "Build: Gradle/npm 빌드, Docker 이미지 생성",
            "Unit Test: JUnit 5, pytest, Vitest 실행",
            "Security Scan: OWASP Dependency Check, Trivy",
            "SonarQube: 코드 품질 분석, 기술 부채 측정",
        ]),
        ("2. CD 단계 상세", [
            "Docker Image: 멀티스테이지 빌드, 경량화",
            "Push Registry: GitLab Container Registry 저장",
            "Staging Deploy: SSH로 docker compose 실행",
            "Manual Approve: 담당자 승인 필수",
            "Production Deploy: 무중단 배포 (Rolling Update)",
        ]),
        ("3. 품질 게이트 상세", [
            "테스트 커버리지: Backend 80%+, Frontend 70%+",
            "보안 취약점: High/Critical 0건 필수",
            "SonarQube: Bugs 0건, Code Smell A등급",
            "코드 리뷰: 최소 1명 Approve 필수",
        ]),
        ("4. 배포 전략 상세", [
            "Rolling Update: 컨테이너 순차 교체",
            "Health Check: 30초 주기 확인 후 트래픽 전환",
            "Rollback: 실패 시 이전 버전으로 자동 복구",
            "배포 시간: 전체 파이프라인 약 15분 소요",
        ]),
    ])

    return slide


def create_phase2_slide(prs):
    """2단계 이관"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "2단계 구축 사업 이관 항목 (Phase 2 Items)")

    # 이관 문서
    add_layer_box(slide, 0.5, 3, 9.5, 3, "📋 이관 대상 문서")
    docs = [
        "• 성능/확장성 설계서 (2-3주) - K8s 마이그레이션 시",
        "• 재해복구 설계서 (1-2주) - 고가용성 요구 시",
        "• 캐싱 전략 상세화 (0.5일) - Redis 클러스터 시",
    ]
    for i, doc in enumerate(docs):
        tb = slide.shapes.add_textbox(Cm(1), Cm(3.7 + i * 0.6), Cm(9), Cm(0.5))
        tf = tb.text_frame
        tf.text = doc
        tp = tf.paragraphs[0]
        tp.font.size = Pt(8)
        tp.font.color.rgb = THEME['text']

    # Medium/Low Priority
    add_layer_box(slide, 10.5, 3, 9, 3, "⚡ Medium/Low Priority 기능")
    items = [
        "• MFA 설계 (0.5일) - 2단계 인증 요구 시",
        "• 데이터 거버넌스 (1일) - 규정 준수 시",
        "• PWA 설계 (0.5일) - 모바일 요구 시",
    ]
    for i, item in enumerate(items):
        tb = slide.shapes.add_textbox(Cm(11), Cm(3.7 + i * 0.6), Cm(8.5), Cm(0.5))
        tf = tb.text_frame
        tf.text = item
        tp = tf.paragraphs[0]
        tp.font.size = Pt(8)
        tp.font.color.rgb = THEME['text']

    # 로드맵
    add_layer_box(slide, 0.5, 6.3, 19, 3.5, "🚀 점진적 확장 로드맵")
    phases = [
        ("Phase 1\n(현재)", "Docker Compose\n$14K/년", THEME['success'], 1),
        ("Phase 2", "Redis Sentinel\n$20K/년", THEME['info'], 5.5),
        ("Phase 3", "Docker Swarm\n$40K/년", THEME['highlight'], 10),
        ("Phase 4", "Kubernetes\n$100K/년", THEME['danger'], 14.5),
    ]
    for title, desc, color, x in phases:
        add_box(slide, x, 7, 4, 2.3, f"{title}\n\n{desc}", color, font_size=9)
        if x < 14:
            add_arrow_text(slide, x + 4, 7.8, "→")

    # 트리거
    trigger = slide.shapes.add_textbox(Cm(0.5), Cm(10), Cm(19), Cm(0.6))
    ttf = trigger.text_frame
    ttf.text = "🎯 트리거: Phase 2 (동시 100+명) | Phase 3 (동시 500+명) | Phase 4 (동시 1,000+명)"
    ttp = ttf.paragraphs[0]
    ttp.font.size = Pt(9)
    ttp.font.bold = True
    ttp.font.color.rgb = THEME['dark']

    # 우측 상세 설명
    add_detailed_description(slide, [
        ("1. 1단계 집중 전략", [
            "핵심 기능 구현에 집중하여 빠른 가치 전달",
            "과도한 초기 설계 지양, MVP 우선 접근",
            "실제 사용량 기반 의사결정",
            "비용 효율성 최우선 고려",
        ]),
        ("2. 이관 문서 필요 시점", [
            "성능/확장성: K8s 마이그레이션 결정 시",
            "재해복구: 99.9% 이상 SLA 요구 시",
            "캐싱 전략: Redis 클러스터 구성 시",
            "원칙: 필요 시점에 작성, 미리 작성 지양",
        ]),
        ("3. 마이그레이션 조건", [
            "Phase 2: 동시 사용자 100+, 일일 요청 10,000+",
            "Phase 3: 동시 사용자 500+, 일일 요청 50,000+",
            "Phase 4: 동시 사용자 1,000+, 일일 요청 200,000+",
            "추가 조건: 가용성 99.9%+ 요구 시 Phase 3 이상",
        ]),
        ("4. 비용 최적화 원칙", [
            "Phase 1: 최소 비용으로 시작 ($14K/년)",
            "점진적 투자: 실제 트래픽 증가에 따라 확장",
            "ROI: 초기 86% 비용 절감 효과",
            "전략: 필요 시점에 맞춰 투자 확대",
        ]),
    ])

    return slide


def create_tech_stack_slide(prs):
    """기술 스택"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "기술 스택 요약 (Tech Stack)")

    # Frontend
    add_layer_box(slide, 0.5, 3, 4.8, 4, "🖥️ Frontend")
    fe_items = ["React 18.3+", "TypeScript 5.4+", "Vite 5.x", "MUI v5", "Redux Toolkit 2.x", "React Query 5.x"]
    for i, item in enumerate(fe_items):
        tb = slide.shapes.add_textbox(Cm(0.8), Cm(3.7 + i * 0.5), Cm(4.2), Cm(0.4))
        tf = tb.text_frame
        tf.text = f"• {item}"
        tp = tf.paragraphs[0]
        tp.font.size = Pt(8)
        tp.font.color.rgb = THEME['text']

    # Backend
    add_layer_box(slide, 5.5, 3, 4.8, 4, "⚙️ Backend")
    be_items = ["Spring Boot 3.2+", "Spring Security 6.x", "Spring Data JPA 3.x", "Resilience4j 2.x", "Gradle 8.x", "Java 17+"]
    for i, item in enumerate(be_items):
        tb = slide.shapes.add_textbox(Cm(5.8), Cm(3.7 + i * 0.5), Cm(4.2), Cm(0.4))
        tf = tb.text_frame
        tf.text = f"• {item}"
        tp = tf.paragraphs[0]
        tp.font.size = Pt(8)
        tp.font.color.rgb = THEME['text']

    # AI Service
    add_layer_box(slide, 10.5, 3, 4.8, 4, "🧠 AI Service")
    ai_items = ["Python 3.11+", "FastAPI 0.110+", "LangGraph 1.0+", "Docling 2.x", "BGE-M3", "DeepSeek V3.2"]
    for i, item in enumerate(ai_items):
        tb = slide.shapes.add_textbox(Cm(10.8), Cm(3.7 + i * 0.5), Cm(4.2), Cm(0.4))
        tf = tb.text_frame
        tf.text = f"• {item}"
        tp = tf.paragraphs[0]
        tp.font.size = Pt(8)
        tp.font.color.rgb = THEME['text']

    # Infra
    add_layer_box(slide, 15.5, 3, 4, 4, "🏗️ Infrastructure")
    infra_items = ["Docker 24.x", "Compose 2.x", "PostgreSQL 16", "ES 8.x", "Neo4j 5.x", "Redis 7.x"]
    for i, item in enumerate(infra_items):
        tb = slide.shapes.add_textbox(Cm(15.8), Cm(3.7 + i * 0.5), Cm(3.5), Cm(0.4))
        tf = tb.text_frame
        tf.text = f"• {item}"
        tp = tf.paragraphs[0]
        tp.font.size = Pt(8)
        tp.font.color.rgb = THEME['text']

    # 서비스 분리 원칙
    principle = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(0.5), Cm(7.3), Cm(19), Cm(1.5)
    )
    principle.fill.solid()
    principle.fill.fore_color.rgb = THEME['light']
    principle.line.color.rgb = THEME['secondary']
    principle.line.width = Pt(2)

    principle_text = slide.shapes.add_textbox(Cm(0.8), Cm(7.5), Cm(18.5), Cm(1.2))
    ptf = principle_text.text_frame
    ptf.text = "⚠️ 서비스 분리 원칙: SpringBoot는 비즈니스 로직/트랜잭션만 담당, AI Service가 LLM/임베딩/검색 전담"
    ptp = ptf.paragraphs[0]
    ptp.font.size = Pt(10)
    ptp.font.bold = True
    ptp.font.color.rgb = THEME['dark']

    # 선정 사유
    reason = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(0.5), Cm(9), Cm(19), Cm(2.2)
    )
    reason.fill.solid()
    reason.fill.fore_color.rgb = THEME['secondary']
    reason.line.fill.background()

    reason_text = slide.shapes.add_textbox(Cm(0.8), Cm(9.2), Cm(18.5), Cm(2))
    rtf = reason_text.text_frame
    rtf.word_wrap = True
    rtf.text = "📌 기술 선정 사유\n• React: 가장 큰 생태계, TypeScript 친화적 | Spring Boot: 기업 환경 표준, 안정성 검증\n• FastAPI: 비동기 성능, Python AI 라이브러리 호환 | DeepSeek: 비용 효율성 95%, GPT-4급 품질"
    for p in rtf.paragraphs:
        p.font.size = Pt(9)
        p.font.color.rgb = THEME['white']

    # 우측 상세 설명
    add_detailed_description(slide, [
        ("1. Frontend 기술 선정", [
            "React 18: Concurrent Features, Suspense 지원",
            "TypeScript: 타입 안정성, IDE 자동완성",
            "Vite: 빠른 HMR, ES Module 기반 빌드",
            "MUI v5: Material Design, 풍부한 컴포넌트",
        ]),
        ("2. Backend 기술 선정", [
            "Spring Boot 3.x: Virtual Thread 지원",
            "Spring Security: OAuth 2.0 Resource Server",
            "Resilience4j: Circuit Breaker, Rate Limiter",
            "WebClient: 비동기 HTTP 클라이언트",
        ]),
        ("3. AI Service 기술 선정", [
            "FastAPI: 비동기 처리, OpenAPI 자동 생성",
            "LangGraph: 워크플로우 기반 AI 파이프라인",
            "BGE-M3: 다국어 지원, 1024차원 임베딩",
            "DeepSeek: GPT-4급 품질, 95% 비용 절감",
        ]),
        ("4. Database 선정", [
            "PostgreSQL: SSOT, ACID 트랜잭션 보장",
            "Elasticsearch: Vector Search, knn 검색",
            "Neo4j: Graph DB, Cypher 쿼리 지원",
            "Redis: 인메모리 캐시, 세션 관리",
        ]),
    ])

    return slide


def create_cost_slide(prs):
    """비용 분석"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "예상 비용 분석 (Cost Analysis)")

    # 비용 테이블
    table = slide.shapes.add_table(7, 4, Cm(0.5), Cm(3), Cm(19), Cm(5)).table

    headers = ["항목", "월간 비용", "연간 비용", "비고"]
    data = [
        ["서버 (Production)", "$800", "$9,600", "32C/128GB/1TB"],
        ["서버 (Staging)", "$200", "$2,400", "16C/64GB/500GB"],
        ["DeepSeek API", "$150", "$1,800", "예상 사용량"],
        ["도메인/SSL", "$20", "$240", "Let's Encrypt"],
        ["합계", "$1,170", "$14,040", ""],
        ["GPT-4 사용 시", "$8,500", "$102,000", "비교 기준"],
    ]

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME['secondary']
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = THEME['white']
        p.alignment = PP_ALIGN.CENTER

    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = val
            if r == 4:
                cell.fill.solid()
                cell.fill.fore_color.rgb = THEME['success']
                p = cell.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.color.rgb = THEME['white']
            elif r == 5:
                cell.fill.solid()
                cell.fill.fore_color.rgb = THEME['danger']
                p = cell.text_frame.paragraphs[0]
                p.font.color.rgb = THEME['white']
            else:
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = THEME['light']
                p = cell.text_frame.paragraphs[0]
                p.font.color.rgb = THEME['text']
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.alignment = PP_ALIGN.CENTER

    # 비용 절감
    savings = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(0.5), Cm(8.3), Cm(19), Cm(1.5)
    )
    savings.fill.solid()
    savings.fill.fore_color.rgb = THEME['success']
    savings.line.fill.background()

    savings_text = slide.shapes.add_textbox(Cm(0.8), Cm(8.5), Cm(18.5), Cm(1.2))
    stf = savings_text.text_frame
    stf.text = "💰 총 비용 절감: $102,000 - $14,040 = $87,960/년 (86% 절감)"
    stp = stf.paragraphs[0]
    stp.font.size = Pt(14)
    stp.font.bold = True
    stp.font.color.rgb = THEME['white']
    stp.alignment = PP_ALIGN.CENTER

    # ROI
    roi = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(0.5), Cm(10), Cm(19), Cm(1.3)
    )
    roi.fill.solid()
    roi.fill.fore_color.rgb = THEME['highlight']
    roi.line.fill.background()

    roi_text = slide.shapes.add_textbox(Cm(0.8), Cm(10.2), Cm(18.5), Cm(1))
    rtf = roi_text.text_frame
    rtf.text = "📈 ROI: 1년차 투자 대비 7.3배 비용 절감 효과"
    rtp = rtf.paragraphs[0]
    rtp.font.size = Pt(12)
    rtp.font.bold = True
    rtp.font.color.rgb = THEME['dark']
    rtp.alignment = PP_ALIGN.CENTER

    # 우측 상세 설명
    add_detailed_description(slide, [
        ("1. 서버 비용 상세", [
            "Production: 32코어/128GB는 AI 워크로드 고려 사양",
            "Elasticsearch: 8GB RAM 권장 (벡터 인덱스)",
            "AI Service: 16GB RAM (임베딩 모델 로드)",
            "클라우드 렌탈: AWS/GCP 동급 사양 기준",
        ]),
        ("2. LLM API 비용 상세", [
            "DeepSeek: $0.14/1M input, $0.28/1M output",
            "GPT-4: $30/1M input, $60/1M output",
            "예상 사용량: 월 10M input, 5M output 토큰",
            "DeepSeek: $1.4 + $1.4 = $2.8/월 × 50배 버퍼",
        ]),
        ("3. K8s 대비 절감 상세", [
            "K8s 클러스터: 13대 서버, 전문 인력 필요",
            "Docker Compose: 1대 서버, 간단한 운영",
            "인프라 비용: $100K → $14K (86% 절감)",
            "운영 복잡도: High → Low",
        ]),
        ("4. 확장 시 비용 증가", [
            "Phase 2: +$6K (Redis Sentinel)",
            "Phase 3: +$26K (Docker Swarm 3노드)",
            "Phase 4: +$86K (K8s 풀 클러스터)",
            "원칙: 트래픽 기반 점진적 투자",
        ]),
    ])

    return slide


def create_closing_slide(prs):
    """마무리"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['dark']
    bg.line.fill.background()

    deco = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(5), Inches(5)
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = THEME['primary']
    deco.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    tf.text = "감사합니다"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1.5))
    sf = sub_box.text_frame
    p1 = sf.paragraphs[0]
    p1.text = "Hybrid RAG Knowledge Platform"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = THEME['accent']
    p1.alignment = PP_ALIGN.CENTER

    p2 = sf.add_paragraph()
    p2.text = '"기업 지식의 80%가 잠들어 있는 문서에서, AI가 즉시 답을 찾아드립니다"'
    p2.font.size = Pt(14)
    p2.font.color.rgb = THEME['light']
    p2.alignment = PP_ALIGN.CENTER

    date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(0.5))
    df = date_box.text_frame
    df.text = "2026-01-16  |  Version 1.0  |  통합 상세 설계서"
    dp = df.paragraphs[0]
    dp.font.size = Pt(11)
    dp.font.color.rgb = THEME['accent']
    dp.alignment = PP_ALIGN.CENTER

    return slide


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("🎨 프레젠테이션 생성 v3 (아키텍처 도식화 개선, 설명 상세화)...")

    # 1. 표지
    create_title_slide(prs)
    print("  ✓ 표지")

    # 2. Executive Summary
    create_section_slide(prs, 1, "Executive Summary")
    create_executive_summary_slide(prs)
    create_key_metrics_slide(prs)
    print("  ✓ Executive Summary")

    # 3. 시스템 아키텍처
    create_section_slide(prs, 2, "시스템 아키텍처")
    create_architecture_slide(prs)
    create_advantages_slide(prs)
    print("  ✓ 시스템 아키텍처")

    # 4. 플랫폼 핵심 설계
    create_section_slide(prs, 3, "플랫폼 핵심 설계")
    create_vip_architecture_slide(prs)
    create_hybrid_search_slide(prs)
    create_data_flow_slide(prs)
    print("  ✓ 플랫폼 핵심 설계")

    # 5. 보안 설계
    create_section_slide(prs, 4, "인증 및 보안 설계")
    create_security_slide(prs)
    print("  ✓ 보안 설계")

    # 6. 인프라/DevOps
    create_section_slide(prs, 5, "인프라 및 DevOps")
    create_infra_slide(prs)
    create_cicd_slide(prs)
    print("  ✓ 인프라/DevOps")

    # 7. 2단계 이관
    create_section_slide(prs, 6, "2단계 구축 이관 항목")
    create_phase2_slide(prs)
    print("  ✓ 2단계 이관")

    # 8. 기술 스택/비용
    create_section_slide(prs, 7, "기술 스택 및 비용")
    create_tech_stack_slide(prs)
    create_cost_slide(prs)
    print("  ✓ 기술 스택/비용")

    # 9. 마무리
    create_closing_slide(prs)
    print("  ✓ 마무리")

    # 저장
    output_path = "/mnt/d/Users/KTDS/Documents/06.과제/hybrid-rag-knowledge-ops/docs/Hybrid_RAG_Platform_Design_v3.pptx"
    prs.save(output_path)

    print(f"\n✅ 저장 완료: {output_path}")
    print(f"   슬라이드: {len(prs.slides)}장")
    print("   개선: 레이어 박스 구분, 상세 설명 (각 4-6개 불렛)")


if __name__ == "__main__":
    main()
