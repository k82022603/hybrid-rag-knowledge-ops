#!/usr/bin/env python3
"""
Hybrid RAG Knowledge Platform 통합 설계서 프레젠테이션 생성기 v2
Brown Theme 적용 - 레이아웃 개선, 내용 보강
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os

# ============================================================
# Brown Theme Color Palette
# ============================================================
THEME = {
    'primary': RGBColor(101, 67, 33),      # #654321 - Deep Brown
    'secondary': RGBColor(139, 90, 43),    # #8B5A2B - Sienna
    'accent': RGBColor(205, 133, 63),      # #CD853F - Peru
    'dark': RGBColor(62, 39, 35),          # #3E2723 - Very Dark Brown
    'light': RGBColor(245, 235, 220),      # #F5EBDC - Light Beige
    'text': RGBColor(62, 39, 35),          # #3E2723 - Dark Brown
    'white': RGBColor(255, 255, 255),      # #FFFFFF
    'highlight': RGBColor(255, 152, 0),    # #FF9800 - Orange
    'success': RGBColor(76, 175, 80),      # #4CAF50 - Green
    'info': RGBColor(33, 150, 243),        # #2196F3 - Blue
    'danger': RGBColor(244, 67, 54),       # #F44336 - Red
    'desc_bg': RGBColor(232, 245, 233),    # #E8F5E9 - Light Green
}

# 슬라이드 치수 (16:9)
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)
HEADER_HEIGHT = Cm(2.2)
CONTENT_TOP = Cm(2.8)  # 헤더 아래 여백


def add_header(slide, title):
    """공통 헤더 추가"""
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, HEADER_HEIGHT
    )
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(30), Cm(1.2))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"


def add_box(slide, left, top, width, height, text, fill_color,
            text_color=None, font_size=11, bold=False, align_center=True):
    """박스 추가"""
    if text_color is None:
        text_color = THEME['white']

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(left), Cm(top), Cm(width), Cm(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = "맑은 고딕"
    p.font.bold = bold
    if align_center:
        p.alignment = PP_ALIGN.CENTER

    return shape


def add_description_box(slide, descriptions):
    """우측 설명 박스 추가 (상세)"""
    # 배경 박스
    desc_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(22), Cm(3), Cm(11), Cm(14)
    )
    desc_bg.fill.solid()
    desc_bg.fill.fore_color.rgb = THEME['desc_bg']
    desc_bg.line.fill.background()

    # 텍스트
    desc_text = slide.shapes.add_textbox(Cm(22.5), Cm(3.3), Cm(10), Cm(13.4))
    tf = desc_text.text_frame
    tf.word_wrap = True

    for i, (title, desc) in enumerate(descriptions):
        if i > 0:
            p = tf.add_paragraph()
            p.text = ""
            p.space_before = Pt(6)

        # 제목
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = THEME['dark']
        p.font.name = "맑은 고딕"

        # 설명
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(9)
        p2.font.color.rgb = THEME['text']
        p2.font.name = "맑은 고딕"
        p2.space_after = Pt(4)


def create_title_slide(prs):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['dark']
    bg.line.fill.background()

    # 장식 원
    deco1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(5), Inches(5)
    )
    deco1.fill.solid()
    deco1.fill.fore_color.rgb = THEME['primary']
    deco1.line.fill.background()

    # 장식 사각형
    deco2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(9), Inches(5.5), Inches(5), Inches(2.5)
    )
    deco2.fill.solid()
    deco2.fill.fore_color.rgb = THEME['secondary']
    deco2.line.fill.background()

    # 좌측 강조선
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.8), Inches(0.08), Inches(1.8)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = THEME['accent']
    line.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Hybrid RAG Knowledge Platform"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    p2 = tf.add_paragraph()
    p2.text = "통합 상세 설계서"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = THEME['accent']
    p2.font.name = "맑은 고딕"

    # 부제목
    sub_box = slide.shapes.add_textbox(Inches(1.1), Inches(4.5), Inches(8), Inches(0.8))
    sf = sub_box.text_frame
    sf.text = "Vector + Graph 기반 차세대 지식 검색 시스템"
    sp = sf.paragraphs[0]
    sp.font.size = Pt(16)
    sp.font.color.rgb = THEME['light']
    sp.font.name = "맑은 고딕"

    # 날짜
    date_box = slide.shapes.add_textbox(Inches(1.1), Inches(6.5), Inches(5), Inches(0.4))
    df = date_box.text_frame
    df.text = "2026-01-16  |  Version 1.0  |  Final Draft"
    dp = df.paragraphs[0]
    dp.font.size = Pt(11)
    dp.font.color.rgb = THEME['light']
    dp.font.name = "맑은 고딕"

    return slide


def create_section_slide(prs, section_number, section_title):
    """섹션 구분 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['primary']
    bg.line.fill.background()

    # 섹션 번호
    num_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(2), Inches(1))
    nf = num_box.text_frame
    nf.text = f"{section_number:02d}"
    np = nf.paragraphs[0]
    np.font.size = Pt(60)
    np.font.bold = True
    np.font.color.rgb = THEME['accent']
    np.font.name = "맑은 고딕"

    # 제목
    title_box = slide.shapes.add_textbox(Inches(3.2), Inches(3), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = section_title
    tp = tf.paragraphs[0]
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = THEME['white']
    tp.font.name = "맑은 고딕"

    return slide


def create_executive_summary_slide(prs):
    """Executive Summary 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Executive Summary - 프로젝트 비전")

    # 비전 문구
    vision_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(3.2), Cm(20), Cm(2)
    )
    vision_box.fill.solid()
    vision_box.fill.fore_color.rgb = THEME['secondary']
    vision_box.line.fill.background()

    vision_text = slide.shapes.add_textbox(Cm(1.5), Cm(3.5), Cm(19), Cm(1.5))
    vf = vision_text.text_frame
    vf.text = '"기업 지식의 80%가 잠들어 있는 문서에서, AI가 즉시 답을 찾아드립니다"'
    vp = vf.paragraphs[0]
    vp.font.size = Pt(14)
    vp.font.bold = True
    vp.font.color.rgb = THEME['white']
    vp.font.name = "맑은 고딕"
    vp.alignment = PP_ALIGN.CENTER

    # 핵심 가치 4개
    values = [
        ("💰 비용 효율성", "DeepSeek 단일 모델\nLLM 비용 95% 절감\n연간 ~$14,000", THEME['success']),
        ("🎯 정확한 검색", "Vector + Graph 융합\nRRF 재순위 알고리즘\nPrecision@5 > 85%", THEME['info']),
        ("🔐 기업 보안", "OAuth 2.0 + PKCE\nAES-256-GCM 암호화\nRBAC 권한 관리", RGBColor(156, 39, 176)),
        ("📈 확장 가능성", "Docker → K8s 경로\n마이크로서비스 구조\nAPI 기반 통합", THEME['highlight']),
    ]

    x_pos = 1
    for title, desc, color in values:
        add_box(slide, x_pos, 6, 4.8, 1.2, title, color, font_size=11, bold=True)

        desc_box = slide.shapes.add_textbox(Cm(x_pos), Cm(7.5), Cm(4.8), Cm(3))
        df = desc_box.text_frame
        df.word_wrap = True
        for i, line in enumerate(desc.split('\n')):
            p = df.paragraphs[0] if i == 0 else df.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(9)
            p.font.color.rgb = THEME['text']
            p.font.name = "맑은 고딕"

        x_pos += 5.2

    # 우측 설명
    add_description_box(slide, [
        ("1. Hybrid RAG란?", "Vector Search(의미 검색)와 Graph Search(관계 검색)를 결합한 차세대 검색 기술입니다."),
        ("2. 비용 절감 원리", "GPT-4 대비 95% 저렴한 DeepSeek 모델을 사용하면서도 고품질 답변을 생성합니다."),
        ("3. 보안 강화", "Keycloak 기반 OAuth 2.0 인증, 데이터 암호화, 역할 기반 접근 제어를 적용합니다."),
        ("4. 확장 전략", "Docker Compose로 시작하여 트래픽 증가 시 Kubernetes로 자연스럽게 마이그레이션합니다."),
    ])

    return slide


def create_key_metrics_slide(prs):
    """핵심 수치 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "핵심 목표 수치 (Key Metrics)")

    # 메트릭 카드
    metrics = [
        ("응답 시간", "< 3초", "(P95)", THEME['info']),
        ("검색 정확도", "> 85%", "(Precision@5)", THEME['success']),
        ("LLM 비용", "95%↓", "(GPT-4 대비)", THEME['highlight']),
        ("가용성", "99.5%", "(SLA)", THEME['secondary']),
        ("동시 사용자", "100명", "(1단계 목표)", THEME['primary']),
    ]

    x_pos = 1
    for title, value, sub, color in metrics:
        # 카드 배경
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x_pos), Cm(3.5), Cm(3.8), Cm(4)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()

        # 수치
        val_box = slide.shapes.add_textbox(Cm(x_pos), Cm(4), Cm(3.8), Cm(1.5))
        vf = val_box.text_frame
        vf.text = value
        vp = vf.paragraphs[0]
        vp.font.size = Pt(24)
        vp.font.bold = True
        vp.font.color.rgb = THEME['white']
        vp.font.name = "맑은 고딕"
        vp.alignment = PP_ALIGN.CENTER

        # 제목
        title_box = slide.shapes.add_textbox(Cm(x_pos), Cm(5.5), Cm(3.8), Cm(0.8))
        tf = title_box.text_frame
        tf.text = title
        tp = tf.paragraphs[0]
        tp.font.size = Pt(10)
        tp.font.color.rgb = THEME['white']
        tp.font.name = "맑은 고딕"
        tp.alignment = PP_ALIGN.CENTER

        # 부제
        sub_box = slide.shapes.add_textbox(Cm(x_pos), Cm(6.3), Cm(3.8), Cm(0.6))
        sf = sub_box.text_frame
        sf.text = sub
        sp = sf.paragraphs[0]
        sp.font.size = Pt(8)
        sp.font.color.rgb = RGBColor(220, 220, 220)
        sp.font.name = "맑은 고딕"
        sp.alignment = PP_ALIGN.CENTER

        x_pos += 4.1

    # 비용 비교 테이블
    table_title = slide.shapes.add_textbox(Cm(1), Cm(8.5), Cm(10), Cm(0.6))
    ttf = table_title.text_frame
    ttf.text = "📊 연간 비용 비교"
    ttp = ttf.paragraphs[0]
    ttp.font.size = Pt(11)
    ttp.font.bold = True
    ttp.font.color.rgb = THEME['dark']

    # 테이블
    rows = 3
    cols = 4
    table = slide.shapes.add_table(rows, cols, Cm(1), Cm(9.3), Cm(19), Cm(3)).table

    headers = ["항목", "GPT-4 사용 시", "DeepSeek 사용 시", "절감액"]
    data = [
        ["LLM API 비용", "$100,000", "$1,800", "$98,200 (98%)"],
        ["총 운영 비용", "$114,000", "$14,040", "$99,960 (88%)"],
    ]

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME['secondary']
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = THEME['white']
        p.alignment = PP_ALIGN.CENTER

    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = THEME['light']
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.font.color.rgb = THEME['text']
            p.alignment = PP_ALIGN.CENTER

    # 우측 설명
    add_description_box(slide, [
        ("1. 응답 시간 < 3초", "Hybrid 검색(Vector+Graph) + 답변 합성까지 P95 기준 3초 이내 완료를 목표로 합니다."),
        ("2. 검색 정확도 > 85%", "상위 5개 검색 결과 중 적합한 문서가 85% 이상 포함되어야 합니다."),
        ("3. LLM 비용 95% 절감", "DeepSeek V3.2 모델은 GPT-4 대비 토큰당 가격이 95% 저렴합니다."),
        ("4. 비용 구조", "서버: $9,600 + Staging: $2,400 + LLM: $1,800 + 기타: $240 = $14,040/년"),
    ])

    return slide


def create_architecture_slide(prs):
    """시스템 아키텍처 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "시스템 아키텍처 (System Architecture)")

    # External
    add_box(slide, 1, 3.3, 2.5, 1, "👤 사용자\n(Browser)", THEME['accent'], font_size=9)
    add_box(slide, 17.5, 3.3, 2.5, 1, "🤖 DeepSeek\nAPI", RGBColor(26, 115, 232), font_size=9)

    # Gateway Layer
    add_box(slide, 4.5, 3.3, 2.5, 1, "🚪 Nginx\nReverse Proxy", THEME['secondary'], font_size=9)
    add_box(slide, 7.5, 3.3, 3.5, 1, "🔀 API Gateway\nJWT 검증 / Rate Limit", THEME['secondary'], font_size=9)

    # Auth
    add_box(slide, 11.5, 3.3, 2.5, 1, "🔐 Keycloak\nOAuth 2.0 IdP", RGBColor(183, 28, 28), font_size=9)

    # Application Layer 라벨
    app_label = slide.shapes.add_textbox(Cm(1), Cm(4.8), Cm(5), Cm(0.5))
    alf = app_label.text_frame
    alf.text = "Application Layer"
    alp = alf.paragraphs[0]
    alp.font.size = Pt(9)
    alp.font.bold = True
    alp.font.color.rgb = THEME['dark']

    add_box(slide, 1, 5.5, 3, 1.2, "🖥️ Frontend\nReact 18 + TS\nMUI v5", RGBColor(97, 218, 251), THEME['dark'], font_size=9)
    add_box(slide, 4.5, 5.5, 3.5, 1.2, "⚙️ Backend\nSpring Boot 3.x\nResilience4j", RGBColor(109, 179, 63), THEME['white'], font_size=9)
    add_box(slide, 8.5, 5.5, 4, 1.2, "🧠 AI Service\nFastAPI + LangGraph\nBGE-M3 Embedding", RGBColor(0, 150, 136), THEME['white'], font_size=9)

    # Data Layer 라벨
    data_label = slide.shapes.add_textbox(Cm(1), Cm(7.2), Cm(5), Cm(0.5))
    dlf = data_label.text_frame
    dlf.text = "Data Layer"
    dlp = dlf.paragraphs[0]
    dlp.font.size = Pt(9)
    dlp.font.bold = True
    dlp.font.color.rgb = THEME['dark']

    add_box(slide, 1, 8, 2.5, 1.2, "🐘 PostgreSQL\nSSOT (원본)", RGBColor(51, 103, 145), font_size=9)
    add_box(slide, 4, 8, 3, 1.2, "🔍 Elasticsearch\nVector Search\n1024 dims", RGBColor(0, 175, 170), font_size=9)
    add_box(slide, 7.5, 8, 2.5, 1.2, "🕸️ Neo4j\nGraph DB\n엔티티/관계", RGBColor(0, 110, 170), font_size=9)
    add_box(slide, 10.5, 8, 2, 1.2, "💾 Redis\nCache\nSession", RGBColor(220, 63, 66), font_size=9)
    add_box(slide, 13, 8, 2.2, 1.2, "📦 MinIO\nFile Storage\nS3 호환", RGBColor(198, 40, 40), font_size=9)

    # Monitoring Layer 라벨
    mon_label = slide.shapes.add_textbox(Cm(1), Cm(9.7), Cm(5), Cm(0.5))
    mlf = mon_label.text_frame
    mlf.text = "Monitoring Layer"
    mlp = mlf.paragraphs[0]
    mlp.font.size = Pt(9)
    mlp.font.bold = True
    mlp.font.color.rgb = THEME['dark']

    add_box(slide, 1, 10.5, 2.2, 1, "📊 Prometheus\nMetrics", RGBColor(230, 84, 49), font_size=9)
    add_box(slide, 3.7, 10.5, 2.2, 1, "📈 Grafana\nDashboard", RGBColor(241, 144, 13), THEME['dark'], font_size=9)
    add_box(slide, 6.4, 10.5, 2.2, 1, "📝 Loki\nLog", RGBColor(140, 86, 75), font_size=9)

    # 화살표 표시 (텍스트로)
    arrow_box = slide.shapes.add_textbox(Cm(13), Cm(3.5), Cm(4), Cm(0.8))
    af = arrow_box.text_frame
    af.text = "→ LLM API 호출 →"
    ap = af.paragraphs[0]
    ap.font.size = Pt(8)
    ap.font.color.rgb = THEME['text']

    # 우측 설명
    add_description_box(slide, [
        ("1. Gateway Layer", "Nginx가 SSL을 종단하고, API Gateway가 JWT 토큰 검증 및 요청 제한(Rate Limiting)을 수행합니다."),
        ("2. Application Layer", "Frontend(React SPA), Backend(Spring Boot), AI Service(FastAPI)로 구성됩니다. 서비스 간 REST API 통신합니다."),
        ("3. 서비스 분리 원칙", "SpringBoot는 비즈니스 로직/트랜잭션만 담당, AI Service가 LLM 호출/임베딩/검색을 전담합니다."),
        ("4. Data Layer", "PostgreSQL이 SSOT(Single Source of Truth)이며, ES/Neo4j는 검색용 비정규화 인덱스입니다."),
        ("5. Monitoring", "Prometheus(메트릭 수집), Grafana(시각화), Loki(로그 집계)로 통합 모니터링합니다."),
    ])

    return slide


def create_vip_architecture_slide(prs):
    """VIP 3단계 LLM 아키텍처"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "VIP 3단계 LLM 아키텍처 (Value-Intelligent-Planning)")

    # Stage 1: Value
    add_box(slide, 1, 4, 1.5, 0.8, "Stage 1", THEME['dark'], font_size=10, bold=True)
    add_box(slide, 2.8, 3.5, 4, 5.5,
            "📄 문서 업로드\n\n↓\n\n🔧 Docling 파싱\n청크 분할 (500자)\n\n↓\n\n🤖 DeepSeek-Chat\n엔티티/관계 추출\n\n↓\n\n💎 지식 그래프 구축",
            RGBColor(200, 230, 201), THEME['dark'], font_size=9, bold=False)

    # Stage 2: Intelligent
    add_box(slide, 7.5, 4, 1.5, 0.8, "Stage 2", THEME['dark'], font_size=10, bold=True)
    add_box(slide, 9.3, 3.5, 4, 5.5,
            "❓ 사용자 질문\n\n↓\n\n🧠 의도 분석\n질의 유형 분류\n\n↓\n\n🎯 검색 전략 수립\nVector/Graph 가중치\n\n↓\n\n🔍 병렬 검색 실행",
            RGBColor(187, 222, 251), THEME['dark'], font_size=9, bold=False)

    # Stage 3: Planning
    add_box(slide, 14, 4, 1.5, 0.8, "Stage 3", THEME['dark'], font_size=10, bold=True)
    add_box(slide, 15.8, 3.5, 4, 5.5,
            "📊 검색 결과 수신\n\n↓\n\n📝 컨텍스트 구성\nRRF 융합 결과\n\n↓\n\n🤖 DeepSeek-Chat\n답변 합성\n\n↓\n\n💬 최종 답변 + 출처",
            RGBColor(248, 187, 217), THEME['dark'], font_size=9, bold=False)

    # 화살표
    arrow1 = slide.shapes.add_textbox(Cm(6.8), Cm(6), Cm(1), Cm(0.5))
    a1f = arrow1.text_frame
    a1f.text = "→"
    a1p = a1f.paragraphs[0]
    a1p.font.size = Pt(18)
    a1p.font.color.rgb = THEME['accent']

    arrow2 = slide.shapes.add_textbox(Cm(13.3), Cm(6), Cm(1), Cm(0.5))
    a2f = arrow2.text_frame
    a2f.text = "→"
    a2p = a2f.paragraphs[0]
    a2p.font.size = Pt(18)
    a2p.font.color.rgb = THEME['accent']

    # 하단 요약
    summary_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(10), Cm(19), Cm(1.8)
    )
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = THEME['light']
    summary_box.line.fill.background()

    sum_text = slide.shapes.add_textbox(Cm(1.3), Cm(10.2), Cm(18.5), Cm(1.5))
    stf = sum_text.text_frame
    stf.word_wrap = True
    stf.text = "💡 핵심 원리: 문서 업로드 시 지식 추출(Value) → 질의 시 최적 검색 전략 수립(Intelligent) → 검색 결과 기반 답변 합성(Planning)"
    stp = stf.paragraphs[0]
    stp.font.size = Pt(10)
    stp.font.color.rgb = THEME['dark']

    # 우측 설명
    add_description_box(slide, [
        ("1. Value 단계 (오프라인)", "문서 업로드 시 실행됩니다. Docling으로 파싱 후 500자 단위로 청킹하고, DeepSeek으로 엔티티와 관계를 추출합니다."),
        ("2. Intelligent 단계 (온라인)", "사용자 질문 입력 시 실행됩니다. 질의 의도를 분석하고 Vector/Graph 검색 가중치를 결정합니다."),
        ("3. Planning 단계 (온라인)", "검색 결과를 RRF로 융합한 후, 컨텍스트를 구성하여 DeepSeek에 답변 합성을 요청합니다."),
        ("4. 비용 최적화 포인트", "DeepSeek-Chat 단일 모델로 엔티티 추출과 답변 합성을 모두 처리하여 비용을 95% 절감합니다."),
        ("5. 품질 보장", "Hallucination 방지를 위해 반드시 검색된 문서 기반으로만 답변을 생성합니다."),
    ])

    return slide


def create_hybrid_search_slide(prs):
    """Hybrid Search 융합 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Hybrid Search 융합 (Vector + Graph)")

    # 사용자 질의
    add_box(slide, 7.5, 3.3, 5, 1, "❓ 사용자 질의 입력", THEME['dark'], font_size=11, bold=True)

    # 분기 화살표
    arrow_l = slide.shapes.add_textbox(Cm(6), Cm(4.5), Cm(2), Cm(0.5))
    alf = arrow_l.text_frame
    alf.text = "↙"
    alp = alf.paragraphs[0]
    alp.font.size = Pt(16)
    alp.font.color.rgb = THEME['accent']

    arrow_r = slide.shapes.add_textbox(Cm(12), Cm(4.5), Cm(2), Cm(0.5))
    arf = arrow_r.text_frame
    arf.text = "↘"
    arp = arf.paragraphs[0]
    arp.font.size = Pt(16)
    arp.font.color.rgb = THEME['accent']

    # Vector Search
    add_box(slide, 1, 5.5, 6, 3,
            "🔍 Vector Search\n\n• Elasticsearch knn 검색\n• BGE-M3 임베딩 (1024차원)\n• 의미적 유사도 기반\n• Top-K 문서 반환",
            RGBColor(249, 183, 22), THEME['dark'], font_size=10, bold=False, align_center=False)

    # Graph Search
    add_box(slide, 13, 5.5, 6.5, 3,
            "🕸️ Graph Search\n\n• Neo4j Cypher 쿼리\n• 엔티티 관계 탐색\n• 연관 문서 확장\n• 지식 그래프 활용",
            RGBColor(1, 139, 255), THEME['white'], font_size=10, bold=False, align_center=False)

    # 병렬 표시
    parallel_box = slide.shapes.add_textbox(Cm(7.5), Cm(6.5), Cm(5), Cm(1))
    pf = parallel_box.text_frame
    pf.text = "⚡ 병렬 실행"
    pp = pf.paragraphs[0]
    pp.font.size = Pt(12)
    pp.font.bold = True
    pp.font.color.rgb = THEME['highlight']
    pp.alignment = PP_ALIGN.CENTER

    # 수렴 화살표
    conv_l = slide.shapes.add_textbox(Cm(6), Cm(8.5), Cm(2), Cm(0.5))
    clf = conv_l.text_frame
    clf.text = "↘"
    clp = clf.paragraphs[0]
    clp.font.size = Pt(16)
    clp.font.color.rgb = THEME['accent']

    conv_r = slide.shapes.add_textbox(Cm(12), Cm(8.5), Cm(2), Cm(0.5))
    crf = conv_r.text_frame
    crf.text = "↙"
    crp = crf.paragraphs[0]
    crp.font.size = Pt(16)
    crp.font.color.rgb = THEME['accent']

    # RRF 융합
    add_box(slide, 6.5, 9.5, 7, 1.5,
            "⚡ RRF (Reciprocal Rank Fusion)\nscore = Σ 1/(k+rank), k=60",
            RGBColor(255, 107, 107), THEME['white'], font_size=10)

    # 결과
    add_box(slide, 6.5, 11.3, 7, 1, "📊 통합 검색 결과 (Top-K)", THEME['success'], font_size=10, bold=True)

    # 우측 설명
    add_description_box(slide, [
        ("1. Vector Search 특징", "BGE-M3 임베딩으로 질의와 문서의 의미적 유사도를 계산합니다. 동의어, 유사 표현도 검색됩니다."),
        ("2. Graph Search 특징", "엔티티 관계를 따라 연관 문서를 탐색합니다. '프로젝트 A의 담당자' 같은 관계 질의에 강점이 있습니다."),
        ("3. 병렬 실행", "두 검색을 동시에 실행하여 응답 시간을 최소화합니다."),
        ("4. RRF 융합 알고리즘", "각 검색 결과의 순위를 기반으로 점수를 계산합니다. k=60은 순위 차이의 민감도를 조절합니다."),
        ("5. 융합 효과", "Vector만으로 부족한 관계 정보를 Graph가 보완하고, Graph만으로 놓칠 수 있는 유사 문서를 Vector가 보완합니다."),
    ])

    return slide


def create_data_flow_slide(prs):
    """데이터 흐름 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "검색 데이터 흐름 (Search Flow)")

    # 단계별 박스
    steps = [
        ("① 사용자", "질문 입력", THEME['accent'], 1),
        ("② Frontend", "API 요청\n/search/chat", RGBColor(97, 218, 251), 4),
        ("③ Gateway", "JWT 검증\nRate Limit", THEME['secondary'], 7),
        ("④ Backend", "요청 전달\n로깅", RGBColor(109, 179, 63), 10),
        ("⑤ AI Service", "Hybrid 검색\n답변 합성", RGBColor(0, 150, 136), 13),
        ("⑥ DeepSeek", "LLM 답변\n생성", RGBColor(26, 115, 232), 16),
    ]

    for title, desc, color, x in steps:
        add_box(slide, x, 3.5, 2.7, 2.5, f"{title}\n\n{desc}", color, font_size=9)

    # 화살표
    for i in range(5):
        x = 3.7 + i * 3
        arrow = slide.shapes.add_textbox(Cm(x), Cm(4.3), Cm(1), Cm(0.5))
        af = arrow.text_frame
        af.text = "→"
        ap = af.paragraphs[0]
        ap.font.size = Pt(14)
        ap.font.color.rgb = THEME['dark']

    # 병렬 검색 영역
    parallel_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(7), Cm(19), Cm(5)
    )
    parallel_bg.fill.solid()
    parallel_bg.fill.fore_color.rgb = THEME['light']
    parallel_bg.line.color.rgb = THEME['secondary']
    parallel_bg.line.width = Pt(2)

    par_title = slide.shapes.add_textbox(Cm(1.3), Cm(7.2), Cm(10), Cm(0.5))
    ptf = par_title.text_frame
    ptf.text = "⚡ AI Service 내부 병렬 검색"
    ptp = ptf.paragraphs[0]
    ptp.font.size = Pt(10)
    ptp.font.bold = True
    ptp.font.color.rgb = THEME['dark']

    # ES
    add_box(slide, 1.5, 8.3, 5, 1.5,
            "🔍 Elasticsearch\nVector Search (knn)\nBGE-M3 유사도 Top-20",
            RGBColor(249, 183, 22), THEME['dark'], font_size=9)

    # Neo4j
    add_box(slide, 7, 8.3, 5, 1.5,
            "🕸️ Neo4j\nGraph Search (Cypher)\n엔티티 관계 탐색 Top-20",
            RGBColor(1, 139, 255), THEME['white'], font_size=9)

    # RRF
    add_box(slide, 12.5, 8.3, 3.5, 1.5,
            "⚡ RRF 융합\n결과 통합\nTop-10 선정",
            RGBColor(255, 107, 107), THEME['white'], font_size=9)

    # DeepSeek
    add_box(slide, 16.5, 8.3, 3, 1.5,
            "🤖 DeepSeek\n답변 합성\n출처 첨부",
            RGBColor(76, 175, 80), THEME['white'], font_size=9)

    # 흐름 화살표
    arrows_pos = [(6.5, 8.8), (12, 8.8), (16, 8.8)]
    for x, y in arrows_pos:
        arr = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(1), Cm(0.5))
        arf = arr.text_frame
        arf.text = "→"
        arp = arf.paragraphs[0]
        arp.font.size = Pt(12)
        arp.font.color.rgb = THEME['dark']

    # 우측 설명
    add_description_box(slide, [
        ("1. 요청 흐름", "사용자 → Frontend → Gateway → Backend → AI Service → DeepSeek 순서로 요청이 전달됩니다."),
        ("2. Gateway 역할", "JWT 토큰 검증, Rate Limiting(분당 100회), 요청 라우팅을 담당합니다."),
        ("3. 병렬 검색", "AI Service 내에서 ES와 Neo4j를 동시에 조회하여 응답 시간을 단축합니다."),
        ("4. RRF 융합", "두 검색 결과를 Reciprocal Rank Fusion으로 통합하여 Top-10을 선정합니다."),
        ("5. 답변 생성", "선정된 문서를 컨텍스트로 DeepSeek에 전달하여 최종 답변을 생성합니다."),
    ])

    return slide


def create_advantages_slide(prs):
    """아키텍처 장점 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "아키텍처 장점 (Architecture Advantages)")

    advantages = [
        ("💰 비용 효율성 95% 절감", [
            "DeepSeek V3.2: GPT-4 대비 토큰당 95% 저렴",
            "연간 LLM 비용: $100,000 → $1,800",
            "총 운영 비용: $14,040/년 (K8s 대비 86% 절감)",
            "단일 모델로 엔티티 추출 + 답변 합성 처리",
        ], THEME['success']),
        ("🎯 검색 정확도 향상", [
            "Vector Search: 의미적 유사도 기반 검색",
            "Graph Search: 엔티티 관계 기반 확장 검색",
            "RRF 융합: 두 검색의 장점 결합",
            "목표: Precision@5 > 85%, MRR > 0.8",
        ], THEME['info']),
        ("🔐 기업 보안 강화", [
            "Keycloak: OAuth 2.0 + PKCE 인증",
            "AES-256-GCM: 민감 데이터 필드 암호화",
            "RBAC: 역할 기반 접근 제어 (Admin/Manager/User)",
            "TLS 1.3: 전송 구간 암호화",
        ], RGBColor(156, 39, 176)),
        ("📈 점진적 확장 전략", [
            "Phase 1: Docker Compose (현재, $14K/년)",
            "Phase 2: + Redis Sentinel ($20K/년)",
            "Phase 3: Docker Swarm 3노드 ($40K/년)",
            "Phase 4: Kubernetes 클러스터 ($100K/년)",
        ], THEME['highlight']),
    ]

    y_pos = 3.2
    for title, items, color in advantages:
        # 제목 박스
        add_box(slide, 1, y_pos, 6, 0.9, title, color, font_size=11, bold=True)

        # 항목들
        items_box = slide.shapes.add_textbox(Cm(7.5), Cm(y_pos), Cm(12), Cm(2.2))
        itf = items_box.text_frame
        itf.word_wrap = True

        for i, item in enumerate(items):
            p = itf.paragraphs[0] if i == 0 else itf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(9)
            p.font.color.rgb = THEME['text']
            p.font.name = "맑은 고딕"

        y_pos += 2.7

    # 우측 설명
    add_description_box(slide, [
        ("1. 비용 효율성 핵심", "DeepSeek 모델은 품질 대비 가격이 매우 저렴합니다. GPT-4와 유사한 품질을 95% 낮은 비용으로 제공합니다."),
        ("2. Hybrid 검색 장점", "Vector만으로는 관계 정보를 찾기 어렵고, Graph만으로는 유사 표현을 놓칩니다. 두 방식의 융합이 핵심입니다."),
        ("3. 보안 설계 원칙", "Zero Trust 원칙을 적용하여 모든 요청에 대해 인증/인가를 검증합니다."),
        ("4. 확장 트리거 조건", "동시 사용자 100+ → Phase 2, 500+ → Phase 3, 1000+ → Phase 4로 전환합니다."),
    ])

    return slide


def create_security_slide(prs):
    """보안 아키텍처 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "인증 및 보안 설계 (Security Architecture)")

    # 인증 흐름
    add_box(slide, 1, 3.3, 6, 2.5,
            "🔐 인증 (Authentication)\n\n• OAuth 2.0 Authorization Code Flow\n• PKCE (Proof Key for Code Exchange)\n• Keycloak IdP 연동\n• JWT Access Token (RS256 서명)",
            THEME['secondary'], font_size=9, bold=False, align_center=False)

    add_box(slide, 7.5, 3.3, 6, 2.5,
            "👤 권한 (Authorization)\n\n• RBAC (Role-Based Access Control)\n• 3단계 역할: ADMIN / MANAGER / USER\n• API 레벨 권한 검증\n• 리소스 기반 접근 제어",
            RGBColor(156, 39, 176), font_size=9, bold=False, align_center=False)

    add_box(slide, 14, 3.3, 5.5, 2.5,
            "🔒 암호화 (Encryption)\n\n• TLS 1.3 (전송 구간)\n• AES-256-GCM (필드 암호화)\n• bcrypt (비밀번호 해싱)\n• HashiCorp Vault (키 관리)",
            RGBColor(198, 40, 40), font_size=9, bold=False, align_center=False)

    # JWT 토큰 전략
    add_box(slide, 1, 6.3, 9, 2.2,
            "🎫 JWT 토큰 전략\n\n• Access Token: 수명 15분, 메모리 저장, 만료 2분 전 자동 갱신\n• Refresh Token: 수명 7일, HttpOnly Cookie, Rotation 적용\n• 토큰 구조: Header.Payload.Signature (RS256)",
            THEME['primary'], font_size=9, bold=False, align_center=False)

    # 데이터 분류
    add_box(slide, 10.5, 6.3, 9, 2.2,
            "📊 데이터 분류 체계\n\n• Level 4 (극비): 암호화 키, 마스터 비밀번호 → HSM/Vault\n• Level 3 (비밀): API 키, 토큰 → DB 암호화\n• Level 2 (대외비): 이름, 이메일 → 필드 암호화\n• Level 1 (일반): 문서 메타데이터 → TDE",
            RGBColor(255, 152, 0), THEME['dark'], font_size=9, bold=False, align_center=False)

    # RBAC 역할 테이블
    table_title = slide.shapes.add_textbox(Cm(1), Cm(9), Cm(10), Cm(0.5))
    ttf = table_title.text_frame
    ttf.text = "📋 RBAC 권한 매트릭스"
    ttp = ttf.paragraphs[0]
    ttp.font.size = Pt(10)
    ttp.font.bold = True
    ttp.font.color.rgb = THEME['dark']

    table = slide.shapes.add_table(4, 5, Cm(1), Cm(9.8), Cm(18.5), Cm(2.5)).table

    headers = ["역할", "사용자 관리", "지식 CRUD", "지식 조회", "북마크"]
    data = [
        ["ADMIN", "✅", "✅", "✅", "✅"],
        ["MANAGER", "❌", "✅", "✅", "✅"],
        ["USER", "❌", "❌", "✅", "✅"],
    ]

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME['secondary']
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = THEME['white']
        p.alignment = PP_ALIGN.CENTER

    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.font.color.rgb = THEME['text']
            p.alignment = PP_ALIGN.CENTER

    # 우측 설명
    add_description_box(slide, [
        ("1. OAuth 2.0 + PKCE", "Authorization Code Flow에 PKCE를 추가하여 인증 코드 탈취 공격을 방지합니다."),
        ("2. JWT 토큰 보안", "Access Token은 메모리에만 저장하고, Refresh Token은 HttpOnly Cookie로 XSS 공격을 방지합니다."),
        ("3. 암호화 계층", "전송(TLS) → 애플리케이션(AES-256) → 저장(TDE) 3계층 암호화를 적용합니다."),
        ("4. 키 관리", "HashiCorp Vault로 암호화 키를 중앙 집중 관리하고, 키 로테이션을 자동화합니다."),
        ("5. RBAC 원칙", "최소 권한 원칙을 적용하여 역할별로 필요한 권한만 부여합니다."),
    ])

    return slide


def create_infra_slide(prs):
    """인프라 구성 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "인프라 구성 (Docker Compose 기반)")

    # Docker Host 외곽
    host_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(3), Cm(19), Cm(9.5)
    )
    host_bg.fill.solid()
    host_bg.fill.fore_color.rgb = THEME['light']
    host_bg.line.color.rgb = THEME['secondary']
    host_bg.line.width = Pt(2)

    host_title = slide.shapes.add_textbox(Cm(1.3), Cm(3.2), Cm(10), Cm(0.5))
    htf = host_title.text_frame
    htf.text = "🐳 Docker Host (Production: 32C/128GB/1TB SSD)"
    htp = htf.paragraphs[0]
    htp.font.size = Pt(10)
    htp.font.bold = True
    htp.font.color.rgb = THEME['dark']

    # Application Layer
    app_label = slide.shapes.add_textbox(Cm(1.3), Cm(4), Cm(5), Cm(0.4))
    alf = app_label.text_frame
    alf.text = "Application (6개)"
    alp = alf.paragraphs[0]
    alp.font.size = Pt(8)
    alp.font.bold = True
    alp.font.color.rgb = THEME['dark']

    apps = [("nginx", 1.5), ("frontend", 4), ("gateway", 6.5), ("backend", 9), ("ai-service", 11.5), ("keycloak", 14.5)]
    for name, x in apps:
        add_box(slide, x, 4.6, 2.2, 0.9, name, RGBColor(129, 236, 236), THEME['dark'], font_size=8)

    # Database Layer
    db_label = slide.shapes.add_textbox(Cm(1.3), Cm(5.8), Cm(5), Cm(0.4))
    dlf = db_label.text_frame
    dlf.text = "Database (5개)"
    dlp = dlf.paragraphs[0]
    dlp.font.size = Pt(8)
    dlp.font.bold = True
    dlp.font.color.rgb = THEME['dark']

    dbs = [("postgresql", 1.5), ("elasticsearch", 4.5), ("neo4j", 7.5), ("redis", 10), ("minio", 12.5)]
    for name, x in dbs:
        add_box(slide, x, 6.4, 2.5, 0.9, name, RGBColor(162, 155, 254), THEME['white'], font_size=8)

    # Monitoring Layer
    mon_label = slide.shapes.add_textbox(Cm(1.3), Cm(7.6), Cm(5), Cm(0.4))
    mlf = mon_label.text_frame
    mlf.text = "Monitoring (4개)"
    mlp = mlf.paragraphs[0]
    mlp.font.size = Pt(8)
    mlp.font.bold = True
    mlp.font.color.rgb = THEME['dark']

    mons = [("prometheus", 1.5), ("grafana", 4.5), ("loki", 7.5), ("promtail", 10)]
    for name, x in mons:
        add_box(slide, x, 8.2, 2.5, 0.9, name, RGBColor(253, 121, 168), THEME['white'], font_size=8)

    # 스펙 요약
    spec_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(9.8), Cm(19), Cm(2.5)
    )
    spec_box.fill.solid()
    spec_box.fill.fore_color.rgb = THEME['secondary']
    spec_box.line.fill.background()

    spec_text = slide.shapes.add_textbox(Cm(1.3), Cm(10), Cm(18.5), Cm(2.2))
    stf = spec_text.text_frame
    stf.word_wrap = True
    stf.text = "📦 총 15개 컨테이너 | Production: 32C/128GB/1TB SSD ($800/월) | Staging: 16C/64GB/500GB ($200/월) | 연간 총비용: $14,040"
    stp = stf.paragraphs[0]
    stp.font.size = Pt(10)
    stp.font.color.rgb = THEME['white']
    stp.alignment = PP_ALIGN.CENTER

    # 우측 설명
    add_description_box(slide, [
        ("1. 단일 서버 전략", "초기 단계에서는 Docker Compose로 단일 서버에 모든 컨테이너를 배포하여 운영 복잡도를 최소화합니다."),
        ("2. 컨테이너 구성", "Application 6개(nginx, frontend, gateway, backend, ai-service, keycloak) + Database 5개 + Monitoring 4개"),
        ("3. 서버 사양", "Production: 32코어/128GB RAM/1TB SSD, Staging: 16코어/64GB RAM/500GB SSD"),
        ("4. 비용 구조", "서버 비용 $1,000/월 + LLM API $150/월 + 기타 $20/월 = 연간 $14,040"),
        ("5. 확장 경로", "사용자 증가 시 Redis Sentinel → Docker Swarm → Kubernetes 순으로 마이그레이션합니다."),
    ])

    return slide


def create_cicd_slide(prs):
    """CI/CD 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "CI/CD 파이프라인 (GitLab)")

    # CI 단계
    ci_label = slide.shapes.add_textbox(Cm(1), Cm(3.2), Cm(5), Cm(0.5))
    cif = ci_label.text_frame
    cif.text = "CI (Continuous Integration)"
    cip = cif.paragraphs[0]
    cip.font.size = Pt(10)
    cip.font.bold = True
    cip.font.color.rgb = THEME['dark']

    ci_steps = [
        ("Git Push", THEME['accent'], 1),
        ("Build", THEME['info'], 4),
        ("Unit Test", THEME['success'], 7),
        ("Security\nScan", RGBColor(156, 39, 176), 10),
        ("SonarQube", THEME['highlight'], 13),
    ]

    for name, color, x in ci_steps:
        add_box(slide, x, 4, 2.5, 1.2, name, color, font_size=9)

    for i in range(4):
        x = 3.5 + i * 3
        arr = slide.shapes.add_textbox(Cm(x), Cm(4.3), Cm(1), Cm(0.5))
        arf = arr.text_frame
        arf.text = "→"
        arp = arf.paragraphs[0]
        arp.font.size = Pt(12)
        arp.font.color.rgb = THEME['dark']

    # CD 단계
    cd_label = slide.shapes.add_textbox(Cm(1), Cm(5.8), Cm(5), Cm(0.5))
    cdf = cd_label.text_frame
    cdf.text = "CD (Continuous Deployment)"
    cdp = cdf.paragraphs[0]
    cdp.font.size = Pt(10)
    cdp.font.bold = True
    cdp.font.color.rgb = THEME['dark']

    cd_steps = [
        ("Docker\nImage", THEME['info'], 1),
        ("Push\nRegistry", THEME['secondary'], 4),
        ("Staging\nDeploy", THEME['success'], 7),
        ("Manual\nApprove", THEME['highlight'], 10),
        ("Production\nDeploy", THEME['danger'], 13),
    ]

    for name, color, x in cd_steps:
        add_box(slide, x, 6.6, 2.5, 1.2, name, color, font_size=9)

    for i in range(4):
        x = 3.5 + i * 3
        arr = slide.shapes.add_textbox(Cm(x), Cm(6.9), Cm(1), Cm(0.5))
        arf = arr.text_frame
        arf.text = "→"
        arp = arf.paragraphs[0]
        arp.font.size = Pt(12)
        arp.font.color.rgb = THEME['dark']

    # 품질 게이트
    gate_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(8.5), Cm(18.5), Cm(1.8)
    )
    gate_box.fill.solid()
    gate_box.fill.fore_color.rgb = THEME['light']
    gate_box.line.fill.background()

    gate_text = slide.shapes.add_textbox(Cm(1.3), Cm(8.7), Cm(18), Cm(1.5))
    gtf = gate_text.text_frame
    gtf.word_wrap = True
    gtf.text = "🚦 품질 게이트: 테스트 커버리지 80%+ | 보안 취약점 High 0건 | SonarQube Quality Gate Pass | 코드 리뷰 승인"
    gtp = gtf.paragraphs[0]
    gtp.font.size = Pt(10)
    gtp.font.color.rgb = THEME['dark']

    # 브랜치 전략
    branch_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(10.5), Cm(18.5), Cm(2)
    )
    branch_box.fill.solid()
    branch_box.fill.fore_color.rgb = THEME['secondary']
    branch_box.line.fill.background()

    branch_text = slide.shapes.add_textbox(Cm(1.3), Cm(10.7), Cm(18), Cm(1.7))
    btf = branch_text.text_frame
    btf.word_wrap = True
    btf.text = "🌿 브랜치 전략: main (프로덕션 릴리스) | develop (개발 통합) | feature/* (새 기능) | fix/* (버그 수정) | hotfix/* (긴급 수정)"
    btp = btf.paragraphs[0]
    btp.font.size = Pt(10)
    btp.font.color.rgb = THEME['white']

    # 우측 설명
    add_description_box(slide, [
        ("1. CI 단계", "코드 푸시 시 자동으로 빌드, 테스트, 보안 스캔, 정적 분석을 수행합니다."),
        ("2. 보안 스캔", "OWASP Dependency Check로 의존성 취약점을 검사하고, Trivy로 컨테이너 이미지를 스캔합니다."),
        ("3. CD 단계", "Docker 이미지 빌드 후 레지스트리에 푸시하고, Staging → Manual Approve → Production 순서로 배포합니다."),
        ("4. 품질 게이트", "모든 조건을 통과해야 다음 단계로 진행됩니다. 실패 시 파이프라인이 중단됩니다."),
        ("5. 배포 전략", "SSH 기반 배포로 docker compose pull && docker compose up -d 명령을 실행합니다."),
    ])

    return slide


def create_phase2_slide(prs):
    """2단계 이관 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "2단계 구축 사업 이관 항목 (Phase 2 Items)")

    # 이관 문서
    add_box(slide, 1, 3.3, 9, 2.5,
            "📋 이관 대상 문서\n\n• 성능/확장성 설계서: Docker Compose → K8s 마이그레이션 시 필요 (2-3주)\n• 재해복구 설계서: 고가용성 요구사항 발생 시 필요 (1-2주)\n• 캐싱 전략 상세화: Redis 클러스터 구성 시 필요 (0.5일)",
            THEME['danger'], font_size=9, bold=False, align_center=False)

    # Medium/Low Priority
    add_box(slide, 10.5, 3.3, 9, 2.5,
            "⚡ Medium/Low Priority 기능\n\n• MFA 설계 추가 (0.5일) - 2단계 인증 요구 시\n• 데이터 거버넌스 설계서 (1일) - 규정 준수 필요 시\n• PWA 설계 추가 (0.5일) - 모바일 앱 요구 시\n• 통합 테스트 계획서 (1일) - QA 프로세스 강화 시",
            RGBColor(255, 193, 7), THEME['dark'], font_size=9, bold=False, align_center=False)

    # 마이그레이션 로드맵
    roadmap_label = slide.shapes.add_textbox(Cm(1), Cm(6.3), Cm(10), Cm(0.5))
    rlf = roadmap_label.text_frame
    rlf.text = "🚀 점진적 확장 로드맵"
    rlp = rlf.paragraphs[0]
    rlp.font.size = Pt(10)
    rlp.font.bold = True
    rlp.font.color.rgb = THEME['dark']

    phases = [
        ("Phase 1\n(현재)", "Docker Compose\n단일 서버\n$14K/년", THEME['success'], 1),
        ("Phase 2", "Redis Sentinel\n캐시 이중화\n$20K/년", THEME['info'], 5.5),
        ("Phase 3", "Docker Swarm\n3노드 클러스터\n$40K/년", THEME['highlight'], 10),
        ("Phase 4", "Kubernetes\n풀 클러스터\n$100K/년", THEME['danger'], 14.5),
    ]

    for title, desc, color, x in phases:
        add_box(slide, x, 7, 4, 2.8, f"{title}\n\n{desc}", color, font_size=9)

    for i in range(3):
        x = 5 + i * 4.5
        arr = slide.shapes.add_textbox(Cm(x), Cm(8), Cm(1), Cm(0.5))
        arf = arr.text_frame
        arf.text = "→"
        arp = arf.paragraphs[0]
        arp.font.size = Pt(14)
        arp.font.color.rgb = THEME['dark']

    # 트리거 조건
    trigger_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(10.3), Cm(18.5), Cm(2)
    )
    trigger_box.fill.solid()
    trigger_box.fill.fore_color.rgb = THEME['light']
    trigger_box.line.fill.background()

    trigger_text = slide.shapes.add_textbox(Cm(1.3), Cm(10.5), Cm(18), Cm(1.7))
    ttf = trigger_text.text_frame
    ttf.word_wrap = True
    ttf.text = "🎯 마이그레이션 트리거: Phase 2 (동시 사용자 100+, 일일 요청 10,000+) | Phase 3 (500+, 50,000+) | Phase 4 (1,000+, 200,000+)"
    ttp = ttf.paragraphs[0]
    ttp.font.size = Pt(10)
    ttp.font.color.rgb = THEME['dark']

    # 우측 설명
    add_description_box(slide, [
        ("1. 1단계 집중 전략", "핵심 기능 구현에 집중하여 빠른 가치 전달을 우선합니다. 과도한 설계는 지양합니다."),
        ("2. 이관 문서 필요 시점", "트래픽 증가로 확장이 필요하거나, 규정 준수 요구사항이 발생할 때 작성합니다."),
        ("3. 점진적 확장 원칙", "사용량에 맞춰 단계적으로 인프라를 확장하여 비용 효율성을 유지합니다."),
        ("4. 마이그레이션 경로", "Docker Compose → Redis Sentinel → Docker Swarm → K8s 순서로 자연스럽게 확장됩니다."),
        ("5. 비용 최적화", "필요 시점에 맞춰 투자하여 초기 비용을 86% 절감합니다."),
    ])

    return slide


def create_tech_stack_slide(prs):
    """기술 스택 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "기술 스택 요약 (Tech Stack)")

    # Frontend
    add_box(slide, 1, 3.3, 4.5, 3.5,
            "🖥️ Frontend\n\n• React 18.3+\n• TypeScript 5.4+\n• Vite 5.x\n• MUI v5\n• Redux Toolkit 2.x\n• React Query 5.x\n• React Hook Form 7.x",
            RGBColor(97, 218, 251), THEME['dark'], font_size=9, bold=False, align_center=False)

    # Backend
    add_box(slide, 6, 3.3, 4.5, 3.5,
            "⚙️ Backend\n\n• Spring Boot 3.2+\n• Spring Security 6.x\n• Spring Data JPA 3.x\n• WebClient (HTTP)\n• Resilience4j 2.x\n• Gradle 8.x\n• Java 17+",
            RGBColor(109, 179, 63), THEME['white'], font_size=9, bold=False, align_center=False)

    # AI Service
    add_box(slide, 11, 3.3, 4.5, 3.5,
            "🧠 AI Service\n\n• Python 3.11+\n• FastAPI 0.110+\n• LangChain 1.2+\n• LangGraph 1.0+\n• Docling 2.x\n• BGE-M3 (Embedding)\n• DeepSeek V3.2 (LLM)",
            RGBColor(0, 150, 136), THEME['white'], font_size=9, bold=False, align_center=False)

    # Infrastructure
    add_box(slide, 16, 3.3, 3.5, 3.5,
            "🏗️ Infra\n\n• Docker 24.x\n• Compose 2.x\n• Nginx 1.25+\n• PostgreSQL 16\n• ES 8.x\n• Neo4j 5.x\n• Redis 7.x",
            THEME['secondary'], font_size=9, bold=False, align_center=False)

    # 서비스 분리 원칙
    principle_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(7.3), Cm(18.5), Cm(2)
    )
    principle_box.fill.solid()
    principle_box.fill.fore_color.rgb = THEME['light']
    principle_box.line.color.rgb = THEME['secondary']
    principle_box.line.width = Pt(2)

    principle_text = slide.shapes.add_textbox(Cm(1.3), Cm(7.5), Cm(18), Cm(1.7))
    ptf = principle_text.text_frame
    ptf.word_wrap = True
    ptf.text = "⚠️ 서비스 분리 원칙: SpringBoot는 비즈니스 로직/트랜잭션만 담당, AI Service가 LLM/임베딩/검색 전담. SpringBoot는 LLM과 직접 연동하지 않음."
    ptp = ptf.paragraphs[0]
    ptp.font.size = Pt(10)
    ptp.font.bold = True
    ptp.font.color.rgb = THEME['dark']

    # 선정 사유
    reason_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(9.8), Cm(18.5), Cm(2.5)
    )
    reason_box.fill.solid()
    reason_box.fill.fore_color.rgb = THEME['secondary']
    reason_box.line.fill.background()

    reason_text = slide.shapes.add_textbox(Cm(1.3), Cm(10), Cm(18), Cm(2.2))
    rtf = reason_text.text_frame
    rtf.word_wrap = True
    rtf.text = "📌 기술 선정 사유:\n• React: 가장 큰 생태계, TypeScript 친화적 | Spring Boot: 기업 환경 표준, 안정성 검증\n• FastAPI: 비동기 성능, Python AI 라이브러리 호환 | DeepSeek: 비용 효율성 95%, GPT-4급 품질"
    for i, p in enumerate(rtf.paragraphs):
        p.font.size = Pt(9)
        p.font.color.rgb = THEME['white']

    # 우측 설명
    add_description_box(slide, [
        ("1. Frontend 기술", "React 18의 Concurrent Features와 TypeScript로 타입 안정성을 확보합니다."),
        ("2. Backend 기술", "Spring Boot 3.x의 Virtual Thread 지원과 Resilience4j로 장애 복구를 처리합니다."),
        ("3. AI Service 기술", "FastAPI의 비동기 처리와 LangGraph의 워크플로우 관리로 AI 파이프라인을 구현합니다."),
        ("4. 데이터베이스", "PostgreSQL(SSOT) + ES(Vector) + Neo4j(Graph)로 Zero-Join 아키텍처를 구현합니다."),
        ("5. 비용 최적화", "DeepSeek V3.2는 GPT-4 대비 95% 저렴하면서 유사한 품질을 제공합니다."),
    ])

    return slide


def create_cost_slide(prs):
    """비용 분석 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "예상 비용 분석 (Cost Analysis)")

    # 비용 테이블
    table = slide.shapes.add_table(7, 4, Cm(1), Cm(3.3), Cm(18.5), Cm(5)).table

    headers = ["항목", "월간 비용", "연간 비용", "비고"]
    data = [
        ["서버 (Production)", "$800", "$9,600", "32C/128GB/1TB SSD"],
        ["서버 (Staging)", "$200", "$2,400", "16C/64GB/500GB"],
        ["DeepSeek API", "$150", "$1,800", "예상 사용량 기준"],
        ["도메인/SSL", "$20", "$240", "Let's Encrypt 무료"],
        ["합계", "$1,170", "$14,040", ""],
        ["GPT-4 사용 시", "$8,500", "$102,000", "LLM만 $100K"],
    ]

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME['secondary']
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = THEME['white']
        p.alignment = PP_ALIGN.CENTER

    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = val
            if r == 4:  # 합계 행
                cell.fill.solid()
                cell.fill.fore_color.rgb = THEME['success']
                p = cell.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.color.rgb = THEME['white']
            elif r == 5:  # GPT-4 행
                cell.fill.solid()
                cell.fill.fore_color.rgb = THEME['danger']
                p = cell.text_frame.paragraphs[0]
                p.font.color.rgb = THEME['white']
            else:
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = THEME['light']
                p = cell.text_frame.paragraphs[0]
                p.font.color.rgb = THEME['text']
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER

    # 비용 절감 강조
    savings_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(8.8), Cm(18.5), Cm(1.8)
    )
    savings_box.fill.solid()
    savings_box.fill.fore_color.rgb = THEME['success']
    savings_box.line.fill.background()

    savings_text = slide.shapes.add_textbox(Cm(1.3), Cm(9), Cm(18), Cm(1.5))
    stf = savings_text.text_frame
    stf.text = "💰 총 비용 절감: $102,000 - $14,040 = $87,960/년 (86% 절감)"
    stp = stf.paragraphs[0]
    stp.font.size = Pt(16)
    stp.font.bold = True
    stp.font.color.rgb = THEME['white']
    stp.alignment = PP_ALIGN.CENTER

    # ROI
    roi_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(11), Cm(18.5), Cm(1.5)
    )
    roi_box.fill.solid()
    roi_box.fill.fore_color.rgb = THEME['highlight']
    roi_box.line.fill.background()

    roi_text = slide.shapes.add_textbox(Cm(1.3), Cm(11.2), Cm(18), Cm(1.2))
    rtf = roi_text.text_frame
    rtf.text = "📈 ROI: 1년차 투자 대비 6.3배 비용 절감 효과 (GPT-4 대비)"
    rtp = rtf.paragraphs[0]
    rtp.font.size = Pt(14)
    rtp.font.bold = True
    rtp.font.color.rgb = THEME['dark']
    rtp.alignment = PP_ALIGN.CENTER

    # 우측 설명
    add_description_box(slide, [
        ("1. 서버 비용", "클라우드 서버 렌탈 비용입니다. Production 32코어/128GB는 AI 워크로드를 고려한 사양입니다."),
        ("2. LLM API 비용", "DeepSeek V3.2는 토큰당 가격이 GPT-4의 5% 수준입니다. 동일 사용량 기준 $100K → $1.8K로 절감됩니다."),
        ("3. K8s 대비 절감", "Kubernetes 클러스터 운영 시 연간 $100K 이상 소요됩니다. Docker Compose로 86% 절감됩니다."),
        ("4. 확장 시 비용", "Phase 2: $20K, Phase 3: $40K, Phase 4: $100K로 점진적 증가합니다."),
        ("5. 투자 효율성", "초기 비용을 최소화하고, 실제 트래픽 증가에 따라 투자를 확대하는 전략입니다."),
    ])

    return slide


def create_closing_slide(prs):
    """마무리 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['dark']
    bg.line.fill.background()

    # 장식
    deco = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(5), Inches(5)
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = THEME['primary']
    deco.line.fill.background()

    # 감사합니다
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    tf.text = "감사합니다"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    # 비전
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1.5))
    sf = sub_box.text_frame
    p1 = sf.paragraphs[0]
    p1.text = "Hybrid RAG Knowledge Platform"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = THEME['accent']
    p1.alignment = PP_ALIGN.CENTER

    p2 = sf.add_paragraph()
    p2.text = '"기업 지식의 80%가 잠들어 있는 문서에서, AI가 즉시 답을 찾아드립니다"'
    p2.font.size = Pt(14)
    p2.font.color.rgb = THEME['light']
    p2.alignment = PP_ALIGN.CENTER

    # 날짜
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(0.5))
    df = date_box.text_frame
    df.text = "2026-01-16  |  Version 1.0  |  통합 상세 설계서"
    dp = df.paragraphs[0]
    dp.font.size = Pt(11)
    dp.font.color.rgb = THEME['accent']
    dp.alignment = PP_ALIGN.CENTER

    return slide


def main():
    """메인 함수"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("🎨 프레젠테이션 생성 시작 (v2 - 레이아웃 개선)...")

    # 1. 표지
    create_title_slide(prs)
    print("  ✓ 표지 슬라이드")

    # 2. Executive Summary
    create_section_slide(prs, 1, "Executive Summary")
    create_executive_summary_slide(prs)
    create_key_metrics_slide(prs)
    print("  ✓ Executive Summary")

    # 3. 시스템 아키텍처
    create_section_slide(prs, 2, "시스템 아키텍처")
    create_architecture_slide(prs)
    create_advantages_slide(prs)
    print("  ✓ 시스템 아키텍처")

    # 4. 플랫폼 핵심 설계
    create_section_slide(prs, 3, "플랫폼 핵심 설계")
    create_vip_architecture_slide(prs)
    create_hybrid_search_slide(prs)
    create_data_flow_slide(prs)
    print("  ✓ 플랫폼 핵심 설계")

    # 5. 보안 설계
    create_section_slide(prs, 4, "인증 및 보안 설계")
    create_security_slide(prs)
    print("  ✓ 보안 설계")

    # 6. 인프라 및 DevOps
    create_section_slide(prs, 5, "인프라 및 DevOps 설계")
    create_infra_slide(prs)
    create_cicd_slide(prs)
    print("  ✓ 인프라 및 DevOps")

    # 7. 2단계 이관
    create_section_slide(prs, 6, "2단계 구축 이관 항목")
    create_phase2_slide(prs)
    print("  ✓ 2단계 이관 항목")

    # 8. 기술 스택 및 비용
    create_section_slide(prs, 7, "기술 스택 및 비용 분석")
    create_tech_stack_slide(prs)
    create_cost_slide(prs)
    print("  ✓ 기술 스택 및 비용")

    # 9. 마무리
    create_closing_slide(prs)
    print("  ✓ 마무리 슬라이드")

    # 저장
    output_path = "/mnt/d/Users/KTDS/Documents/06.과제/hybrid-rag-knowledge-ops/docs/Hybrid_RAG_Platform_Design_Brown.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

    print(f"\n✅ 프레젠테이션 저장 완료: {output_path}")
    print(f"   총 슬라이드 수: {len(prs.slides)}장")
    print(f"   테마: Brown (브라운)")
    print(f"   개선 사항: 헤더 겹침 해결, 내용 보강, 우측 설명 상세화, 폰트 크기 축소")


if __name__ == "__main__":
    main()
