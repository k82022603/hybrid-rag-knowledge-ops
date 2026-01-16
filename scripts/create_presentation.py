#!/usr/bin/env python3
"""
Hybrid RAG Knowledge Platform 통합 설계서 프레젠테이션 생성기
Brown Theme 적용
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os

# ============================================================
# Brown Theme Color Palette
# ============================================================
BROWN_THEME = {
    'primary': RGBColor(101, 67, 33),      # #654321 - Deep Brown
    'secondary': RGBColor(139, 90, 43),    # #8B5A2B - Sienna
    'accent': RGBColor(205, 133, 63),      # #CD853F - Peru
    'dark': RGBColor(62, 39, 35),          # #3E2723 - Very Dark Brown
    'light': RGBColor(245, 235, 220),      # #F5EBDC - Light Beige
    'text': RGBColor(62, 39, 35),          # #3E2723 - Dark Brown
    'white': RGBColor(255, 255, 255),      # #FFFFFF
    'highlight': RGBColor(255, 152, 0),    # #FF9800 - Orange accent
    'success': RGBColor(76, 175, 80),      # #4CAF50 - Green
}

THEME = BROWN_THEME


def add_rounded_box(slide, left, top, width, height, text,
                    fill_color, text_color=None, font_size=14,
                    font_name="맑은 고딕", bold=True):
    """둥근 모서리 박스 추가"""
    if text_color is None:
        text_color = THEME['white']

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(left), Cm(top),
        Cm(width), Cm(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = font_name
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER

    return shape


def add_arrow(slide, x1, y1, x2, y2, color=None, width=2):
    """화살표 추가"""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Cm(x1), Cm(y1),
        Cm(x2), Cm(y2)
    )
    connector.line.color.rgb = color or THEME['accent']
    connector.line.width = Pt(width)
    return connector


def create_title_slide(prs, title, subtitle):
    """표지 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['dark']
    bg.line.color.rgb = THEME['dark']

    # 장식 요소 - 우측 상단 원
    deco1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(9.5), Inches(-0.5), Inches(4), Inches(4)
    )
    deco1.fill.solid()
    deco1.fill.fore_color.rgb = THEME['primary']
    deco1.line.color.rgb = THEME['primary']

    # 장식 요소 - 우측 하단 사각형
    deco2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8), Inches(5), Inches(5.5), Inches(3)
    )
    deco2.fill.solid()
    deco2.fill.fore_color.rgb = THEME['secondary']
    deco2.line.color.rgb = THEME['secondary']

    # 좌측 장식 라인
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(3.2), Inches(0.1), Inches(1.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = THEME['accent']
    line.line.color.rgb = THEME['accent']

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # 부제목
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8), Inches(0.8))
    sf = sub_box.text_frame
    sf.text = subtitle
    sp = sf.paragraphs[0]
    sp.font.size = Pt(20)
    sp.font.color.rgb = THEME['light']
    sp.font.name = "맑은 고딕"

    # 날짜 및 버전
    date_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(5), Inches(0.5))
    df = date_box.text_frame
    df.text = "2026-01-16  |  Version 1.0  |  Final Draft"
    dp = df.paragraphs[0]
    dp.font.size = Pt(12)
    dp.font.color.rgb = THEME['light']
    dp.font.name = "맑은 고딕"

    return slide


def create_section_slide(prs, section_title, section_number):
    """섹션 구분 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['primary']
    bg.line.color.rgb = THEME['primary']

    # 섹션 번호
    num_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(2), Inches(1))
    nf = num_box.text_frame
    nf.text = f"{section_number:02d}"
    np = nf.paragraphs[0]
    np.font.size = Pt(72)
    np.font.bold = True
    np.font.color.rgb = THEME['accent']
    np.font.name = "맑은 고딕"

    # 섹션 제목
    title_box = slide.shapes.add_textbox(Inches(3.5), Inches(3), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = section_title
    tp = tf.paragraphs[0]
    tp.font.size = Pt(40)
    tp.font.bold = True
    tp.font.color.rgb = THEME['white']
    tp.font.name = "맑은 고딕"

    return slide


def create_content_slide(prs, title, bullet_points, add_header=True):
    """일반 콘텐츠 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if add_header:
        # 헤더 바
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.33), Inches(0.9)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = THEME['primary']
        header.line.color.rgb = THEME['primary']

        # 제목
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = THEME['white']
        p.font.name = "맑은 고딕"

    # 내용
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.8))
    cf = content_box.text_frame
    cf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = cf.paragraphs[0]
        else:
            p = cf.add_paragraph()

        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = THEME['text']
        p.font.name = "맑은 고딕"
        p.space_after = Pt(14)
        p.level = 0

    return slide


def create_table_slide(prs, title, headers, rows):
    """테이블 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더 바
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.9)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = THEME['primary']
    header_bar.line.color.rgb = THEME['primary']

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # 테이블 추가
    cols = len(headers)
    table_rows = len(rows) + 1

    table_width = Inches(12)
    table_height = Inches(0.5 * table_rows)
    table_left = Inches(0.67)
    table_top = Inches(1.3)

    table = slide.shapes.add_table(table_rows, cols, table_left, table_top, table_width, table_height).table

    # 헤더 행
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME['secondary']
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = THEME['white']
        p.font.size = Pt(14)
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.CENTER

    # 데이터 행
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = THEME['light']
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = THEME['text']
            p.font.name = "맑은 고딕"

    return slide


def create_architecture_slide(prs):
    """시스템 아키텍처 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.color.rgb = THEME['primary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "시스템 아키텍처 (System Architecture)"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # 외부 영역 - 사용자
    add_rounded_box(slide, 1, 1.5, 2.5, 1.2, "👤 사용자", THEME['accent'], font_size=14)

    # Gateway Layer
    add_rounded_box(slide, 5, 1.5, 2.5, 1.2, "🚪 Nginx\nSSL Termination", THEME['secondary'], font_size=12)
    add_rounded_box(slide, 8, 1.5, 3.5, 1.2, "🔀 API Gateway\nJWT 검증 / Rate Limit", THEME['secondary'], font_size=12)

    # Application Layer - Frontend
    add_rounded_box(slide, 1, 3.2, 3, 1.2, "🖥️ Frontend\nReact 18 + TypeScript", RGBColor(97, 218, 251), THEME['dark'], font_size=12)

    # Application Layer - Backend
    add_rounded_box(slide, 4.5, 3.2, 3, 1.2, "⚙️ Backend\nSpring Boot 3.x", RGBColor(109, 179, 63), THEME['white'], font_size=12)

    # Application Layer - AI Service
    add_rounded_box(slide, 8, 3.2, 3.5, 1.2, "🧠 AI Service\nFastAPI + LangGraph", RGBColor(0, 150, 136), THEME['white'], font_size=12)

    # Auth Layer - Keycloak
    add_rounded_box(slide, 12, 3.2, 2.5, 1.2, "🔐 Keycloak\nOAuth 2.0", RGBColor(183, 28, 28), THEME['white'], font_size=12)

    # External - DeepSeek
    add_rounded_box(slide, 12, 1.5, 2.5, 1.2, "🤖 DeepSeek\nAPI", RGBColor(26, 115, 232), THEME['white'], font_size=12)

    # Data Layer
    add_rounded_box(slide, 1, 5, 2.3, 1.2, "🐘 PostgreSQL\nSSOT", RGBColor(51, 103, 145), THEME['white'], font_size=12)
    add_rounded_box(slide, 3.8, 5, 2.5, 1.2, "🔍 Elasticsearch\nVector Search", RGBColor(0, 175, 170), THEME['white'], font_size=12)
    add_rounded_box(slide, 6.8, 5, 2.3, 1.2, "🕸️ Neo4j\nGraph DB", RGBColor(0, 110, 170), THEME['white'], font_size=12)
    add_rounded_box(slide, 9.6, 5, 2, 1.2, "💾 Redis\nCache", RGBColor(220, 63, 66), THEME['white'], font_size=12)
    add_rounded_box(slide, 12, 5, 2.3, 1.2, "📦 MinIO\nFile Storage", RGBColor(198, 40, 40), THEME['white'], font_size=12)

    # 설명 박스
    desc_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(19), Cm(4), Cm(8.5), Cm(10)
    )
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = RGBColor(226, 240, 217)
    desc_box.line.color.rgb = RGBColor(226, 240, 217)

    desc_text = slide.shapes.add_textbox(Cm(19.3), Cm(4.3), Cm(8), Cm(9.5))
    dtf = desc_text.text_frame
    dtf.word_wrap = True

    descriptions = [
        ("1. Gateway Layer", "Nginx가 SSL을 처리하고, API Gateway가 JWT 검증 및 Rate Limiting 수행"),
        ("2. Application Layer", "React SPA, Spring Boot 백엔드, FastAPI AI 서비스로 구성"),
        ("3. 서비스 분리 원칙", "SpringBoot는 비즈니스 로직, AI Service는 LLM/임베딩 처리 담당"),
        ("4. Data Layer", "PostgreSQL(SSOT) + 비정규화된 검색 인덱스(ES, Neo4j)"),
    ]

    for i, (title, desc) in enumerate(descriptions):
        if i == 0:
            p = dtf.paragraphs[0]
        else:
            p = dtf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = THEME['dark']
        p.font.name = "맑은 고딕"
        p.space_before = Pt(8) if i > 0 else Pt(0)

        p2 = dtf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = THEME['text']
        p2.font.name = "맑은 고딕"
        p2.space_after = Pt(6)

    return slide


def create_vip_architecture_slide(prs):
    """VIP 3단계 LLM 아키텍처 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.color.rgb = THEME['primary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "VIP 3단계 LLM 아키텍처 (Value-Intelligent-Planning)"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # Stage 1: Value
    add_rounded_box(slide, 0.5, 2.5, 4.5, 4,
                    "Stage 1: Value\n\n📄 문서\n     ↓\n🤖 DeepSeek-Chat\n     ↓\n💎 엔티티/관계 추출",
                    RGBColor(200, 230, 201), THEME['dark'], font_size=14)

    # Stage 2: Intelligent
    add_rounded_box(slide, 5.5, 2.5, 4.5, 4,
                    "Stage 2: Intelligent\n\n❓ 사용자 질문\n     ↓\n🧠 의도 분석\n     ↓\n🎯 검색 전략 수립",
                    RGBColor(187, 222, 251), THEME['dark'], font_size=14)

    # Stage 3: Planning
    add_rounded_box(slide, 10.5, 2.5, 4.5, 4,
                    "Stage 3: Planning\n\n📊 검색 결과\n     ↓\n📝 컨텍스트 구성\n     ↓\n💬 답변 합성",
                    RGBColor(248, 187, 217), THEME['dark'], font_size=14)

    # 설명 박스
    desc_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(19), Cm(4), Cm(8.5), Cm(10)
    )
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = RGBColor(226, 240, 217)
    desc_box.line.color.rgb = RGBColor(226, 240, 217)

    desc_text = slide.shapes.add_textbox(Cm(19.3), Cm(4.3), Cm(8), Cm(9.5))
    dtf = desc_text.text_frame
    dtf.word_wrap = True

    descriptions = [
        ("1. Value 단계", "문서에서 엔티티와 관계를 추출하여 지식 그래프 구축"),
        ("2. Intelligent 단계", "사용자 질의 의도를 분석하고 최적 검색 전략 수립"),
        ("3. Planning 단계", "검색 결과를 기반으로 정확한 답변 생성"),
        ("4. 핵심 장점", "단일 모델(DeepSeek)로 95% 비용 절감하면서 고품질 유지"),
    ]

    for i, (title, desc) in enumerate(descriptions):
        if i == 0:
            p = dtf.paragraphs[0]
        else:
            p = dtf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = THEME['dark']
        p.font.name = "맑은 고딕"
        p.space_before = Pt(8) if i > 0 else Pt(0)

        p2 = dtf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = THEME['text']
        p2.font.name = "맑은 고딕"
        p2.space_after = Pt(6)

    return slide


def create_hybrid_search_slide(prs):
    """Hybrid Search 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.color.rgb = THEME['primary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "Hybrid Search 융합 (Vector + Graph)"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # 사용자 질의
    add_rounded_box(slide, 6.5, 2, 4, 1, "❓ 사용자 질의", THEME['dark'], font_size=16)

    # Vector Search
    add_rounded_box(slide, 2, 4, 5, 1.5, "🔍 Vector Search\nElasticsearch knn (BGE-M3)", RGBColor(249, 183, 22), THEME['dark'], font_size=14)

    # Graph Search
    add_rounded_box(slide, 10, 4, 5, 1.5, "🕸️ Graph Search\nNeo4j Cypher Query", RGBColor(1, 139, 255), THEME['white'], font_size=14)

    # RRF 융합
    add_rounded_box(slide, 6, 6.5, 5, 1.5, "⚡ RRF 융합 알고리즘\nReciprocal Rank Fusion", RGBColor(255, 107, 107), THEME['white'], font_size=14)

    # 통합 결과
    add_rounded_box(slide, 6, 9, 5, 1, "📊 통합 검색 결과 (Top-K)", THEME['success'], THEME['white'], font_size=14)

    # 설명 박스
    desc_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(19), Cm(4), Cm(8.5), Cm(10)
    )
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = RGBColor(226, 240, 217)
    desc_box.line.color.rgb = RGBColor(226, 240, 217)

    desc_text = slide.shapes.add_textbox(Cm(19.3), Cm(4.3), Cm(8), Cm(9.5))
    dtf = desc_text.text_frame
    dtf.word_wrap = True

    descriptions = [
        ("1. Vector Search", "의미적 유사도 기반 검색, BGE-M3 1024차원 임베딩 사용"),
        ("2. Graph Search", "엔티티 관계 기반 검색, 연관 지식 탐색에 강점"),
        ("3. RRF 융합", "score = Σ 1/(k+rank), k=60, 두 검색 결과 통합"),
        ("4. 장점", "의미 유사도 + 관계 기반 검색 결합으로 정확도 향상"),
    ]

    for i, (title, desc) in enumerate(descriptions):
        if i == 0:
            p = dtf.paragraphs[0]
        else:
            p = dtf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = THEME['dark']
        p.font.name = "맑은 고딕"
        p.space_before = Pt(8) if i > 0 else Pt(0)

        p2 = dtf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = THEME['text']
        p2.font.name = "맑은 고딕"
        p2.space_after = Pt(6)

    return slide


def create_advantages_slide(prs):
    """아키텍처 장점 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.color.rgb = THEME['primary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "아키텍처 장점 (Architecture Advantages)"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    advantages = [
        ("💰 비용 효율성", "DeepSeek 단일 모델로 LLM 비용 95% 절감\n연간 예상 비용: ~$14,000 (GPT-4 대비 86% 절감)", THEME['success']),
        ("🎯 검색 정확도", "Vector + Graph 융합으로 Precision@5 > 85%\nRRF 알고리즘으로 두 검색 방식의 장점 결합", RGBColor(33, 150, 243)),
        ("🔐 기업 보안", "Keycloak OAuth 2.0 + PKCE 인증\nAES-256-GCM 암호화, RBAC 권한 관리", RGBColor(156, 39, 176)),
        ("📈 확장 가능성", "Docker Compose → K8s 마이그레이션 경로 확보\n마이크로서비스 구조로 독립적 확장 가능", THEME['highlight']),
    ]

    y_pos = 1.5
    for title, desc, color in advantages:
        # 아이콘 박스
        add_rounded_box(slide, 0.5, y_pos, 5, 2.3, title, color, font_size=18)

        # 설명
        desc_box = slide.shapes.add_textbox(Cm(6.5), Cm(y_pos + 0.3), Cm(12), Cm(2))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        p = dtf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = THEME['text']
        p.font.name = "맑은 고딕"

        y_pos += 3

    return slide


def create_data_flow_slide(prs):
    """데이터 흐름 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.color.rgb = THEME['primary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "데이터 흐름 (Search Flow)"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # 검색 흐름 박스들
    steps = [
        ("1", "사용자", "질문 입력", 0.5),
        ("2", "Frontend", "API 요청", 2.8),
        ("3", "Gateway", "JWT 검증", 5.1),
        ("4", "Backend", "요청 전달", 7.4),
        ("5", "AI Service", "Hybrid 검색", 9.7),
        ("6", "DeepSeek", "답변 합성", 12),
    ]

    for num, name, action, x in steps:
        # 번호
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Cm(x), Cm(2.5), Cm(1), Cm(1)
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = THEME['accent']
        num_box.line.color.rgb = THEME['accent']

        num_text = slide.shapes.add_textbox(Cm(x), Cm(2.5), Cm(1), Cm(1))
        ntf = num_text.text_frame
        ntf.text = num
        ntf.paragraphs[0].font.size = Pt(14)
        ntf.paragraphs[0].font.bold = True
        ntf.paragraphs[0].font.color.rgb = THEME['white']
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 이름
        add_rounded_box(slide, x - 0.3, 4, 2.5, 0.8, name, THEME['secondary'], font_size=12)

        # 액션
        action_box = slide.shapes.add_textbox(Cm(x - 0.3), Cm(5.2), Cm(2.5), Cm(0.8))
        atf = action_box.text_frame
        atf.text = action
        atf.paragraphs[0].font.size = Pt(10)
        atf.paragraphs[0].font.color.rgb = THEME['text']
        atf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 병렬 검색 영역
    parallel_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(0.5), Cm(7), Cm(16), Cm(3.5)
    )
    parallel_box.fill.solid()
    parallel_box.fill.fore_color.rgb = THEME['light']
    parallel_box.line.color.rgb = THEME['secondary']
    parallel_box.line.width = Pt(2)

    # 병렬 검색 제목
    par_title = slide.shapes.add_textbox(Cm(0.8), Cm(7.2), Cm(5), Cm(0.6))
    ptf = par_title.text_frame
    ptf.text = "⚡ 병렬 검색 실행"
    ptf.paragraphs[0].font.bold = True
    ptf.paragraphs[0].font.size = Pt(14)
    ptf.paragraphs[0].font.color.rgb = THEME['dark']

    # ES 검색
    add_rounded_box(slide, 1.5, 8.5, 4, 1.2, "Elasticsearch\nknn Vector Search", RGBColor(249, 183, 22), THEME['dark'], font_size=12)

    # Neo4j 검색
    add_rounded_box(slide, 7, 8.5, 4, 1.2, "Neo4j\nCypher Graph Search", RGBColor(1, 139, 255), THEME['white'], font_size=12)

    # RRF 융합
    add_rounded_box(slide, 12.5, 8.5, 3, 1.2, "RRF 융합", RGBColor(255, 107, 107), THEME['white'], font_size=12)

    # 설명 박스
    desc_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(19), Cm(4), Cm(8.5), Cm(10)
    )
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = RGBColor(226, 240, 217)
    desc_box.line.color.rgb = RGBColor(226, 240, 217)

    desc_text = slide.shapes.add_textbox(Cm(19.3), Cm(4.3), Cm(8), Cm(9.5))
    dtf = desc_text.text_frame
    dtf.word_wrap = True

    descriptions = [
        ("1. 요청 흐름", "사용자 → Frontend → Gateway → Backend → AI Service"),
        ("2. Gateway 역할", "JWT 토큰 검증, Rate Limiting, 라우팅"),
        ("3. 병렬 검색", "Vector Search와 Graph Search 동시 실행"),
        ("4. 결과 융합", "RRF 알고리즘으로 두 검색 결과 통합 후 답변 생성"),
    ]

    for i, (title, desc) in enumerate(descriptions):
        if i == 0:
            p = dtf.paragraphs[0]
        else:
            p = dtf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = THEME['dark']
        p.font.name = "맑은 고딕"
        p.space_before = Pt(8) if i > 0 else Pt(0)

        p2 = dtf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = THEME['text']
        p2.font.name = "맑은 고딕"
        p2.space_after = Pt(6)

    return slide


def create_security_slide(prs):
    """보안 아키텍처 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.color.rgb = THEME['primary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "보안 아키텍처 (Security Architecture)"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # 인증 레이어
    add_rounded_box(slide, 0.5, 2, 5, 2.5, "🔐 인증 (Authentication)\n\n• OAuth 2.0 + PKCE\n• Keycloak IdP\n• JWT (RS256)", THEME['secondary'], font_size=12, bold=False)

    # 권한 레이어
    add_rounded_box(slide, 6, 2, 5, 2.5, "👤 권한 (Authorization)\n\n• RBAC 모델\n• 3단계 역할 (Admin/Manager/User)\n• API 레벨 권한 검증", RGBColor(156, 39, 176), font_size=12, bold=False)

    # 암호화 레이어
    add_rounded_box(slide, 11.5, 2, 5, 2.5, "🔒 암호화 (Encryption)\n\n• TLS 1.3 (전송)\n• AES-256-GCM (필드)\n• TDE (저장)", RGBColor(198, 40, 40), font_size=12, bold=False)

    # 토큰 전략
    add_rounded_box(slide, 0.5, 5.5, 7.5, 2, "🎫 JWT 토큰 전략\n\n• Access Token: 15분, 메모리 저장\n• Refresh Token: 7일, HttpOnly Cookie\n• 만료 2분 전 자동 갱신", THEME['primary'], font_size=11, bold=False)

    # 데이터 분류
    add_rounded_box(slide, 9, 5.5, 7.5, 2, "📊 데이터 분류\n\n• Level 4 (극비): HSM/Vault 암호화\n• Level 3 (비밀): DB 암호화\n• Level 2 (대외비): 필드 암호화", RGBColor(255, 152, 0), THEME['dark'], font_size=11, bold=False)

    # 설명 박스
    desc_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(19), Cm(4), Cm(8.5), Cm(10)
    )
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = RGBColor(226, 240, 217)
    desc_box.line.color.rgb = RGBColor(226, 240, 217)

    desc_text = slide.shapes.add_textbox(Cm(19.3), Cm(4.3), Cm(8), Cm(9.5))
    dtf = desc_text.text_frame
    dtf.word_wrap = True

    descriptions = [
        ("1. Zero Trust 원칙", "모든 요청에 대해 인증/인가 검증 수행"),
        ("2. PKCE 적용", "Authorization Code 탈취 공격 방지"),
        ("3. 계층적 암호화", "전송/애플리케이션/저장 레벨별 암호화"),
        ("4. HashiCorp Vault", "암호화 키 중앙 집중 관리"),
    ]

    for i, (title, desc) in enumerate(descriptions):
        if i == 0:
            p = dtf.paragraphs[0]
        else:
            p = dtf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = THEME['dark']
        p.font.name = "맑은 고딕"
        p.space_before = Pt(8) if i > 0 else Pt(0)

        p2 = dtf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = THEME['text']
        p2.font.name = "맑은 고딕"
        p2.space_after = Pt(6)

    return slide


def create_infra_slide(prs):
    """인프라 구성 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.color.rgb = THEME['primary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "인프라 구성 (Docker Compose 기반)"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # Docker Host
    host_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(0.5), Cm(3), Cm(17), Cm(11)
    )
    host_box.fill.solid()
    host_box.fill.fore_color.rgb = THEME['light']
    host_box.line.color.rgb = THEME['secondary']
    host_box.line.width = Pt(3)

    # Host 제목
    host_title = slide.shapes.add_textbox(Cm(0.8), Cm(3.2), Cm(6), Cm(0.6))
    htf = host_title.text_frame
    htf.text = "🐳 Docker Host (Single Server)"
    htf.paragraphs[0].font.bold = True
    htf.paragraphs[0].font.size = Pt(14)
    htf.paragraphs[0].font.color.rgb = THEME['dark']

    # Application 컨테이너
    apps = [
        ("nginx", 1, 4.5),
        ("frontend", 3.5, 4.5),
        ("gateway", 6, 4.5),
        ("backend", 8.5, 4.5),
        ("ai-service", 11, 4.5),
        ("keycloak", 14, 4.5),
    ]

    for name, x, y in apps:
        add_rounded_box(slide, x, y, 2.2, 1, name, RGBColor(129, 236, 236), THEME['dark'], font_size=10)

    # Database 컨테이너
    dbs = [
        ("postgres", 1, 7),
        ("elastic", 3.5, 7),
        ("neo4j", 6, 7),
        ("redis", 8.5, 7),
        ("minio", 11, 7),
    ]

    for name, x, y in dbs:
        add_rounded_box(slide, x, y, 2.2, 1, name, RGBColor(162, 155, 254), THEME['white'], font_size=10)

    # Monitoring 컨테이너
    monitors = [
        ("prometheus", 1, 9.5),
        ("grafana", 3.5, 9.5),
        ("loki", 6, 9.5),
        ("promtail", 8.5, 9.5),
    ]

    for name, x, y in monitors:
        add_rounded_box(slide, x, y, 2.2, 1, name, RGBColor(253, 121, 168), THEME['white'], font_size=10)

    # 범례
    legend_box = slide.shapes.add_textbox(Cm(11), Cm(9.5), Cm(6), Cm(2))
    ltf = legend_box.text_frame
    ltf.text = "📦 총 15개 컨테이너\n• Application: 6개\n• Database: 5개\n• Monitoring: 4개"
    for para in ltf.paragraphs:
        para.font.size = Pt(10)
        para.font.color.rgb = THEME['text']

    # 설명 박스
    desc_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(19), Cm(4), Cm(8.5), Cm(10)
    )
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = RGBColor(226, 240, 217)
    desc_box.line.color.rgb = RGBColor(226, 240, 217)

    desc_text = slide.shapes.add_textbox(Cm(19.3), Cm(4.3), Cm(8), Cm(9.5))
    dtf = desc_text.text_frame
    dtf.word_wrap = True

    descriptions = [
        ("1. 단일 서버 배포", "Docker Compose로 15개 컨테이너 통합 관리"),
        ("2. 서버 사양", "Production: 32C/128GB RAM, 1TB SSD"),
        ("3. 예상 비용", "연간 ~$14,000 (K8s 대비 86% 절감)"),
        ("4. 확장 경로", "Phase 2: Redis Sentinel → Phase 3: Docker Swarm → Phase 4: K8s"),
    ]

    for i, (title, desc) in enumerate(descriptions):
        if i == 0:
            p = dtf.paragraphs[0]
        else:
            p = dtf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = THEME['dark']
        p.font.name = "맑은 고딕"
        p.space_before = Pt(8) if i > 0 else Pt(0)

        p2 = dtf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = THEME['text']
        p2.font.name = "맑은 고딕"
        p2.space_after = Pt(6)

    return slide


def create_phase2_slide(prs):
    """2단계 이관 항목 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = THEME['primary']
    header.line.color.rgb = THEME['primary']

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "2단계 구축 사업 이관 항목 (Phase 2 Items)"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"

    # 이관 문서
    add_rounded_box(slide, 0.5, 2, 8, 2, "📋 이관 대상 문서\n\n• 성능/확장성 설계서 (2-3주)\n• 재해복구 설계서 (1-2주)", RGBColor(255, 107, 107), font_size=12, bold=False)

    # Medium Priority
    add_rounded_box(slide, 9.5, 2, 7.5, 2, "⚡ Medium Priority\n\n• 캐싱 전략 상세화 (0.5일)\n• MFA 설계 추가 (0.5일)", RGBColor(255, 193, 7), THEME['dark'], font_size=12, bold=False)

    # Low Priority
    add_rounded_box(slide, 0.5, 5, 8, 2.5, "📌 Low Priority\n\n• 데이터 거버넌스 설계서 (1일)\n• 모델 버전 관리 (0.5일)\n• PWA 설계 추가 (0.5일)\n• 통합 테스트 계획서 (1일)", RGBColor(158, 158, 158), font_size=12, bold=False)

    # 마이그레이션 로드맵
    roadmap_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(10), Cm(6.5), Cm(7.5), Cm(5)
    )
    roadmap_box.fill.solid()
    roadmap_box.fill.fore_color.rgb = THEME['light']
    roadmap_box.line.color.rgb = THEME['secondary']

    roadmap_title = slide.shapes.add_textbox(Cm(10.3), Cm(6.8), Cm(7), Cm(0.6))
    rtf = roadmap_title.text_frame
    rtf.text = "🚀 확장 로드맵"
    rtf.paragraphs[0].font.bold = True
    rtf.paragraphs[0].font.size = Pt(12)
    rtf.paragraphs[0].font.color.rgb = THEME['dark']

    roadmap_text = slide.shapes.add_textbox(Cm(10.3), Cm(7.8), Cm(7), Cm(3.5))
    rttf = roadmap_text.text_frame
    rttf.word_wrap = True

    phases = [
        "Phase 1 (현재): Docker Compose - $14K/년",
        "Phase 2: + Redis Sentinel - $20K/년",
        "Phase 3: Docker Swarm 3노드 - $40K/년",
        "Phase 4: Kubernetes 풀 클러스터 - $100K/년"
    ]

    for i, phase in enumerate(phases):
        if i == 0:
            p = rttf.paragraphs[0]
        else:
            p = rttf.add_paragraph()
        p.text = f"• {phase}"
        p.font.size = Pt(10)
        p.font.color.rgb = THEME['text']
        p.font.name = "맑은 고딕"

    # 설명 박스
    desc_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(19), Cm(4), Cm(8.5), Cm(10)
    )
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = RGBColor(226, 240, 217)
    desc_box.line.color.rgb = RGBColor(226, 240, 217)

    desc_text = slide.shapes.add_textbox(Cm(19.3), Cm(4.3), Cm(8), Cm(9.5))
    dtf = desc_text.text_frame
    dtf.word_wrap = True

    descriptions = [
        ("1. 1단계 집중", "핵심 기능 구현에 집중하여 빠른 가치 전달"),
        ("2. 점진적 확장", "사용량 증가에 따라 단계적 인프라 확장"),
        ("3. 마이그레이션 조건", "동시 사용자 100+ → Phase 2\n동시 사용자 500+ → Phase 3"),
        ("4. 비용 최적화", "필요 시점에 맞춰 투자로 비용 효율성 유지"),
    ]

    for i, (title, desc) in enumerate(descriptions):
        if i == 0:
            p = dtf.paragraphs[0]
        else:
            p = dtf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = THEME['dark']
        p.font.name = "맑은 고딕"
        p.space_before = Pt(8) if i > 0 else Pt(0)

        p2 = dtf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = THEME['text']
        p2.font.name = "맑은 고딕"
        p2.space_after = Pt(6)

    return slide


def create_closing_slide(prs):
    """마무리 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME['dark']
    bg.line.color.rgb = THEME['dark']

    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    tf.text = "감사합니다"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = THEME['white']
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
    sf = sub_box.text_frame
    sf.text = "Hybrid RAG Knowledge Platform\n기업 지식의 80%가 잠들어 있는 문서에서, AI가 즉시 답을 찾아드립니다"
    for para in sf.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = THEME['light']
        para.font.name = "맑은 고딕"
        para.alignment = PP_ALIGN.CENTER

    # 연락처
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(0.5))
    cf = contact_box.text_frame
    cf.text = "2026-01-16  |  Version 1.0"
    cp = cf.paragraphs[0]
    cp.font.size = Pt(14)
    cp.font.color.rgb = THEME['accent']
    cp.font.name = "맑은 고딕"
    cp.alignment = PP_ALIGN.CENTER

    return slide


def main():
    """메인 함수"""
    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 비율
    prs.slide_height = Inches(7.5)

    print("프레젠테이션 생성 시작...")

    # 1. 표지
    create_title_slide(prs,
                       "Hybrid RAG Knowledge Platform\n통합 상세 설계서",
                       "Vector + Graph 기반 차세대 지식 검색 시스템")
    print("  - 표지 슬라이드 생성")

    # 2. Executive Summary
    create_section_slide(prs, "Executive Summary", 1)
    create_content_slide(prs, "프로젝트 비전", [
        '"기업 지식의 80%가 잠들어 있는 문서에서, AI가 즉시 답을 찾아드립니다"',
        "Vector Search + Graph Search 결합한 차세대 Hybrid RAG 시스템",
        "DeepSeek 단일 모델로 LLM 비용 95% 절감 (GPT-4 대비)",
        "기업 보안 요구사항 충족: OAuth 2.0, AES-256, RBAC",
    ])
    print("  - Executive Summary 생성")

    # 3. 핵심 수치
    create_table_slide(prs, "핵심 목표 수치 (Key Metrics)",
                       ["항목", "목표", "비고"],
                       [
                           ["응답 시간", "< 3초 (P95)", "Hybrid 검색 + 답변 생성"],
                           ["검색 정확도", "> 85% (Precision@5)", "상위 5개 결과 기준"],
                           ["LLM 비용", "95% 절감", "GPT-4 대비 DeepSeek"],
                           ["가용성", "99.5%", "SLA 기준"],
                           ["동시 사용자", "100명", "1단계 목표"],
                       ])
    print("  - 핵심 수치 슬라이드 생성")

    # 4. 시스템 아키텍처
    create_section_slide(prs, "시스템 아키텍처", 2)
    create_architecture_slide(prs)
    print("  - 시스템 아키텍처 슬라이드 생성")

    # 5. 아키텍처 장점
    create_advantages_slide(prs)
    print("  - 아키텍처 장점 슬라이드 생성")

    # 6. 플랫폼 핵심 설계
    create_section_slide(prs, "플랫폼 핵심 설계", 3)
    create_vip_architecture_slide(prs)
    create_hybrid_search_slide(prs)
    print("  - VIP 아키텍처, Hybrid Search 슬라이드 생성")

    # 7. 데이터 흐름
    create_data_flow_slide(prs)
    print("  - 데이터 흐름 슬라이드 생성")

    # 8. 백엔드 설계
    create_section_slide(prs, "백엔드 설계", 4)
    create_content_slide(prs, "백엔드 아키텍처 (Spring Boot 3.x)", [
        "Layered Architecture: Presentation → Application → Domain → Infrastructure",
        "모듈 구조: platform-common, platform-domain, platform-gateway, platform-api",
        "AI Service 연동: WebClient + Resilience4j (Circuit Breaker)",
        "장애 복구: failureRateThreshold 50%, waitDurationInOpenState 30초",
        "서비스 분리 원칙: SpringBoot는 비즈니스 로직, AI Service는 LLM/임베딩",
    ])
    print("  - 백엔드 설계 슬라이드 생성")

    # 9. 프론트엔드 설계
    create_section_slide(prs, "프론트엔드 설계", 5)
    create_content_slide(prs, "프론트엔드 아키텍처 (React 18 + TypeScript)", [
        "컴포넌트 계층: Pages → Features → Components (3계층)",
        "상태 관리: Redux Toolkit (UI 상태) + React Query (서버 상태)",
        "폼 관리: React Hook Form + Zod 스키마 검증",
        "UI 라이브러리: MUI v5 기반 커스텀 디자인 시스템",
        "핵심 기능: 검색 모드, 채팅 모드, 지식 관리, 북마크, 내보내기",
    ])
    print("  - 프론트엔드 설계 슬라이드 생성")

    # 10. 보안 설계
    create_section_slide(prs, "인증 및 보안 설계", 6)
    create_security_slide(prs)
    print("  - 보안 설계 슬라이드 생성")

    # 11. 인프라 설계
    create_section_slide(prs, "인프라 및 DevOps 설계", 7)
    create_infra_slide(prs)
    print("  - 인프라 설계 슬라이드 생성")

    # 12. CI/CD
    create_content_slide(prs, "CI/CD 파이프라인 (GitLab)", [
        "CI 단계: Git Push → Build → Test → Security Scan → SonarQube",
        "CD 단계: Docker Image → Staging Deploy → Manual Approve → Production",
        "브랜치 전략: main (프로덕션), develop (개발), feature/*, fix/*, hotfix/*",
        "품질 게이트: 테스트 커버리지 80%+, 보안 취약점 High 0건",
        "모니터링: Prometheus + Grafana + Loki (메트릭/대시보드/로그)",
    ])
    print("  - CI/CD 슬라이드 생성")

    # 13. API 설계
    create_section_slide(prs, "API 통합 설계", 8)
    create_table_slide(prs, "주요 API 엔드포인트",
                       ["API", "메서드", "경로", "설명"],
                       [
                           ["External", "POST", "/api/v1/search/chat", "채팅 검색 (RAG)"],
                           ["External", "POST", "/api/v1/knowledge", "지식 등록"],
                           ["External", "GET", "/api/v1/knowledge/{id}", "지식 상세"],
                           ["Internal", "POST", "/internal/v1/search/hybrid", "Hybrid 검색"],
                           ["Internal", "POST", "/internal/v1/extract/entities", "엔티티 추출"],
                       ])
    print("  - API 설계 슬라이드 생성")

    # 14. 품질 보증
    create_section_slide(prs, "품질 보증 전략", 9)
    create_content_slide(prs, "테스트 및 품질 관리", [
        "테스트 피라미드: Unit(60%) → Integration(30%) → E2E(10%)",
        "커버리지 목표: Backend 80%+, AI Service 75%+, Frontend 70%+",
        "코드 품질: SonarQube, ESLint, Black/isort, Checkstyle",
        "RAG 평가: Precision@K, Recall@K, MRR, Faithfulness, TTFB",
        "성능 목표: P95 Latency < 3초, 검색 정확도 > 85%",
    ])
    print("  - 품질 보증 슬라이드 생성")

    # 15. 2단계 이관 항목
    create_section_slide(prs, "2단계 구축 이관 항목", 10)
    create_phase2_slide(prs)
    print("  - 2단계 이관 슬라이드 생성")

    # 16. 기술 스택
    create_section_slide(prs, "기술 스택 요약", 11)
    create_table_slide(prs, "기술 스택 (Tech Stack)",
                       ["영역", "기술", "버전", "용도"],
                       [
                           ["Frontend", "React + TypeScript", "18.3+ / 5.4+", "UI 개발"],
                           ["Backend", "Spring Boot", "3.2+", "비즈니스 로직"],
                           ["AI Service", "FastAPI + LangGraph", "0.110+ / 1.0+", "AI 파이프라인"],
                           ["Database", "PostgreSQL + ES + Neo4j", "16+ / 8.x / 5.x", "데이터 저장"],
                           ["Auth", "Keycloak", "24.x", "OAuth 2.0 IdP"],
                           ["Infra", "Docker Compose", "2.x", "컨테이너 오케스트레이션"],
                       ])
    print("  - 기술 스택 슬라이드 생성")

    # 17. 비용
    create_table_slide(prs, "예상 비용 (Annual Cost)",
                       ["항목", "월간 비용", "연간 비용", "비고"],
                       [
                           ["서버 (Production)", "$800", "$9,600", "32C/128G"],
                           ["서버 (Staging)", "$200", "$2,400", "16C/64G"],
                           ["DeepSeek API", "$150", "$1,800", "예상 사용량"],
                           ["도메인/SSL", "$20", "$240", ""],
                           ["합계", "$1,170", "$14,040", "GPT-4 대비 86% 절감"],
                       ])
    print("  - 비용 슬라이드 생성")

    # 18. 마무리
    create_closing_slide(prs)
    print("  - 마무리 슬라이드 생성")

    # 저장
    output_path = "/mnt/d/Users/KTDS/Documents/06.과제/hybrid-rag-knowledge-ops/docs/Hybrid_RAG_Platform_Design_Brown.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

    print(f"\n✅ 프레젠테이션 저장 완료: {output_path}")
    print(f"   총 슬라이드 수: {len(prs.slides)}장")
    print(f"   테마: Brown (브라운)")


if __name__ == "__main__":
    main()
